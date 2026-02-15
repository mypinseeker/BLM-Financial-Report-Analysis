#!/usr/bin/env python3
"""Xavier Niel Consolidated Group Report Generator.

Generates a comprehensive BLM strategic assessment covering all 18 operators
across NJJ Holding, Iliad Group, and Millicom/Tigo.

Outputs:
  - MD:   data/output/group/blm_niel_group_consolidated_cq4_2025.md
  - PDF:  data/output/group/blm_niel_group_consolidated_cq4_2025.pdf
  - PDF (zh): data/output/group/blm_niel_group_consolidated_cq4_2025_zh.pdf
  - PPTX: data/output/group/blm_niel_group_consolidated_cq4_2025.pptx

Usage:
    python3 generate_niel_group_report.py
"""
from __future__ import annotations

import sys
import tempfile
from pathlib import Path

_project_root = Path(__file__).resolve().parent
if str(_project_root) not in sys.path:
    sys.path.insert(0, str(_project_root))

# ---------------------------------------------------------------------------
# Constants — the 18 Xavier Niel operators across 3 groups
# ---------------------------------------------------------------------------

PERIOD = "CQ4_2025"
OUTPUT_DIR = Path("data/output/group")

# (operator_id, market_id, display_name, country)
NJJ_OPERATORS = [
    ("salt_ch", "switzerland", "Salt", "Switzerland"),
    ("eir_ie", "ireland", "eir", "Ireland"),
    ("lifecell_ua", "ukraine", "lifecell", "Ukraine"),
    ("epic_cy", "cyprus", "Epic Cyprus", "Cyprus"),
    ("epic_mt", "malta", "Epic Malta", "Malta"),
]

ILIAD_OPERATORS = [
    ("free_fr", "france", "Free", "France"),
    ("iliad_it", "italy", "Iliad Italia", "Italy"),
    ("play_pl", "poland", "Play", "Poland"),
]

MILLICOM_OPERATORS = [
    ("tigo_guatemala", "guatemala", "Tigo Guatemala", "Guatemala"),
    ("tigo_colombia", "colombia", "Tigo Colombia", "Colombia"),
    ("tigo_honduras", "honduras", "Tigo Honduras", "Honduras"),
    ("tigo_paraguay", "paraguay", "Tigo Paraguay", "Paraguay"),
    ("tigo_bolivia", "bolivia", "Tigo Bolivia", "Bolivia"),
    ("tigo_el_salvador", "el_salvador", "Tigo El Salvador", "El Salvador"),
    ("tigo_panama", "panama", "Tigo Panama", "Panama"),
    ("tigo_nicaragua", "nicaragua", "Tigo Nicaragua", "Nicaragua"),
    ("tigo_ecuador", "ecuador", "Tigo Ecuador", "Ecuador"),
    ("tigo_uruguay", "uruguay", "Tigo Uruguay", "Uruguay"),
]

ALL_OPERATORS = NJJ_OPERATORS + ILIAD_OPERATORS + MILLICOM_OPERATORS

GROUPS = [
    ("njj_holding", "NJJ Holding", NJJ_OPERATORS),
    ("iliad_group", "Iliad Group", ILIAD_OPERATORS),
    ("millicom", "Millicom / Tigo", MILLICOM_OPERATORS),
]


# =========================================================================
# Step 1: Seed & Analyze
# =========================================================================

def seed_and_analyze():
    """Seed all 22 markets and run Five Looks for 18 target operators."""
    from src.database.seed_orchestrator import seed_all_markets
    from src.blm.engine import BLMAnalysisEngine
    from src.web.services.group_summary import GroupSummaryGenerator
    from src.database.operator_directory import OPERATOR_GROUPS
    from src.output.strategic_diagnosis import StrategicDiagnosisComputer
    from src.models.market_configs import get_market_config

    print("=" * 70)
    print("  Xavier Niel Consolidated Group Report Generator")
    print(f"  Operators: {len(ALL_OPERATORS)} across 3 groups")
    print(f"  Period:    {PERIOD}")
    print("=" * 70)

    # Seed
    print("\n[Phase 1] Seeding all 22 markets into in-memory SQLite...")
    db = seed_all_markets(":memory:")

    # Run Five Looks for each operator
    print(f"\n[Phase 2] Running Five Looks for {len(ALL_OPERATORS)} operators...")
    all_results = {}  # operator_id -> FiveLooksResult
    all_diagnoses = {}  # operator_id -> StrategicDiagnosis

    for i, (op_id, market_id, display, country) in enumerate(ALL_OPERATORS, 1):
        print(f"  [{i:2d}/{len(ALL_OPERATORS)}] {display} ({country})...", end=" ", flush=True)
        try:
            engine = BLMAnalysisEngine(
                db=db, target_operator=op_id, market=market_id,
                target_period=PERIOD, n_quarters=8,
            )
            result = engine.run_five_looks()
            all_results[op_id] = result

            config = get_market_config(market_id)
            diag = StrategicDiagnosisComputer(result, config).compute()
            all_diagnoses[op_id] = diag
            print("OK")
        except Exception as e:
            print(f"FAILED: {e}")

    # Group summaries
    print(f"\n[Phase 3] Generating group summaries...")
    gen = GroupSummaryGenerator()
    group_summaries = {}
    for group_id, group_name, operators in GROUPS:
        group_results = {
            mkt: all_results[op] for op, mkt, _, _ in operators if op in all_results
        }
        if group_results:
            group_info = OPERATOR_GROUPS.get(group_id, {"group_id": group_id, "group_name": group_name})
            group_summaries[group_id] = gen.generate(group_results, group_info)
            print(f"  {group_name}: {len(group_results)} markets")

    db.close()
    return all_results, all_diagnoses, group_summaries


# =========================================================================
# Step 2: Generate Consolidated MD Report
# =========================================================================

def _safe_financial(result):
    """Extract key financial metrics from a FiveLooksResult."""
    sa = result.self_analysis
    if sa is None:
        return {}
    fh = getattr(sa, "financial_health", None) or {}
    if isinstance(fh, dict):
        return fh
    # Might be a dataclass — convert
    return fh if isinstance(fh, dict) else {}


def _fmt(val, suffix="", prefix="", decimals=1):
    """Format a numeric value."""
    if val is None:
        return "N/A"
    try:
        n = float(val)
        if suffix == "%" and prefix == "+":
            return f"{n:+.{decimals}f}%"
        if suffix == "%":
            return f"{n:.{decimals}f}%"
        if abs(n) >= 1000:
            return f"{prefix}{n:,.0f}{suffix}"
        return f"{prefix}{n:,.{decimals}f}{suffix}"
    except (TypeError, ValueError):
        return str(val)


def _get_metric(fh, *keys):
    """Try multiple keys to get a metric from financial_health dict."""
    for k in keys:
        v = fh.get(k)
        if v is not None:
            return v
    return None


def generate_md(all_results, all_diagnoses, group_summaries):
    """Build the consolidated Markdown report."""
    lines = []

    # ── Header ──
    lines.append("# Xavier Niel 电信帝国 — BLM战略评估综合报告 (CQ4 2025)")
    lines.append("")
    lines.append(f"> **报告日期**: 2026-02-15 | **分析期间**: CQ4 2025 (October–December 2025)")
    lines.append(f"> **覆盖范围**: 18个运营商 × 3大集团 × 15个国家")
    lines.append(f"> **分析框架**: 华为BLM五看方法论 (PEST + $APPEALS + Porter's Five Forces + BMC + SPAN)")
    lines.append("")

    # ── Table of Contents ──
    lines.append("## 目录")
    lines.append("")
    lines.append("1. [第一部分：全局总览](#第一部分全局总览)")
    lines.append("2. [第二部分：NJJ Holding（5个市场）](#第二部分njj-holding5个市场)")
    lines.append("3. [第三部分：Iliad Group（3个市场）](#第三部分iliad-group3个市场)")
    lines.append("4. [第四部分：Millicom/Tigo（10个市场）](#第四部分millicomtigo10个市场)")
    lines.append("5. [第五部分：跨集团战略建议](#第五部分跨集团战略建议)")
    lines.append("")
    lines.append("---")
    lines.append("")

    # ══════════════════════════════════════════════════════════════════════
    # Part 1: Global Overview
    # ══════════════════════════════════════════════════════════════════════
    lines.append("## 第一部分：全局总览")
    lines.append("")
    lines.append("### 1.1 18个市场总览表")
    lines.append("")

    # Build overview table
    headers = ["集团", "运营商", "国家", "营收 (M)", "EBITDA利润率", "移动用户 (K)", "市场排名", "健康评级"]
    rows = []
    total_revenue = 0
    total_mobile_k = 0

    for group_id, group_name, operators in GROUPS:
        for op_id, market_id, display, country in operators:
            if op_id not in all_results:
                rows.append([group_name, display, country, "N/A", "N/A", "N/A", "N/A", "N/A"])
                continue
            result = all_results[op_id]
            fh = _safe_financial(result)
            diag = all_diagnoses.get(op_id)

            rev = _get_metric(fh, "total_revenue", "revenue", "quarterly_revenue")
            margin = _get_metric(fh, "ebitda_margin_pct", "ebitda_margin")
            mobile_k = _get_metric(fh, "mobile_total_k", "mobile_subscribers_k")
            health = getattr(result.self_analysis, "health_rating", "N/A") if result.self_analysis else "N/A"
            rank = diag.operator_rank if diag else 0

            if rev is not None:
                try:
                    total_revenue += float(rev)
                except (TypeError, ValueError):
                    pass
            if mobile_k is not None:
                try:
                    total_mobile_k += float(mobile_k)
                except (TypeError, ValueError):
                    pass

            rows.append([
                group_name,
                display,
                country,
                _fmt(rev, "M"),
                _fmt(margin, "%"),
                _fmt(mobile_k, "K", decimals=0),
                f"#{rank}" if rank else "N/A",
                health,
            ])

    lines.append(_md_table(headers, rows))
    lines.append("")

    # Aggregate financials
    lines.append("### 1.2 集团合计财务数据")
    lines.append("")
    lines.append(f"- **合并营收 (季度)**: {total_revenue:,.0f}M (所有18个运营商)")
    lines.append(f"- **合并移动用户**: {total_mobile_k:,.0f}K")
    lines.append(f"- **覆盖国家**: 15")
    lines.append(f"- **分析完成率**: {len(all_results)}/{len(ALL_OPERATORS)} operators")
    lines.append("")

    # Per-group aggregate
    for group_id, group_name, operators in GROUPS:
        g_rev = 0
        g_mobile = 0
        g_count = 0
        for op_id, market_id, _, _ in operators:
            if op_id in all_results:
                fh = _safe_financial(all_results[op_id])
                r = _get_metric(fh, "total_revenue", "revenue", "quarterly_revenue")
                m = _get_metric(fh, "mobile_total_k", "mobile_subscribers_k")
                if r: g_rev += float(r)
                if m: g_mobile += float(m)
                g_count += 1
        lines.append(f"**{group_name}**: {g_count} markets, Revenue {g_rev:,.0f}M, Mobile {g_mobile:,.0f}K")
        lines.append("")

    # Cross-group opportunities and threats
    lines.append("### 1.3 跨集团共性机会 & 威胁")
    lines.append("")

    # Aggregate common opportunities
    all_opps = set()
    all_threats = set()
    for gid, gsummary in group_summaries.items():
        for opp in gsummary.get("common_opportunities", [])[:5]:
            all_opps.add(opp)
        for threat in gsummary.get("common_threats", [])[:5]:
            all_threats.add(threat)

    lines.append("**共性机会 (Opportunities):**")
    for opp in sorted(all_opps)[:10]:
        lines.append(f"- {opp}")
    lines.append("")
    lines.append("**共性威胁 (Threats):**")
    for threat in sorted(all_threats)[:10]:
        lines.append(f"- {threat}")
    lines.append("")

    # Strategic diagnosis summary
    lines.append("### 1.4 战略诊断：整体投资组合评估")
    lines.append("")
    diagnosis_counts = {}
    for op_id, diag in all_diagnoses.items():
        label = diag.central_diagnosis_label
        diagnosis_counts.setdefault(label, []).append(op_id)

    for label, ops in sorted(diagnosis_counts.items(), key=lambda x: -len(x[1])):
        op_displays = []
        for op in ops:
            for _, _, display, country in ALL_OPERATORS:
                # Match by finding the entry in ALL_OPERATORS
                pass
            for o_id, _, disp, ctry in ALL_OPERATORS:
                if o_id == op:
                    op_displays.append(f"{disp} ({ctry})")
                    break
        lines.append(f"- **{label}** ({len(ops)}): {', '.join(op_displays)}")
    lines.append("")

    lines.append("---")
    lines.append("")

    # ══════════════════════════════════════════════════════════════════════
    # Parts 2-4: Per-Group Sections
    # ══════════════════════════════════════════════════════════════════════
    part_names = {
        "njj_holding": ("第二部分：NJJ Holding（5个市场）", "NJJ Holding"),
        "iliad_group": ("第三部分：Iliad Group（3个市场）", "Iliad Group"),
        "millicom": ("第四部分：Millicom/Tigo（10个市场）", "Millicom / Tigo"),
    }

    for group_id, group_name, operators in GROUPS:
        part_title, short_name = part_names[group_id]
        gsummary = group_summaries.get(group_id, {})

        lines.append(f"## {part_title}")
        lines.append("")

        # Group summary table
        lines.append(f"### {short_name} 集团小结")
        lines.append("")

        # Revenue comparison
        rev_comp = gsummary.get("revenue_comparison", {})
        if rev_comp:
            headers = ["市场", "营收 (M)", "营收增长", "EBITDA利润率"]
            rows = []
            for market, data in sorted(rev_comp.items()):
                rows.append([
                    market.replace("_", " ").title(),
                    _fmt(data.get("total_revenue"), "M"),
                    _fmt(data.get("revenue_growth_pct"), "%", "+"),
                    _fmt(data.get("ebitda_margin_pct"), "%"),
                ])
            lines.append("**营收对比:**")
            lines.append("")
            lines.append(_md_table(headers, rows))
            lines.append("")

        # Subscriber comparison
        sub_comp = gsummary.get("subscriber_comparison", {})
        if sub_comp:
            headers = ["市场", "移动用户 (K)", "宽带用户 (K)", "ARPU"]
            rows = []
            for market, data in sorted(sub_comp.items()):
                rows.append([
                    market.replace("_", " ").title(),
                    _fmt(data.get("mobile_subs_k"), "K", decimals=0),
                    _fmt(data.get("broadband_subs_k"), "K", decimals=0),
                    _fmt(data.get("mobile_arpu")),
                ])
            lines.append("**用户对比:**")
            lines.append("")
            lines.append(_md_table(headers, rows))
            lines.append("")

        # Key findings
        findings = gsummary.get("key_findings", [])
        if findings:
            lines.append("**关键发现:**")
            for f in findings:
                lines.append(f"- {f}")
            lines.append("")

        # Per-market breakdowns
        lines.append(f"### {short_name} 各市场摘要")
        lines.append("")

        for op_id, market_id, display, country in operators:
            lines.append(f"#### {display} ({country})")
            lines.append("")

            if op_id not in all_results:
                lines.append("*分析失败 — 数据不可用*")
                lines.append("")
                continue

            result = all_results[op_id]
            diag = all_diagnoses.get(op_id)
            fh = _safe_financial(result)

            # One-line verdict
            if diag and diag.one_line_verdict:
                lines.append(f"> {diag.one_line_verdict}")
                lines.append("")

            # Key metrics
            rev = _get_metric(fh, "total_revenue", "revenue", "quarterly_revenue")
            margin = _get_metric(fh, "ebitda_margin_pct", "ebitda_margin")
            mobile_k = _get_metric(fh, "mobile_total_k", "mobile_subscribers_k")
            share = _get_metric(fh, "market_share_pct", "revenue_share_pct")
            health = getattr(result.self_analysis, "health_rating", "N/A") if result.self_analysis else "N/A"

            lines.append(f"- **营收**: {_fmt(rev, 'M')} | **EBITDA利润率**: {_fmt(margin, '%')} | **移动用户**: {_fmt(mobile_k, 'K', decimals=0)}")
            lines.append(f"- **市场份额**: {_fmt(share, '%')} | **健康评级**: {health} | **市场排名**: #{diag.operator_rank if diag else 'N/A'}")
            lines.append(f"- **战略诊断**: {diag.central_diagnosis_label if diag else 'N/A'}")
            lines.append(f"- **竞争态势**: {diag.competitive_stance if diag else 'N/A'}")
            lines.append("")

            # Top 3 priorities
            if diag and diag.priorities:
                lines.append("**Top 3 战略优先级:**")
                for p in diag.priorities[:3]:
                    lines.append(f"1. **{p.get('name', '')}** ({p.get('priority', '')}) — {p.get('description', '')[:100]}")
                lines.append("")

            # Key risks from SWOT threats
            if result.swot:
                threats = getattr(result.swot, "threats", [])
                if threats:
                    lines.append("**关键风险:**")
                    for t in threats[:3]:
                        lines.append(f"- {str(t)[:120]}")
                    lines.append("")

            # Momentum phase
            if result.self_analysis:
                momentum = getattr(result.self_analysis, "momentum_phase", None)
                if momentum:
                    lines.append(f"**动量阶段 (Momentum)**: {momentum}")
                    lines.append("")

        lines.append("---")
        lines.append("")

    # ══════════════════════════════════════════════════════════════════════
    # Part 5: Cross-Group Strategic Recommendations
    # ══════════════════════════════════════════════════════════════════════
    lines.append("## 第五部分：跨集团战略建议")
    lines.append("")

    lines.append("### 5.1 协同效应 & 资源共享机会")
    lines.append("")
    lines.append("1. **采购协同 (Procurement Synergies)**: 18个运营商的网络设备、终端和IT系统联合采购，可降低单位成本15-25%")
    lines.append("2. **技术平台共享 (Shared Technology Platforms)**: BSS/OSS平台统一化，减少重复投资")
    lines.append("3. **最佳实践转移 (Best Practice Transfer)**: Iliad在法国的低成本运营模式可推广至NJJ和Millicom市场")
    lines.append("4. **漫游 & 互联 (Roaming & Interconnect)**: 跨集团优惠漫游费率，提升客户体验")
    lines.append("5. **数字服务平台 (Digital Services)**: 共享移动金融(Tigo Money)、企业云和安全服务平台")
    lines.append("")

    lines.append("### 5.2 风险矩阵")
    lines.append("")
    risk_headers = ["风险类别", "描述", "影响程度", "涉及市场"]
    risk_rows = [
        ["监管风险", "各国电信监管政策变化、频谱拍卖条件", "高", "全部18个市场"],
        ["汇率风险", "LATAM货币波动、UAH贬值", "高", "Millicom(10) + Ukraine"],
        ["地缘政治", "乌克兰冲突持续、LATAM政治不稳定", "中-高", "lifecell + LATAM"],
        ["竞争加剧", "价格战、新进入者、OTT替代", "中", "France, Italy, Poland"],
        ["技术转型", "5G投资回报不确定、网络升级成本", "中", "欧洲市场为主"],
        ["债务/融资", "高利率环境下的再融资压力", "中", "Millicom, eir"],
    ]
    lines.append(_md_table(risk_headers, risk_rows))
    lines.append("")

    lines.append("### 5.3 优先级排序")
    lines.append("")
    lines.append("**P0 — 立即行动 (0-6个月):**")
    lines.append("1. 稳定Millicom LATAM市场的EBITDA利润率")
    lines.append("2. 加速Iliad Italia的用户增长和市场份额扩大")
    lines.append("3. lifecell乌克兰的Datagroup-Volia整合和固移融合")
    lines.append("")
    lines.append("**P1 — 中期计划 (6-18个月):**")
    lines.append("1. NJJ欧洲市场的5G覆盖加速（Salt/eir/Epic）")
    lines.append("2. Millicom B2B和数字服务扩展（Tigo Money, 企业云）")
    lines.append("3. Play波兰的固网宽带渗透率提升")
    lines.append("")
    lines.append("**P2 — 长期战略 (18-36个月):**")
    lines.append("1. 跨集团采购协同落地")
    lines.append("2. 统一BSS/OSS技术平台")
    lines.append("3. 新市场进入评估（非洲、东南亚）")
    lines.append("")

    lines.append("---")
    lines.append("")
    lines.append("*报告由 BLM 分析引擎自动生成 | 数据来源: 公开财报、监管数据、行业报告 | 置信度: 中-高*")

    return "\n".join(lines)


def _md_table(headers, rows):
    """Build a Markdown table."""
    lines = []
    lines.append("| " + " | ".join(str(h) for h in headers) + " |")
    lines.append("| " + " | ".join("---" for _ in headers) + " |")
    for row in rows:
        cells = list(row) + [""] * (len(headers) - len(row))
        lines.append("| " + " | ".join(str(c) for c in cells[:len(headers)]) + " |")
    return "\n".join(lines)


# =========================================================================
# Step 3: Convert to PDF
# =========================================================================

def generate_pdf(md_content: str, output_path: Path, chinese: bool = True):
    """Convert Markdown to PDF using weasyprint."""
    import markdown
    from weasyprint import HTML

    # Convert MD to HTML
    html_body = markdown.markdown(
        md_content,
        extensions=["tables", "toc", "fenced_code"],
    )

    # CSS for professional styling
    font_family = "'Noto Sans CJK SC', 'Noto Sans', Arial, sans-serif" if chinese else "'Noto Sans', Arial, sans-serif"

    css = f"""
    @page {{
        size: A4;
        margin: 2cm 1.5cm;
        @bottom-center {{
            content: counter(page) " / " counter(pages);
            font-size: 9pt;
            color: #888;
        }}
    }}
    body {{
        font-family: {font_family};
        font-size: 10pt;
        line-height: 1.6;
        color: #333;
    }}
    h1 {{
        font-size: 22pt;
        color: #003366;
        border-bottom: 3px solid #003366;
        padding-bottom: 8px;
        page-break-before: avoid;
    }}
    h2 {{
        font-size: 16pt;
        color: #003366;
        border-bottom: 1px solid #ccc;
        padding-bottom: 4px;
        page-break-before: always;
        margin-top: 1em;
    }}
    h2:first-of-type {{
        page-break-before: avoid;
    }}
    h3 {{
        font-size: 13pt;
        color: #1a5276;
        margin-top: 0.8em;
    }}
    h4 {{
        font-size: 11pt;
        color: #2c3e50;
        margin-top: 0.6em;
    }}
    table {{
        width: 100%;
        border-collapse: collapse;
        margin: 0.8em 0;
        font-size: 8.5pt;
    }}
    th {{
        background-color: #003366;
        color: white;
        padding: 6px 8px;
        text-align: left;
        font-weight: bold;
    }}
    td {{
        border: 1px solid #ddd;
        padding: 5px 8px;
    }}
    tr:nth-child(even) {{
        background-color: #f8f9fa;
    }}
    blockquote {{
        border-left: 4px solid #003366;
        margin: 0.5em 0;
        padding: 0.5em 1em;
        background-color: #f0f4f8;
        font-style: italic;
    }}
    ul, ol {{
        margin: 0.3em 0;
        padding-left: 1.5em;
    }}
    li {{
        margin-bottom: 0.2em;
    }}
    strong {{
        color: #1a5276;
    }}
    hr {{
        border: none;
        border-top: 1px solid #ccc;
        margin: 1.5em 0;
    }}
    """

    full_html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>{css}</style>
</head>
<body>
{html_body}
</body>
</html>"""

    HTML(string=full_html).write_pdf(str(output_path))


# =========================================================================
# Step 4: Generate PPT
# =========================================================================

def generate_ppt(all_results, all_diagnoses, group_summaries, output_path: Path):
    """Generate consolidated group PPT deck."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color scheme
    NAVY = (0, 51, 102)
    WHITE = (255, 255, 255)
    LIGHT_GRAY = (245, 245, 245)
    DARK_GRAY = (51, 51, 51)
    ACCENT_BLUE = (26, 82, 118)
    TEAL = (0, 128, 128)

    slide_num = [0]  # Mutable counter

    def new_slide():
        slide_num[0] += 1
        return prs.slides.add_slide(prs.slide_layouts[6])

    def add_shape(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color)
        shape.line.fill.background()
        return shape

    def add_text(slide, left, top, width, height, text, size=14,
                 color=DARK_GRAY, bold=False, align="left"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(*color)
        p.font.bold = bold
        p.font.name = "Arial"
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        return txBox

    def add_bullet_text(slide, left, top, width, height, items, size=12, color=DARK_GRAY):
        """Add a multi-line bulleted text box."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(size)
            p.font.color.rgb = RGBColor(*color)
            p.font.name = "Arial"
            p.space_after = Pt(4)
        return txBox

    def add_header(slide, title, subtitle=""):
        add_shape(slide, 0, 0, Inches(13.333), Inches(0.08), NAVY)
        add_text(slide, Inches(0.5), Inches(0.25), Inches(10), Inches(0.6),
                 title, size=24, color=DARK_GRAY, bold=True)
        if subtitle:
            add_text(slide, Inches(0.5), Inches(0.85), Inches(10), Inches(0.35),
                     subtitle.upper(), size=12, color=NAVY, bold=True)
        add_text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.35),
                 str(slide_num[0]), size=10, color=(128, 128, 128), align="right")

    def add_key_message(slide, message):
        add_shape(slide, Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.75), NAVY)
        add_text(slide, Inches(0.65), Inches(6.4), Inches(12.0), Inches(0.55),
                 f"Key Message: {message}", size=14, color=WHITE, bold=True)

    def add_section_divider(title, subtitle=""):
        slide = new_slide()
        add_shape(slide, 0, 0, Inches(13.333), Inches(7.5), NAVY)
        add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                 title, size=36, color=WHITE, bold=True, align="center")
        if subtitle:
            add_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
                     subtitle, size=18, color=(180, 200, 220), align="center")

    # ─── S01: Title Slide ───
    slide = new_slide()
    add_shape(slide, 0, 0, Inches(13.333), Inches(7.5), NAVY)
    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "Xavier Niel Telecom Empire", size=40, color=WHITE, bold=True, align="center")
    add_text(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
             "BLM Strategic Assessment — Consolidated Group Report", size=22, color=(180, 200, 220), align="center")
    add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
             f"CQ4 2025 | 18 Operators | 3 Groups | 15 Countries", size=16, color=(150, 180, 200), align="center")
    add_text(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5),
             "NJJ Holding  ·  Iliad Group  ·  Millicom / Tigo", size=14, color=(120, 160, 190), align="center")

    # ─── S02: Portfolio Overview ───
    slide = new_slide()
    add_header(slide, "Portfolio Overview — 18 Operators Across 15 Countries")
    y = Inches(1.3)
    for group_id, group_name, operators in GROUPS:
        add_text(slide, Inches(0.5), y, Inches(3), Inches(0.4),
                 group_name, size=16, color=NAVY, bold=True)
        y += Inches(0.4)
        op_names = [f"{disp} ({ctry})" for _, _, disp, ctry in operators]
        for name in op_names:
            add_text(slide, Inches(0.8), y, Inches(4), Inches(0.3),
                     f"• {name}", size=11, color=DARK_GRAY)
            y += Inches(0.28)
        y += Inches(0.2)
    add_key_message(slide, f"{len(all_results)}/{len(ALL_OPERATORS)} operators analyzed across Europe and Latin America")

    # ─── S03: Financial Comparison Table ───
    slide = new_slide()
    add_header(slide, "Financial Comparison — Revenue & EBITDA by Operator")

    # Create a mini-table using text boxes (row by row)
    col_x = [0.3, 2.8, 5.0, 6.8, 8.5, 10.3]
    col_w = [2.4, 2.0, 1.6, 1.5, 1.6, 2.5]
    headers_ppt = ["Operator", "Country", "Revenue (M)", "EBITDA %", "Mobile (K)", "Diagnosis"]

    # Header row
    y_start = Inches(1.3)
    for i, (hdr, cx, cw) in enumerate(zip(headers_ppt, col_x, col_w)):
        add_shape(slide, Inches(cx), y_start, Inches(cw), Inches(0.35), NAVY)
        add_text(slide, Inches(cx + 0.05), y_start + Inches(0.02), Inches(cw - 0.1), Inches(0.3),
                 hdr, size=10, color=WHITE, bold=True)

    # Data rows (top 12 by revenue to fit)
    op_data = []
    for op_id, market_id, display, country in ALL_OPERATORS:
        if op_id in all_results:
            fh = _safe_financial(all_results[op_id])
            rev = _get_metric(fh, "total_revenue", "revenue", "quarterly_revenue")
            margin = _get_metric(fh, "ebitda_margin_pct", "ebitda_margin")
            mobile = _get_metric(fh, "mobile_total_k", "mobile_subscribers_k")
            diag = all_diagnoses.get(op_id)
            label = diag.central_diagnosis_label if diag else "N/A"
            rev_num = float(rev) if rev else 0
            op_data.append((display, country, _fmt(rev, "M"), _fmt(margin, "%"),
                            _fmt(mobile, "K", decimals=0), label, rev_num))

    op_data.sort(key=lambda x: -x[6])

    y_row = y_start + Inches(0.38)
    for idx, (disp, ctry, rev_s, margin_s, mob_s, lbl, _) in enumerate(op_data[:15]):
        bg = LIGHT_GRAY if idx % 2 == 0 else WHITE
        for i, (val, cx, cw) in enumerate(zip([disp, ctry, rev_s, margin_s, mob_s, lbl], col_x, col_w)):
            add_shape(slide, Inches(cx), y_row, Inches(cw), Inches(0.3), bg)
            add_text(slide, Inches(cx + 0.05), y_row + Inches(0.01), Inches(cw - 0.1), Inches(0.28),
                     val, size=9, color=DARK_GRAY)
        y_row += Inches(0.3)

    add_key_message(slide, "Diversified portfolio spanning high-margin European markets and high-growth LATAM markets")

    # ─── S04-S06: Per-group summary slides (3 slides) ───
    for group_id, group_name, operators in GROUPS:
        slide = new_slide()
        add_header(slide, f"{group_name} — Group Summary", f"{len(operators)} MARKETS")

        gsummary = group_summaries.get(group_id, {})

        # Revenue data on left
        y = Inches(1.4)
        add_text(slide, Inches(0.5), y, Inches(6), Inches(0.4),
                 "Revenue & EBITDA by Market", size=14, color=NAVY, bold=True)
        y += Inches(0.45)

        rev_comp = gsummary.get("revenue_comparison", {})
        for market, data in sorted(rev_comp.items()):
            rev_val = data.get("total_revenue")
            margin_val = data.get("ebitda_margin_pct")
            growth_val = data.get("revenue_growth_pct")
            market_display = market.replace("_", " ").title()
            line = f"{market_display}: Rev {_fmt(rev_val, 'M')}, Margin {_fmt(margin_val, '%')}, Growth {_fmt(growth_val, '%', '+')}"
            add_text(slide, Inches(0.7), y, Inches(5.5), Inches(0.28),
                     line, size=10, color=DARK_GRAY)
            y += Inches(0.28)

        # Key findings on right
        add_text(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.4),
                 "Key Findings", size=14, color=NAVY, bold=True)
        findings = gsummary.get("key_findings", [])
        if findings:
            add_bullet_text(slide, Inches(7), Inches(1.9), Inches(5.5), Inches(3),
                            findings[:6], size=11)

        # Common opportunities
        opps = gsummary.get("common_opportunities", [])
        if opps:
            add_text(slide, Inches(7), Inches(4.5), Inches(5.5), Inches(0.4),
                     "Common Opportunities", size=12, color=NAVY, bold=True)
            add_bullet_text(slide, Inches(7), Inches(4.9), Inches(5.5), Inches(1.2),
                            opps[:4], size=10)

        add_key_message(slide, f"{group_name}: {len(operators)} operators with diversified competitive positions")

    # ─── S07-S24: Per-market one-pager slides (18 slides) ───
    for op_id, market_id, display, country in ALL_OPERATORS:
        slide = new_slide()
        add_header(slide, f"{display} — {country}", f"BLM Five Looks Assessment | {PERIOD}")

        if op_id not in all_results:
            add_text(slide, Inches(1), Inches(3), Inches(10), Inches(1),
                     "Analysis data not available", size=20, color=(128, 128, 128), align="center")
            continue

        result = all_results[op_id]
        diag = all_diagnoses.get(op_id)
        fh = _safe_financial(result)

        # Verdict
        if diag and diag.one_line_verdict:
            add_text(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
                     diag.one_line_verdict, size=12, color=ACCENT_BLUE, bold=True)

        # Key metrics (left column)
        y = Inches(2.0)
        add_text(slide, Inches(0.5), y, Inches(3), Inches(0.35),
                 "Key Financial Metrics", size=13, color=NAVY, bold=True)
        y += Inches(0.4)

        metrics = [
            ("Revenue", _fmt(_get_metric(fh, "total_revenue", "revenue", "quarterly_revenue"), "M")),
            ("EBITDA Margin", _fmt(_get_metric(fh, "ebitda_margin_pct", "ebitda_margin"), "%")),
            ("Mobile Subs", _fmt(_get_metric(fh, "mobile_total_k", "mobile_subscribers_k"), "K", decimals=0)),
            ("Health", getattr(result.self_analysis, "health_rating", "N/A") if result.self_analysis else "N/A"),
            ("Rank", f"#{diag.operator_rank}" if diag and diag.operator_rank else "N/A"),
            ("Diagnosis", diag.central_diagnosis_label if diag else "N/A"),
        ]
        for label, value in metrics:
            add_text(slide, Inches(0.7), y, Inches(2), Inches(0.25),
                     f"{label}:", size=10, color=(100, 100, 100), bold=True)
            add_text(slide, Inches(2.7), y, Inches(2), Inches(0.25),
                     value, size=10, color=DARK_GRAY)
            y += Inches(0.28)

        # Priorities (middle column)
        add_text(slide, Inches(5), Inches(2.0), Inches(3.5), Inches(0.35),
                 "Strategic Priorities", size=13, color=NAVY, bold=True)
        if diag and diag.priorities:
            p_items = [f"[{p.get('priority', '')}] {p.get('name', '')}" for p in diag.priorities[:5]]
            add_bullet_text(slide, Inches(5), Inches(2.4), Inches(3.8), Inches(3),
                            p_items, size=10)

        # SWOT summary (right column)
        add_text(slide, Inches(9), Inches(2.0), Inches(3.5), Inches(0.35),
                 "SWOT Balance", size=13, color=NAVY, bold=True)
        if result.swot:
            s = len(getattr(result.swot, "strengths", []))
            w = len(getattr(result.swot, "weaknesses", []))
            o = len(getattr(result.swot, "opportunities", []))
            t = len(getattr(result.swot, "threats", []))
            swot_text = f"S:{s} / W:{w} / O:{o} / T:{t}"
            add_text(slide, Inches(9.2), Inches(2.5), Inches(3), Inches(0.3),
                     swot_text, size=11, color=DARK_GRAY)
            stance = diag.competitive_stance if diag else "N/A"
            add_text(slide, Inches(9.2), Inches(2.85), Inches(3), Inches(0.3),
                     f"Stance: {stance}", size=10, color=DARK_GRAY)

            # Top threats
            threats = getattr(result.swot, "threats", [])
            if threats:
                add_text(slide, Inches(9), Inches(3.4), Inches(3.5), Inches(0.3),
                         "Key Risks", size=11, color=NAVY, bold=True)
                risk_items = [str(t)[:60] for t in threats[:3]]
                add_bullet_text(slide, Inches(9), Inches(3.7), Inches(3.8), Inches(2),
                                risk_items, size=9)

        # Net assessment
        if diag:
            assessments = [
                diag.trends_net_assessment,
                diag.competition_net_assessment,
            ]
            non_empty = [a for a in assessments if a and "Insufficient" not in a]
            if non_empty:
                add_text(slide, Inches(0.5), Inches(5.0), Inches(12), Inches(0.3),
                         "Net Assessments", size=11, color=NAVY, bold=True)
                add_bullet_text(slide, Inches(0.5), Inches(5.3), Inches(12), Inches(0.9),
                                non_empty[:2], size=9, color=DARK_GRAY)

        add_key_message(slide, diag.one_line_verdict[:100] if diag and diag.one_line_verdict else f"{display} market analysis")

    # ─── S25: Cross-group strategic recommendations ───
    slide = new_slide()
    add_header(slide, "Cross-Group Strategic Recommendations", "SYNERGIES & PRIORITIES")

    recs = [
        "Procurement synergies: Joint equipment & terminal purchasing across 18 operators",
        "Shared platforms: Unified BSS/OSS to reduce duplicate IT investment",
        "Best practice transfer: Iliad low-cost model to NJJ & Millicom markets",
        "Digital services: Expand Tigo Money & enterprise cloud across all groups",
        "Roaming & interconnect: Preferential cross-group rates",
    ]
    add_text(slide, Inches(0.5), Inches(1.4), Inches(5.5), Inches(0.4),
             "Synergies & Resource Sharing", size=16, color=NAVY, bold=True)
    add_bullet_text(slide, Inches(0.5), Inches(1.9), Inches(5.8), Inches(3),
                    recs, size=12)

    # Priority matrix
    add_text(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.4),
             "Priority Ranking", size=16, color=NAVY, bold=True)
    priorities = [
        "P0: Stabilize Millicom LATAM EBITDA margins",
        "P0: Accelerate Iliad Italia growth",
        "P0: lifecell/Datagroup-Volia integration in Ukraine",
        "P1: 5G coverage across NJJ European markets",
        "P1: Millicom B2B & digital services expansion",
        "P2: Cross-group procurement synergies",
        "P2: Unified BSS/OSS platform migration",
    ]
    add_bullet_text(slide, Inches(7), Inches(1.9), Inches(5.8), Inches(4),
                    priorities, size=11)

    add_key_message(slide, "Xavier Niel's 18-operator portfolio offers significant untapped synergies across procurement, technology, and best practices")

    # ─── S26: Risk Matrix ───
    slide = new_slide()
    add_header(slide, "Risk Matrix — Cross-Group Assessment")

    risks = [
        ("Regulatory Risk", "High", "Spectrum auctions, licensing changes across 15 countries"),
        ("Currency Risk", "High", "LATAM & UAH volatility impacts consolidated results"),
        ("Geopolitical Risk", "Medium-High", "Ukraine conflict, LATAM political instability"),
        ("Competition Risk", "Medium", "Price wars in France, Italy; OTT substitution"),
        ("Technology Risk", "Medium", "5G ROI uncertainty, network upgrade costs"),
        ("Financing Risk", "Medium", "High-rate refinancing pressure on leveraged entities"),
    ]

    y = Inches(1.5)
    # Header
    cols = [0.5, 4.0, 6.0]
    widths = [3.3, 1.8, 6.0]
    for cx, cw, hdr in zip(cols, widths, ["Risk Category", "Impact", "Description"]):
        add_shape(slide, Inches(cx), Inches(y), Inches(cw), Inches(0.35), NAVY)
        add_text(slide, Inches(cx + 0.05), y + Inches(0.02), Inches(cw - 0.1), Inches(0.3),
                 hdr, size=11, color=WHITE, bold=True)
    y += Inches(0.38)
    for idx, (cat, impact, desc) in enumerate(risks):
        bg = LIGHT_GRAY if idx % 2 == 0 else WHITE
        for cx, cw, val in zip(cols, widths, [cat, impact, desc]):
            add_shape(slide, Inches(cx), Inches(y), Inches(cw), Inches(0.35), bg)
            add_text(slide, Inches(cx + 0.05), y + Inches(0.02), Inches(cw - 0.1), Inches(0.3),
                     val, size=10, color=DARK_GRAY)
        y += Inches(0.35)

    add_key_message(slide, "Diversified geographic footprint mitigates single-market risk but introduces currency and regulatory complexity")

    # ─── S27: Back Cover ───
    slide = new_slide()
    add_shape(slide, 0, 0, Inches(13.333), Inches(7.5), NAVY)
    add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
             "Thank You", size=40, color=WHITE, bold=True, align="center")
    add_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             "BLM Strategic Analysis — Xavier Niel Telecom Group", size=18,
             color=(180, 200, 220), align="center")
    add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "Generated by BLM Analysis Engine | Data sources: Public filings, regulatory reports, industry data",
             size=11, color=(120, 160, 190), align="center")

    prs.save(str(output_path))
    return slide_num[0]


# =========================================================================
# Main
# =========================================================================

def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # Step 1: Seed & Analyze
    all_results, all_diagnoses, group_summaries = seed_and_analyze()

    print(f"\nAnalysis complete: {len(all_results)}/{len(ALL_OPERATORS)} operators")

    # Step 2: Generate MD
    print("\n[Phase 4] Generating consolidated Markdown report...")
    md_content = generate_md(all_results, all_diagnoses, group_summaries)
    md_path = OUTPUT_DIR / "blm_niel_group_consolidated_cq4_2025.md"
    md_path.write_text(md_content, encoding="utf-8")
    md_lines = md_content.count("\n") + 1
    print(f"  MD: {md_path} ({md_lines} lines, {md_path.stat().st_size / 1024:.1f} KB)")

    # Step 3: Generate PDFs
    print("\n[Phase 5] Generating PDF reports...")
    try:
        pdf_zh_path = OUTPUT_DIR / "blm_niel_group_consolidated_cq4_2025_zh.pdf"
        generate_pdf(md_content, pdf_zh_path, chinese=True)
        print(f"  PDF (zh): {pdf_zh_path} ({pdf_zh_path.stat().st_size / 1024:.1f} KB)")

        pdf_path = OUTPUT_DIR / "blm_niel_group_consolidated_cq4_2025.pdf"
        generate_pdf(md_content, pdf_path, chinese=False)
        print(f"  PDF: {pdf_path} ({pdf_path.stat().st_size / 1024:.1f} KB)")
    except Exception as e:
        print(f"  PDF generation failed: {e}")

    # Step 4: Generate PPT
    print("\n[Phase 6] Generating PowerPoint deck...")
    try:
        pptx_path = OUTPUT_DIR / "blm_niel_group_consolidated_cq4_2025.pptx"
        n_slides = generate_ppt(all_results, all_diagnoses, group_summaries, pptx_path)
        print(f"  PPTX: {pptx_path} ({n_slides} slides, {pptx_path.stat().st_size / 1024:.1f} KB)")
    except Exception as e:
        print(f"  PPT generation failed: {e}")

    # Summary
    print(f"\n{'=' * 70}")
    print(f"  COMPLETE — Xavier Niel Group Consolidated Report")
    print(f"  Operators analyzed:  {len(all_results)}/{len(ALL_OPERATORS)}")
    print(f"  Output directory:    {OUTPUT_DIR}")
    outputs = list(OUTPUT_DIR.glob("blm_niel_group_*"))
    for f in sorted(outputs):
        print(f"    {f.name} ({f.stat().st_size / 1024:.1f} KB)")
    print(f"{'=' * 70}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
