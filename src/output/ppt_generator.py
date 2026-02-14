"""BLM PPT Generator — transforms FiveLooksResult into a branded PowerPoint deck.

Generates 35-42 slides (Draft) or 20-25 (Final mode with user_decisions filter).
Every content slide ends with a Key Message bar at the bottom.

Usage:
    from src.output.ppt_styles import get_style
    from src.output.ppt_generator import BLMPPTGenerator

    style = get_style("vodafone")
    gen = BLMPPTGenerator(style=style, operator_id="vodafone_germany")
    ppt_path = gen.generate(five_looks_result, mode="draft")
"""

from __future__ import annotations

import tempfile
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from src.output.ppt_styles import PPTStyle, get_style, DEFAULT_STYLE
from src.output.ppt_charts import BLMChartGenerator


@dataclass
class SlideSpec:
    """Specification for a slide in the deck."""
    slide_id: str
    title: str
    section: str = ""
    required: bool = True  # False = optional deep-dive slide


class BLMPPTGenerator:
    """Generate PowerPoint presentations from FiveLooksResult.

    Takes a FiveLooksResult and PPTStyle, produces a branded 16:9 deck
    with ~38 slides (Draft) or a filtered subset (Final).
    """

    # Slide dimensions: 16:9 widescreen
    SLIDE_WIDTH = 13.333
    SLIDE_HEIGHT = 7.5

    # Key message bar position
    KM_LEFT = 0.5
    KM_TOP = 6.3
    KM_WIDTH = 12.333
    KM_HEIGHT = 0.75

    def __init__(
        self,
        style: Optional[PPTStyle] = None,
        operator_id: str = "",
        output_dir: Optional[str] = None,
        chart_dpi: int = 150,
    ):
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for PPT generation. "
                "Install it with: pip install python-pptx"
            )
        self.style = style or DEFAULT_STYLE
        self.operator_id = operator_id
        self.output_dir = Path(output_dir) if output_dir else Path(
            tempfile.mkdtemp())
        self.output_dir.mkdir(parents=True, exist_ok=True)

        self._chart_dir = Path(tempfile.mkdtemp())
        self.chart_gen = BLMChartGenerator(
            style=self.style, output_dir=str(self._chart_dir), dpi=chart_dpi)

        self.prs = None
        self.slide_num = 0
        self._slide_specs: list[SlideSpec] = []

    # =========================================================================
    # Public API
    # =========================================================================

    def generate(
        self,
        result,
        mode: str = "draft",
        user_decisions: Optional[dict] = None,
        filename: Optional[str] = None,
    ) -> str:
        """Generate the PPT deck.

        Args:
            result: FiveLooksResult from the analysis engine
            mode: "draft" (full ~38 slides) or "final" (filtered by user_decisions)
            user_decisions: dict mapping slide_id -> "keep"/"remove"/"merge"
            filename: Output filename (auto-generated if None)

        Returns:
            Path to generated .pptx file
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.SLIDE_WIDTH)
        self.prs.slide_height = Inches(self.SLIDE_HEIGHT)
        self.slide_num = 0
        self._slide_specs = []

        # Build all slides
        self._build_all_slides(result)

        # In final mode, remove slides the user declined
        if mode == "final" and user_decisions:
            self._apply_user_decisions(user_decisions)

        # Save
        if filename is None:
            safe_name = self._sanitize_name(result.target_operator)
            filename = f"blm_{safe_name}_{mode}.pptx"
        output_path = self.output_dir / filename
        self.prs.save(str(output_path))
        return str(output_path)

    # =========================================================================
    # Slide orchestration
    # =========================================================================

    def _build_all_slides(self, result):
        """Build the complete slide deck."""
        # S01 Cover
        self._add_cover_slide(result)
        # S02 TOC
        self._add_toc_slide(result)
        # S03 Data Quality
        self._add_data_quality_slide(result)
        # S04 Executive Summary
        self._add_executive_summary(result)

        # --- Look 1: Trends ---
        self._add_section_divider("01 Look at Trends", "PEST Framework — Macro Environment")
        self._add_pest_dashboard(result.trends)
        self._add_industry_environment(result.trends)
        self._add_trend_deep_dive(result.trends)

        # --- Look 2: Market/Customer ---
        self._add_section_divider("02 Look at Market", "$APPEALS — Market Changes & Customer Needs")
        self._add_market_change_panorama(result.market_customer)
        self._add_customer_segments(result.market_customer)
        self._add_appeals_assessment(result.market_customer)
        self._add_market_deep_dive(result.market_customer)

        # --- Tariff Analysis ---
        self._add_tariff_slides(result)

        # --- Look 3: Competition ---
        self._add_section_divider("03 Look at Competition", "Porter's Five Forces — Competitive Landscape")
        self._add_five_forces(result.competition)
        self._add_competitor_deep_dives(result.competition)
        self._add_competition_summary(result.competition)
        self._add_gap_analysis(result.competition)
        self._add_new_entrants(result.competition)

        # --- Look 4: Self ---
        self._add_section_divider("04 Look at Self", "BMC + Capability Assessment")
        self._add_health_check(result.self_analysis)
        self._add_revenue_composition_trend(result.self_analysis)
        self._add_revenue_comparison(result.self_analysis, result.competition)
        self._add_segment_deep_dives(result.self_analysis)
        self._add_network_analysis(result.self_analysis)
        self._add_bmc_canvas(result.self_analysis)
        self._add_org_talent(result.self_analysis)
        self._add_strengths_weaknesses_exposure(result.self_analysis)

        # --- SWOT ---
        self._add_section_divider("SWOT Synthesis", "Strengths, Weaknesses, Opportunities & Threats")
        self._add_swot_matrix(result.swot)

        # --- Look 5: Opportunities ---
        self._add_section_divider("05 Look at Opportunities", "SPAN Matrix — Opportunity Selection")
        self._add_span_bubble(result.opportunities)
        self._add_priority_table(result.opportunities)
        self._add_opportunity_deep_dive(result.opportunities)

        # --- Three Decisions (BLM Phase 2) ---
        self._add_three_decisions_slides(result)

        # --- Summary ---
        self._add_summary_slide(result)
        self._add_provenance_appendix(result)
        self._add_back_cover()

    def _apply_user_decisions(self, decisions: dict):
        """Remove slides that the user marked as 'remove' in Final mode."""
        indices_to_remove = []
        for i, spec in enumerate(self._slide_specs):
            action = decisions.get(spec.slide_id, "keep")
            if action == "remove":
                indices_to_remove.append(i)

        # Remove in reverse order to preserve indices
        for idx in reversed(indices_to_remove):
            xml_slides = self.prs.slides._sldIdLst
            slides_list = list(xml_slides)
            if idx < len(slides_list):
                xml_slides.remove(slides_list[idx])

    # =========================================================================
    # Utility methods
    # =========================================================================

    def _new_slide(self, slide_id: str = "", title: str = "",
                   section: str = "", required: bool = True):
        """Add a blank slide and register its spec."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        self.slide_num += 1
        spec = SlideSpec(
            slide_id=slide_id or f"S{self.slide_num:02d}",
            title=title, section=section, required=required)
        self._slide_specs.append(spec)
        return slide

    def _add_header(self, slide, title: str, subtitle: str = ""):
        """Standard slide header with accent bar and UPPERCASE subtitle."""
        # Top accent bar
        self._add_shape(slide, 0, 0, Inches(self.SLIDE_WIDTH), Inches(0.08),
                        self.style.primary_color)
        # Title
        self._add_text_box(slide, Inches(0.5), Inches(0.25), Inches(10), Inches(0.6),
                           title, font_size=self.style.subtitle_size,
                           font_color=self.style.text_color, bold=True)
        if subtitle:
            self._add_text_box(slide, Inches(0.5), Inches(0.85), Inches(10), Inches(0.35),
                               subtitle.upper(), font_size=12,
                               font_color=self.style.primary_color, bold=True)
        # Page number
        self._add_text_box(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.35),
                           str(self.slide_num), font_size=10,
                           font_color=self.style.light_text_color, align="right")

    def _add_key_message_bar(self, slide, message: str):
        """Add Key Message bar at the bottom of a content slide."""
        if not message:
            return
        shape = self._add_shape(
            slide, Inches(self.KM_LEFT), Inches(self.KM_TOP),
            Inches(self.KM_WIDTH), Inches(self.KM_HEIGHT),
            self.style.key_message_bg)

        txBox = slide.shapes.add_textbox(
            Inches(self.KM_LEFT + 0.15), Inches(self.KM_TOP + 0.1),
            Inches(self.KM_WIDTH - 0.3), Inches(self.KM_HEIGHT - 0.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Key Message: {self._strip_cjk(message)}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*self.style.key_message_text)
        p.font.bold = True
        p.font.name = self.style.body_font
        p.alignment = PP_ALIGN.LEFT

    def _add_shape(self, slide, left, top, width, height, fill_color,
                   shape_type=None):
        """Add a colored shape."""
        if shape_type is None:
            shape_type = MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_color)
        shape.line.fill.background()
        return shape

    def _add_text_box(self, slide, left, top, width, height, text: str,
                      font_size: int = 14, font_color: tuple = (0, 0, 0),
                      bold: bool = False, align: str = "left"):
        """Add a text box."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = self._strip_cjk(str(text))
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*font_color)
        p.font.bold = bold
        p.font.name = self.style.body_font
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        return txBox

    def _add_bullet_list(self, slide, left, top, width, height,
                         items: list, font_size: int = 12,
                         font_color: tuple = None, max_items: int = 8):
        """Add a bulleted list text box."""
        if font_color is None:
            font_color = self.style.text_color
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items[:max_items]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {self._strip_cjk(str(item))}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(*font_color)
            p.font.name = self.style.body_font
            p.space_after = Pt(4)
        return txBox

    def _add_image(self, slide, image_path: str, left, top, width, height=None):
        """Add an image to the slide."""
        if height:
            slide.shapes.add_picture(image_path, left, top, width, height)
        else:
            slide.shapes.add_picture(image_path, left, top, width)

    def _add_metric_cards(self, slide, metrics: list[tuple], top, left_start=0.5,
                          card_width=2.8, card_height=1.2, gap=0.2):
        """Add a row of metric cards (label, value, change_indicator)."""
        x = Inches(left_start)
        for label, value, change in metrics:
            self._add_shape(slide, x, Inches(top), Inches(card_width),
                            Inches(card_height), (245, 245, 245))
            self._add_text_box(slide, x + Inches(0.15), Inches(top + 0.1),
                               Inches(card_width - 0.3), Inches(0.35),
                               label, font_size=10,
                               font_color=self.style.light_text_color)
            self._add_text_box(slide, x + Inches(0.15), Inches(top + 0.45),
                               Inches(card_width - 0.3), Inches(0.4),
                               str(value), font_size=18,
                               font_color=self.style.primary_color, bold=True)
            if change:
                is_positive = str(change).startswith('+') or (
                    isinstance(change, (int, float)) and change > 0)
                change_color = self.style.positive_color if is_positive else self.style.negative_color
                self._add_text_box(slide, x + Inches(0.15), Inches(top + 0.85),
                                   Inches(card_width - 0.3), Inches(0.25),
                                   str(change), font_size=10,
                                   font_color=change_color)
            x += Inches(card_width + gap)

    def _sanitize_name(self, name: str) -> str:
        return name.lower().replace(" ", "_").replace("/", "_").replace("&", "and")

    @staticmethod
    def _strip_cjk(text: str) -> str:
        """Remove CJK characters and clean up resulting whitespace."""
        import re
        cleaned = re.sub(r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]+', '', text)
        cleaned = re.sub(r'\s{2,}', ' ', cleaned).strip()
        # Remove trailing punctuation artifacts
        cleaned = re.sub(r'[（）\(\)]\s*$', '', cleaned).strip()
        return cleaned

    # =========================================================================
    # S01: Cover Slide
    # =========================================================================

    def _add_cover_slide(self, result):
        slide = self._new_slide("cover", "Cover")

        # Brand header bar
        self._add_shape(slide, 0, 0, Inches(self.SLIDE_WIDTH), Inches(2.8),
                        self.style.primary_color)

        # Title
        title_text = f"{self.style.display_name} BLM Strategic Analysis"
        self._add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(1.2),
                           title_text, font_size=self.style.title_size,
                           font_color=(255, 255, 255), bold=True)

        # Subtitle
        self._add_text_box(slide, Inches(0.5), Inches(1.7), Inches(9), Inches(0.6),
                           "BLM (Business Leadership Model) Five Looks Analysis",
                           font_size=20, font_color=(255, 220, 220))

        # Info cards
        info_items = [
            ("Target Operator", result.target_operator.replace("_", " ").title()),
            ("Market", result.market.replace("_", " ").title()),
            ("Analysis Period", result.analysis_period),
            ("Report Date", datetime.now().strftime('%Y-%m-%d')),
        ]
        x = Inches(0.5)
        for label, value in info_items:
            self._add_shape(slide, x, Inches(3.5), Inches(2.9), Inches(1.3),
                            (245, 245, 245))
            self._add_text_box(slide, x + Inches(0.15), Inches(3.6),
                               Inches(2.6), Inches(0.3), label,
                               font_size=10, font_color=self.style.light_text_color)
            self._add_text_box(slide, x + Inches(0.15), Inches(3.9),
                               Inches(2.6), Inches(0.7), value,
                               font_size=16, font_color=self.style.primary_color,
                               bold=True)
            x += Inches(3.1)

        # Framework overview
        framework = (
            "01 Trends (PEST) | 02 Market/Customer ($APPEALS) | "
            "03 Competition (Porter) | 04 Self (BMC) | "
            "SWOT Synthesis | 05 Opportunities (SPAN)"
        )
        self._add_text_box(slide, Inches(0.5), Inches(5.3), Inches(12), Inches(0.6),
                           framework, font_size=12,
                           font_color=self.style.light_text_color)

    # =========================================================================
    # S02: Table of Contents
    # =========================================================================

    def _add_toc_slide(self, result):
        slide = self._new_slide("toc", "Table of Contents")
        self._add_header(slide, "Report Contents", "Table of Contents")

        sections = [
            ("01", "Executive Summary", "Key findings and recommendations"),
            ("02", "Look at Trends", "PEST Framework — macro environment"),
            ("03", "Look at Market", "$APPEALS — market changes & customer needs"),
            ("T", "Tariff Analysis", "Pricing comparison — mobile, fixed & convergence"),
            ("04", "Look at Competition", "Porter's Five Forces — competitive landscape"),
            ("05", "Look at Self", "BMC + Capability — internal assessment"),
            ("06", "SWOT Synthesis", "Strengths, Weaknesses, Opportunities, Threats"),
            ("07", "Look at Opportunities", "SPAN Matrix — opportunity selection"),
            ("08", "Summary & Provenance", "Conclusions and data sources"),
        ]

        y = Inches(1.5)
        for num, title, desc in sections:
            self._add_shape(slide, Inches(0.8), y, Inches(0.5), Inches(0.5),
                            self.style.primary_color, MSO_SHAPE.OVAL)
            self._add_text_box(slide, Inches(0.8), y + Inches(0.08), Inches(0.5),
                               Inches(0.35), num, font_size=12,
                               font_color=(255, 255, 255), bold=True, align="center")
            self._add_text_box(slide, Inches(1.5), y + Inches(0.02), Inches(4),
                               Inches(0.3), title, font_size=14,
                               font_color=self.style.text_color, bold=True)
            self._add_text_box(slide, Inches(1.5), y + Inches(0.3), Inches(8),
                               Inches(0.3), desc, font_size=11,
                               font_color=self.style.light_text_color)
            y += Inches(0.7)

    # =========================================================================
    # S03: Data Quality
    # =========================================================================

    def _add_data_quality_slide(self, result):
        slide = self._new_slide("data_quality", "Data Quality")
        self._add_header(slide, "Data Quality Overview", "Data Quality")

        prov = result.provenance
        report = prov.quality_report() if prov else {}

        total = report.get("total_data_points", 0)
        high = report.get("high_confidence", 0)
        medium = report.get("medium_confidence", 0)
        low = report.get("low_confidence", 0)
        estimated = report.get("estimated", 0)
        conflicts = report.get("with_conflicts", 0)
        sources = report.get("unique_sources", 0)

        metrics = [
            ("Total Data Points", str(total), ""),
            ("High Confidence", str(high), f"{high/max(total,1)*100:.0f}%"),
            ("Medium Confidence", str(medium), ""),
            ("Low/Estimated", str(low + estimated), ""),
        ]
        self._add_metric_cards(slide, metrics, top=1.5)

        # Additional info
        info_text = (
            f"Unique Sources: {sources}\n"
            f"Data Conflicts: {conflicts}\n"
            f"Coverage: {'Good' if high > total * 0.5 else 'Partial'}"
        )
        self._add_text_box(slide, Inches(0.5), Inches(3.2), Inches(5), Inches(2),
                           info_text, font_size=14,
                           font_color=self.style.text_color)

        km = f"{total} data points tracked, {high} high-confidence ({high/max(total,1)*100:.0f}%)"
        self._add_key_message_bar(slide, km)

    # =========================================================================
    # S04: Executive Summary
    # =========================================================================

    def _add_executive_summary(self, result):
        slide = self._new_slide("exec_summary", "Executive Summary")
        self._add_header(slide, "Executive Summary")

        y = Inches(1.4)
        summaries = [
            ("Trends", _get_key_message(result.trends)),
            ("Market/Customer", _get_key_message(result.market_customer)),
            ("Competition", _get_key_message(result.competition)),
            ("Self Assessment", _get_key_message(result.self_analysis)),
            ("SWOT", _get_key_message(result.swot)),
            ("Opportunities", _get_key_message(result.opportunities)),
        ]

        for label, msg in summaries:
            # Label
            self._add_shape(slide, Inches(0.5), y, Inches(2.2), Inches(0.6),
                            (245, 245, 245))
            self._add_text_box(slide, Inches(0.6), y + Inches(0.1),
                               Inches(2), Inches(0.4), label,
                               font_size=12, font_color=self.style.primary_color,
                               bold=True)
            # Message
            self._add_text_box(slide, Inches(2.9), y + Inches(0.1),
                               Inches(9.5), Inches(0.5), msg or "Analysis pending",
                               font_size=11, font_color=self.style.text_color)
            y += Inches(0.7)

        self._add_key_message_bar(slide, "Five Looks analysis completed — see details in following sections")

    # =========================================================================
    # Section divider
    # =========================================================================

    def _add_section_divider(self, title: str, subtitle: str = ""):
        slide = self._new_slide(f"section_{self.slide_num}", title)

        self._add_shape(slide, 0, 0, Inches(self.SLIDE_WIDTH), Inches(self.SLIDE_HEIGHT),
                        self.style.primary_color)
        self._add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                           title, font_size=self.style.title_size,
                           font_color=(255, 255, 255), bold=True)
        if subtitle:
            self._add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(1),
                               subtitle, font_size=self.style.heading_size,
                               font_color=(255, 220, 220))

    # =========================================================================
    # Look 1: Trends
    # =========================================================================

    def _add_pest_dashboard(self, trends):
        if trends is None:
            return
        pest = trends.pest if hasattr(trends, 'pest') else None
        if pest is None:
            return

        slide = self._new_slide("pest_dashboard", "PEST Dashboard")
        self._add_header(slide, "PEST Analysis Dashboard", "Macro Environment")

        chart_path = self.chart_gen.create_pest_dashboard(pest, filename="pest_dashboard.png")
        self._add_image(slide, chart_path, Inches(0.5), Inches(1.3),
                        Inches(12), Inches(4.8))

        km = getattr(pest, 'key_message', '') or "PEST analysis of macro environment factors"
        self._add_key_message_bar(slide, km)

    def _add_industry_environment(self, trends):
        if trends is None:
            return

        slide = self._new_slide("industry_env", "Industry Environment")
        self._add_header(slide, "Industry Environment", "Industry Analysis")

        # Top: metric cards for industry overview
        metrics = []
        if trends.industry_market_size:
            metrics.append(("Market Size", str(trends.industry_market_size), ""))
        if trends.industry_growth_rate:
            metrics.append(("Growth Rate", str(trends.industry_growth_rate), ""))
        if trends.industry_concentration:
            metrics.append(("Concentration", str(trends.industry_concentration), ""))
        if trends.industry_lifecycle_stage:
            metrics.append(("Lifecycle", str(trends.industry_lifecycle_stage), ""))
        if metrics:
            self._add_metric_cards(slide, metrics[:4], top=1.4)

        # Left column: PEST factor severity chart (if pest data available)
        pest = trends.pest if hasattr(trends, 'pest') else None
        has_pest_chart = False
        if pest is not None:
            dims = [
                ('Political', pest.political_factors if hasattr(pest, 'political_factors') else []),
                ('Economic', pest.economic_factors if hasattr(pest, 'economic_factors') else []),
                ('Society', pest.society_factors if hasattr(pest, 'society_factors') else []),
                ('Technology', pest.technology_factors if hasattr(pest, 'technology_factors') else []),
            ]
            # Count high-severity factors per dimension
            categories = []
            high_counts = []
            total_counts = []
            for dim_name, factors in dims:
                if factors:
                    categories.append(dim_name)
                    high = sum(1 for f in factors
                               if (getattr(f, 'severity', 'medium') == 'high'
                                   if hasattr(f, 'severity')
                                   else (f.get('severity', 'medium') == 'high'
                                         if isinstance(f, dict) else False)))
                    high_counts.append(high)
                    total_counts.append(len(factors))
            if categories and any(total_counts):
                chart_path = self.chart_gen.create_segment_comparison(
                    categories,
                    {'High Severity': high_counts, 'Total Factors': total_counts},
                    title="PEST Factor Distribution",
                    y_label="Count",
                    filename="pest_factor_dist.png")
                self._add_image(slide, chart_path, Inches(0.3), Inches(2.9),
                                Inches(6), Inches(3.1))
                has_pest_chart = True

        # Right column: business models + tech revolution
        right_items = []
        if trends.new_business_models:
            right_items.append("New Business Models:")
            right_items.extend(f"  {m}" for m in trends.new_business_models[:3])
        if trends.technology_revolution:
            right_items.append("Technology Revolution:")
            right_items.extend(f"  {t}" for t in trends.technology_revolution[:3])
        right_x = Inches(6.5) if has_pest_chart else Inches(6.5)
        right_y = Inches(2.9) if has_pest_chart else Inches(2.9)
        if right_items:
            self._add_bullet_list(slide, right_x, right_y, Inches(6),
                                  Inches(2), right_items, font_size=12)

        # Key success factors (below chart area)
        if trends.key_success_factors and not has_pest_chart:
            self._add_text_box(slide, Inches(0.5), Inches(4.2), Inches(5), Inches(0.3),
                               "Key Success Factors:", font_size=13,
                               font_color=self.style.primary_color, bold=True)
            self._add_bullet_list(slide, Inches(0.5), Inches(4.5), Inches(12),
                                  Inches(1.5), trends.key_success_factors[:5],
                                  font_size=11)

        km = trends.key_message or "Industry environment analysis"
        self._add_key_message_bar(slide, km)

    def _add_trend_deep_dive(self, trends):
        if trends is None:
            return
        if not trends.value_transfer_trends:
            return

        slide = self._new_slide("trend_deep_dive", "Trend Deep Dive", required=False)
        self._add_header(slide, "Value Transfer Trends", "Value Migration")

        self._add_bullet_list(slide, Inches(0.5), Inches(1.4), Inches(12),
                              Inches(4.5), trends.value_transfer_trends[:8],
                              font_size=13)

        km = trends.key_message or "Value transfer trends in the industry"
        self._add_key_message_bar(slide, km)

    # =========================================================================
    # Look 2: Market/Customer
    # =========================================================================

    def _add_market_change_panorama(self, mci):
        if mci is None:
            return

        slide = self._new_slide("market_changes", "Market Changes")
        self._add_header(slide, "Market Change Panorama", "Market Overview")

        # Snapshot metrics
        snapshot = mci.market_snapshot or {}
        if snapshot:
            metrics = []
            key_map = [
                ('total_revenue', 'Total Revenue (€M)'),
                ('total_mobile_subscribers_k', 'Mobile Subs (K)'),
                ('total_broadband_subscribers_k', 'Broadband Subs (K)'),
                ('operator_count', 'Operators'),
            ]
            for key, label in key_map:
                if key in snapshot:
                    v = snapshot[key]
                    if isinstance(v, dict):
                        v = ", ".join(f"{sk}: {sv}" for sk, sv in v.items())
                    elif isinstance(v, float) and v > 1000:
                        v = f"{v:,.0f}"
                    metrics.append((label, str(v), ""))
            if metrics:
                self._add_metric_cards(slide, metrics[:4], top=1.4)

        # Market share chart (right side)
        market_shares = snapshot.get('market_shares', {})
        if market_shares:
            categories = [k.replace('_', ' ').title() for k in market_shares.keys()]
            values = list(market_shares.values())
            chart_path = self.chart_gen.create_horizontal_bar_chart(
                categories, values, target_category='Vodafone Germany',
                title="Market Share (%)", x_label="%",
                filename="market_share.png", value_suffix="%")
            self._add_image(slide, chart_path, Inches(6.5), Inches(2.8),
                            Inches(6), Inches(3.2))

        # Changes (left side, below metrics)
        changes = mci.changes or []
        if changes:
            y = Inches(3.0)
            self._add_text_box(slide, Inches(0.5), y, Inches(5.5), Inches(0.3),
                               "Market Changes:", font_size=13,
                               font_color=self.style.primary_color, bold=True)
            items = []
            for c in changes[:4]:
                desc = c.description if hasattr(c, 'description') else str(c)
                items.append(desc)
            self._add_bullet_list(slide, Inches(0.5), y + Inches(0.3), Inches(5.5),
                                  Inches(2.5), items, font_size=11)

        km = mci.key_message or "Market dynamics overview"
        self._add_key_message_bar(slide, km)

    def _add_customer_segments(self, mci):
        if mci is None:
            return
        if not mci.customer_segments:
            return

        slide = self._new_slide("customer_segments", "Customer Segments")
        self._add_header(slide, "Customer Segments", "Segmentation")

        # Left: segment cards (compact)
        y = Inches(1.4)
        for seg in mci.customer_segments[:4]:
            self._add_shape(slide, Inches(0.5), y, Inches(6.5), Inches(0.85),
                            (248, 248, 248))
            name = seg.segment_name if hasattr(seg, 'segment_name') else str(seg)
            self._add_text_box(slide, Inches(0.7), y + Inches(0.05), Inches(3),
                               Inches(0.3), name, font_size=13,
                               font_color=self.style.primary_color, bold=True)
            details = []
            if hasattr(seg, 'size_estimate') and seg.size_estimate:
                details.append(f"Size: {seg.size_estimate}")
            if hasattr(seg, 'growth_trend'):
                details.append(f"Trend: {seg.growth_trend}")
            if hasattr(seg, 'our_share') and seg.our_share:
                details.append(f"Our Share: {seg.our_share}")
            detail_text = " | ".join(details)
            self._add_text_box(slide, Inches(0.7), y + Inches(0.35), Inches(5.8),
                               Inches(0.25), detail_text, font_size=11,
                               font_color=self.style.text_color)
            if hasattr(seg, 'unmet_needs') and seg.unmet_needs:
                needs_text = "Needs: " + ", ".join(seg.unmet_needs[:2])
                self._add_text_box(slide, Inches(0.7), y + Inches(0.55), Inches(5.8),
                                   Inches(0.25), needs_text, font_size=10,
                                   font_color=self.style.light_text_color)
            y += Inches(0.95)

        # Right: segment comparison chart (if $APPEALS has multi-operator data)
        appeals = mci.appeals_assessment or []
        if appeals and len(appeals) >= 3:
            dimensions = [a.dimension_name if hasattr(a, 'dimension_name') else str(a)
                          for a in appeals[:8]]
            our_scores = [a.our_score if hasattr(a, 'our_score') else 0
                          for a in appeals[:8]]
            if any(s > 0 for s in our_scores):
                chart_path = self.chart_gen.create_bar_chart(
                    dimensions, our_scores,
                    title="$APPEALS Scores",
                    y_label="Score (1-5)",
                    filename="appeals_bar.png")
                self._add_image(slide, chart_path, Inches(7.3), Inches(1.3),
                                Inches(5.5), Inches(4.8))

        km = mci.key_message or "Customer segmentation analysis"
        self._add_key_message_bar(slide, km)

    def _add_appeals_assessment(self, mci):
        if mci is None:
            return
        if not mci.appeals_assessment:
            return

        slide = self._new_slide("appeals", "$APPEALS Assessment")
        self._add_header(slide, "$APPEALS Assessment", "Customer Needs")

        chart_path = self.chart_gen.create_appeals_radar(
            mci.appeals_assessment,
            target_operator=self.operator_id,
            filename="appeals_radar.png")
        self._add_image(slide, chart_path, Inches(0.5), Inches(1.3),
                        Inches(11.5), Inches(4.8))

        km = mci.key_message or "$APPEALS customer needs assessment"
        self._add_key_message_bar(slide, km)

    def _add_market_deep_dive(self, mci):
        if mci is None:
            return
        if not mci.customer_value_migration:
            return

        slide = self._new_slide("market_deep_dive", "Market Deep Dive", required=False)
        self._add_header(slide, "Customer Value Migration", "Value Migration")

        self._add_text_box(slide, Inches(0.5), Inches(1.4), Inches(12), Inches(3),
                           mci.customer_value_migration, font_size=13,
                           font_color=self.style.text_color)

        opps = mci.opportunities or []
        threats = mci.threats or []
        if opps:
            self._add_text_box(slide, Inches(0.5), Inches(3.5), Inches(5.5), Inches(0.3),
                               "Opportunities:", font_size=12,
                               font_color=self.style.positive_color, bold=True)
            self._add_bullet_list(slide, Inches(0.5), Inches(3.8), Inches(5.5),
                                  Inches(2), [o.description for o in opps[:4]
                                              if hasattr(o, 'description')],
                                  font_size=10)
        if threats:
            self._add_text_box(slide, Inches(6.5), Inches(3.5), Inches(5.5), Inches(0.3),
                               "Threats:", font_size=12,
                               font_color=self.style.negative_color, bold=True)
            self._add_bullet_list(slide, Inches(6.5), Inches(3.8), Inches(5.5),
                                  Inches(2), [t.description for t in threats[:4]
                                              if hasattr(t, 'description')],
                                  font_size=10)

        km = mci.key_message or "Market deep dive analysis"
        self._add_key_message_bar(slide, km)

    # =========================================================================
    # Tariff Analysis (4 slides)
    # =========================================================================

    def _add_tariff_slides(self, result):
        """Add tariff analysis section with 3-4 slides."""
        ta = getattr(result, 'tariff_analysis', None)
        if ta is None:
            return

        self._add_section_divider(
            "Tariff Analysis",
            "Mobile, Fixed & Convergence Pricing Landscape"
        )
        self._add_tariff_mobile_postpaid(ta)
        self._add_tariff_value_and_5g(ta)
        self._add_tariff_price_evolution(ta)
        self._add_tariff_fixed_and_fmc(ta)

    def _add_tariff_mobile_postpaid(self, ta):
        """Slide: Mobile Postpaid Price Comparison — grouped bar chart."""
        comparison = ta.get("mobile_postpaid_comparison", [])
        if not comparison:
            return

        slide = self._new_slide("tariff_mobile", "Mobile Postpaid Price Comparison")
        self._add_header(slide, "Mobile Postpaid Price Comparison",
                         "Cross-Operator Tier Benchmarking")

        # Build grouped bar data: tiers as x-axis, operators as series
        tiers = []
        operator_prices: dict[str, list[float]] = {}
        for tier_data in comparison:
            tier_label = tier_data["tier"].upper()
            tiers.append(tier_label)
            for op in tier_data.get("operators", []):
                name = op["display_name"]
                if name not in operator_prices:
                    operator_prices[name] = []
            # Fill values per operator for this tier
            tier_ops = {op["display_name"]: op.get("price", 0)
                        for op in tier_data.get("operators", [])}
            for name in operator_prices:
                operator_prices[name].append(tier_ops.get(name, 0) or 0)

        if tiers and operator_prices:
            our_name = self.operator_id.replace('_', ' ').title()
            chart_path = self.chart_gen.create_segment_comparison(
                tiers, operator_prices,
                target_operator=our_name,
                title="Monthly Price by Tier (EUR)",
                y_label="EUR/month",
                filename="tariff_mobile_postpaid.png",
            )
            self._add_image(slide, chart_path, Inches(0.3), Inches(1.3),
                            Inches(8.5), Inches(4.8))

        # Right column: price gap annotations
        insights = ta.get("strategic_insights", [])
        tier_insights = [i for i in insights if i.startswith("Tier ")]
        if tier_insights:
            self._add_text_box(slide, Inches(9.0), Inches(1.4), Inches(3.8), Inches(0.3),
                               "Price Gaps:", font_size=12,
                               font_color=self.style.primary_color, bold=True)
            self._add_bullet_list(slide, Inches(9.0), Inches(1.8), Inches(3.8),
                                  Inches(4.2), tier_insights[:4], font_size=10)

        km = "Mobile postpaid tier comparison across all operators"
        self._add_key_message_bar(slide, km)

    def _add_tariff_value_and_5g(self, ta):
        """Slide: EUR/GB Value Analysis & 5G Premium Erosion — split layout."""
        value_list = ta.get("value_per_gb", [])
        erosion = ta.get("five_g_erosion", [])
        if not value_list and not erosion:
            return

        slide = self._new_slide("tariff_value_5g",
                                "Value Analysis & 5G Premium")
        self._add_header(slide, "EUR/GB Value Analysis & 5G Premium Erosion",
                         "Pricing Intelligence")

        # Left half: horizontal bar chart EUR/GB ranking
        if value_list:
            labels = [f"{v['operator']} {v['plan']}" for v in value_list[:12]]
            values = [v["eur_per_gb"] for v in value_list[:12]]
            # Truncate long labels
            labels = [l[:30] for l in labels]

            chart_path = self.chart_gen.create_horizontal_bar_chart(
                labels, values,
                title="EUR/GB Ranking (lower = better value)",
                x_label="EUR/GB",
                filename="tariff_eur_per_gb.png",
            )
            self._add_image(slide, chart_path, Inches(0.2), Inches(1.3),
                            Inches(6.3), Inches(4.8))

        # Right half: 5G premium erosion trend
        if erosion and len(erosion) >= 2:
            # Only include snapshots where we have premium data
            snapshots_with_data = [
                e for e in erosion if e.get("premium_pct") is not None
            ]
            if snapshots_with_data:
                x_labels = [e["snapshot"] for e in snapshots_with_data]
                y_values = [e["premium_pct"] for e in snapshots_with_data]
                chart_path2 = self.chart_gen.create_multi_line_trend(
                    x_labels,
                    {"5G Premium %": y_values},
                    title="5G Premium Erosion Over Time",
                    y_label="Premium %",
                    y_format="percent",
                    filename="tariff_5g_erosion.png",
                )
                self._add_image(slide, chart_path2, Inches(6.8), Inches(1.3),
                                Inches(6), Inches(4.8))
            elif erosion:
                # All plans include 5G now — show text
                self._add_text_box(
                    slide, Inches(7.0), Inches(3.0), Inches(5.5), Inches(1),
                    "All mobile postpaid plans now include 5G — premium fully eroded",
                    font_size=14, font_color=self.style.positive_color,
                    bold=True, align="center",
                )

        km_parts = []
        if value_list:
            best = value_list[0]
            km_parts.append(
                f"Best value: {best['operator']} at EUR {best['eur_per_gb']:.2f}/GB"
            )
        erosion_insight = [
            i for i in ta.get("strategic_insights", []) if "5G" in i
        ]
        if erosion_insight:
            km_parts.append(erosion_insight[0])
        km = " | ".join(km_parts) if km_parts else "Value and 5G premium analysis"
        self._add_key_message_bar(slide, km)

    def _add_tariff_price_evolution(self, ta):
        """Slide: Price Evolution table — all operators/tiers across snapshots."""
        evolution = ta.get("price_evolution", {})
        if not evolution:
            return

        slide = self._new_slide("tariff_evolution",
                                "Price Evolution — Mobile Postpaid")
        self._add_header(slide, "Price Evolution — Mobile Postpaid",
                         "H1_2023 to H1_2026")

        # Build a table as image: rows = operator+tier, cols = snapshots
        # Collect all snapshots
        all_snaps = set()
        for op, timeline in evolution.items():
            for entry in timeline:
                all_snaps.add(entry["snapshot"])
        snapshots_sorted = sorted(all_snaps)

        # Build table data
        row_labels = []
        col_labels = snapshots_sorted
        operators_sorted = sorted(evolution.keys())
        tier_order = ["s", "m", "l", "xl"]

        cell_data = {}  # operator -> list of values per snapshot
        for op in operators_sorted:
            op_label = op.replace('_', ' ').title()
            timeline = evolution[op]
            snap_map = {e["snapshot"]: e for e in timeline}
            for tier in tier_order:
                row_key = f"{op_label} {tier.upper()}"
                values = []
                for snap in snapshots_sorted:
                    entry = snap_map.get(snap, {})
                    price = entry.get(tier)
                    values.append(f"EUR {price:.0f}" if price else "—")
                row_labels.append(row_key)
                cell_data[row_key] = values

        if row_labels and col_labels:
            chart_path = self.chart_gen.create_kpi_table_chart(
                col_labels, cell_data,
                target_operator=self.operator_id.replace('_', ' ').title(),
                title="Mobile Postpaid Monthly Prices (EUR)",
                filename="tariff_evolution_table.png",
            )
            self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                            Inches(12.5), Inches(5.0))

        km = "Price evolution across 7 snapshots — H1_2023 to H1_2026"
        self._add_key_message_bar(slide, km)

    def _add_tariff_fixed_and_fmc(self, ta):
        """Slide: Fixed & Convergence Pricing Landscape — tables."""
        fixed = ta.get("fixed_comparison", {})
        fmc = ta.get("fmc_comparison", [])
        if not fixed and not fmc:
            return

        slide = self._new_slide("tariff_fixed_fmc",
                                "Fixed & Convergence Pricing")
        self._add_header(slide, "Fixed & Convergence Pricing Landscape",
                         "DSL / Cable / Fiber / FMC Bundles")

        y = Inches(1.3)

        # Fixed broadband comparison: one section per technology
        for plan_type, label in [
            ("fixed_dsl", "DSL"),
            ("fixed_cable", "Cable"),
            ("fixed_fiber", "Fiber"),
        ]:
            entries = fixed.get(plan_type, [])
            if not entries:
                continue

            self._add_text_box(slide, Inches(0.5), y, Inches(2), Inches(0.3),
                               f"{label}:", font_size=11,
                               font_color=self.style.primary_color, bold=True)

            items = []
            for e in entries[:4]:
                speed = f"{e['speed_mbps']}Mbps" if e.get('speed_mbps') else ""
                items.append(
                    f"{e['display_name']} {e['plan_name']}: "
                    f"EUR {e['price']:.0f}/mo {speed}"
                )
            self._add_bullet_list(slide, Inches(2.5), y, Inches(5), Inches(0.8),
                                  items, font_size=9)
            y += Inches(0.8 + min(len(items), 4) * 0.15)

        # FMC bundles (right column)
        if fmc:
            self._add_text_box(slide, Inches(8.0), Inches(1.3), Inches(4.8), Inches(0.3),
                               "FMC Bundles:", font_size=12,
                               font_color=self.style.primary_color, bold=True)
            fmc_items = []
            for f in fmc:
                fmc_items.append(
                    f"{f['display_name']} {f['plan_name']}: EUR {f['price']:.0f}/mo"
                )
            self._add_bullet_list(slide, Inches(8.0), Inches(1.7), Inches(4.8),
                                  Inches(4.5), fmc_items[:8], font_size=10)

        # Key message
        fmc_insight = [
            i for i in ta.get("strategic_insights", []) if "FMC" in i
        ]
        km = fmc_insight[0] if fmc_insight else "Fixed and convergence pricing overview"
        self._add_key_message_bar(slide, km)

    # =========================================================================
    # Look 3: Competition
    # =========================================================================

    def _add_five_forces(self, comp):
        if comp is None:
            return

        slide = self._new_slide("five_forces", "Porter's Five Forces")
        self._add_header(slide, "Porter's Five Forces", "Competitive Forces")

        forces = comp.five_forces if hasattr(comp, 'five_forces') else {}
        if forces:
            chart_path = self.chart_gen.create_porter_five_forces(
                forces, filename="porter_forces.png")
            self._add_image(slide, chart_path, Inches(3), Inches(1.2),
                            Inches(7), Inches(5))

        # Overall intensity label
        intensity = getattr(comp, 'overall_competition_intensity', 'medium')
        self._add_text_box(slide, Inches(0.5), Inches(1.5), Inches(2.5), Inches(0.4),
                           f"Overall Intensity: {intensity.upper()}",
                           font_size=13, font_color=self.style.primary_color, bold=True)

        km = comp.key_message or "Competitive forces analysis"
        self._add_key_message_bar(slide, km)

    def _add_competitor_deep_dives(self, comp):
        if comp is None:
            return
        analyses = comp.competitor_analyses if hasattr(comp, 'competitor_analyses') else {}
        for operator_id, deep_dive in list(analyses.items())[:3]:
            slide = self._new_slide(f"competitor_{operator_id}", f"Competitor: {operator_id}")
            op_name = operator_id.replace('_', ' ').title()
            self._add_header(slide, f"Competitor Analysis: {op_name}",
                             "Competitor Deep Dive")

            # Financial health — formatted KPI cards
            fh = deep_dive.financial_health if hasattr(deep_dive, 'financial_health') else {}
            if fh:
                kpi_map = [
                    ('revenue', 'Revenue', '€{:,.0f}M'),
                    ('ebitda', 'EBITDA', '€{:,.0f}M'),
                    ('ebitda_margin_pct', 'Margin', '{:.1f}%'),
                    ('capex_to_revenue_pct', 'Capex/Rev', '{:.1f}%'),
                ]
                metrics = []
                for key, label, fmt in kpi_map:
                    v = fh.get(key)
                    if v is not None and not isinstance(v, str):
                        growth_key = f'{key.replace("_pct","")}_growth_pct'
                        change = ''
                        g = fh.get(growth_key) or fh.get('ebitda_growth_pct')
                        if key == 'revenue' and fh.get('service_revenue_growth_pct') is not None:
                            change = f"{fh['service_revenue_growth_pct']:+.1f}% YoY"
                        metrics.append((label, fmt.format(v), change))
                self._add_metric_cards(slide, metrics[:4], top=1.4)

            # Strengths / Weaknesses (left side)
            y = Inches(2.8)
            cols = [
                ("Strengths", getattr(deep_dive, 'strengths', []), self.style.positive_color),
                ("Weaknesses", getattr(deep_dive, 'weaknesses', []), self.style.negative_color),
            ]
            x = Inches(0.5)
            for col_title, items, color in cols:
                self._add_text_box(slide, x, y, Inches(5.5), Inches(0.3),
                                   col_title, font_size=13, font_color=color, bold=True)
                self._add_bullet_list(slide, x, y + Inches(0.3), Inches(5.5),
                                      Inches(2.5), items[:4], font_size=11)
                x += Inches(6)

            # Growth strategy
            strategy = getattr(deep_dive, 'growth_strategy', '')
            if strategy:
                self._add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.5),
                                   f"Strategy: {strategy}", font_size=11,
                                   font_color=self.style.text_color)

            self._add_key_message_bar(slide, f"Competitor analysis: {op_name}")

    def _add_competition_summary(self, comp):
        if comp is None:
            return

        slide = self._new_slide("competition_summary", "Competition Summary")
        self._add_header(slide, "Competitive Landscape Summary", "Competition")

        # Comparison table
        table = comp.comparison_table if hasattr(comp, 'comparison_table') else {}
        if table:
            metrics = list(table.keys())[:6]
            operators = set()
            for m in metrics:
                operators.update(table[m].keys())
            operators = sorted(operators)[:5]

            data = {op: [str(table[m].get(op, 'N/A')) for m in metrics]
                    for op in operators}
            if data:
                chart_path = self.chart_gen.create_kpi_table_chart(
                    metrics, data, target_operator=self.operator_id,
                    title="Competitive Comparison", filename="comp_table.png")
                self._add_image(slide, chart_path, Inches(0.5), Inches(1.3),
                                Inches(12), Inches(4.5))

        landscape = getattr(comp, 'competitive_landscape', '')
        if landscape and not table:
            self._add_text_box(slide, Inches(0.5), Inches(1.4), Inches(12),
                               Inches(4), landscape, font_size=13,
                               font_color=self.style.text_color)

        km = comp.key_message or "Competitive landscape overview"
        self._add_key_message_bar(slide, km)

    def _add_gap_analysis(self, comp):
        """Add a gap analysis chart comparing target operator vs market leader."""
        if comp is None:
            return
        table = comp.comparison_table if hasattr(comp, 'comparison_table') else {}
        if not table:
            return

        # Find the market leader (highest revenue that isn't us)
        rev_data = table.get('revenue', table.get('total_revenue', {}))
        if not rev_data:
            return

        our_rev = rev_data.get(self.operator_id, 0)
        leader_id = None
        leader_rev = 0
        for op_id, rev in rev_data.items():
            if op_id != self.operator_id and isinstance(rev, (int, float)):
                if rev > leader_rev:
                    leader_rev = rev
                    leader_id = op_id
        if leader_id is None:
            return

        # Collect numeric metrics for gap analysis (try both naming conventions)
        numeric_metrics = ['revenue', 'ebitda_margin', 'subscribers', 'arpu',
                           '5g_coverage', 'total_revenue', 'ebitda_margin_pct',
                           'mobile_total_k']
        dimensions = []
        target_scores = []
        leader_scores = []
        for metric in numeric_metrics:
            metric_data = table.get(metric, {})
            our_val = metric_data.get(self.operator_id)
            their_val = metric_data.get(leader_id)
            if our_val is not None and their_val is not None:
                try:
                    our_f = float(our_val) if not isinstance(our_val, (int, float)) else our_val
                    their_f = float(their_val) if not isinstance(their_val, (int, float)) else their_val
                    if their_f != 0:
                        # Normalize to 0-100 scale relative to leader
                        dimensions.append(metric.replace('_', ' ').replace('total ', '').title())
                        target_scores.append(our_f / their_f * 100)
                        leader_scores.append(100)
                except (ValueError, TypeError):
                    continue

        if len(dimensions) < 3:
            return

        slide = self._new_slide("gap_analysis", "Gap Analysis")
        our_name = self.operator_id.replace('_', ' ').title()
        leader_name = leader_id.replace('_', ' ').title()
        self._add_header(slide, f"Gap Analysis: {our_name} vs {leader_name}",
                         "Competitive Gap")

        chart_path = self.chart_gen.create_gap_analysis_chart(
            dimensions, target_scores, leader_scores,
            target_name=our_name, leader_name=leader_name,
            title=f"Competitive Gap (Indexed to {leader_name} = 100)",
            filename="gap_analysis.png")
        self._add_image(slide, chart_path, Inches(0.5), Inches(1.3),
                        Inches(12), Inches(4.8))

        km = f"Gap analysis vs market leader {leader_name}"
        self._add_key_message_bar(slide, km)

    def _add_new_entrants(self, comp):
        if comp is None:
            return
        forces = comp.five_forces if hasattr(comp, 'five_forces') else {}
        ne_force = forces.get('new_entrants')
        if ne_force is None:
            return
        if not (hasattr(ne_force, 'key_factors') and ne_force.key_factors):
            return

        slide = self._new_slide("new_entrants", "New Entrants", required=False)
        self._add_header(slide, "New Entrants Threat", "Market Entry")

        items = []
        for factor in ne_force.key_factors[:6]:
            if isinstance(factor, dict):
                items.append(f"{factor.get('name', '')}: {factor.get('description', '')}")
            else:
                items.append(str(factor))
        self._add_bullet_list(slide, Inches(0.5), Inches(1.4), Inches(12),
                              Inches(4), items, font_size=13)

        level = getattr(ne_force, 'force_level', 'medium')
        km = f"New entrants threat level: {level}"
        self._add_key_message_bar(slide, km)

    # =========================================================================
    # Look 4: Self
    # =========================================================================

    def _add_health_check(self, self_analysis):
        if self_analysis is None:
            return

        slide = self._new_slide("health_check", "Health Check")
        self._add_header(slide, "Operating Health Check", "Financial Health")

        # Financial health metrics — show formatted top-line KPIs
        fh = self_analysis.financial_health or {}
        if fh:
            kpi_map = [
                ('total_revenue', 'Revenue', '€{:,.0f}M'),
                ('ebitda', 'EBITDA', '€{:,.0f}M'),
                ('ebitda_margin_pct', 'EBITDA Margin', '{:.1f}%'),
                ('capex_to_revenue_pct', 'Capex/Rev', '{:.1f}%'),
            ]
            metrics = []
            for key, label, fmt in kpi_map:
                v = fh.get(key)
                if v is not None:
                    change = ''
                    yoy = fh.get(f'{key.replace("_pct","")}_yoy_pct') or fh.get('revenue_yoy_pct')
                    if key == 'total_revenue' and fh.get('revenue_yoy_pct') is not None:
                        change = f"+{fh['revenue_yoy_pct']}% YoY"
                    elif key == 'ebitda' and fh.get('ebitda_growth_pct') is not None:
                        change = f"+{fh['ebitda_growth_pct']}% YoY"
                    metrics.append((label, fmt.format(v), change))
            self._add_metric_cards(slide, metrics[:4], top=1.4)

        # Revenue trend chart (left)
        rev_trend = fh.get('revenue_trend', [])
        if rev_trend and len(rev_trend) >= 4:
            quarters = [f'Q{i+1}' for i in range(len(rev_trend))]
            chart_path = self.chart_gen.create_multi_line_trend(
                quarters, {self.operator_id.replace('_', ' ').title(): rev_trend},
                title="Revenue Trend (€M)", y_label="€M",
                target_operator=self.operator_id.replace('_', ' ').title(),
                filename="revenue_trend.png")
            self._add_image(slide, chart_path, Inches(0.3), Inches(2.9),
                            Inches(6), Inches(3.1))

        # Revenue breakdown bar chart (right)
        rb = self_analysis.revenue_breakdown or {}
        rb_items = {k: v for k, v in rb.items()
                    if isinstance(v, dict) and 'value' in v}
        if rb_items:
            labels = [k.replace('_revenue', '').replace('_', ' ').title()
                      for k in rb_items.keys()]
            values = [v['value'] for v in rb_items.values()]
            chart_path = self.chart_gen.create_bar_chart(
                labels, values, title="Revenue Breakdown (€M)",
                y_label="€M", filename="revenue_breakdown.png")
            self._add_image(slide, chart_path, Inches(6.5), Inches(2.9),
                            Inches(6.3), Inches(3.1))

        km = self_analysis.key_message or "Operating health check summary"
        self._add_key_message_bar(slide, km)

    def _add_revenue_composition_trend(self, self_analysis):
        """Stacked bar chart showing revenue composition by segment over 8 quarters."""
        if self_analysis is None:
            return
        segments = self_analysis.segment_analyses or []
        # Need at least 2 segments with revenue trend data
        seg_data = {}
        for seg in segments:
            rev = (seg.trend_data or {}).get('revenue', [])
            if rev and len(rev) >= 4:
                seg_data[seg.segment_name] = rev
        if len(seg_data) < 2:
            return

        slide = self._new_slide("revenue_composition", "Revenue Composition Trend")
        self._add_header(slide, "Revenue Composition Trend",
                         "Segment Revenue Mix Over Time")

        # Build quarter labels matching data length
        n = len(next(iter(seg_data.values())))
        x_labels = [f"Q{i+1}" for i in range(n)]

        chart_path = self.chart_gen.create_stacked_bar(
            x_labels, seg_data,
            title="Revenue by Segment (€M)",
            y_label="€M",
            filename="revenue_composition_trend.png")
        self._add_image(slide, chart_path, Inches(0.5), Inches(1.5),
                        Inches(12), Inches(4.5))

        # Key message: identify growing vs declining segments
        growing = [n for n, v in seg_data.items() if v[-1] > v[0]]
        declining = [n for n, v in seg_data.items() if v[-1] < v[0]]
        parts = []
        if growing:
            parts.append(f"Growing: {', '.join(growing)}")
        if declining:
            parts.append(f"Declining: {', '.join(declining)}")
        km = " | ".join(parts) if parts else "Revenue composition over time"
        self._add_key_message_bar(slide, km)

    def _add_revenue_comparison(self, self_analysis, competition):
        """Add a multi-operator revenue comparison slide with trend + bar charts."""
        if self_analysis is None or competition is None:
            return
        table = competition.comparison_table if hasattr(competition, 'comparison_table') else {}
        rev_data = table.get('revenue', table.get('total_revenue', {}))
        if not rev_data or len(rev_data) < 2:
            return

        slide = self._new_slide("revenue_comparison", "Revenue Comparison")
        self._add_header(slide, "Revenue Comparison", "Competitive Benchmarking")

        # Left: horizontal bar chart of revenue across operators
        categories = [k.replace('_', ' ').title() for k in rev_data.keys()]
        values = []
        for v in rev_data.values():
            try:
                values.append(float(v))
            except (TypeError, ValueError):
                values.append(0)

        our_name = self.operator_id.replace('_', ' ').title()
        chart_path = self.chart_gen.create_horizontal_bar_chart(
            categories, values, target_category=our_name,
            title="Revenue by Operator (€M)", x_label="€M",
            filename="revenue_comparison.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.3),
                        Inches(6.3), Inches(4.8))

        # Right: margin comparison (if available)
        margin_data = table.get('ebitda_margin', table.get('ebitda_margin_pct', {}))
        if margin_data and len(margin_data) >= 2:
            m_cats = [k.replace('_', ' ').title() for k in margin_data.keys()]
            m_vals = []
            for v in margin_data.values():
                try:
                    m_vals.append(float(v))
                except (TypeError, ValueError):
                    m_vals.append(0)
            chart_path2 = self.chart_gen.create_bar_chart(
                m_cats, m_vals, target_category=our_name,
                title="EBITDA Margin (%)", y_label="%",
                filename="margin_comparison.png")
            self._add_image(slide, chart_path2, Inches(6.8), Inches(1.3),
                            Inches(6), Inches(4.8))

        km = "Revenue and margin benchmarking across operators"
        self._add_key_message_bar(slide, km)

    def _add_segment_deep_dives(self, self_analysis):
        if self_analysis is None:
            return
        segments = self_analysis.segment_analyses or []
        for seg_idx, seg in enumerate(segments[:4]):
            slide = self._new_slide(
                f"segment_{seg.segment_id}" if hasattr(seg, 'segment_id') else f"segment_{self.slide_num}",
                f"Segment: {seg.segment_name}")
            self._add_header(slide, f"Segment Analysis: {seg.segment_name}",
                             "Segment Performance")

            # Key metrics — formatted with proper labels
            km_data = seg.key_metrics if hasattr(seg, 'key_metrics') else {}
            if km_data:
                formatted_metrics = []
                for k, v in list(km_data.items())[:4]:
                    label = k.replace('_', ' ').replace('mobile ', '').replace('fixed ', '').replace('broadband ', '').replace('b2b ', '').replace('tv ', '').title()
                    if isinstance(v, float):
                        if 'pct' in k:
                            formatted_metrics.append((label, f"{v:.1f}%", ""))
                        elif 'revenue' in k or v > 100:
                            formatted_metrics.append((label, f"€{v:,.0f}M", ""))
                        elif 'arpu' in k:
                            formatted_metrics.append((label, f"€{v:.1f}", ""))
                        else:
                            formatted_metrics.append((label, f"{v:,.0f}", ""))
                    else:
                        formatted_metrics.append((label, str(v), ""))
                self._add_metric_cards(slide, formatted_metrics[:4], top=1.4)

            # Bar chart of key numeric metrics (right side)
            numeric_items = {k: v for k, v in km_data.items()
                            if isinstance(v, (int, float)) and 'pct' not in k and v > 0}
            if len(numeric_items) >= 2:
                labels = [k.replace('_', ' ').replace('mobile ', '').replace('fixed ', '').replace('broadband ', '').replace('b2b ', '').replace('tv ', '').title()
                         for k in numeric_items.keys()]
                values = list(numeric_items.values())
                chart_path = self.chart_gen.create_bar_chart(
                    labels[:6], values[:6],
                    title=f"{seg.segment_name} Metrics",
                    filename=f"segment_{seg_idx}_metrics.png")
                self._add_image(slide, chart_path, Inches(6.5), Inches(2.8),
                                Inches(6), Inches(3.2))

            # Health status + Changes (left side)
            y = Inches(2.8)
            health = getattr(seg, 'health_status', 'stable')
            self._add_text_box(slide, Inches(0.5), y, Inches(3), Inches(0.3),
                               f"Health: {health.upper()}", font_size=12,
                               font_color=self.style.primary_color, bold=True)

            changes = getattr(seg, 'changes', [])
            attributions = getattr(seg, 'attributions', [])
            combined_items = []
            for c in changes[:3]:
                if hasattr(c, 'metric'):
                    combined_items.append(f"{c.metric}: {getattr(c, 'direction', '')}")
                else:
                    combined_items.append(str(c))
            for a in attributions[:3]:
                if hasattr(a, 'description'):
                    combined_items.append(a.description)
            if combined_items:
                self._add_bullet_list(slide, Inches(0.5), y + Inches(0.4),
                                      Inches(5.5), Inches(2.5), combined_items, font_size=11)

            seg_km = getattr(seg, 'key_message', '') or f"{seg.segment_name} segment analysis"
            self._add_key_message_bar(slide, seg_km)

    def _add_network_analysis(self, self_analysis):
        if self_analysis is None:
            return
        network = getattr(self_analysis, 'network', None)
        if network is None:
            return

        slide = self._new_slide("network", "Network Analysis")
        self._add_header(slide, "Network Analysis", "Infrastructure")

        y = Inches(1.4)
        # Coverage gauges
        coverage = network.coverage or {}
        if coverage:
            labels = list(coverage.keys())[:4]
            values = [float(coverage[k]) for k in labels]
            chart_path = self.chart_gen.create_donut_gauges(
                labels, values, title="Network Coverage", filename="coverage_gauges.png")
            self._add_image(slide, chart_path, Inches(0.5), Inches(1.3),
                            Inches(12), Inches(2.5))
            y = Inches(4.0)

        # Tech mix
        tech = network.technology_mix or {}
        if tech:
            items = []
            for k, v in tech.items():
                label = k.replace('_', ' ').title()
                if isinstance(v, dict):
                    formatted = ", ".join(f"{sk}: {sv}" for sk, sv in v.items())
                    items.append(f"{label}: {formatted}")
                elif isinstance(v, (int, float)) and 'pct' in k:
                    items.append(f"{label}: {v}%")
                else:
                    items.append(f"{label}: {v}")
            self._add_text_box(slide, Inches(0.5), y, Inches(3), Inches(0.3),
                               "Technology Mix:", font_size=12,
                               font_color=self.style.primary_color, bold=True)
            self._add_bullet_list(slide, Inches(0.5), y + Inches(0.3), Inches(5),
                                  Inches(1.5), items, font_size=11)

        # Investment direction
        inv = network.investment_direction or ''
        if inv:
            self._add_text_box(slide, Inches(6.5), y, Inches(5.5), Inches(2),
                               f"Investment Direction: {inv}", font_size=12,
                               font_color=self.style.text_color)

        km = "Network infrastructure and coverage assessment"
        self._add_key_message_bar(slide, km)

    def _add_bmc_canvas(self, self_analysis):
        if self_analysis is None:
            return
        bmc = getattr(self_analysis, 'bmc', None)
        if bmc is None:
            return

        slide = self._new_slide("bmc", "Business Model Canvas")
        self._add_header(slide, "Business Model Canvas", "BMC")

        chart_path = self.chart_gen.create_bmc_canvas(bmc, filename="bmc_canvas.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(12.5), Inches(4.9))

        km = "Business model canvas overview"
        self._add_key_message_bar(slide, km)

    def _add_org_talent(self, self_analysis):
        if self_analysis is None:
            return

        slide = self._new_slide("org_talent", "Organization & Talent")
        self._add_header(slide, "Organization & Talent", "People & Culture")

        y = Inches(1.4)
        # Top: KPI cards for workforce metrics
        fh = self_analysis.financial_health or {}
        org_metrics = []
        if fh.get('employees'):
            org_metrics.append(("Employees", f"{fh['employees']:,.0f}", ""))
        talent = self_analysis.talent_assessment or {}
        for k, v in list(talent.items())[:3]:
            label = k.replace('_', ' ').title()
            org_metrics.append((label, str(v), ""))
        if org_metrics:
            self._add_metric_cards(slide, org_metrics[:4], top=1.4)
            y = Inches(2.9)

        # Left: Org culture
        culture = getattr(self_analysis, 'org_culture', '')
        if culture:
            self._add_text_box(slide, Inches(0.5), y, Inches(5.5), Inches(0.3),
                               "Organization Culture:", font_size=12,
                               font_color=self.style.primary_color, bold=True)
            self._add_text_box(slide, Inches(0.5), y + Inches(0.35), Inches(5.5),
                               Inches(1.2), culture, font_size=11,
                               font_color=self.style.text_color)

        # Right: Talent score bar chart (if talent has numeric values)
        numeric_talent = {k: v for k, v in talent.items()
                          if isinstance(v, (int, float))}
        if len(numeric_talent) >= 2:
            labels = [k.replace('_', ' ').title() for k in numeric_talent.keys()]
            values = list(numeric_talent.values())
            chart_path = self.chart_gen.create_horizontal_bar_chart(
                labels, values,
                title="Talent Assessment Scores",
                filename="talent_scores.png")
            self._add_image(slide, chart_path, Inches(6.5), y,
                            Inches(6), Inches(3))
        elif talent:
            # Text fallback
            items = [f"{k}: {v}" for k, v in talent.items()]
            self._add_text_box(slide, Inches(6.5), y, Inches(5.5), Inches(0.3),
                               "Talent Assessment:", font_size=12,
                               font_color=self.style.primary_color, bold=True)
            self._add_bullet_list(slide, Inches(6.5), y + Inches(0.35), Inches(5.5),
                                  Inches(2), items, font_size=11)

        # Leadership changes
        changes = self_analysis.leadership_changes or []
        if changes:
            change_y = Inches(5.0) if org_metrics else Inches(4.0)
            self._add_text_box(slide, Inches(0.5), change_y, Inches(12), Inches(0.3),
                               "Leadership Changes:", font_size=12,
                               font_color=self.style.primary_color, bold=True)
            items = [str(c) if not isinstance(c, dict) else
                     f"{c.get('name', '')}: {c.get('role', '')} ({c.get('date', '')})"
                     for c in changes[:4]]
            self._add_bullet_list(slide, Inches(0.5), change_y + Inches(0.3), Inches(12),
                                  Inches(1), items, font_size=11)

        km = "Organization and talent assessment"
        self._add_key_message_bar(slide, km)

    def _add_strengths_weaknesses_exposure(self, self_analysis):
        if self_analysis is None:
            return

        slide = self._new_slide("sw_exposure", "Strengths/Weaknesses/Exposure")
        self._add_header(slide, "Strengths, Weaknesses & Exposure Points",
                         "Self Assessment")

        # Three columns
        cols = [
            ("Strengths", getattr(self_analysis, 'strengths', []),
             self.style.positive_color, Inches(0.5)),
            ("Weaknesses", getattr(self_analysis, 'weaknesses', []),
             self.style.negative_color, Inches(4.5)),
            ("Exposure Points", [], self.style.warning_color, Inches(8.5)),
        ]
        # Build exposure items
        exposure = getattr(self_analysis, 'exposure_points', [])
        exposure_items = []
        for ep in exposure[:5]:
            if hasattr(ep, 'trigger_action'):
                exposure_items.append(f"{ep.trigger_action}: {ep.side_effect}")
            else:
                exposure_items.append(str(ep))
        cols[2] = ("Exposure Points", exposure_items, self.style.warning_color, Inches(8.5))

        for title, items, color, x in cols:
            self._add_text_box(slide, x, Inches(1.4), Inches(3.8), Inches(0.3),
                               title, font_size=13, font_color=color, bold=True)
            self._add_bullet_list(slide, x, Inches(1.8), Inches(3.8), Inches(4),
                                  items[:6], font_size=11)

        km = self_analysis.key_message or "Self-assessment summary"
        self._add_key_message_bar(slide, km)

    # =========================================================================
    # SWOT Synthesis
    # =========================================================================

    def _add_swot_matrix(self, swot):
        if swot is None:
            return

        slide = self._new_slide("swot", "SWOT Matrix")
        self._add_header(slide, "SWOT Analysis", "Strategic Synthesis")

        chart_path = self.chart_gen.create_swot_matrix(swot, filename="swot_matrix.png")
        self._add_image(slide, chart_path, Inches(0.5), Inches(1.2),
                        Inches(7.5), Inches(4.9))

        # Strategy quadrants on the right
        strategies = [
            ("SO", getattr(swot, 'so_strategies', []), self.style.positive_color),
            ("WO", getattr(swot, 'wo_strategies', []), (70, 130, 180)),
            ("ST", getattr(swot, 'st_strategies', []), self.style.warning_color),
            ("WT", getattr(swot, 'wt_strategies', []), self.style.negative_color),
        ]
        y = Inches(1.4)
        for label, items, color in strategies:
            self._add_text_box(slide, Inches(8.5), y, Inches(1), Inches(0.25),
                               f"{label}:", font_size=10, font_color=color, bold=True)
            if items:
                text = items[0] if len(items[0]) <= 60 else items[0][:57] + "..."
                self._add_text_box(slide, Inches(9.3), y, Inches(3.5), Inches(0.6),
                                   text, font_size=9, font_color=self.style.text_color)
            y += Inches(0.7)

        km = swot.key_message or "SWOT synthesis and strategic implications"
        self._add_key_message_bar(slide, km)

    # =========================================================================
    # Look 5: Opportunities
    # =========================================================================

    def _add_span_bubble(self, opp):
        if opp is None:
            return

        slide = self._new_slide("span_bubble", "SPAN Matrix")
        self._add_header(slide, "SPAN Matrix — Opportunity Positioning", "Opportunity Map")

        positions = opp.span_positions if hasattr(opp, 'span_positions') else []
        if positions:
            chart_path = self.chart_gen.create_span_bubble_chart(
                positions, filename="span_bubble.png")
            self._add_image(slide, chart_path, Inches(0.5), Inches(1.2),
                            Inches(11.5), Inches(4.8))
        else:
            self._add_text_box(slide, Inches(2), Inches(3), Inches(9), Inches(1),
                               "No SPAN positions available",
                               font_size=14, font_color=self.style.light_text_color,
                               align="center")

        km = opp.key_message or "Opportunity positioning on SPAN matrix"
        self._add_key_message_bar(slide, km)

    def _add_priority_table(self, opp):
        if opp is None:
            return
        opportunities = getattr(opp, 'opportunities', [])
        if not opportunities:
            return

        slide = self._new_slide("priority_table", "Opportunity Priorities")
        self._add_header(slide, "Opportunity Priority Ranking", "Priorities")

        items = [o.name if hasattr(o, 'name') else str(o) for o in opportunities[:8]]
        priorities = [getattr(o, 'priority', 'P1') for o in opportunities[:8]]

        chart_path = self.chart_gen.create_priority_chart(
            items, priorities, title="", filename="priority_table.png")
        self._add_image(slide, chart_path, Inches(0.5), Inches(1.3),
                        Inches(12), Inches(4.5))

        km = opp.key_message or "Prioritized opportunities"
        self._add_key_message_bar(slide, km)

    def _add_opportunity_deep_dive(self, opp):
        if opp is None:
            return
        opportunities = getattr(opp, 'opportunities', [])
        if not opportunities:
            return

        slide = self._new_slide("opp_deep_dive", "Opportunity Details", required=False)
        self._add_header(slide, "Opportunity Details", "Deep Dive")

        y = Inches(1.4)
        for item in opportunities[:4]:
            name = item.name if hasattr(item, 'name') else str(item)
            desc = getattr(item, 'description', '')
            priority = getattr(item, 'priority', 'P1')
            market = getattr(item, 'addressable_market', 'N/A')
            window = getattr(item, 'time_window', '')

            self._add_shape(slide, Inches(0.5), y, Inches(12), Inches(1.0),
                            (248, 248, 248))
            # Priority badge
            p_colors = {'P0': self.style.primary_color,
                        'P1': self.style.warning_color,
                        'P2': self.style.positive_color}
            self._add_shape(slide, Inches(0.6), y + Inches(0.1), Inches(0.5), Inches(0.35),
                            p_colors.get(priority, self.style.neutral_color))
            self._add_text_box(slide, Inches(0.6), y + Inches(0.12), Inches(0.5),
                               Inches(0.3), priority, font_size=10,
                               font_color=(255, 255, 255), bold=True, align="center")

            self._add_text_box(slide, Inches(1.3), y + Inches(0.05), Inches(5),
                               Inches(0.3), name, font_size=13,
                               font_color=self.style.primary_color, bold=True)
            detail = f"Market: {market} | Window: {window}"
            self._add_text_box(slide, Inches(1.3), y + Inches(0.35), Inches(10),
                               Inches(0.25), detail, font_size=10,
                               font_color=self.style.light_text_color)
            if desc:
                self._add_text_box(slide, Inches(1.3), y + Inches(0.6), Inches(10),
                                   Inches(0.35), desc, font_size=10,
                                   font_color=self.style.text_color)
            y += Inches(1.15)

        km = opp.key_message or "Detailed opportunity analysis"
        self._add_key_message_bar(slide, km)

    # =========================================================================
    # Three Decisions slides (BLM Phase 2)
    # =========================================================================

    def _add_three_decisions_slides(self, result):
        """Generate Three Decisions slides from FiveLooksResult."""
        try:
            from src.output.strategic_diagnosis import StrategicDiagnosisComputer
            from src.blm.three_decisions_engine import ThreeDecisionsComputer
            from src.models.market_configs import get_market_config

            config = get_market_config(result.market) if result.market else None
            diag = StrategicDiagnosisComputer(result, config).compute()
            decisions = ThreeDecisionsComputer(result, diag, config).compute()
        except Exception:
            return  # Graceful fallback — skip if engine fails

        self._add_section_divider("Three Decisions",
                                  "Strategy — Key Tasks — Execution")
        self._add_strategy_slide(decisions)
        self._add_key_tasks_slide(decisions)
        self._add_execution_slide(decisions)

    def _add_strategy_slide(self, decisions):
        """Decision 1: Strategic Direction with 4 pillars."""
        slide = self._new_slide("strategy_decision", "Define Strategy")
        self._add_header(slide, "Decision 1: Define Strategy",
                         decisions.strategy.competitive_posture)

        # Overall direction as subtitle
        self._add_text_box(
            slide, Inches(0.5), Inches(1.3), Inches(12), Inches(0.6),
            decisions.strategy.overall_direction,
            font_size=13, bold=True,
            font_color=self.style.accent_color)

        # 4 strategic pillars as colored cards
        pillars = decisions.strategy.pillars
        card_w = Inches(3.0)
        card_h = Inches(3.5)
        gap = Inches(0.15)
        left_start = Inches(0.35)
        top = Inches(2.1)

        priority_colors = {"P0": (0xC0, 0x39, 0x2B),
                           "P1": (0xE6, 0x7E, 0x22),
                           "P2": (0x27, 0xAE, 0x60)}

        for i, pillar in enumerate(pillars[:4]):
            left = left_start + i * (card_w + gap)
            # Card background
            fill = (0xF8, 0xF9, 0xFA) if i % 2 == 0 else (0xEB, 0xF5, 0xFB)
            self._add_shape(slide, left, top, card_w, card_h, fill_color=fill)

            # Priority badge background
            p_color = priority_colors.get(pillar.priority, (0x95, 0xA5, 0xA6))
            self._add_shape(slide, left + Inches(0.08), top + Inches(0.08),
                            Inches(0.45), Inches(0.28), fill_color=p_color)
            self._add_text_box(slide, left + Inches(0.12), top + Inches(0.09),
                               Inches(0.4), Inches(0.25), pillar.priority,
                               font_size=9, bold=True, font_color=(0xFF, 0xFF, 0xFF))

            # Pillar name
            self._add_text_box(slide, left + Inches(0.1), top + Inches(0.45),
                               card_w - Inches(0.2), Inches(0.35), pillar.name,
                               font_size=12, bold=True,
                               font_color=self.style.text_color)

            # Direction text
            self._add_text_box(slide, left + Inches(0.1), top + Inches(0.85),
                               card_w - Inches(0.2), Inches(1.3), pillar.direction,
                               font_size=10, font_color=(0x2C, 0x3E, 0x50))

            # KPIs
            kpi_text = "\n".join(f"- {k}" for k in pillar.kpis[:3])
            self._add_text_box(slide, left + Inches(0.1), top + Inches(2.3),
                               card_w - Inches(0.2), Inches(1.0), kpi_text,
                               font_size=9, font_color=(0x7F, 0x8C, 0x8D))

        self._add_key_message_bar(slide, decisions.narrative[:120])

    def _add_key_tasks_slide(self, decisions):
        """Decision 2: Key Tasks prioritized by domain."""
        slide = self._new_slide("key_tasks_decision", "Define Key Tasks")
        self._add_header(slide, "Decision 2: Define Key Tasks",
                         decisions.key_tasks.resource_implication)

        tasks = decisions.key_tasks.tasks
        if not tasks:
            return

        domain_colors = {
            "Network": (0x29, 0x80, 0xB9),
            "Business": (0x27, 0xAE, 0x60),
            "Customer": (0x8E, 0x44, 0xAD),
            "Efficiency": (0xE6, 0x7E, 0x22),
        }
        priority_fills = {
            "P0": (0xFD, 0xED, 0xEC),
            "P1": (0xFE, 0xF5, 0xE7),
            "P2": (0xEA, 0xFB, 0xF0),
        }
        priority_badge = {
            "P0": (0xC0, 0x39, 0x2B),
            "P1": (0xE6, 0x7E, 0x22),
            "P2": (0x27, 0xAE, 0x60),
        }

        row_h = Inches(0.65)
        top = Inches(1.5)

        for i, task in enumerate(tasks[:8]):
            y = top + i * row_h
            fill = priority_fills.get(task.priority, (0xF5, 0xF5, 0xF5))
            self._add_shape(slide, Inches(0.3), y, Inches(12.4),
                            row_h - Inches(0.05), fill_color=fill)

            # Priority badge
            self._add_shape(slide, Inches(0.4), y + Inches(0.15),
                            Inches(0.45), Inches(0.3),
                            fill_color=priority_badge.get(task.priority, (0x95, 0xA5, 0xA6)))
            self._add_text_box(slide, Inches(0.42), y + Inches(0.16),
                               Inches(0.42), Inches(0.28), task.priority,
                               font_size=9, bold=True, font_color=(0xFF, 0xFF, 0xFF))

            # Domain tag
            d_color = domain_colors.get(task.domain, (0x7F, 0x8C, 0x8D))
            self._add_text_box(slide, Inches(1.0), y + Inches(0.16),
                               Inches(1.2), Inches(0.3), task.domain,
                               font_size=9, bold=True, font_color=d_color)

            # Task name + description
            self._add_text_box(slide, Inches(2.3), y + Inches(0.08),
                               Inches(5.5), Inches(0.25), task.name,
                               font_size=11, bold=True, font_color=(0x2C, 0x3E, 0x50))
            self._add_text_box(slide, Inches(2.3), y + Inches(0.35),
                               Inches(5.5), Inches(0.25), task.description[:80],
                               font_size=8, font_color=(0x7F, 0x8C, 0x8D))

            # KPIs on right
            kpi_str = " | ".join(task.kpis[:2])
            self._add_text_box(slide, Inches(8.0), y + Inches(0.15),
                               Inches(4.5), Inches(0.35), kpi_str,
                               font_size=8, font_color=(0x56, 0x6C, 0x73))

        km = f"{len(tasks)} critical tasks: {decisions.key_tasks.resource_implication}"
        self._add_key_message_bar(slide, km)

    def _add_execution_slide(self, decisions):
        """Decision 3: Execution timeline + governance."""
        slide = self._new_slide("execution_decision", "Define Execution")
        self._add_header(slide, "Decision 3: Define Execution",
                         "Quarterly Roadmap & Governance")

        # Use timeline chart for milestones
        milestones_data = []
        for ms in decisions.execution.milestones:
            delivs = "\n".join(d[:40] for d in ms.deliverables[:2])
            milestones_data.append({
                "date": ms.quarter,
                "name": f"{ms.name}\n{delivs}",
                "priority": ms.priority,
            })

        if milestones_data:
            chart_path = self.chart_gen.create_timeline_chart(
                milestones_data,
                title="Execution Roadmap",
                filename="execution_timeline.png")
            self._add_image(slide, chart_path, Inches(0.3), Inches(1.3),
                            Inches(12.5), Inches(2.8))

        # Governance section
        gov = decisions.execution.governance
        if gov:
            gov_top = Inches(4.4)
            self._add_text_box(slide, Inches(0.5), gov_top,
                               Inches(4), Inches(0.3), "Governance",
                               font_size=12, bold=True,
                               font_color=self.style.text_color)
            for j, g in enumerate(gov[:3]):
                self._add_text_box(
                    slide, Inches(0.5), gov_top + Inches(0.4) + j * Inches(0.4),
                    Inches(5.5), Inches(0.35),
                    f"- {g.mechanism} ({g.cadence})",
                    font_size=9, font_color=(0x2C, 0x3E, 0x50))

        # Traps to avoid section
        traps = decisions.execution.traps_to_avoid
        if traps:
            trap_top = Inches(4.4)
            self._add_text_box(slide, Inches(7.0), trap_top,
                               Inches(5), Inches(0.3), "Strategic Traps to Avoid",
                               font_size=12, bold=True,
                               font_color=(0xC0, 0x39, 0x2B))
            for j, trap in enumerate(traps[:3]):
                trap_text = trap.get("trap", "") if isinstance(trap, dict) else str(trap)
                self._add_text_box(
                    slide, Inches(7.0), trap_top + Inches(0.4) + j * Inches(0.4),
                    Inches(5.5), Inches(0.35),
                    f"x {trap_text}",
                    font_size=9, font_color=(0xC0, 0x39, 0x2B))

        # Risks
        risks = decisions.execution.risk_mitigation
        if risks:
            risk_top = Inches(5.8)
            self._add_text_box(slide, Inches(0.5), risk_top,
                               Inches(12), Inches(0.3), "Key Risks & Mitigation",
                               font_size=11, bold=True,
                               font_color=self.style.text_color)
            for j, risk in enumerate(risks[:3]):
                self._add_text_box(
                    slide, Inches(0.5), risk_top + Inches(0.35) + j * Inches(0.3),
                    Inches(12), Inches(0.28),
                    f"! {risk.get('risk', '')} -- {risk.get('mitigation', '')}",
                    font_size=8, font_color=(0x56, 0x6C, 0x73))

        self._add_key_message_bar(slide,
            "Execute with discipline: monthly reviews, quarterly checkpoints, annual refresh")

    # =========================================================================
    # Summary slides
    # =========================================================================

    def _add_summary_slide(self, result):
        slide = self._new_slide("summary", "Summary")
        self._add_header(slide, "Analysis Summary", "Key Takeaways")

        sections = [
            ("01 Trends", _get_key_message(result.trends)),
            ("02 Market", _get_key_message(result.market_customer)),
            ("03 Competition", _get_key_message(result.competition)),
            ("04 Self", _get_key_message(result.self_analysis)),
            ("SWOT", _get_key_message(result.swot)),
            ("05 Opportunities", _get_key_message(result.opportunities)),
        ]

        y = Inches(1.4)
        for label, msg in sections:
            self._add_shape(slide, Inches(0.5), y, Inches(0.08), Inches(0.6),
                            self.style.primary_color)
            self._add_text_box(slide, Inches(0.8), y + Inches(0.02), Inches(2.2),
                               Inches(0.3), label, font_size=12,
                               font_color=self.style.primary_color, bold=True)
            self._add_text_box(slide, Inches(3.2), y + Inches(0.02), Inches(9.5),
                               Inches(0.55), msg or "—", font_size=11,
                               font_color=self.style.text_color)
            y += Inches(0.7)

        self._add_key_message_bar(slide, "Complete Five Looks analysis — see opportunity priorities for next steps")

    def _add_provenance_appendix(self, result):
        slide = self._new_slide("provenance", "Provenance Appendix")
        self._add_header(slide, "Data Provenance", "Sources & Methodology")

        prov = result.provenance
        footnotes = prov.to_footnotes() if prov else []
        if footnotes:
            self._add_bullet_list(slide, Inches(0.5), Inches(1.4), Inches(12),
                                  Inches(4.5), footnotes[:15], font_size=10)
        else:
            self._add_text_box(slide, Inches(0.5), Inches(2), Inches(12), Inches(1),
                               "No provenance data recorded for this analysis",
                               font_size=13, font_color=self.style.light_text_color)

        self._add_key_message_bar(slide, "All data points are traceable to original sources")

    def _add_back_cover(self):
        slide = self._new_slide("back_cover", "Back Cover")
        self._add_shape(slide, 0, 0, Inches(self.SLIDE_WIDTH), Inches(self.SLIDE_HEIGHT),
                        self.style.primary_color)
        self._add_text_box(slide, Inches(2), Inches(2.5), Inches(9), Inches(1.5),
                           "Thank You", font_size=48,
                           font_color=(255, 255, 255), bold=True, align="center")
        self._add_text_box(slide, Inches(2), Inches(4.0), Inches(9), Inches(1),
                           "BLM Five Looks Strategic Analysis",
                           font_size=20, font_color=(255, 220, 220), align="center")


# =============================================================================
# Helpers
# =============================================================================

def _get_key_message(obj) -> str:
    """Safely extract key_message from a data model object."""
    if obj is None:
        return ""
    return getattr(obj, 'key_message', '') or ""
