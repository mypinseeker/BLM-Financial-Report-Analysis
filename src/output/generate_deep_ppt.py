"""Generate a comprehensive deep-analysis PPT from the 8 MD files.

Inherits BLMPPTGenerator for utility methods and BLMChartGenerator for
all 16 chart types. Does NOT modify existing files.

Usage:
    python3 -m src.output.generate_deep_ppt
"""

from __future__ import annotations

import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from src.output.ppt_styles import PPTStyle, get_style, VODAFONE_STYLE
from src.output.ppt_charts import BLMChartGenerator
from src.output.ppt_generator import BLMPPTGenerator
from src.output.deep_data_extractor import (
    DeepAnalysisData,
    load_deep_analysis,
)


class DeepAnalysisPPTGenerator(BLMPPTGenerator):
    """Generate a 46-48 slide PPT from the deep analysis MD content.

    Inherits all utility methods from BLMPPTGenerator:
      _new_slide, _add_header, _add_key_message_bar,
      _add_shape, _add_text_box, _add_bullet_list,
      _add_image, _add_metric_cards, _add_section_divider
    """

    def __init__(self, data_dir: str = "data/output", **kwargs):
        kwargs.setdefault('style', VODAFONE_STYLE)
        kwargs.setdefault('operator_id', 'vodafone_germany')
        kwargs.setdefault('output_dir', 'data/output')
        kwargs.setdefault('chart_dpi', 150)
        super().__init__(**kwargs)
        self.data_dir = data_dir

    def generate_deep(self, filename: str = None) -> str:
        """Generate the deep analysis PPT.

        Returns:
            Path to generated .pptx file.
        """
        # Load data
        data = load_deep_analysis(self.data_dir)

        # Init presentation
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.SLIDE_WIDTH)
        self.prs.slide_height = Inches(self.SLIDE_HEIGHT)
        self.slide_num = 0
        self._slide_specs = []

        # Build all slides
        self._build_deep_slides(data)

        # Save
        if filename is None:
            filename = "blm_vodafone_germany_deep_analysis.pptx"
        output_path = self.output_dir / filename
        self.prs.save(str(output_path))
        print(f"Generated: {output_path} ({self.slide_num} slides)")
        return str(output_path)

    # =====================================================================
    # Master slide orchestration
    # =====================================================================

    def _build_deep_slides(self, data: DeepAnalysisData):
        """Build all 46-48 slides."""
        d = data

        # --- 1-3: Cover, TOC, Data Note ---
        self._slide_cover(d)
        self._slide_toc()
        self._slide_data_note(d)

        # --- 4-8: Executive Summary (5 slides) ---
        self._add_section_divider("Executive Summary",
                                  "Vodafone Germany BLM Strategic Assessment")
        self._slide_situation_glance(d)
        self._slide_vf_position(d)
        self._slide_squeezed_middle(d)
        self._slide_exec_priorities(d)

        # --- 9-13: Look 1: Trends / PEST (5 slides) ---
        self._add_section_divider("01 Look at Trends",
                                  "PEST Framework — Macro Environment")
        self._slide_pest_dashboard(d)
        self._slide_pest_factors(d)
        self._slide_industry_landscape(d)
        self._slide_value_migration(d)

        # --- 14-18: Look 2: Market/Customer (5 slides) ---
        self._add_section_divider("02 Look at Market",
                                  "$APPEALS — Market Changes & Customer Needs")
        self._slide_market_snapshot(d)
        self._slide_customer_segments(d)
        self._slide_appeals_radar(d)
        self._slide_market_events(d)

        # --- 19-23: Tariff Analysis (5 slides) ---
        self._add_section_divider("Tariff Analysis",
                                  "Mobile, Fixed & Convergence Pricing Landscape")
        self._slide_tariff_postpaid(d)
        self._slide_tariff_eur_gb(d)
        self._slide_tariff_evolution(d)
        self._slide_tariff_fixed_fmc(d)

        # --- 24-29: Look 3: Competition (6 slides) ---
        self._add_section_divider("03 Look at Competition",
                                  "Porter's Five Forces — Competitive Landscape")
        self._slide_five_forces(d)
        self._slide_competitor_dt(d)
        self._slide_competitor_o2(d)
        self._slide_competitor_1and1(d)
        self._slide_competition_benchmark(d)

        # --- 30-35: Look 4: Self (6 slides) ---
        self._add_section_divider("04 Look at Self",
                                  "BMC + Capability Assessment")
        self._slide_self_health(d)
        self._slide_self_revenue_mix(d)
        self._slide_self_segments(d)
        self._slide_self_network(d)
        self._slide_self_strengths_weaknesses(d)

        # --- 36-40: SWOT Synthesis (5 slides) ---
        self._add_section_divider("SWOT Synthesis",
                                  "Strengths, Weaknesses, Opportunities & Threats")
        self._slide_swot_matrix(d)
        self._slide_swot_strategies(d)
        self._slide_swot_squeezed(d)
        self._slide_swot_priorities(d)

        # --- 41-45: Look 5: Opportunities (5 slides) ---
        self._add_section_divider("05 Look at Opportunities",
                                  "SPAN Matrix — Opportunity Selection")
        self._slide_span_bubble(d)
        self._slide_opp_priority_table(d)
        self._slide_opp_financial(d)
        self._slide_opp_timeline(d)

        # --- 46-48: Summary slides ---
        self._slide_risk_reward(d)
        self._slide_kpi_dashboard(d)
        self._slide_back_cover()

    # =====================================================================
    # Helper: table-to-chart data conversion
    # =====================================================================

    def _table_col(self, rows: list[dict], col: str) -> list[str]:
        """Extract a column from parsed table rows."""
        return [r.get(col, '') for r in rows]

    def _table_numeric_col(self, rows: list[dict], col: str) -> list[float]:
        """Extract numeric column, stripping €, %, etc."""
        import re
        vals = []
        for r in rows:
            raw = r.get(col, '0')
            cleaned = re.sub(r'[€$%,KMBkmb\s+]', '', str(raw))
            try:
                vals.append(float(cleaned))
            except ValueError:
                vals.append(0)
        return vals

    def _rows_to_bullets(self, rows: list[dict], key_col: str,
                         val_col: str = "", max_items: int = 8) -> list[str]:
        """Convert table rows to bullet strings."""
        items = []
        for r in rows[:max_items]:
            text = r.get(key_col, '')
            if val_col and r.get(val_col):
                text = f"{text}: {r[val_col]}"
            if text:
                items.append(text)
        return items

    def _safe_items(self, data_list: list, key: str, max_items: int = 6) -> list[str]:
        """Safely extract text items from a list of dicts."""
        if not data_list:
            return []
        items = []
        for item in data_list[:max_items]:
            if isinstance(item, dict):
                items.append(item.get(key, str(item)))
            else:
                items.append(str(item))
        return items

    # =====================================================================
    # Slides 1-3: Cover, TOC, Data Note
    # =====================================================================

    def _slide_cover(self, d: DeepAnalysisData):
        slide = self._new_slide("cover", "Cover")
        self._add_shape(slide, 0, 0, Inches(self.SLIDE_WIDTH), Inches(2.8),
                        self.style.primary_color)
        self._add_text_box(
            slide, Inches(0.5), Inches(0.5), Inches(11), Inches(1.2),
            "Vodafone Germany — BLM Strategic Assessment",
            font_size=self.style.title_size,
            font_color=(255, 255, 255), bold=True)
        self._add_text_box(
            slide, Inches(0.5), Inches(1.7), Inches(11), Inches(0.6),
            "Deep Analysis | Five Looks + SWOT + SPAN | CQ4 2025",
            font_size=20, font_color=(255, 220, 220))

        info_items = [
            ("Target Operator", "Vodafone Germany"),
            ("Market", "German Telecom (4-Op Oligopoly)"),
            ("Period", "CQ4 2025 / Q3 FY26"),
            ("Report Date", datetime.now().strftime('%Y-%m-%d')),
        ]
        x = Inches(0.5)
        for label, value in info_items:
            self._add_shape(slide, x, Inches(3.5), Inches(2.9), Inches(1.3),
                            (245, 245, 245))
            self._add_text_box(slide, x + Inches(0.15), Inches(3.6),
                               Inches(2.6), Inches(0.3), label,
                               font_size=10, font_color=self.style.light_text_color)
            self._add_text_box(slide, x + Inches(0.15), Inches(3.9),
                               Inches(2.6), Inches(0.7), value,
                               font_size=16, font_color=self.style.primary_color,
                               bold=True)
            x += Inches(3.1)

        framework = (
            "01 Trends (PEST) | 02 Market ($APPEALS) | "
            "Tariff | 03 Competition (Porter) | 04 Self (BMC) | "
            "SWOT | 05 Opportunities (SPAN)"
        )
        self._add_text_box(slide, Inches(0.5), Inches(5.3), Inches(12), Inches(0.6),
                           framework, font_size=12,
                           font_color=self.style.light_text_color)

    def _slide_toc(self):
        slide = self._new_slide("toc", "Table of Contents")
        self._add_header(slide, "Report Contents", "Table of Contents")

        sections = [
            ("ES", "Executive Summary", "Key findings, position, priorities", "5"),
            ("01", "Look at Trends", "PEST macro environment analysis", "5"),
            ("02", "Look at Market", "$APPEALS, segments, events", "5"),
            ("T", "Tariff Analysis", "Mobile, fixed, FMC pricing", "5"),
            ("03", "Look at Competition", "Porter, 3 competitor deep dives", "6"),
            ("04", "Look at Self", "Health, revenue, network, BMC", "6"),
            ("SW", "SWOT Synthesis", "Matrix, strategies, priorities", "5"),
            ("05", "Look at Opportunities", "SPAN, priorities, financials", "5"),
            ("S", "Summary", "Risk/reward, KPI dashboard", "3"),
        ]
        y = Inches(1.4)
        for num, title, desc, slides_str in sections:
            self._add_shape(slide, Inches(0.8), y, Inches(0.5), Inches(0.45),
                            self.style.primary_color, MSO_SHAPE.OVAL)
            self._add_text_box(slide, Inches(0.8), y + Inches(0.07),
                               Inches(0.5), Inches(0.3), num,
                               font_size=10, font_color=(255, 255, 255),
                               bold=True, align="center")
            self._add_text_box(slide, Inches(1.5), y, Inches(4), Inches(0.3),
                               title, font_size=13,
                               font_color=self.style.text_color, bold=True)
            self._add_text_box(slide, Inches(1.5), y + Inches(0.25), Inches(8),
                               Inches(0.25), f"{desc}  ({slides_str} slides)",
                               font_size=10,
                               font_color=self.style.light_text_color)
            y += Inches(0.6)

    def _slide_data_note(self, d: DeepAnalysisData):
        slide = self._new_slide("data_note", "Data Sources")
        self._add_header(slide, "Data Sources & Methodology")

        metrics = [
            ("Data Points", "~2,700 lines", ""),
            ("Tariff Records", "466", ""),
            ("Analysis Documents", "8", ""),
            ("Time Span", "8 Quarters", ""),
        ]
        self._add_metric_cards(slide, metrics, top=1.4)

        sources = [
            "CQ4_2025 Market Snapshot (4 operators, all KPIs)",
            "466 Tariff Records (H1_2023 - H1_2026, 7 snapshots)",
            "8-Quarter Financial Trend (Q1-Q8 by segment)",
            "10 Market Events (ranked by strategic impact)",
            "$APPEALS Assessment (8 dimensions x 4 operators)",
            "Porter Five Forces + SWOT + SPAN Matrix",
            "Management Commentary (Q3 FY26 Earnings Call)",
        ]
        self._add_bullet_list(slide, Inches(0.5), Inches(3.0), Inches(12),
                              Inches(3), sources, font_size=12)
        self._add_key_message_bar(
            slide, "All data traceable to original sources; 8 deep analysis documents totaling ~2,700 lines")

    # =====================================================================
    # Slides 4-8: Executive Summary
    # =====================================================================

    def _slide_situation_glance(self, d: DeepAnalysisData):
        """Slide 5: Situation at a Glance — 8 metric cards + revenue bar chart."""
        slide = self._new_slide("situation", "Situation at a Glance")
        self._add_header(slide, "Situation at a Glance", "Executive Summary")

        # Top row: market metrics
        top_metrics = [
            ("Market", "EUR 12.3B/q", "~EUR 49B ann."),
            ("Growth", "+0.3% YoY", "Near-zero real"),
            ("HHI", "~3,500+", "Highly concentrated"),
            ("Lifecycle", "Mature", "Zero-sum game"),
        ]
        self._add_metric_cards(slide, top_metrics, top=1.3, card_width=2.7,
                               card_height=1.0, gap=0.15)

        # Bottom row: VF metrics
        bot_metrics = [
            ("Revenue", "EUR 3,092M #2", "+0.7% YoY"),
            ("EBITDA", "36.2%", "+0.8pp over 8Q"),
            ("Mobile", "32,500K #3", "+80K/q"),
            ("Churn", "1.05%", "Worst in market"),
        ]
        self._add_metric_cards(slide, bot_metrics, top=2.6, card_width=2.7,
                               card_height=1.0, gap=0.15)

        # Revenue bar chart: 4 operators
        cats = ['DT', 'Vodafone', 'O2', '1&1']
        vals = [6200, 3092, 2000, 1035]
        chart_path = self.chart_gen.create_bar_chart(
            cats, vals, target_category='Vodafone',
            title="Quarterly Revenue (EUR M)",
            y_label="EUR M", filename="deep_revenue_bar.png")
        self._add_image(slide, chart_path, Inches(0.5), Inches(3.9),
                        Inches(5.5), Inches(2.3))

        # Right: key bullets
        bullets = [
            "DT alone > other 3 combined (EUR 6,200M vs EUR 6,127M)",
            "VF growing +0.7% vs DT +1.1% — gap widening",
            "O2 winning subs but losing revenue (-3.4%)",
            "1&1 at 8.4% share, still in MVNO-to-MNO build phase",
        ]
        self._add_bullet_list(slide, Inches(6.3), Inches(4.0), Inches(6.3),
                              Inches(2.0), bullets, font_size=11)

        self._add_key_message_bar(
            slide, "EUR 12.3B market growing at near-zero real — every euro must come from competitors")

    def _slide_vf_position(self, d: DeepAnalysisData):
        """Slide 6: VF Position — headline numbers."""
        slide = self._new_slide("vf_position", "Vodafone Position")
        self._add_header(slide, "Vodafone Germany — Headline Numbers",
                         "Executive Summary")

        headlines = [
            ("Revenue", "EUR 3,092M/q (+0.7% YoY)", "Growing, barely"),
            ("EBITDA", "EUR 1,120M/q (36.2% margin)", "Healthy"),
            ("Mobile", "+80K/q net adds", "Weakest in market (DT +340K)"),
            ("Broadband", "-63K/q net loss", "ONLY operator losing BB subs"),
            ("B2B", "EUR 520M/q (+8.5% YoY)", "Strongest growth engine"),
            ("Wholesale", "EUR 380M/q (+90% YoY)", "Temporary 1&1 windfall"),
            ("FMC", "4,900K subs (49.3% pen.)", "Key retention metric"),
        ]
        y = Inches(1.4)
        for label, value, note in headlines:
            self._add_shape(slide, Inches(0.5), y, Inches(2.0), Inches(0.5),
                            (245, 245, 245))
            self._add_text_box(slide, Inches(0.6), y + Inches(0.08),
                               Inches(1.8), Inches(0.35), label,
                               font_size=12, font_color=self.style.primary_color,
                               bold=True)
            self._add_text_box(slide, Inches(2.7), y + Inches(0.08),
                               Inches(5), Inches(0.35), value,
                               font_size=12, font_color=self.style.text_color,
                               bold=True)
            self._add_text_box(slide, Inches(8.0), y + Inches(0.08),
                               Inches(4.8), Inches(0.35), note,
                               font_size=11,
                               font_color=self.style.light_text_color)
            y += Inches(0.58)

        # Donut gauges for key %
        labels = ['EBITDA Margin', 'FMC Pen.', '5G Coverage']
        values = [36.2, 49.3, 92.0]
        chart_path = self.chart_gen.create_donut_gauges(
            labels, values, title="Key Ratios",
            filename="deep_donut_kpis.png")
        self._add_image(slide, chart_path, Inches(0.5), Inches(5.0),
                        Inches(12), Inches(1.2))

        self._add_key_message_bar(
            slide, "Stable #2 but growth engines (B2B, wholesale) mask flat consumer and declining broadband")

    def _slide_squeezed_middle(self, d: DeepAnalysisData):
        """Slide 7: Squeezed Middle diagnosis."""
        slide = self._new_slide("squeezed_middle", "Squeezed Middle Diagnosis")
        self._add_header(slide, 'The "Squeezed Middle" — Central Diagnosis',
                         "Executive Summary")

        # KPI comparison table: DT / VF / O2-1&1
        cols = ['DT', 'Vodafone', 'O2/1&1']
        rows_data = {
            'Mobile pricing': ['EUR 40-85 (never discounts)', 'EUR 25-55 (reactive cuts)', 'EUR 15-60 (aggressive)'],
            'Network quality': ['97% 5G, SA ready, 8.5M fiber', '92% 5G, no SA, 1.5M fiber', '98% 5G (O2), 25% (1&1)'],
            'Brand perception': ['Premium, trusted (90/100)', 'Strong not premium (82)', 'Functional (65-75)'],
            'EUR/GB value': ['EUR 0.60-2.00 (brand justified)', 'EUR 0.80-3.57 (not justified)', 'EUR 0.17-0.50 (best)'],
            'Customer service': ['Above average (82)', 'Below average (70)', 'Average (65-72)'],
            'Churn': ['0.80% (lowest)', '1.05% (highest)', '0.90-0.95%'],
            'BB sub trend': ['+80K/quarter', '-63K/quarter', '+8K to -5K'],
        }
        metrics = list(rows_data.keys())
        # Transpose: rows_data[metric][i] -> operator_data[operator][metric_idx]
        operator_data = {c: [rows_data[m][i] for m in metrics] for i, c in enumerate(cols)}
        chart_path = self.chart_gen.create_kpi_table_chart(
            metrics, operator_data,
            target_operator='Vodafone',
            title="Squeezed Middle — 7 Dimension Comparison",
            filename="deep_squeezed_table.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(12.5), Inches(4.0))

        # Escape routes
        escapes = [
            "Convergence moat: Cable+mobile FMC unique — neither O2 nor 1&1 can match",
            "B2B pivot: Enterprise less price-sensitive; Skaylink adds cloud/security",
            "Infrastructure: FibreCo JV + DOCSIS 4.0 to defend fixed broadband",
        ]
        self._add_text_box(slide, Inches(0.5), Inches(5.3), Inches(3), Inches(0.3),
                           "Escape Routes:", font_size=12,
                           font_color=self.style.primary_color, bold=True)
        self._add_bullet_list(slide, Inches(3.3), Inches(5.3), Inches(9.5),
                              Inches(0.8), escapes, font_size=10)

        self._add_key_message_bar(
            slide, "Customers wanting premium choose DT; wanting value choose O2/1&1 — VF must differentiate via FMC, B2B, infra")

    def _slide_exec_priorities(self, d: DeepAnalysisData):
        """Slide 8: 5 strategic priorities."""
        slide = self._new_slide("exec_priorities", "Strategic Priorities")
        self._add_header(slide, "5 Strategic Priorities — Consolidated",
                         "Executive Summary")

        priorities = [
            ("P1", "Defend Fixed Broadband", "EXISTENTIAL", "FibreCo JV (7M FTTH) + DOCSIS 4.0"),
            ("P2", "Maximize FMC Penetration", "STRATEGIC MOAT", "49% to 60% — cable+mobile unique"),
            ("P3", "Fix Customer Service", "HIGHEST ROI", "Score 70 to 74+ — EUR 200M/yr retained"),
            ("P4", "Accelerate B2B via Skaylink", "GROWTH ENGINE", "EUR 520M to EUR 800M/q by 5yr"),
            ("P5", "Manage Wholesale Transition", "DEFEND REVENUE", "EUR 380M/q — 3-5 year window"),
        ]

        items = [f"{p}: {desc} [{tag}] — {detail}"
                 for p, desc, tag, detail in priorities]
        prios = [p for p, _, _, _ in priorities]

        chart_path = self.chart_gen.create_priority_chart(
            items, prios, title="Strategic Priorities",
            filename="deep_priorities.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(12.5), Inches(4.8))

        self._add_key_message_bar(
            slide, "Five priorities across all analyses — P1 (BB defense) is existential, P3 (service) is highest ROI")

    # =====================================================================
    # Slides 9-13: Look 1 — Trends (PEST)
    # =====================================================================

    def _slide_pest_dashboard(self, d: DeepAnalysisData):
        """PEST 4-quadrant dashboard."""
        slide = self._new_slide("pest_dash", "PEST Dashboard")
        self._add_header(slide, "PEST Analysis Dashboard", "Macro Environment")

        # Build PEST data for chart
        pest_obj = {
            'political_factors': [],
            'economic_factors': [],
            'society_factors': [],
            'technology_factors': [],
        }
        factor_map = {
            'Political': ('political_factors', [
                ('BNetzA Regulatory Environment', 'both', 'high'),
                ('Gigabit Funding EUR 1.2B', 'opportunity', 'high'),
                ('Nebenkostenprivileg Abolition', 'threat', 'high'),
                ('Spectrum Extension (5 years)', 'opportunity', 'high'),
                ('MVNO Access Requirements', 'threat', 'medium'),
            ]),
            'Economic': ('economic_factors', [
                ('GDP Growth 0.8%', 'threat', 'medium'),
                ('Inflation 2.1%', 'neutral', 'low'),
            ]),
            'Society': ('society_factors', [
                ('5G Adoption at 40%', 'opportunity', 'high'),
                ('Fiber Penetration at 18%', 'both', 'high'),
            ]),
            'Technology': ('technology_factors', [
                ('5G SA Monetization Phase', 'opportunity', 'high'),
                ('Fiber/FTTH Deployment Race', 'both', 'high'),
                ('Cable Frequency Restructuring', 'opportunity', 'medium'),
                ('VF-Altice FibreCo JV', 'opportunity', 'high'),
            ]),
        }
        for dim, (key, factors) in factor_map.items():
            for name, impact, severity in factors:
                pest_obj[key].append({
                    'factor_name': name,
                    'impact_type': impact,
                    'severity': severity,
                })

        chart_path = self.chart_gen.create_pest_dashboard(
            pest_obj, title="PEST Analysis Dashboard",
            filename="deep_pest_dashboard.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(12.5), Inches(4.8))

        self._add_key_message_bar(
            slide, "Net favorable (9 opps vs 3 threats) — but fiber displacement and 5G SA gap are existential risks")

    def _slide_pest_factors(self, d: DeepAnalysisData):
        """PEST factor ranking horizontal bar."""
        slide = self._new_slide("pest_factors", "PEST Factor Ranking")
        self._add_header(slide, "PEST Factor Priority Ranking — VF-Specific",
                         "Macro Impact")

        categories = [
            'Fiber penetration growth',
            'Nebenkostenprivileg abolition',
            '5G SA readiness gap',
            'Gigabit Funding EUR 1.2B',
            'Spectrum extension (5yr)',
            'Low GDP growth (0.8%)',
            '5G adoption at 40%',
            'DOCSIS 4.0 readiness',
        ]
        severity_scores = [10, 9, 9, 8, 8, 5, 7, 8]

        chart_path = self.chart_gen.create_horizontal_bar_chart(
            categories, severity_scores,
            target_category='Fiber penetration growth',
            title="PEST Factor Severity (Vodafone-Specific Impact)",
            x_label="Impact Score",
            filename="deep_pest_factors.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(12.5), Inches(4.8))

        self._add_key_message_bar(
            slide, "Fiber displacement is #1 threat; Gigabit Fund and spectrum extension are #1 opportunities")

    def _slide_industry_landscape(self, d: DeepAnalysisData):
        """Industry landscape + KSF."""
        slide = self._new_slide("industry", "Industry Landscape")
        self._add_header(slide, "Industry Landscape & Key Success Factors",
                         "Industry Analysis")

        metrics = [
            ("Market Size", "EUR 49B/yr", "4th globally"),
            ("Penetration", "Mobile 170% / BB 37%", ""),
            ("CR4", "100%", "Pure oligopoly"),
            ("Avg Margin", "~36.6%", "Healthy"),
        ]
        self._add_metric_cards(slide, metrics, top=1.3)

        # KSF radar
        dims = ['Network Quality', 'FMC Bundling', 'B2B/ICT', 'Op. Efficiency', 'Brand/CX']
        scores = {
            'Vodafone': [60, 80, 70, 60, 50],
            'DT': [90, 75, 85, 80, 85],
        }
        chart_path = self.chart_gen.create_radar_chart(
            dims, scores, target_operator='Vodafone',
            title="Key Success Factor Scores",
            filename="deep_ksf_radar.png")
        self._add_image(slide, chart_path, Inches(0.5), Inches(2.9),
                        Inches(5.5), Inches(3.3))

        # Right: lifecycle bullets
        items = [
            "Volume growth exhausted — mobile pen. >130%",
            "Competition shifts to retention and upsell",
            "ARPU defense > subscriber growth",
            "FMC and B2B are remaining growth pockets",
            "Regulatory and tech shifts are primary disruptors",
        ]
        self._add_text_box(slide, Inches(6.5), Inches(2.9), Inches(6), Inches(0.3),
                           "Mature Market Implications:", font_size=12,
                           font_color=self.style.primary_color, bold=True)
        self._add_bullet_list(slide, Inches(6.5), Inches(3.3), Inches(6),
                              Inches(2.8), items, font_size=11)

        self._add_key_message_bar(
            slide, "Mature market — growth only from stealing share; FMC and B2B are the growth pockets")

    def _slide_value_migration(self, d: DeepAnalysisData):
        """Value migration map."""
        slide = self._new_slide("value_migration", "Value Migration")
        self._add_header(slide, "Value Transfer & New Business Models",
                         "Value Migration")

        transfers = [
            ("Voice/SMS -> Data/OTT", "Complete"),
            ("Linear TV -> Streaming", "Accelerating"),
            ("Consumer -> B2B Enterprise", "Growing"),
            ("DSL/Copper -> Fiber", "Accelerating"),
            ("Cable BB -> Fiber", "Beginning — VF main risk"),
            ("Traditional capex -> NaaS", "Emerging"),
            ("Mobile-only -> FMC", "Mature — proven retention"),
        ]

        y = Inches(1.5)
        for transfer, status in transfers:
            self._add_shape(slide, Inches(0.5), y, Inches(6.5), Inches(0.5),
                            (248, 248, 248))
            self._add_text_box(slide, Inches(0.7), y + Inches(0.08),
                               Inches(4.5), Inches(0.35), transfer,
                               font_size=12, font_color=self.style.text_color)
            is_risk = 'risk' in status.lower() or 'accelerating' in status.lower()
            color = self.style.negative_color if is_risk else self.style.positive_color
            self._add_text_box(slide, Inches(5.3), y + Inches(0.08),
                               Inches(1.7), Inches(0.35), status,
                               font_size=10, font_color=color, bold=True)
            y += Inches(0.55)

        # Right: emerging models
        models = [
            "FWA (5G as BB substitute) — moderate for VF",
            "Network-as-a-Service (slicing) — high, needs 5G SA",
            "Wholesale/MVNO — EUR 380M/q, structural decline",
            "IoT platform — 3M connections, B2B synergy",
        ]
        self._add_text_box(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.3),
                           "Emerging Models:", font_size=12,
                           font_color=self.style.primary_color, bold=True)
        self._add_bullet_list(slide, Inches(7.5), Inches(1.9), Inches(5),
                              Inches(3), models, font_size=11)

        self._add_key_message_bar(
            slide, "Cable-to-fiber migration is VF's #1 value risk; FMC convergence is proven value creation")

    # =====================================================================
    # Slides 14-18: Look 2 — Market/Customer
    # =====================================================================

    def _slide_market_snapshot(self, d: DeepAnalysisData):
        slide = self._new_slide("mkt_snapshot", "Market Snapshot")
        self._add_header(slide, "Market Snapshot — CQ4 2025", "Market Overview")

        metrics = [
            ("Revenue", "EUR 12,327M/q", "~EUR 49B/yr"),
            ("Mobile Subs", "143,100K", "170% pen."),
            ("BB Subs", "31,420K", "37% pen."),
            ("Operators", "4", "CR4=100%"),
        ]
        self._add_metric_cards(slide, metrics, top=1.3)

        # Market share bar
        cats = ['DT 50.3%', 'VF 25.1%', 'O2 16.2%', '1&1 8.4%']
        vals = [50.3, 25.1, 16.2, 8.4]
        chart_path = self.chart_gen.create_horizontal_bar_chart(
            cats, vals, target_category='VF 25.1%',
            title="Revenue Market Share (%)",
            x_label="%", filename="deep_mkt_share.png",
            value_suffix="%")
        self._add_image(slide, chart_path, Inches(0.3), Inches(2.9),
                        Inches(6), Inches(3.2))

        # Right column: monetization insight
        insights = [
            "DT monetizes best: 36.5% mobile subs -> 50.3% revenue",
            "O2 monetizes worst: 32.1% subs -> 16.2% revenue",
            "VF in between: revenue share (25.1%) > mobile sub share (22.7%)",
            "Revenue share tells the real competitive story",
        ]
        self._add_text_box(slide, Inches(6.5), Inches(2.9), Inches(6), Inches(0.3),
                           "Monetization Insight:", font_size=12,
                           font_color=self.style.primary_color, bold=True)
        self._add_bullet_list(slide, Inches(6.5), Inches(3.3), Inches(6),
                              Inches(2.5), insights, font_size=11)

        self._add_key_message_bar(
            slide, "EUR 12.3B market — DT captures more revenue than the other 3 combined")

    def _slide_customer_segments(self, d: DeepAnalysisData):
        slide = self._new_slide("segments", "Customer Segments")
        self._add_header(slide, "7 Customer Segments", "Segmentation")

        segments = [
            ("Consumer High-End", "~13.9M", "Growing", "P1"),
            ("Consumer Mainstream", "~51.2M", "Stable", "P1"),
            ("Consumer Price-Sensitive", "~35.1M", "Shrinking", "P3"),
            ("Consumer Youth", "~17.2M", "Growing", "P2"),
            ("Enterprise Large", "~47K", "Growing", "P1"),
            ("Enterprise SME", "~896K", "Stable", "P2"),
            ("Wholesale/MVNO", "~17.2M", "Stable", "P2"),
        ]

        y = Inches(1.4)
        for name, size, growth, priority in segments:
            self._add_shape(slide, Inches(0.5), y, Inches(12), Inches(0.55),
                            (248, 248, 248))
            # Priority badge
            p_colors = {'P1': self.style.primary_color,
                        'P2': self.style.warning_color,
                        'P3': self.style.neutral_color}
            self._add_shape(slide, Inches(0.6), y + Inches(0.08),
                            Inches(0.45), Inches(0.35),
                            p_colors.get(priority, self.style.neutral_color))
            self._add_text_box(slide, Inches(0.6), y + Inches(0.1),
                               Inches(0.45), Inches(0.3), priority,
                               font_size=9, font_color=(255, 255, 255),
                               bold=True, align="center")
            self._add_text_box(slide, Inches(1.2), y + Inches(0.1),
                               Inches(3.5), Inches(0.3), name,
                               font_size=12, font_color=self.style.text_color,
                               bold=True)
            self._add_text_box(slide, Inches(5.0), y + Inches(0.1),
                               Inches(1.5), Inches(0.3), size,
                               font_size=11, font_color=self.style.text_color)
            g_color = self.style.positive_color if growth == 'Growing' else (
                self.style.negative_color if growth == 'Shrinking' else self.style.light_text_color)
            self._add_text_box(slide, Inches(6.8), y + Inches(0.1),
                               Inches(1.5), Inches(0.3), growth,
                               font_size=11, font_color=g_color)
            y += Inches(0.6)

        self._add_key_message_bar(
            slide, "Focus on High-End + Mainstream + Enterprise; deprioritize price-sensitive (35M subs)")

    def _slide_appeals_radar(self, d: DeepAnalysisData):
        """$APPEALS radar chart — 4 operators, 8 dimensions."""
        slide = self._new_slide("appeals_radar", "$APPEALS Assessment")
        self._add_header(slide, "$APPEALS Assessment — Customer Needs",
                         "Customer Perception")

        dims = ['Price', 'Availability', 'Packaging', 'Performance',
                'Ease of Use', 'Assurances', 'Lifecycle Cost', 'Social/Brand']
        scores = {
            'Vodafone': [3.6, 4.0, 3.8, 3.8, 3.6, 4.1, 3.6, 4.0],
            'DT': [3.2, 4.8, 4.4, 4.7, 4.2, 4.6, 3.2, 4.3],
            'O2': [4.2, 4.2, 3.5, 4.2, 3.5, 3.4, 4.2, 3.8],
            '1&1': [4.5, 2.2, 3.6, 2.4, 3.5, 2.8, 4.5, 3.6],
        }
        # Scale to 0-100 for radar
        scaled = {op: [s * 20 for s in sc] for op, sc in scores.items()}
        chart_path = self.chart_gen.create_radar_chart(
            dims, scaled, target_operator='Vodafone',
            title="$APPEALS Scores (scaled 0-100)",
            filename="deep_appeals_radar.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.1),
                        Inches(7.5), Inches(5.0))

        # Right: gap text
        gaps = [
            "Performance: VF 3.8 vs DT 4.7 (gap -0.9pt) — CRITICAL",
            "Price: VF 3.6 vs 1&1 4.5 (gap -0.9pt) — CRITICAL",
            "Availability: VF 4.0 vs DT 4.8 (gap -0.8pt) — CRITICAL",
            "VF beats DT on: NOTHING",
            "VF best score: Assurances (4.1) — trust/reliability",
            "Triple weakness: Price, Ease of Use, Lifecycle Cost (all 3.6)",
        ]
        self._add_text_box(slide, Inches(8.0), Inches(1.3), Inches(5), Inches(0.3),
                           "Key Gaps:", font_size=12,
                           font_color=self.style.primary_color, bold=True)
        self._add_bullet_list(slide, Inches(8.0), Inches(1.7), Inches(5),
                              Inches(4.5), gaps, font_size=10)

        self._add_key_message_bar(
            slide, "VF beats DT on nothing — squeezed middle confirmed in customer perception")

    def _slide_market_events(self, d: DeepAnalysisData):
        slide = self._new_slide("mkt_events", "Market Events")
        self._add_header(slide, "Top Market Events — CQ4 2025",
                         "Competitive Intelligence")

        events = [
            ("1", "DT reboots fiber, targets VF cable footprint", "High"),
            ("2", "VF completes 1&1 migration: 12M users on network", "High"),
            ("3", "1&1 completes OpenRAN, reaches 25% coverage", "High"),
            ("4", "VF-Altice FibreCo JV: 7M FTTH homes", "High"),
            ("5", "VF acquires Skaylink for EUR 175M", "High"),
            ("6", "BNetzA extends spectrum by 5 years", "High"),
            ("7", "DT Q3: revenue +1.5%, profit +14.3%", "Medium"),
            ("8", "VF restructures cable TV frequencies", "Medium"),
        ]

        y = Inches(1.4)
        for num, event, severity in events:
            self._add_shape(slide, Inches(0.5), y, Inches(0.4), Inches(0.4),
                            self.style.primary_color, MSO_SHAPE.OVAL)
            self._add_text_box(slide, Inches(0.5), y + Inches(0.05),
                               Inches(0.4), Inches(0.3), num,
                               font_size=10, font_color=(255, 255, 255),
                               bold=True, align="center")
            self._add_text_box(slide, Inches(1.1), y + Inches(0.05),
                               Inches(9.5), Inches(0.3), event,
                               font_size=11, font_color=self.style.text_color)
            sev_color = self.style.negative_color if severity == 'High' else self.style.warning_color
            self._add_text_box(slide, Inches(11.0), y + Inches(0.05),
                               Inches(1.5), Inches(0.3), severity,
                               font_size=10, font_color=sev_color, bold=True)
            y += Inches(0.5)

        # Three battles summary
        self._add_text_box(slide, Inches(0.5), Inches(5.6), Inches(12), Inches(0.3),
                           "Three Battles: 1) Fixed BB Infrastructure  |  2) Wholesale Economics  |  3) B2B Transformation",
                           font_size=12, font_color=self.style.primary_color, bold=True)

        self._add_key_message_bar(
            slide, "10 market events — DT fiber overbuild (#1) and 1&1 network build (#3) are the strategic threats")

    # =====================================================================
    # Slides 19-23: Tariff Analysis
    # =====================================================================

    def _slide_tariff_postpaid(self, d: DeepAnalysisData):
        """Mobile postpaid comparison — grouped bar."""
        slide = self._new_slide("tariff_postpaid", "Mobile Postpaid Comparison")
        self._add_header(slide, "Mobile Postpaid Price Comparison — H1 2026",
                         "Cross-Operator Benchmarking")

        tiers = ['S', 'M', 'L', 'XL']
        prices = {
            'DT': [40, 50, 60, 85],
            'Vodafone': [25, 33, 40, 55],
            'O2': [15, 20, 25, 60],
            '1&1': [15, 20, 30, 40],
        }
        chart_path = self.chart_gen.create_segment_comparison(
            tiers, prices, target_operator='Vodafone',
            title="Monthly Price by Tier (EUR)",
            y_label="EUR/month",
            filename="deep_tariff_postpaid.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(8.5), Inches(4.8))

        # Right: premium gaps
        gaps = [
            "Tier S: DT premium 167% vs O2/1&1",
            "Tier M: DT premium 150%",
            "Tier L: DT premium 140%",
            "Tier XL: DT premium 112%",
            "",
            "VF sits 65-75% above O2/1&1",
            "VF sits 20-35% below DT",
            "= Classic 'squeezed middle'",
        ]
        self._add_text_box(slide, Inches(9.0), Inches(1.3), Inches(4), Inches(0.3),
                           "Price Gaps:", font_size=12,
                           font_color=self.style.primary_color, bold=True)
        self._add_bullet_list(slide, Inches(9.0), Inches(1.7), Inches(4),
                              Inches(4.5), [g for g in gaps if g], font_size=10)

        self._add_key_message_bar(
            slide, "DT commands 112-167% premium; VF caught in the middle — too expensive for value, not premium enough for DT")

    def _slide_tariff_eur_gb(self, d: DeepAnalysisData):
        """EUR/GB ranking."""
        slide = self._new_slide("tariff_eur_gb", "EUR/GB Value Analysis")
        self._add_header(slide, "EUR/GB Value Analysis — Mobile Postpaid H1 2026",
                         "Pricing Intelligence")

        labels = [
            'O2 Mobile L EUR 0.17', '1&1 L EUR 0.25',
            'O2 Mobile M EUR 0.33', '1&1 M EUR 0.33',
            '1&1 S EUR 0.50', 'DT L EUR 0.60',
            'VF GigaMobil L EUR 0.80', 'DT M EUR 1.25',
            'VF GigaMobil M EUR 1.32', 'O2 S EUR 1.50',
            'DT S EUR 2.00', 'VF GigaMobil S EUR 3.57',
        ]
        vals = [0.17, 0.25, 0.33, 0.33, 0.50, 0.60, 0.80, 1.25, 1.32, 1.50, 2.00, 3.57]

        chart_path = self.chart_gen.create_horizontal_bar_chart(
            labels, vals,
            target_category='VF GigaMobil S EUR 3.57',
            title="EUR/GB Ranking (lower = better value)",
            x_label="EUR/GB",
            filename="deep_eur_gb.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(8.5), Inches(4.8))

        insights = [
            "VF GigaMobil S at EUR 3.57/GB is worst in market",
            "That's 21x worse than O2 Mobile L (EUR 0.17/GB)",
            "O2 best value: 1/5 of VF at L tier",
            "1&1 consistently 2nd-best value",
            "",
            "Options: increase data to 15-20GB at same price",
            "or cut price to EUR 20 to close O2/1&1 gap",
        ]
        self._add_bullet_list(slide, Inches(9.0), Inches(1.4), Inches(4),
                              Inches(4.5), [i for i in insights if i], font_size=10)

        self._add_key_message_bar(
            slide, "VF GigaMobil S worst value in market (EUR 3.57/GB) — 21x worse than O2; repricing or data rebalance needed")

    def _slide_tariff_evolution(self, d: DeepAnalysisData):
        """Price evolution trend."""
        slide = self._new_slide("tariff_evolution", "Price Evolution")
        self._add_header(slide, "Mobile Postpaid Price Evolution — H1 2023 to H1 2026",
                         "3-Year Pricing Trends")

        # Trend line: Tier M price over time
        x_labels = ['H1_23', 'H2_23', 'H1_24', 'H2_24', 'H1_25', 'H2_25', 'H1_26']
        series = {
            'DT': [50, 50, 50, 50, 50, 50, 50],
            'Vodafone': [40, 40, 40, 33, 33, 33, 33],
            'O2': [25, 23, 23, 20, 20, 20, 20],
            '1&1': [15, 15, 20, 20, 20, 20, 20],
        }
        chart_path = self.chart_gen.create_multi_line_trend(
            x_labels, series,
            title="Tier M Price Evolution (EUR/month)",
            y_label="EUR/month",
            target_operator='Vodafone',
            filename="deep_tariff_evolution.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(8), Inches(4.8))

        # Right: strategy summary
        items = [
            "DT: Never discounts — doubled data instead",
            "VF: One major reset H2_2024 (-17-20%), then frozen",
            "O2: Most aggressive — two disruptions in 18mo",
            "1&1: MVNO-to-MNO price increase justified by own network",
            "",
            "5G premium eroded from +123% to 0% in 18 months",
            "Differentiation must come from FMC, SA, B2B",
        ]
        self._add_bullet_list(slide, Inches(8.5), Inches(1.4), Inches(4.3),
                              Inches(4.5), [i for i in items if i], font_size=10)

        self._add_key_message_bar(
            slide, "DT never discounts; O2 most aggressive; VF reactive — 5G premium fully eroded to 0%")

    def _slide_tariff_fixed_fmc(self, d: DeepAnalysisData):
        """Fixed & FMC comparison."""
        slide = self._new_slide("tariff_fixed", "Fixed & FMC Pricing")
        self._add_header(slide, "Fixed Broadband & FMC Bundle Pricing",
                         "Cable, DSL, Fiber, FMC")

        # Cable: VF's crown jewel
        cable_data = {
            'VF Cable 250': ['EUR 35', '250 Mbps'],
            'VF Cable 500': ['EUR 40', '500 Mbps'],
            'VF Cable 1000': ['EUR 50', '1000 Mbps'],
            'DT Cable 250': ['EUR 50', '250 Mbps'],
        }
        chart_path = self.chart_gen.create_kpi_table_chart(
            ['Price', 'Speed'], cable_data,
            target_operator='VF Cable 1000',
            title="Cable Broadband — VF vs DT",
            filename="deep_cable_compare.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(6), Inches(2.5))

        # FMC bundles
        fmc_data = {
            'O2': ['EUR 35', 'EUR 60'],
            '1&1': ['EUR 35', 'EUR 60'],
            'Vodafone': ['EUR 50', 'EUR 80'],
            'DT': ['EUR 65', 'EUR 120'],
        }
        chart_path2 = self.chart_gen.create_kpi_table_chart(
            ['Basic FMC', 'Premium FMC'], fmc_data,
            target_operator='Vodafone',
            title="FMC Bundle Comparison",
            filename="deep_fmc_compare.png")
        self._add_image(slide, chart_path2, Inches(0.3), Inches(3.9),
                        Inches(6), Inches(2.3))

        # Right: insights
        insights = [
            "Cable is VF's crown jewel:",
            "  1000 Mbps at EUR 50 vs DT 250 Mbps at EUR 50",
            "  = 4x speed at same price",
            "",
            "FMC narrows the price gap:",
            "  VF Basic EUR 50 is 43% above O2/1&1",
            "  (vs 65-167% gap in mobile alone)",
            "",
            "FMC should be the lead product —",
            "  cable speed advantage = sustainable moat",
        ]
        self._add_bullet_list(slide, Inches(6.5), Inches(1.4), Inches(6),
                              Inches(4.5), [i for i in insights if i], font_size=11)

        self._add_key_message_bar(
            slide, "Cable is VF's only clear pricing win (4x speed at same price); FMC narrows the competitive gap")

    # =====================================================================
    # Slides 24-29: Look 3 — Competition
    # =====================================================================

    def _slide_five_forces(self, d: DeepAnalysisData):
        slide = self._new_slide("five_forces", "Porter's Five Forces")
        self._add_header(slide, "Porter's Five Forces", "Competitive Forces")

        forces = {
            'existing_competitors': 'high',
            'new_entrants': 'medium',
            'substitutes': 'high',
            'supplier_power': 'medium',
            'buyer_power': 'medium',
        }
        chart_path = self.chart_gen.create_porter_five_forces(
            forces, filename="deep_porter.png")
        self._add_image(slide, chart_path, Inches(2.5), Inches(1.1),
                        Inches(7), Inches(5.2))

        # Left annotations
        self._add_text_box(slide, Inches(0.3), Inches(1.5), Inches(2.2), Inches(0.3),
                           "Overall: MEDIUM", font_size=13,
                           font_color=self.style.primary_color, bold=True)
        notes = [
            "Price competition fierce at value end",
            "DT insulated in premium bubble",
            "1&1 OpenRAN is structural change",
            "OTT substitution 90% complete",
            "VF has highest churn (1.05%)",
        ]
        self._add_bullet_list(slide, Inches(0.3), Inches(1.9), Inches(2.3),
                              Inches(4), notes, font_size=9)

        self._add_key_message_bar(
            slide, "Medium overall intensity — but masks fierce price competition at value end while DT operates in premium bubble")

    def _slide_competitor_dt(self, d: DeepAnalysisData):
        """DT competitor card."""
        slide = self._new_slide("comp_dt", "Competitor: Deutsche Telekom")
        self._add_header(slide, "Deutsche Telekom — The Unassailable Incumbent",
                         "Competitor Deep Dive")

        metrics = [
            ("Revenue", "EUR 6,200M/q", "2.0x VF"),
            ("EBITDA Margin", "42.1%", "+5.9pp vs VF"),
            ("Capex/Rev", "19.4%", "Most efficient"),
            ("Mobile Subs", "52,200K", "1.6x VF"),
        ]
        self._add_metric_cards(slide, metrics, top=1.3, card_width=2.7, card_height=1.0)

        # Strategy and threat
        strategy = [
            "Never discount — doubles data instead",
            "Fiber-first: 8,500K homepass (4x VF)",
            "Commands 112-167% premium across all tiers",
            "Only operator growing revenue + margin + subs",
            "Actively targeting VF cable footprint with fiber",
        ]
        self._add_text_box(slide, Inches(0.5), Inches(2.6), Inches(5.5), Inches(0.3),
                           'Strategy: "Premium fortress with fiber moat"',
                           font_size=12, font_color=self.style.primary_color, bold=True)
        self._add_bullet_list(slide, Inches(0.5), Inches(3.0), Inches(5.5),
                              Inches(2.5), strategy, font_size=11)

        # Network gap chart
        dims_net = ['5G Coverage', 'Fiber Homepass', 'Mobile Subs', 'BB Subs', 'EBITDA Margin']
        dt_norm = [100, 100, 100, 100, 100]
        vf_norm = [94.8, 17.6, 62.3, 65.3, 86.0]
        chart_path = self.chart_gen.create_gap_analysis_chart(
            dims_net, vf_norm, dt_norm,
            target_name='Vodafone', leader_name='DT',
            title="VF vs DT Gap (DT=100)",
            filename="deep_dt_gap.png")
        self._add_image(slide, chart_path, Inches(6.3), Inches(2.5),
                        Inches(6.5), Inches(3.5))

        self._add_key_message_bar(
            slide, "DT is 2x VF on revenue with 6pp higher margins — gap widening; fiber overbuild is existential threat")

    def _slide_competitor_o2(self, d: DeepAnalysisData):
        slide = self._new_slide("comp_o2", "Competitor: O2")
        self._add_header(slide, "Telefonica O2 — The Value Disruptor",
                         "Competitor Deep Dive")

        metrics = [
            ("Revenue", "EUR 2,000M/q", "0.65x VF"),
            ("EBITDA Margin", "32.5%", "-3.7pp vs VF"),
            ("Revenue Growth", "-3.4%", "Declining!"),
            ("Mobile Subs", "45,900K", "1.4x VF"),
        ]
        self._add_metric_cards(slide, metrics, top=1.3, card_width=2.7, card_height=1.0)

        strategy = [
            "Most aggressive pricing: two major restructures in 18mo",
            "5G democratization — first to make 5G standard on all plans",
            '"Grow" innovation: auto-increasing data at same price',
            "98% 5G but only 2,400K BB subs — no convergence play",
            "Revenue paradox: winning subs (+170K/q) but losing revenue",
        ]
        self._add_text_box(slide, Inches(0.5), Inches(2.6), Inches(6), Inches(0.3),
                           'Strategy: "More for less — volume at any cost"',
                           font_size=12, font_color=self.style.primary_color, bold=True)
        self._add_bullet_list(slide, Inches(0.5), Inches(3.0), Inches(6),
                              Inches(2.5), strategy, font_size=11)

        threat_opp = [
            "THREAT: O2 EUR 20/60GB vs VF EUR 33/25GB — stark value gap",
            "O2 'Grow' feature is a retention innovation VF lacks",
            "OPPORTUNITY: O2 has no fixed BB (2,400K vs VF 9,940K)",
            "FMC bundling is VF's counter to O2 price aggression",
        ]
        self._add_text_box(slide, Inches(6.8), Inches(2.6), Inches(5.8), Inches(0.3),
                           "Threat & Opportunity:", font_size=12,
                           font_color=self.style.primary_color, bold=True)
        self._add_bullet_list(slide, Inches(6.8), Inches(3.0), Inches(5.8),
                              Inches(3), threat_opp, font_size=10)

        self._add_key_message_bar(
            slide, "O2 winning subs but destroying value (-3.4% revenue); FMC is VF's counter — O2 has no fixed broadband")

    def _slide_competitor_1and1(self, d: DeepAnalysisData):
        slide = self._new_slide("comp_1and1", "Competitor: 1&1")
        self._add_header(slide, "1&1 AG — The Insurgent Network Builder",
                         "Competitor Deep Dive")

        metrics = [
            ("Revenue", "EUR 1,035M/q", "0.33x VF"),
            ("EBITDA Margin", "12.2%", "-24pp vs VF"),
            ("Capex/Rev", "38.6%", "Building network"),
            ("5G Coverage", "55%", "vs VF 92%"),
        ]
        self._add_metric_cards(slide, metrics, top=1.3, card_width=2.7, card_height=1.0)

        strategy = [
            "MVNO to MNO transformation — Europe's first OpenRAN 5G",
            "Successfully raised prices 14-50% at network launch",
            "Cheapest unlimited in Germany: EUR 40 (VF EUR 55, DT EUR 85)",
            "12.2% EBITDA reflects 'invest now, monetize later'",
            "Coverage gap (55%) is critical weakness",
        ]
        self._add_text_box(slide, Inches(0.5), Inches(2.6), Inches(6), Inches(0.3),
                           'Strategy: "From MVNO underdog to OpenRAN pioneer"',
                           font_size=12, font_color=self.style.primary_color, bold=True)
        self._add_bullet_list(slide, Inches(0.5), Inches(3.0), Inches(6),
                              Inches(2.5), strategy, font_size=11)

        threat_opp = [
            "THREAT (short-term): competes on price at bottom tier",
            "THREAT (long-term): 12M hosted subs will migrate off VF network",
            "At 50% coverage (est. 2027), migration accelerates",
            "Wholesale revenue EUR 380M/q at risk over 3-5 years",
            "OPPORTUNITY: OpenRAN quality untested — may drive subs back",
        ]
        self._add_text_box(slide, Inches(6.8), Inches(2.6), Inches(5.8), Inches(0.3),
                           "Wholesale Risk:", font_size=12,
                           font_color=self.style.primary_color, bold=True)
        self._add_bullet_list(slide, Inches(6.8), Inches(3.0), Inches(5.8),
                              Inches(3), threat_opp, font_size=10)

        self._add_key_message_bar(
            slide, "1&1 at 25% coverage with cheapest unlimited (EUR 40) — wholesale revenue (EUR 380M/q) is VF's main risk")

    def _slide_competition_benchmark(self, d: DeepAnalysisData):
        """Competitive radar + benchmark."""
        slide = self._new_slide("comp_radar", "Competitive Benchmark")
        self._add_header(slide, "Competitive Strength Radar",
                         "Cross-Operator Comparison")

        dims = ['Brand', 'Network', '5G', 'Pricing', 'Enterprise',
                'Service', 'Digital', 'Innovation']
        scores = {
            'DT': [90, 92, 95, 60, 92, 82, 85, 88],
            'Vodafone': [82, 80, 85, 72, 82, 70, 75, 75],
            'O2': [75, 78, 88, 85, 70, 72, 70, 80],
            '1&1': [65, 45, 55, 90, 40, 65, 60, 72],
        }
        chart_path = self.chart_gen.create_radar_chart(
            dims, scores, target_operator='Vodafone',
            title="Competitive Strength (0-100)",
            filename="deep_comp_radar.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.1),
                        Inches(7.5), Inches(5.2))

        # Right: key insight
        insight = [
            "DT outscores VF on 8 of 9 dimensions",
            "VF's only relative advantage: Pricing (72 vs DT 60)",
            "But O2 (85) and 1&1 (90) outscore VF on Pricing too",
            "= Squeezed middle quantified",
            "",
            "VF wins on:",
            "  Cable broadband (~24M homepass)",
            "  Convergence/FMC (4,900K subs)",
            "  TV leadership (7,740K — 1.8x DT)",
        ]
        self._add_bullet_list(slide, Inches(8.0), Inches(1.5), Inches(5),
                              Inches(4.5), [i for i in insight if i], font_size=11)

        self._add_key_message_bar(
            slide, "DT outscores VF on 8/9 dimensions; VF's advantages are cable, FMC, and TV — must defend and extend")

    # =====================================================================
    # Slides 30-35: Look 4 — Self
    # =====================================================================

    def _slide_self_health(self, d: DeepAnalysisData):
        slide = self._new_slide("self_health", "Operating Health Check")
        self._add_header(slide, "Operating Health Check", "Financial Health")

        metrics = [
            ("Revenue", "EUR 3,092M/q", "+0.7% YoY"),
            ("EBITDA", "EUR 1,120M/q", "36.2% margin"),
            ("Capex", "EUR 800M/q", "25.9% intensity"),
            ("Employees", "14,500", ""),
        ]
        self._add_metric_cards(slide, metrics, top=1.3)

        # Revenue trend line
        quarters = ['Q1', 'Q2', 'Q3', 'Q4', 'Q5', 'Q6', 'Q7', 'Q8']
        trend = {
            'Revenue': [3050, 3080, 3060, 3070, 3080, 3090, 3090, 3092],
            'EBITDA': [1080, 1090, 1090, 1100, 1100, 1110, 1110, 1120],
        }
        chart_path = self.chart_gen.create_multi_line_trend(
            quarters, trend,
            title="Revenue & EBITDA Trend (EUR M)",
            y_label="EUR M",
            target_operator='Revenue',
            filename="deep_rev_trend.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(2.9),
                        Inches(7), Inches(3.2))

        # Right: diagnosis
        items = [
            "Revenue flat: EUR 3,050M -> EUR 3,092M over 8Q (+1.4%)",
            "Margin expanding: 35.4% -> 36.2% (+0.8pp)",
            "Cost optimization doing the heavy lifting, not growth",
            "Capex/Rev 25.9% vs DT 19.4% — less efficient",
            '"Managed optimization" playbook, not a growth story',
        ]
        self._add_bullet_list(slide, Inches(7.5), Inches(3.0), Inches(5.3),
                              Inches(3), items, font_size=11)

        self._add_key_message_bar(
            slide, "Slow growth mode — revenue +0.7% but margin improving; cost optimization doing heavy lifting")

    def _slide_self_revenue_mix(self, d: DeepAnalysisData):
        slide = self._new_slide("self_rev_mix", "Revenue Mix")
        self._add_header(slide, "Revenue Breakdown & Segment Portfolio",
                         "Revenue Analysis")

        # Bar chart: segment revenue
        segments = ['Mobile\n49.2%', 'Fixed BB\n25.7%', 'B2B\n16.8%',
                    'Wholesale\n12.3%', 'TV\n8.7%']
        values = [1520, 795, 520, 380, 268]
        chart_path = self.chart_gen.create_bar_chart(
            segments, values, target_category='B2B\n16.8%',
            title="Segment Revenue (EUR M/q)",
            y_label="EUR M",
            filename="deep_rev_mix.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(7), Inches(3.0))

        # Growth arrows
        growth_items = [
            ("Mobile +2.8%", self.style.positive_color),
            ("Fixed BB -1.1%", self.style.negative_color),
            ("B2B +8.5%", self.style.positive_color),
            ("Wholesale +90%", self.style.positive_color),
            ("TV -3.6%", self.style.negative_color),
        ]
        x = Inches(7.5)
        y = Inches(1.4)
        self._add_text_box(slide, x, y, Inches(5), Inches(0.3),
                           "YoY Growth:", font_size=12,
                           font_color=self.style.primary_color, bold=True)
        y += Inches(0.35)
        for label, color in growth_items:
            self._add_text_box(slide, x, y, Inches(5), Inches(0.3),
                               label, font_size=12, font_color=color, bold=True)
            y += Inches(0.35)

        # Wholesale trend
        x_labels = ['Q1', 'Q2', 'Q3', 'Q4', 'Q5', 'Q6', 'Q7', 'Q8']
        ws_trend = {'Wholesale': [180, 185, 190, 200, 220, 260, 320, 380]}
        chart_path2 = self.chart_gen.create_multi_line_trend(
            x_labels, ws_trend,
            title="Wholesale Revenue Ramp (EUR M)",
            y_label="EUR M",
            filename="deep_wholesale_trend.png")
        self._add_image(slide, chart_path2, Inches(0.3), Inches(4.4),
                        Inches(6), Inches(1.8))

        warn = [
            "Wholesale EUR 180M -> EUR 380M in 8Q (1&1 driven)",
            "Peak reached — structural decline as 1&1 builds coverage",
            "EUR 100-200M/q at risk by 2028-29",
        ]
        self._add_bullet_list(slide, Inches(6.5), Inches(4.5), Inches(6),
                              Inches(1.7), warn, font_size=10)

        self._add_key_message_bar(
            slide, "B2B (+8.5%) and Wholesale (+90%) mask flat consumer; wholesale is temporary windfall")

    def _slide_self_segments(self, d: DeepAnalysisData):
        """Key segment metrics: cable erosion + fiber growth + FMC."""
        slide = self._new_slide("self_segments", "Segment Deep Dive")
        self._add_header(slide, "Fixed Broadband — The Critical Battleground",
                         "Segment Analysis")

        # Cable erosion + fiber growth trend
        x_labels = ['Q1', 'Q2', 'Q3', 'Q4', 'Q5', 'Q6', 'Q7', 'Q8']
        series = {
            'Cable': [7500, 7420, 7340, 7250, 7170, 7090, 7030, 6980],
            'Fiber': [850, 930, 1020, 1100, 1180, 1280, 1380, 1480],
            'DSL': [1480, 1480, 1480, 1480, 1480, 1480, 1480, 1480],
        }
        chart_path = self.chart_gen.create_multi_line_trend(
            x_labels, series,
            title="Fixed BB Subscribers by Technology (K)",
            y_label="Subscribers (K)",
            target_operator='Cable',
            filename="deep_bb_tech.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(7.5), Inches(3.0))

        # Right: key numbers
        bb_items = [
            "Cable: 7,500K -> 6,980K (-6.9%) HEMORRHAGING",
            "Fiber: 850K -> 1,480K (+74.1%) BRIGHT SPOT",
            "DSL: 1,480K stable (legacy)",
            "Net: -63K/q — ONLY operator losing BB subs",
            "",
            "At current pace, cable < 6,000K by end-2027",
            "FibreCo JV + DOCSIS 4.0 is the dual-track answer",
        ]
        self._add_bullet_list(slide, Inches(8.0), Inches(1.4), Inches(5),
                              Inches(3), [i for i in bb_items if i], font_size=11)

        # FMC growth
        fmc_series = {'FMC Subs': [4200, 4300, 4400, 4500, 4600, 4700, 4800, 4900]}
        chart_path2 = self.chart_gen.create_multi_line_trend(
            x_labels, fmc_series,
            title="FMC Subscribers (K)",
            y_label="K",
            filename="deep_fmc_trend.png")
        self._add_image(slide, chart_path2, Inches(0.3), Inches(4.4),
                        Inches(6), Inches(1.8))

        fmc_items = [
            "FMC: 4,200K -> 4,900K (+16.7%)",
            "Penetration: 40.7% -> 49.3% (+8.6pp)",
            "FMC customers churn 30-40% less",
        ]
        self._add_bullet_list(slide, Inches(6.5), Inches(4.5), Inches(6),
                              Inches(1.7), fmc_items, font_size=10)

        self._add_key_message_bar(
            slide, "Cable hemorrhaging (-6.9% over 8Q); fiber growing (+74%); FMC at 49% is the retention engine")

    def _slide_self_network(self, d: DeepAnalysisData):
        slide = self._new_slide("self_network", "Network Assessment")
        self._add_header(slide, "Network Infrastructure & Technology Gap",
                         "Network")

        # Donut gauges: coverage
        labels = ['5G 92%', '4G 99.5%', 'Virtual. 45%']
        values = [92, 99.5, 45]
        chart_path = self.chart_gen.create_donut_gauges(
            labels, values, title="Network Coverage",
            filename="deep_coverage.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(12), Inches(2.5))

        # Three critical gaps
        gaps = [
            "5G coverage: 92% vs DT 97%, O2 98% — now #3 network",
            "5G SA: DT commercial since Sep 2024; VF still testing — 12-18mo behind",
            "Fiber homepass: 1,500K vs DT 8,500K (5.7x gap)",
        ]
        self._add_text_box(slide, Inches(0.5), Inches(4.0), Inches(5), Inches(0.3),
                           "Three Critical Gaps:", font_size=12,
                           font_color=self.style.negative_color, bold=True)
        self._add_bullet_list(slide, Inches(0.5), Inches(4.4), Inches(5.5),
                              Inches(2), gaps, font_size=11)

        # Unique advantage
        strength = [
            "Cable network (~24M homepass) — unique differentiator",
            "1000 Mbps cable at competitive prices",
            "DOCSIS 4.0 -> 10Gbps symmetric (matches fiber)",
        ]
        self._add_text_box(slide, Inches(6.5), Inches(4.0), Inches(6), Inches(0.3),
                           "Unique Advantage:", font_size=12,
                           font_color=self.style.positive_color, bold=True)
        self._add_bullet_list(slide, Inches(6.5), Inches(4.4), Inches(6),
                              Inches(2), strength, font_size=11)

        self._add_key_message_bar(
            slide, "5G #3, fiber 5.7x behind DT, SA 12-18mo late — but cable 24M homepass is unchallengeable moat")

    def _slide_self_strengths_weaknesses(self, d: DeepAnalysisData):
        slide = self._new_slide("self_sw", "Strengths / Weaknesses / Exposure")
        self._add_header(slide, "Strengths, Weaknesses & Exposure Points",
                         "Self Assessment")

        # Three columns
        cols = [
            ("Strengths", [
                "Cable BB network (24M homes)",
                "Brand strength (82/100)",
                "Enterprise solutions (82/100)",
                "EBITDA margin (36.2%)",
                "TV leadership (7,740K)",
                "Revenue growth trajectory",
            ], self.style.positive_color, Inches(0.3)),
            ("Weaknesses", [
                "Highest churn (1.05%)",
                "Customer service (70/100)",
                "Pricing competitiveness (72)",
                "5G coverage gap (92%)",
                "Fiber homepass gap (1.5M)",
                "BB subscriber losses (-63K/q)",
            ], self.style.negative_color, Inches(4.5)),
            ("Exposure Points", [
                "1&1 wholesale decline (EUR 100-200M/q)",
                "DT fiber overbuild in cable footprint",
                "O2 further price disruption",
                "5G SA delay beyond 2026",
                "Cable obsolescence perception",
            ], self.style.warning_color, Inches(8.7)),
        ]

        for title, items, color, x in cols:
            self._add_text_box(slide, x, Inches(1.4), Inches(3.8), Inches(0.3),
                               title, font_size=13, font_color=color, bold=True)
            self._add_bullet_list(slide, x, Inches(1.8), Inches(3.8),
                                  Inches(4.2), items, font_size=10)

        self._add_key_message_bar(
            slide, "Stable but diverging — B2B and wholesale masking declines in fixed and competitive slippage")

    # =====================================================================
    # Slides 36-40: SWOT Synthesis
    # =====================================================================

    def _slide_swot_matrix(self, d: DeepAnalysisData):
        slide = self._new_slide("swot_matrix", "SWOT Matrix")
        self._add_header(slide, "SWOT Analysis", "Strategic Synthesis")

        swot_data = {
            'strengths': [
                "Brand strength (82/100)",
                "Enterprise solutions (82)",
                "EBITDA margin 36.2%",
                "Cable network (24M homes)",
                "TV leadership (7,740K)",
                "Market position #2 (25.1%)",
            ],
            'weaknesses': [
                "Customer service (70 vs avg 74)",
                "Pricing competitiveness (72 vs 78)",
                "Wholesale revenue risk (1&1)",
            ],
            'opportunities': [
                "FibreCo JV (7M FTTH)",
                "1&1 Wholesale (EUR 380M/q)",
                "Skaylink B2B (EUR 175M)",
                "Spectrum extension (5yr)",
                "Gigabit Funding (EUR 1.2B)",
                "FMC convergence growth",
            ],
            'threats': [
                "DT fiber overbuild",
                "Nebenkostenprivileg",
                "Competitor pressure (O2/1&1)",
                "Substitutes (OTT, FWA)",
            ],
        }
        chart_path = self.chart_gen.create_swot_matrix(
            swot_data, title="SWOT Analysis",
            filename="deep_swot_matrix.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.1),
                        Inches(8), Inches(5.2))

        # Posture
        self._add_text_box(slide, Inches(8.5), Inches(1.3), Inches(4.3), Inches(0.3),
                           "Posture: OFFENSIVE", font_size=14,
                           font_color=self.style.positive_color, bold=True)
        posture_items = [
            "S=7 > W=3",
            "O=11 > T=4 (ratio ~3:1)",
            "SO-dominant strategy",
            "Aggressive exploitation,",
            "not defensive retrenchment",
        ]
        self._add_bullet_list(slide, Inches(8.5), Inches(1.7), Inches(4.3),
                              Inches(4), posture_items, font_size=11)

        self._add_key_message_bar(
            slide, "Offensive posture — strengths > weaknesses, opportunities 3:1 vs threats; SO-dominant strategy")

    def _slide_swot_strategies(self, d: DeepAnalysisData):
        """Strategy 2x2 matrix."""
        slide = self._new_slide("swot_strategies", "Strategy Matrix")
        self._add_header(slide, "SWOT Strategy Matrix — Cross-Quadrant",
                         "Strategic Responses")

        # Build as KPI table
        strategies = {
            'SO (Exploit)': [
                'Use EBITDA to fund infra',
                'Enterprise -> digital subsidies',
                'Brand -> regulatory position',
                'Network -> spectrum value',
            ],
            'WO (Improve)': [
                'Fix Service 70->74+',
                'Reframe pricing to value',
                'Manage wholesale transition',
                '',
            ],
            'ST (Defend)': [
                'Network vs competitors',
                'Enterprise vs Nebenkostenprivileg',
                'Brand vs regulation',
                '',
            ],
            'WT (Mitigate)': [
                'Break weakness spiral',
                'Pricing + Nebenkostenprivileg',
                'Wholesale + competitor defense',
                '',
            ],
        }
        row_labels = ['Strategy 1', 'Strategy 2', 'Strategy 3', 'Strategy 4']
        chart_path = self.chart_gen.create_kpi_table_chart(
            row_labels, strategies,
            target_operator='SO (Exploit)',
            title="SWOT Strategy Matrix",
            filename="deep_swot_strategies.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(12.5), Inches(4.0))

        # Key message
        self._add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.5),
                           "SO Mega-Strategy: Use financial strength and enterprise capabilities "
                           "to fund and capture the infrastructure transformation opportunity.",
                           font_size=12, font_color=self.style.primary_color, bold=True)

        self._add_key_message_bar(
            slide, "SO-dominant: use EBITDA + enterprise to fund infra transformation; WO: fix service as highest ROI")

    def _slide_swot_squeezed(self, d: DeepAnalysisData):
        slide = self._new_slide("swot_squeezed", "Squeezed Middle — SWOT View")
        self._add_header(slide, 'The "Squeezed Middle" — SWOT Evidence',
                         "Strategic Diagnosis")

        # Gap analysis: VF vs DT on SWOT strengths
        dims = ['Brand', 'Enterprise', 'Network', 'EBITDA Margin', '5G Coverage', 'Revenue Share']
        vf_scores = [82, 82, 80, 36.2, 92, 25.1]
        dt_scores = [85, 85, 88, 40.0, 97, 50.3]
        # Normalize to DT=100
        vf_norm = [v / d * 100 for v, d in zip(vf_scores, dt_scores)]
        dt_norm = [100] * len(dims)

        chart_path = self.chart_gen.create_gap_analysis_chart(
            dims, vf_norm, dt_norm,
            target_name='Vodafone', leader_name='DT',
            title="VF Strengths vs DT (DT=100)",
            filename="deep_swot_gap.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(7.5), Inches(4.0))

        # Key insight
        insight = [
            "VF strengths are above market average",
            "But BELOW DT on every single dimension",
            "This is the fundamental positioning problem:",
            "Not weakness — insufficiency of strength vs leader",
            "",
            "VF charges near-premium (closer to DT)",
            "But delivers mid-range (closer to O2)",
        ]
        self._add_bullet_list(slide, Inches(8.2), Inches(1.4), Inches(4.5),
                              Inches(4.5), [i for i in insight if i], font_size=11)

        self._add_key_message_bar(
            slide, "VF's strengths are above average but below DT on every dimension — not weak, but insufficient vs leader")

    def _slide_swot_priorities(self, d: DeepAnalysisData):
        slide = self._new_slide("swot_priorities", "SWOT Priorities")
        self._add_header(slide, "SWOT-Derived Strategic Priorities",
                         "Priority Ranking")

        items = [
            "P1: Defend Fixed BB via FibreCo JV + DOCSIS 4.0 (ST) — 2-6yr, HIGH invest",
            "P2: Fix Customer Service 70->74+ (WO) — 12-18mo, MEDIUM invest",
            "P3: Maximize FMC 49%->60% (SO) — 12-24mo, MEDIUM invest",
            "P4: Accelerate B2B via Skaylink (SO+WO) — 12-36mo, HIGH invest",
            "P5: Manage Wholesale Transition (WO+WT) — 3-5yr, LOW invest",
        ]
        prios = ['P1', 'P2', 'P3', 'P4', 'P5']
        chart_path = self.chart_gen.create_priority_chart(
            items, prios, title="SWOT-Derived Priorities",
            filename="deep_swot_prios.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(12.5), Inches(3.0))

        # What NOT to do
        traps = [
            "Price wars — destroys 36.2% margin without winning share",
            "Premium repositioning — behind DT on every dimension, 5+ years",
            "Standalone TV investment — declining; value is as FMC anchor",
            "Delay fiber — every quarter = ~63K BB subs lost to DT",
        ]
        self._add_text_box(slide, Inches(0.5), Inches(4.5), Inches(3), Inches(0.3),
                           "What NOT To Do:", font_size=12,
                           font_color=self.style.negative_color, bold=True)
        self._add_bullet_list(slide, Inches(0.5), Inches(4.9), Inches(12),
                              Inches(1.5), traps, font_size=10)

        self._add_key_message_bar(
            slide, "5 priorities: defend BB (existential), fix service (highest ROI), grow FMC, accelerate B2B, manage wholesale")

    # =====================================================================
    # Slides 41-45: Look 5 — Opportunities
    # =====================================================================

    def _slide_span_bubble(self, d: DeepAnalysisData):
        slide = self._new_slide("span_bubble", "SPAN Matrix")
        self._add_header(slide, "SPAN Matrix — 21 Opportunities",
                         "Opportunity Positioning")

        # Build positions from extracted data
        positions = []
        for opp in d.opportunities.opportunities:
            positions.append({
                'competitive_position': opp.span_cp if opp.span_cp > 0 else 5.5,
                'market_attractiveness': opp.span_ma if opp.span_ma > 0 else 5.5,
                'bubble_size': 2.0 if opp.priority == 'P0' else 1.5 if opp.priority == 'P1' else 1.0,
                'opportunity_name': opp.name[:20],
                'quadrant': opp.quadrant,
            })

        # Add some default positions if extraction was sparse
        if len(positions) < 10:
            defaults = [
                ('FibreCo JV', 5.25, 5.5, 'grow_invest', 2.5),
                ('Spectrum Ext.', 5.5, 5.5, 'grow_invest', 2.0),
                ('Gigabit Fund', 5.25, 5.5, 'grow_invest', 2.0),
                ('1&1 Wholesale', 5.5, 5.5, 'grow_invest', 2.0),
                ('Skaylink B2B', 5.25, 5.5, 'grow_invest', 2.0),
                ('FMC Expansion', 5.75, 5.5, 'grow_invest', 2.0),
                ('Cable Restructure', 5.5, 5.0, 'grow_invest', 1.5),
                ('5G Enterprise', 5.25, 5.25, 'grow_invest', 1.5),
                ('AI/ML Network', 5.25, 5.25, 'grow_invest', 1.0),
                ('SO-1 Brand', 5.5, 5.6, 'grow_invest', 1.5),
                ('SO-2 Enterprise', 5.5, 5.6, 'grow_invest', 1.5),
                ('SO-3 Network', 5.5, 5.6, 'grow_invest', 1.5),
                ('SO-4 EBITDA', 5.5, 5.6, 'grow_invest', 1.5),
                ('WO-1 Service', 4.75, 5.2, 'acquire_skills', 1.5),
                ('WO-2 Pricing', 4.75, 5.2, 'acquire_skills', 1.0),
                ('WO-3 Wholesale', 4.75, 5.2, 'acquire_skills', 1.0),
            ]
            existing_names = {p['opportunity_name'] for p in positions}
            for name, cp, ma, quad, size in defaults:
                if name[:20] not in existing_names:
                    positions.append({
                        'competitive_position': cp,
                        'market_attractiveness': ma,
                        'bubble_size': size,
                        'opportunity_name': name,
                        'quadrant': quad,
                    })

        chart_path = self.chart_gen.create_span_bubble_chart(
            positions, title="SPAN Matrix — 21 Opportunities",
            filename="deep_span_bubble.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.1),
                        Inches(8.5), Inches(5.2))

        # Right: quadrant summary
        summary = [
            "Grow/Invest: 18 (86%)",
            "Acquire Skills: 3 (14%)",
            "Harvest: 0",
            "Avoid/Exit: 0",
            "",
            "86% in Grow/Invest confirms",
            "offensive strategic posture",
            "",
            "Challenge is not finding opportunities",
            "but prioritizing and sequencing them",
        ]
        self._add_bullet_list(slide, Inches(9.0), Inches(1.5), Inches(4),
                              Inches(5), [s for s in summary if s], font_size=11)

        self._add_key_message_bar(
            slide, "86% of opportunities in Grow/Invest — market offering favorable hand; challenge is execution speed")

    def _slide_opp_priority_table(self, d: DeepAnalysisData):
        slide = self._new_slide("opp_priority", "Opportunity Priorities")
        self._add_header(slide, "Opportunity Priority Ranking",
                         "Resource Allocation")

        items = [
            "P0: FibreCo JV — 7M FTTH homes, defend BB (existential)",
            "P0: 1&1 Wholesale — maximize EUR 380M/q before decline",
            "P0: FMC Expansion — 49% to 60%, cable+mobile unique",
            "P0: Skaylink B2B — EUR 1B 5yr target, cloud/security",
            "P0: Spectrum Extension — 5yr certainty, avoid re-auction",
            "P1: Customer Service — score 70 to 74+, highest ROI",
            "P1: Gigabit Fund — EUR 1.2B co-funding for fiber",
            "P1: Cable Restructuring — DOCSIS 4.0 enabler",
            "P2: 5G SA Enterprise — network slicing, needs SA launch",
            "P2: AI/ML Network — 10-20% OpEx reduction",
        ]
        prios = ['P0', 'P0', 'P0', 'P0', 'P0', 'P1', 'P1', 'P1', 'P2', 'P2']

        chart_path = self.chart_gen.create_priority_chart(
            items, prios, title="Opportunity Priorities",
            filename="deep_opp_priorities.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(12.5), Inches(4.8))

        self._add_key_message_bar(
            slide, "5 P0 (must-do), 3 P1 (should-do), 2 P2 (explore) — 60/30/10% resource allocation")

    def _slide_opp_financial(self, d: DeepAnalysisData):
        """Financial impact assessment."""
        slide = self._new_slide("opp_financial", "Financial Impact")
        self._add_header(slide, "Financial Impact Assessment",
                         "Revenue & Investment")

        # Revenue impact table
        rev_data = {
            '1&1 Wholesale': ['EUR 400-800M', 'Immediate, declining', 'High'],
            'FMC Expansion': ['+EUR 600-900M', '12-24 months', 'Med-High'],
            'Skaylink B2B': ['+EUR 200M/yr', '24-36 months', 'Medium'],
            'FibreCo JV': ['Defend EUR 3.2B', '3-6 years', 'Medium'],
            'Service Fix': ['+EUR 150-250M', '12-18 months', 'Medium'],
            '5G Enterprise': ['+EUR 50-100M', '3-5 years', 'Low'],
        }
        chart_path = self.chart_gen.create_kpi_table_chart(
            ['Annual Revenue', 'Timeline', 'Confidence'], rev_data,
            title="Revenue Impact by Opportunity (Annual, at Maturity)",
            filename="deep_fin_impact.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(12.5), Inches(3.0))

        # Net value comparison
        scenarios = ['Execute P0', 'Execute P0+P1', 'Do Nothing']
        rev_vals = [6.5, 8.0, -3.0]
        inv_vals = [5.0, 6.0, 0]
        net_vals = [1.5, 3.0, -3.0]

        chart_path2 = self.chart_gen.create_segment_comparison(
            scenarios, {
                '5yr Revenue (EUR B)': rev_vals,
                '5yr Investment (EUR B)': inv_vals,
                'Net Value (EUR B)': net_vals,
            },
            title="5-Year Scenario Comparison (EUR B)",
            y_label="EUR B",
            filename="deep_scenario_compare.png")
        self._add_image(slide, chart_path2, Inches(0.3), Inches(4.4),
                        Inches(8), Inches(2.3))

        insight = [
            "Cost of inaction: -EUR 2-4B",
            "Net value of full execution: +EUR 2-4B",
            "Asymmetry is clear: invest or decline",
        ]
        self._add_bullet_list(slide, Inches(8.5), Inches(4.5), Inches(4),
                              Inches(2), insight, font_size=11)

        self._add_key_message_bar(
            slide, "Cost of inaction (EUR 2-4B loss) exceeds cost of action (EUR 1-3B net gain) — VF must invest")

    def _slide_opp_timeline(self, d: DeepAnalysisData):
        """Execution timeline."""
        slide = self._new_slide("opp_timeline", "Execution Timeline")
        self._add_header(slide, "Opportunity Sequencing — 2025 to 2029+",
                         "Execution Plan")

        milestones = [
            {'date': 'Q4 2025', 'name': 'Lock wholesale\nLaunch FMC push', 'priority': 'P0'},
            {'date': 'H1 2026', 'name': 'Skaylink integration\nDOCSIS 4.0 start', 'priority': 'P0'},
            {'date': 'H2 2026', 'name': '5G SA launch\nService score 74', 'priority': 'P1'},
            {'date': '2027', 'name': 'FMC 60%\nB2B EUR 200M+/yr', 'priority': 'P0'},
            {'date': '2028', 'name': 'FibreCo ramp\nBB stabilize', 'priority': 'P0'},
            {'date': '2029+', 'name': 'Full fiber\nWholesale replaced', 'priority': 'P1'},
        ]
        chart_path = self.chart_gen.create_timeline_chart(
            milestones, title="Strategic Execution Timeline",
            filename="deep_timeline.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(12.5), Inches(3.0))

        # Resource allocation
        alloc = [
            "P0 (Must Execute): 60% of strategic investment",
            "  FibreCo JV, Wholesale Max, FMC, Skaylink, Spectrum",
            "P1 (Should Execute): 30% of strategic investment",
            "  Customer Service, Gigabit Fund, Cable Restructuring",
            "P2 (Explore): 10% of strategic investment",
            "  5G SA/Enterprise, AI/ML, Open RAN",
        ]
        self._add_bullet_list(slide, Inches(0.5), Inches(4.5), Inches(12),
                              Inches(1.8), alloc, font_size=11)

        self._add_key_message_bar(
            slide, "3-5 year window to execute — FibreCo JV and wholesale revenue fund the transformation")

    # =====================================================================
    # Slides 46-48: Summary
    # =====================================================================

    def _slide_risk_reward(self, d: DeepAnalysisData):
        """Risk/Reward comparison."""
        slide = self._new_slide("risk_reward", "Risk / Reward Summary")
        self._add_header(slide, "Risk / Reward Summary", "Strategic Outlook")

        # Bull vs Bear comparison
        bull_data = {
            'Bull Case': ['EUR 3,400-3,600M/q', '10,500K+', '65%+', 'EUR 700-800M/q', '<0.85%', '8,500K', '90%+'],
            'Bear Case': ['EUR 2,700-2,900M/q', '<8,000K', '<45%', 'EUR 500M/q', '>1.2%', '<3,000K', '<50%'],
            'Current': ['EUR 3,092M', '9,940K', '49%', 'EUR 520M/q', '1.05%', '1,500K', '0%'],
        }
        row_labels = ['Revenue', 'BB Subs', 'FMC Pen.', 'B2B Rev', 'Churn', 'Fiber HP', '5G SA']
        chart_path = self.chart_gen.create_kpi_table_chart(
            row_labels, bull_data,
            target_operator='Current',
            title="5-Year Scenario Comparison",
            filename="deep_risk_reward.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(12.5), Inches(3.5))

        # Net assessment
        self._add_text_box(slide, Inches(0.5), Inches(5.0), Inches(12), Inches(0.8),
                           "The asymmetry is clear: downside of inaction (EUR 2-4B value destruction) "
                           "exceeds net cost of action (EUR 1-3B net value creation after investment). "
                           "Vodafone must invest.",
                           font_size=13, font_color=self.style.primary_color, bold=True)

        self._add_key_message_bar(
            slide, "Execute: +EUR 1-3B net value | Do nothing: -EUR 2-4B loss — the cost of inaction exceeds the cost of action")

    def _slide_kpi_dashboard(self, d: DeepAnalysisData):
        """KPI dashboard — 9 rows x 4 columns."""
        slide = self._new_slide("kpi_dashboard", "KPI Dashboard")
        self._add_header(slide, "Success Metrics Dashboard", "KPI Tracking")

        kpi_data = {
            'Current': ['+0.7%', '36.2%', '-63K/q', '49%', '1.05%', '70', 'EUR 520M/q', '1,500K', '0%'],
            '12-Month': ['+1.0%', '36.5%', '-30K/q', '55%', '0.95%', '74', 'EUR 560M/q', '2,500K', '30%'],
            '3-Year': ['+2.0%', '37.0%', '0K/q', '60%', '0.88%', '78', 'EUR 650M/q', '5,000K', '70%'],
            '5-Year': ['+2.5%', '38.0%', '+20K/q', '65%', '0.82%', '80', 'EUR 800M/q', '8,500K', '90%+'],
        }
        row_labels = [
            'Revenue growth', 'EBITDA margin', 'BB net adds',
            'FMC penetration', 'Monthly churn', 'Service score',
            'B2B revenue', 'Fiber homepass', '5G SA coverage',
        ]
        chart_path = self.chart_gen.create_kpi_table_chart(
            row_labels, kpi_data,
            target_operator='Current',
            title="Success Metrics Dashboard",
            filename="deep_kpi_dashboard.png")
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2),
                        Inches(12.5), Inches(5.0))

        self._add_key_message_bar(
            slide, "9 KPIs tracked across 4 horizons — BB stabilization, churn reduction, and FMC growth are the leading indicators")

    def _slide_back_cover(self):
        slide = self._new_slide("back_cover", "Back Cover")
        self._add_shape(slide, 0, 0, Inches(self.SLIDE_WIDTH),
                        Inches(self.SLIDE_HEIGHT), self.style.primary_color)
        self._add_text_box(slide, Inches(2), Inches(2.0), Inches(9), Inches(1.5),
                           "Thank You", font_size=48,
                           font_color=(255, 255, 255), bold=True, align="center")
        self._add_text_box(slide, Inches(2), Inches(3.5), Inches(9), Inches(1),
                           "Vodafone Germany — BLM Strategic Assessment",
                           font_size=20, font_color=(255, 220, 220), align="center")
        self._add_text_box(slide, Inches(2), Inches(4.5), Inches(9), Inches(1),
                           "Deep Analysis | CQ4 2025 | Five Looks + SWOT + SPAN",
                           font_size=16, font_color=(255, 220, 220), align="center")
        self._add_text_box(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.5),
                           f"Generated: {datetime.now().strftime('%Y-%m-%d')}",
                           font_size=12, font_color=(200, 150, 150), align="center")


# =====================================================================
# CLI entry point
# =====================================================================

def main():
    data_dir = "data/output"
    if len(sys.argv) > 1:
        data_dir = sys.argv[1]

    gen = DeepAnalysisPPTGenerator(data_dir=data_dir)
    output = gen.generate_deep()
    print(f"\nDone! Output: {output}")

    # Quick verification
    from pptx import Presentation as Prs
    p = Prs(output)
    print(f"Total slides: {len(p.slides)}")


if __name__ == "__main__":
    main()
