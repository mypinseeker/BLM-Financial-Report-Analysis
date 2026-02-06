"""Enhanced BLM PPT Generator with Charts and Visual Elements.

Extends the basic PPT generator with:
- Market share bar charts
- Revenue comparison charts
- Competitive radar charts
- KPI comparison tables
- 5G coverage gauges
- Gap analysis charts
- Strategy priority visualization
- Execution timeline

Follows Huawei 5G template style with rich visual elements.
"""

from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Optional
import tempfile

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from src.blm.five_looks import InsightResult
from src.blm.three_decisions import StrategyResult, StrategyItem
from src.blm.ppt_generator import PPTStyle, HUAWEI_STYLE, VODAFONE_STYLE, STYLES
from src.blm.ppt_charts import PPTChartGenerator


class BLMPPTGeneratorEnhanced:
    """Enhanced PPT generator with chart integration."""

    def __init__(
        self,
        style: str = "huawei",
        output_dir: Optional[str] = None,
        chart_dpi: int = 150,
    ):
        """Initialize enhanced PPT generator.

        Args:
            style: PPT style name ("huawei", "vodafone")
            output_dir: Output directory for generated PPTs
            chart_dpi: DPI for chart images
        """
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for PPT generation. "
                "Install it with: pip install python-pptx"
            )

        if isinstance(style, str):
            self.style = STYLES.get(style, HUAWEI_STYLE)
        else:
            self.style = style

        if output_dir is None:
            output_dir = str(Path(__file__).resolve().parent.parent.parent / "data" / "output")
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # Initialize chart generator
        self.chart_dir = Path(tempfile.mkdtemp())
        self.chart_gen = PPTChartGenerator(output_dir=str(self.chart_dir), dpi=chart_dpi)

        self.prs = None
        self.slide_num = 0

    def generate(
        self,
        five_looks: dict[str, InsightResult],
        three_decisions: dict[str, StrategyResult],
        target_operator: str,
        competitors: list[str] = None,
        financial_data: dict = None,
        competitive_scores: dict = None,
        title: str = None,
        filename: str = None,
    ) -> str:
        """Generate enhanced PPT with charts.

        Args:
            five_looks: Five Looks analysis results
            three_decisions: Three Decisions strategy results
            target_operator: Target operator name
            competitors: List of competitor names
            financial_data: Raw financial data for charts
            competitive_scores: Competitive dimension scores for radar chart
            title: Presentation title
            filename: Output filename

        Returns:
            Path to generated PPT file
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        self.prs.slide_height = Inches(7.5)
        self.slide_num = 0

        competitors = competitors or []
        title = title or f"{target_operator} BLM 战略分析报告"
        filename = filename or f"blm_{self._sanitize_name(target_operator)}_enhanced.pptx"

        # Generate slides
        self._add_title_slide(title, target_operator, competitors)
        self._add_agenda_slide()

        # Executive Summary with key metrics
        self._add_executive_summary_slide(five_looks, three_decisions, target_operator)

        # Section 1: Five Looks with charts
        self._add_section_divider("五看分析", "Five Looks Analysis", "01")

        # Market overview with charts
        if financial_data:
            self._add_market_overview_slide(financial_data, target_operator, competitors)
            self._add_revenue_comparison_slide(financial_data, target_operator)
            self._add_5g_coverage_slide(financial_data, target_operator)

        # Competitive radar chart
        if competitive_scores:
            self._add_competitive_radar_slide(competitive_scores, target_operator)

        # Five Looks insight slides
        for key, insight in five_looks.items():
            self._add_insight_slide(insight, financial_data, target_operator)

        # Gap analysis
        if competitive_scores and financial_data:
            self._add_gap_analysis_slide(competitive_scores, target_operator)

        # Section 2: Three Decisions
        self._add_section_divider("三定策略", "Three Decisions Strategy", "02")

        # Strategy slides with priority visualization
        for key, decision in three_decisions.items():
            self._add_strategy_slide_enhanced(decision)

        # Execution timeline
        execution = three_decisions.get("execution")
        if execution:
            self._add_timeline_slide(execution)

        # KPI Dashboard
        self._add_kpi_dashboard_slide(three_decisions)

        # Conclusion and next steps
        self._add_conclusion_slide(five_looks, three_decisions, target_operator)
        self._add_closing_slide()

        # Save
        output_path = self.output_dir / filename
        self.prs.save(str(output_path))
        return str(output_path)

    def _add_title_slide(
        self,
        title: str,
        target_operator: str,
        competitors: list[str],
    ):
        """Add title slide with visual design."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        # Red header bar
        self._add_shape(
            slide, 0, 0, Inches(13.333), Inches(2.8),
            self.style.primary_color
        )

        # Decorative diagonal stripe
        self._add_shape(
            slide, Inches(10), 0, Inches(3.333), Inches(2.8),
            (180, 0, 10),  # Darker red
        )

        # Title
        self._add_text_box(
            slide,
            Inches(0.5), Inches(0.5), Inches(9), Inches(1.2),
            title,
            font_size=40,
            font_color=(255, 255, 255),
            bold=True,
        )

        # Subtitle
        subtitle = f"基于 BLM (Business Leadership Model) 方法论的战略分析"
        self._add_text_box(
            slide,
            Inches(0.5), Inches(1.7), Inches(9), Inches(0.6),
            subtitle,
            font_size=20,
            font_color=(255, 220, 220),
        )

        # Info cards
        info_items = [
            ("分析对象", target_operator),
            ("竞争对手", ", ".join(competitors) if competitors else "N/A"),
            ("报告日期", datetime.now().strftime('%Y年%m月%d日')),
        ]

        x_pos = Inches(0.5)
        for label, value in info_items:
            # Card background
            self._add_shape(
                slide, x_pos, Inches(3.5), Inches(3.8), Inches(1.5),
                (245, 245, 245),
            )
            # Label
            self._add_text_box(
                slide,
                x_pos + Inches(0.2), Inches(3.6), Inches(3.4), Inches(0.4),
                label,
                font_size=12,
                font_color=self.style.light_text_color,
            )
            # Value
            self._add_text_box(
                slide,
                x_pos + Inches(0.2), Inches(4.0), Inches(3.4), Inches(0.8),
                value,
                font_size=18,
                font_color=self.style.primary_color,
                bold=True,
            )
            x_pos += Inches(4.1)

        # Framework info
        framework_items = [
            ("五看 Five Looks", "看市场 | 看自己 | 看对手 | 看宏观 | 看机会"),
            ("三定 Three Decisions", "定策略 | 定重点工作 | 定执行"),
        ]

        y_pos = Inches(5.5)
        for fw_title, fw_desc in framework_items:
            # Icon circle
            self._add_shape(
                slide,
                Inches(0.5), y_pos, Inches(0.4), Inches(0.4),
                self.style.primary_color,
                shape_type=MSO_SHAPE.OVAL,
            )
            # Text
            self._add_text_box(
                slide,
                Inches(1.1), y_pos, Inches(4), Inches(0.4),
                fw_title,
                font_size=14,
                font_color=self.style.text_color,
                bold=True,
            )
            self._add_text_box(
                slide,
                Inches(1.1), y_pos + Inches(0.35), Inches(6), Inches(0.35),
                fw_desc,
                font_size=11,
                font_color=self.style.light_text_color,
            )
            y_pos += Inches(0.85)

    def _add_agenda_slide(self):
        """Add agenda slide with visual timeline."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "报告目录", "CONTENTS")

        # Vertical timeline line
        self._add_shape(
            slide,
            Inches(1.1), Inches(1.8), Inches(0.02), Inches(5),
            self.style.primary_color,
        )

        agenda_items = [
            ("01", "执行摘要", "Executive Summary - 核心数据与关键结论", "page 3-4"),
            ("02", "五看分析", "Five Looks - 市场/自身/竞争/宏观/机会洞察", "page 5-15"),
            ("03", "三定策略", "Three Decisions - 策略/重点工作/执行计划", "page 16-22"),
            ("04", "总结建议", "Summary - 结论与下一步行动", "page 23-24"),
        ]

        y_pos = Inches(1.8)
        for num, title, subtitle, pages in agenda_items:
            # Number circle
            self._add_shape(
                slide,
                Inches(0.85), y_pos, Inches(0.5), Inches(0.5),
                self.style.primary_color,
                shape_type=MSO_SHAPE.OVAL,
            )
            self._add_text_box(
                slide,
                Inches(0.85), y_pos + Inches(0.08), Inches(0.5), Inches(0.35),
                num,
                font_size=14,
                font_color=(255, 255, 255),
                bold=True,
                align="center",
            )

            # Title and description
            self._add_text_box(
                slide,
                Inches(1.6), y_pos, Inches(8), Inches(0.45),
                title,
                font_size=20,
                font_color=self.style.text_color,
                bold=True,
            )
            self._add_text_box(
                slide,
                Inches(1.6), y_pos + Inches(0.45), Inches(8), Inches(0.35),
                subtitle,
                font_size=12,
                font_color=self.style.light_text_color,
            )

            # Page reference
            self._add_text_box(
                slide,
                Inches(10.5), y_pos + Inches(0.15), Inches(2), Inches(0.3),
                pages,
                font_size=10,
                font_color=self.style.primary_color,
            )

            y_pos += Inches(1.2)

    def _add_executive_summary_slide(
        self,
        five_looks: dict[str, InsightResult],
        three_decisions: dict[str, StrategyResult],
        target_operator: str,
    ):
        """Add executive summary with key metrics."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "执行摘要", "EXECUTIVE SUMMARY")

        # Get key data
        self_insight = five_looks.get("self")
        strategy = three_decisions.get("strategy")

        # Metric cards row
        if self_insight and self_insight.metrics:
            metrics = self_insight.metrics
            metric_cards = [
                ("收入规模", f"€{metrics.get('revenue_eur_billion', 'N/A')}B", "Q2 2025"),
                ("市场排名", f"第{metrics.get('revenue_rank', 'N/A')}位", "德国市场"),
                ("市场份额", f"{metrics.get('market_share_pct', 'N/A')}%", "宽带市场"),
                ("5G覆盖率", f"{metrics.get('5g_coverage_pct', 'N/A')}%", "人口覆盖"),
            ]

            x_pos = Inches(0.5)
            for label, value, sub in metric_cards:
                # Card
                self._add_shape(
                    slide, x_pos, Inches(1.5), Inches(3), Inches(1.3),
                    (250, 250, 250),
                )
                # Top accent
                self._add_shape(
                    slide, x_pos, Inches(1.5), Inches(3), Inches(0.06),
                    self.style.primary_color,
                )
                # Value
                self._add_text_box(
                    slide,
                    x_pos + Inches(0.15), Inches(1.65), Inches(2.7), Inches(0.7),
                    value,
                    font_size=28,
                    font_color=self.style.primary_color,
                    bold=True,
                )
                # Label
                self._add_text_box(
                    slide,
                    x_pos + Inches(0.15), Inches(2.3), Inches(2.7), Inches(0.35),
                    f"{label} ({sub})",
                    font_size=10,
                    font_color=self.style.light_text_color,
                )
                x_pos += Inches(3.2)

        # Strategic priorities section
        self._add_text_box(
            slide,
            Inches(0.5), Inches(3.1), Inches(6), Inches(0.4),
            "战略重点 STRATEGIC PRIORITIES",
            font_size=14,
            font_color=self.style.text_color,
            bold=True,
        )

        # Priority items
        if strategy:
            y_pos = Inches(3.6)
            for item in strategy.items[:4]:
                # Priority badge
                priority_colors = {
                    "P0": self.style.primary_color,
                    "P1": (255, 165, 0),
                    "P2": (34, 139, 34),
                }
                color = priority_colors.get(item.priority, self.style.primary_color)

                self._add_shape(
                    slide,
                    Inches(0.5), y_pos, Inches(0.45), Inches(0.35),
                    color,
                )
                self._add_text_box(
                    slide,
                    Inches(0.5), y_pos + Inches(0.05), Inches(0.45), Inches(0.25),
                    item.priority,
                    font_size=10,
                    font_color=(255, 255, 255),
                    bold=True,
                    align="center",
                )

                # Name and description
                self._add_text_box(
                    slide,
                    Inches(1.05), y_pos, Inches(5), Inches(0.35),
                    item.name,
                    font_size=13,
                    font_color=self.style.text_color,
                    bold=True,
                )
                desc = item.description[:65] + "..." if len(item.description) > 65 else item.description
                self._add_text_box(
                    slide,
                    Inches(1.05), y_pos + Inches(0.32), Inches(5), Inches(0.35),
                    desc,
                    font_size=10,
                    font_color=self.style.light_text_color,
                )

                y_pos += Inches(0.75)

        # Key challenges section
        self._add_text_box(
            slide,
            Inches(6.8), Inches(3.1), Inches(6), Inches(0.4),
            "核心挑战 KEY CHALLENGES",
            font_size=14,
            font_color=self.style.text_color,
            bold=True,
        )

        # Challenge box
        self._add_shape(
            slide,
            Inches(6.8), Inches(3.5), Inches(5.8), Inches(2.8),
            (255, 248, 245),
        )

        challenges = []
        if self_insight:
            for finding in self_insight.findings[:5]:
                if any(kw in finding for kw in ["落后", "低于", "高于", "下滑", "差距"]):
                    challenges.append(f"• {finding[:55]}...")

        challenge_text = "\n".join(challenges[:4]) if challenges else "• 详见五看分析"
        self._add_text_box(
            slide,
            Inches(7.0), Inches(3.7), Inches(5.4), Inches(2.4),
            challenge_text,
            font_size=11,
            font_color=self.style.text_color,
        )

        # Bottom key message
        self._add_shape(
            slide,
            Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.7),
            self.style.primary_color,
        )
        self._add_text_box(
            slide,
            Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
            f"核心定位: 做德国市场「有质量的挑战者」- 以网络和服务质量为核心竞争力",
            font_size=16,
            font_color=(255, 255, 255),
            bold=True,
        )

    def _add_market_overview_slide(
        self,
        financial_data: dict,
        target_operator: str,
        competitors: list[str],
    ):
        """Add market overview with market share chart."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "市场格局", "MARKET OVERVIEW")

        # Market share data
        operators = [target_operator] + competitors
        market_shares = []
        for op in operators:
            if op in financial_data:
                share = financial_data[op].get("market_share_broadband_pct", 0)
                market_shares.append(share)
            else:
                market_shares.append(0)

        # Generate and embed chart
        chart_path = self.chart_gen.create_market_share_bar_chart(
            operators=operators,
            market_shares=market_shares,
            target_operator=target_operator,
            title="德国宽带市场份额对比 (%)",
            filename="market_share_overview.png",
        )

        # Add chart image
        slide.shapes.add_picture(
            chart_path,
            Inches(0.5), Inches(1.5),
            width=Inches(6), height=Inches(3),
        )

        # Key metrics table on right side
        self._add_text_box(
            slide,
            Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
            "市场关键数据 Q2 2025",
            font_size=14,
            font_color=self.style.text_color,
            bold=True,
        )

        # Metrics table
        y_pos = Inches(2.0)
        for op in operators[:4]:
            if op in financial_data:
                data = financial_data[op]

                # Operator row
                is_target = op == target_operator
                bg_color = (255, 240, 240) if is_target else (248, 248, 248)

                self._add_shape(
                    slide,
                    Inches(7), y_pos, Inches(5.5), Inches(1.0),
                    bg_color,
                )

                # Operator name
                self._add_text_box(
                    slide,
                    Inches(7.1), y_pos + Inches(0.05), Inches(2.5), Inches(0.35),
                    op,
                    font_size=12,
                    font_color=self.style.primary_color if is_target else self.style.text_color,
                    bold=True,
                )

                # Metrics
                metrics_text = (
                    f"收入: €{data.get('revenue_eur_billion', 'N/A')}B | "
                    f"EBITDA: {data.get('ebitda_margin_pct', 'N/A')}% | "
                    f"5G: {data.get('5g_coverage_pct', 'N/A')}%"
                )
                self._add_text_box(
                    slide,
                    Inches(7.1), y_pos + Inches(0.4), Inches(5.3), Inches(0.5),
                    metrics_text,
                    font_size=10,
                    font_color=self.style.text_color,
                )

                y_pos += Inches(1.1)

        # Market insight
        self._add_shape(
            slide,
            Inches(0.5), Inches(5.0), Inches(12.333), Inches(2),
            (250, 250, 250),
        )
        self._add_text_box(
            slide,
            Inches(0.7), Inches(5.15), Inches(12), Inches(0.35),
            "市场洞察 MARKET INSIGHT",
            font_size=12,
            font_color=self.style.primary_color,
            bold=True,
        )
        insight_text = (
            "• 德电以40.6%份额稳居第一，连续35季度EBITDA增长，领先优势持续扩大\n"
            "• Vodafone以27.0%份额位居第二，服务收入重回增长(+0.5%)\n"
            "• 1&1作为新进入者正在建设自有网络，短期面临成本压力\n"
            "• 市场竞争加剧，价格战压力持续，需要差异化竞争策略"
        )
        self._add_text_box(
            slide,
            Inches(0.7), Inches(5.55), Inches(12), Inches(1.4),
            insight_text,
            font_size=11,
            font_color=self.style.text_color,
        )

    def _add_revenue_comparison_slide(
        self,
        financial_data: dict,
        target_operator: str,
    ):
        """Add revenue comparison chart slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "收入与利润对比", "REVENUE & PROFITABILITY")

        operators = list(financial_data.keys())
        revenues = [financial_data[op].get("revenue_eur_billion", 0) for op in operators]

        # Revenue chart
        chart_path = self.chart_gen.create_revenue_comparison_chart(
            operators=operators,
            revenues=revenues,
            target_operator=target_operator,
            title="收入规模对比 (€B) - Q2 2025",
            filename="revenue_comparison.png",
        )

        slide.shapes.add_picture(
            chart_path,
            Inches(0.5), Inches(1.5),
            width=Inches(6.3), height=Inches(4),
        )

        # EBITDA comparison
        metrics_data = {
            "EBITDA利润率 (%)": {op: financial_data[op].get("ebitda_margin_pct", 0) for op in operators},
            "ARPU (€)": {op: financial_data[op].get("arpu_eur", 0) for op in operators},
        }

        chart_path2 = self.chart_gen.create_financial_metrics_chart(
            operators=operators,
            metrics_data=metrics_data,
            target_operator=target_operator,
            title="盈利能力指标对比",
            filename="ebitda_arpu.png",
        )

        slide.shapes.add_picture(
            chart_path2,
            Inches(7), Inches(1.5),
            width=Inches(5.8), height=Inches(3.5),
        )

        # Key findings
        self._add_shape(
            slide,
            Inches(0.5), Inches(5.8), Inches(12.333), Inches(1.4),
            (250, 250, 250),
        )

        target_data = financial_data.get(target_operator, {})
        dt_data = financial_data.get("Deutsche Telekom", {})

        findings = [
            f"• {target_operator}收入€{target_data.get('revenue_eur_billion', 'N/A')}B，为德电的{target_data.get('revenue_eur_billion', 0)/dt_data.get('revenue_eur_billion', 1)*100:.0f}%",
            f"• EBITDA利润率{target_data.get('ebitda_margin_pct', 'N/A')}%，低于德电{dt_data.get('ebitda_margin_pct', 0) - target_data.get('ebitda_margin_pct', 0):.1f}pp",
            f"• 需要提升运营效率，向40%利润率目标努力",
        ]

        self._add_text_box(
            slide,
            Inches(0.7), Inches(5.95), Inches(12), Inches(1.1),
            "\n".join(findings),
            font_size=11,
            font_color=self.style.text_color,
        )

    def _add_5g_coverage_slide(
        self,
        financial_data: dict,
        target_operator: str,
    ):
        """Add 5G coverage comparison slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "5G 网络覆盖对比", "5G NETWORK COVERAGE")

        operators = list(financial_data.keys())
        coverage = [financial_data[op].get("5g_coverage_pct", 0) for op in operators]

        # 5G coverage donut charts
        chart_path = self.chart_gen.create_5g_coverage_comparison(
            operators=operators,
            coverage_pct=coverage,
            target_operator=target_operator,
            title="5G人口覆盖率 (%)",
            filename="5g_coverage.png",
        )

        slide.shapes.add_picture(
            chart_path,
            Inches(0.5), Inches(1.5),
            width=Inches(12.333), height=Inches(3.5),
        )

        # Analysis section
        self._add_shape(
            slide,
            Inches(0.5), Inches(5.3), Inches(12.333), Inches(1.9),
            (250, 250, 250),
        )

        self._add_text_box(
            slide,
            Inches(0.7), Inches(5.45), Inches(12), Inches(0.35),
            "5G 网络分析与建议",
            font_size=14,
            font_color=self.style.primary_color,
            bold=True,
        )

        target_5g = financial_data.get(target_operator, {}).get("5g_coverage_pct", 0)
        dt_5g = financial_data.get("Deutsche Telekom", {}).get("5g_coverage_pct", 0)
        o2_5g = financial_data.get("Telefónica O2 Germany", {}).get("5g_coverage_pct", 0)

        analysis = (
            f"• {target_operator} 5G覆盖{target_5g}%，落后德电{dt_5g - target_5g}pp，落后O2 {o2_5g - target_5g}pp\n"
            f"• 5G覆盖差距可能导致高端用户流失，需加速网络建设\n"
            f"• 建议：Q4达到95%覆盖，追平第一梯队水平\n"
            f"• 重点：城市核心区、工业园区、交通枢纽优先覆盖"
        )

        self._add_text_box(
            slide,
            Inches(0.7), Inches(5.9), Inches(12), Inches(1.2),
            analysis,
            font_size=11,
            font_color=self.style.text_color,
        )

    def _add_competitive_radar_slide(
        self,
        competitive_scores: dict,
        target_operator: str,
    ):
        """Add competitive radar chart slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "竞争力雷达图", "COMPETITIVE ANALYSIS")

        # Get dimensions and scores
        if target_operator in competitive_scores:
            dimensions = list(competitive_scores[target_operator].keys())
        else:
            dimensions = ["Network", "Service", "Price", "Brand", "Innovation"]

        scores = {}
        for op, op_scores in competitive_scores.items():
            scores[op] = [op_scores.get(dim, 50) for dim in dimensions]

        # Generate radar chart
        chart_path = self.chart_gen.create_competitive_radar_chart(
            dimensions=dimensions,
            scores=scores,
            target_operator=target_operator,
            title="多维度竞争力对比",
            filename="competitive_radar.png",
        )

        slide.shapes.add_picture(
            chart_path,
            Inches(0.3), Inches(1.3),
            width=Inches(7), height=Inches(5.5),
        )

        # Analysis panel
        self._add_text_box(
            slide,
            Inches(7.5), Inches(1.5), Inches(5.3), Inches(0.4),
            f"{target_operator} 竞争力分析",
            font_size=14,
            font_color=self.style.text_color,
            bold=True,
        )

        target_scores = competitive_scores.get(target_operator, {})
        avg_score = sum(target_scores.values()) / len(target_scores) if target_scores else 0

        strengths = [k for k, v in target_scores.items() if v >= avg_score + 5]
        weaknesses = [k for k, v in target_scores.items() if v <= avg_score - 5]

        # Strengths
        self._add_shape(
            slide,
            Inches(7.5), Inches(2.0), Inches(5.3), Inches(1.8),
            (240, 255, 240),
        )
        self._add_text_box(
            slide,
            Inches(7.6), Inches(2.1), Inches(5.1), Inches(0.3),
            "优势领域 Strengths",
            font_size=11,
            font_color=(34, 139, 34),
            bold=True,
        )
        strength_text = "\n".join([f"• {s}: {target_scores.get(s, 0)}分" for s in strengths[:4]])
        self._add_text_box(
            slide,
            Inches(7.6), Inches(2.45), Inches(5.1), Inches(1.3),
            strength_text or "• 暂无明显优势领域",
            font_size=11,
            font_color=self.style.text_color,
        )

        # Weaknesses
        self._add_shape(
            slide,
            Inches(7.5), Inches(4.0), Inches(5.3), Inches(1.8),
            (255, 240, 240),
        )
        self._add_text_box(
            slide,
            Inches(7.6), Inches(4.1), Inches(5.1), Inches(0.3),
            "改进领域 Weaknesses",
            font_size=11,
            font_color=self.style.primary_color,
            bold=True,
        )
        weakness_text = "\n".join([f"• {w}: {target_scores.get(w, 0)}分" for w in weaknesses[:4]])
        self._add_text_box(
            slide,
            Inches(7.6), Inches(4.45), Inches(5.1), Inches(1.3),
            weakness_text or "• 暂无明显劣势领域",
            font_size=11,
            font_color=self.style.text_color,
        )

        # Score summary
        self._add_text_box(
            slide,
            Inches(7.5), Inches(6.0), Inches(5.3), Inches(0.5),
            f"综合竞争力评分: {avg_score:.1f}/100",
            font_size=14,
            font_color=self.style.primary_color,
            bold=True,
        )

    def _add_insight_slide(
        self,
        insight: InsightResult,
        financial_data: dict = None,
        target_operator: str = None,
    ):
        """Add Five Looks insight slide with metrics."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, insight.title, insight.category.upper())

        # Metrics boxes
        if insight.metrics:
            x_pos = Inches(0.5)
            metrics_items = list(insight.metrics.items())[:4]

            for key, value in metrics_items:
                self._add_shape(
                    slide, x_pos, Inches(1.5), Inches(2.9), Inches(0.9),
                    (248, 248, 248),
                )
                self._add_shape(
                    slide, x_pos, Inches(1.5), Inches(2.9), Inches(0.05),
                    self.style.primary_color,
                )

                label = str(key).replace("_", " ").title()[:25]
                self._add_text_box(
                    slide,
                    x_pos + Inches(0.1), Inches(1.6), Inches(2.7), Inches(0.25),
                    label,
                    font_size=9,
                    font_color=self.style.light_text_color,
                )
                self._add_text_box(
                    slide,
                    x_pos + Inches(0.1), Inches(1.85), Inches(2.7), Inches(0.45),
                    str(value),
                    font_size=18,
                    font_color=self.style.primary_color,
                    bold=True,
                )
                x_pos += Inches(3.1)

        # Findings column
        self._add_text_box(
            slide,
            Inches(0.5), Inches(2.7), Inches(6), Inches(0.35),
            "洞察发现 FINDINGS",
            font_size=12,
            font_color=self.style.text_color,
            bold=True,
        )

        self._add_shape(
            slide,
            Inches(0.5), Inches(3.1), Inches(6), Inches(3.6),
            (252, 252, 252),
        )

        findings_text = ""
        for i, finding in enumerate(insight.findings[:7], 1):
            text = finding[:75] + "..." if len(finding) > 75 else finding
            findings_text += f"{i}. {text}\n"

        self._add_text_box(
            slide,
            Inches(0.7), Inches(3.25), Inches(5.6), Inches(3.3),
            findings_text,
            font_size=10,
            font_color=self.style.text_color,
        )

        # Recommendations column
        self._add_text_box(
            slide,
            Inches(6.8), Inches(2.7), Inches(6), Inches(0.35),
            "行动建议 RECOMMENDATIONS",
            font_size=12,
            font_color=self.style.text_color,
            bold=True,
        )

        self._add_shape(
            slide,
            Inches(6.8), Inches(3.1), Inches(5.8), Inches(3.6),
            (255, 250, 245),
        )

        rec_text = ""
        for i, rec in enumerate(insight.recommendations[:5], 1):
            text = rec[:70] + "..." if len(rec) > 70 else rec
            rec_text += f"→ {text}\n\n"

        self._add_text_box(
            slide,
            Inches(7.0), Inches(3.25), Inches(5.4), Inches(3.3),
            rec_text,
            font_size=10,
            font_color=self.style.text_color,
        )

    def _add_gap_analysis_slide(
        self,
        competitive_scores: dict,
        target_operator: str,
    ):
        """Add gap analysis chart slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "差距分析", "GAP ANALYSIS")

        # Find leader (Deutsche Telekom)
        leader = "Deutsche Telekom"
        if leader not in competitive_scores:
            leader = list(competitive_scores.keys())[0]

        if target_operator in competitive_scores and leader in competitive_scores:
            dimensions = list(competitive_scores[target_operator].keys())[:8]
            target_scores = [competitive_scores[target_operator].get(d, 50) for d in dimensions]
            leader_scores = [competitive_scores[leader].get(d, 50) for d in dimensions]

            chart_path = self.chart_gen.create_gap_analysis_chart(
                dimensions=dimensions,
                target_scores=target_scores,
                leader_scores=leader_scores,
                target_name=target_operator,
                leader_name=leader,
                title=f"与市场领导者 {leader} 差距分析",
                filename="gap_analysis.png",
            )

            slide.shapes.add_picture(
                chart_path,
                Inches(0.5), Inches(1.5),
                width=Inches(8), height=Inches(4.5),
            )

        # Action items
        self._add_text_box(
            slide,
            Inches(9), Inches(1.5), Inches(3.8), Inches(0.4),
            "关键差距与行动",
            font_size=14,
            font_color=self.style.text_color,
            bold=True,
        )

        gaps_info = [
            ("5G覆盖", "-12pp", "加速网络建设"),
            ("网络质量", "-17分", "优化网络性能"),
            ("客户服务", "-14分", "提升服务水平"),
            ("企业方案", "-12分", "强化B2B能力"),
        ]

        y_pos = Inches(2.0)
        for dim, gap, action in gaps_info:
            self._add_shape(
                slide,
                Inches(9), y_pos, Inches(3.8), Inches(1.0),
                (255, 245, 245),
            )
            self._add_text_box(
                slide,
                Inches(9.1), y_pos + Inches(0.1), Inches(2), Inches(0.3),
                dim,
                font_size=11,
                font_color=self.style.text_color,
                bold=True,
            )
            self._add_text_box(
                slide,
                Inches(11), y_pos + Inches(0.1), Inches(1.5), Inches(0.3),
                gap,
                font_size=12,
                font_color=self.style.primary_color,
                bold=True,
            )
            self._add_text_box(
                slide,
                Inches(9.1), y_pos + Inches(0.45), Inches(3.5), Inches(0.4),
                f"→ {action}",
                font_size=10,
                font_color=self.style.light_text_color,
            )
            y_pos += Inches(1.1)

        # Summary
        self._add_shape(
            slide,
            Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.9),
            self.style.primary_color,
        )
        self._add_text_box(
            slide,
            Inches(0.7), Inches(6.45), Inches(12), Inches(0.6),
            "核心策略: 聚焦5G网络和客户服务两大短板，缩小与领导者差距",
            font_size=14,
            font_color=(255, 255, 255),
            bold=True,
        )

    def _add_section_divider(self, title: str, subtitle: str, number: str):
        """Add section divider slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        # Full red background
        self._add_shape(
            slide, 0, 0, Inches(13.333), Inches(7.5),
            self.style.primary_color
        )

        # Decorative elements
        self._add_shape(
            slide, Inches(10), 0, Inches(3.333), Inches(7.5),
            (170, 0, 10),
        )

        # Section number
        self._add_text_box(
            slide,
            Inches(0.5), Inches(2), Inches(2), Inches(1.2),
            number,
            font_size=80,
            font_color=(255, 255, 255),
            bold=True,
        )

        # Line
        self._add_shape(
            slide,
            Inches(0.5), Inches(3.4), Inches(5), Inches(0.03),
            (255, 255, 255),
        )

        # Title
        self._add_text_box(
            slide,
            Inches(0.5), Inches(3.7), Inches(9), Inches(1),
            title,
            font_size=52,
            font_color=(255, 255, 255),
            bold=True,
        )

        # Subtitle
        self._add_text_box(
            slide,
            Inches(0.5), Inches(4.8), Inches(9), Inches(0.6),
            subtitle,
            font_size=24,
            font_color=(255, 200, 200),
        )

    def _add_strategy_slide_enhanced(self, decision: StrategyResult):
        """Add enhanced strategy slide with priority visualization."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, decision.title, decision.decision_type.upper())

        # Summary section
        self._add_shape(
            slide,
            Inches(0.5), Inches(1.4), Inches(12.333), Inches(1.3),
            (245, 250, 255),
        )

        summary_lines = decision.summary.split("\n")[:3]
        summary_text = "\n".join(summary_lines)
        self._add_text_box(
            slide,
            Inches(0.7), Inches(1.5), Inches(12), Inches(1.1),
            summary_text,
            font_size=11,
            font_color=self.style.text_color,
        )

        # Strategy items with visual priority
        y_pos = Inches(2.9)

        for item in decision.items[:6]:
            # Priority colors
            priority_colors = {
                "P0": (self.style.primary_color, (255, 240, 240)),
                "P1": ((255, 165, 0), (255, 250, 240)),
                "P2": ((34, 139, 34), (240, 255, 240)),
            }
            badge_color, bg_color = priority_colors.get(
                item.priority,
                (self.style.primary_color, (248, 248, 248))
            )

            # Row background
            self._add_shape(
                slide,
                Inches(0.5), y_pos, Inches(12.333), Inches(0.7),
                bg_color,
            )

            # Priority badge
            self._add_shape(
                slide,
                Inches(0.6), y_pos + Inches(0.15), Inches(0.45), Inches(0.4),
                badge_color,
            )
            self._add_text_box(
                slide,
                Inches(0.6), y_pos + Inches(0.22), Inches(0.45), Inches(0.25),
                item.priority,
                font_size=10,
                font_color=(255, 255, 255),
                bold=True,
                align="center",
            )

            # Name
            self._add_text_box(
                slide,
                Inches(1.2), y_pos + Inches(0.1), Inches(4), Inches(0.35),
                item.name,
                font_size=13,
                font_color=self.style.text_color,
                bold=True,
            )

            # Description
            desc = item.description[:80] + "..." if len(item.description) > 80 else item.description
            self._add_text_box(
                slide,
                Inches(1.2), y_pos + Inches(0.4), Inches(5), Inches(0.3),
                desc,
                font_size=9,
                font_color=self.style.light_text_color,
            )

            # KPIs
            if item.kpis:
                kpi_text = " | ".join(item.kpis[:2])
                self._add_text_box(
                    slide,
                    Inches(6.5), y_pos + Inches(0.2), Inches(6), Inches(0.4),
                    f"KPI: {kpi_text}",
                    font_size=9,
                    font_color=self.style.primary_color,
                )

            y_pos += Inches(0.75)

    def _add_timeline_slide(self, execution: StrategyResult):
        """Add execution timeline slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "执行时间线", "EXECUTION TIMELINE")

        # Get milestones
        milestones = [
            {
                "date": item.timeline,
                "name": item.name,
                "priority": item.priority,
            }
            for item in execution.items
            if item.category == "milestone" and item.timeline
        ][:6]

        if milestones:
            chart_path = self.chart_gen.create_timeline_chart(
                milestones=milestones,
                title="关键里程碑 Key Milestones",
                filename="timeline.png",
            )

            slide.shapes.add_picture(
                chart_path,
                Inches(0.5), Inches(1.5),
                width=Inches(12.333), height=Inches(3),
            )

        # Governance section
        self._add_text_box(
            slide,
            Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.4),
            "治理机制 GOVERNANCE",
            font_size=14,
            font_color=self.style.text_color,
            bold=True,
        )

        governance_items = [item for item in execution.items if item.category == "governance"]

        x_pos = Inches(0.5)
        for item in governance_items[:3]:
            self._add_shape(
                slide,
                x_pos, Inches(5.3), Inches(4), Inches(1.5),
                (248, 248, 248),
            )
            self._add_shape(
                slide,
                x_pos, Inches(5.3), Inches(4), Inches(0.05),
                self.style.primary_color,
            )
            self._add_text_box(
                slide,
                x_pos + Inches(0.15), Inches(5.45), Inches(3.7), Inches(0.4),
                item.name,
                font_size=12,
                font_color=self.style.text_color,
                bold=True,
            )
            desc = item.description[:60] + "..." if len(item.description) > 60 else item.description
            self._add_text_box(
                slide,
                x_pos + Inches(0.15), Inches(5.9), Inches(3.7), Inches(0.8),
                desc,
                font_size=10,
                font_color=self.style.light_text_color,
            )
            x_pos += Inches(4.2)

    def _add_kpi_dashboard_slide(self, three_decisions: dict[str, StrategyResult]):
        """Add KPI dashboard slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "KPI仪表盘", "KEY PERFORMANCE INDICATORS")

        # Collect KPIs
        all_kpis = []
        for key, decision in three_decisions.items():
            for item in decision.items:
                if item.kpis and item.priority in ["P0", "P1"]:
                    for kpi in item.kpis[:2]:
                        all_kpis.append({
                            "priority": item.priority,
                            "area": item.name[:20],
                            "kpi": kpi,
                        })

        # Display KPIs in grid
        x_positions = [Inches(0.5), Inches(4.4), Inches(8.3)]
        y_pos = Inches(1.5)
        col = 0

        for i, kpi_item in enumerate(all_kpis[:12]):
            x = x_positions[col]

            # KPI card
            is_p0 = kpi_item["priority"] == "P0"
            bg_color = (255, 245, 245) if is_p0 else (245, 250, 255)
            accent_color = self.style.primary_color if is_p0 else (50, 100, 150)

            self._add_shape(
                slide,
                x, y_pos, Inches(3.7), Inches(1.2),
                bg_color,
            )
            self._add_shape(
                slide,
                x, y_pos, Inches(0.05), Inches(1.2),
                accent_color,
            )

            # Priority
            self._add_text_box(
                slide,
                x + Inches(0.15), y_pos + Inches(0.1), Inches(0.4), Inches(0.25),
                kpi_item["priority"],
                font_size=9,
                font_color=accent_color,
                bold=True,
            )

            # Area
            self._add_text_box(
                slide,
                x + Inches(0.6), y_pos + Inches(0.1), Inches(3), Inches(0.25),
                kpi_item["area"],
                font_size=9,
                font_color=self.style.light_text_color,
            )

            # KPI
            self._add_text_box(
                slide,
                x + Inches(0.15), y_pos + Inches(0.45), Inches(3.4), Inches(0.65),
                kpi_item["kpi"],
                font_size=12,
                font_color=self.style.text_color,
                bold=True,
            )

            col += 1
            if col >= 3:
                col = 0
                y_pos += Inches(1.35)

    def _add_conclusion_slide(
        self,
        five_looks: dict[str, InsightResult],
        three_decisions: dict[str, StrategyResult],
        target_operator: str,
    ):
        """Add conclusion slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "总结与建议", "SUMMARY & RECOMMENDATIONS")

        # Key strategies
        self._add_text_box(
            slide,
            Inches(0.5), Inches(1.5), Inches(6), Inches(0.4),
            f"{target_operator} 战略重点",
            font_size=16,
            font_color=self.style.primary_color,
            bold=True,
        )

        strategy = three_decisions.get("strategy")
        if strategy:
            y_pos = Inches(2.0)
            for item in strategy.items[:4]:
                if item.priority == "P0":
                    self._add_shape(
                        slide,
                        Inches(0.5), y_pos, Inches(6), Inches(0.7),
                        (255, 248, 248),
                    )
                    self._add_text_box(
                        slide,
                        Inches(0.6), y_pos + Inches(0.1), Inches(5.8), Inches(0.3),
                        f"【{item.priority}】{item.name}",
                        font_size=12,
                        font_color=self.style.text_color,
                        bold=True,
                    )
                    desc = item.description[:50] + "..."
                    self._add_text_box(
                        slide,
                        Inches(0.6), y_pos + Inches(0.4), Inches(5.8), Inches(0.25),
                        desc,
                        font_size=10,
                        font_color=self.style.light_text_color,
                    )
                    y_pos += Inches(0.8)

        # Key milestones
        self._add_text_box(
            slide,
            Inches(6.8), Inches(1.5), Inches(6), Inches(0.4),
            "关键里程碑",
            font_size=16,
            font_color=self.style.primary_color,
            bold=True,
        )

        execution = three_decisions.get("execution")
        if execution:
            milestones = [item for item in execution.items if item.category == "milestone"]
            y_pos = Inches(2.0)
            for item in milestones[:5]:
                self._add_text_box(
                    slide,
                    Inches(6.8), y_pos, Inches(5.5), Inches(0.4),
                    f"• {item.timeline}: {item.name}",
                    font_size=11,
                    font_color=self.style.text_color,
                )
                y_pos += Inches(0.45)

        # Next steps banner
        self._add_shape(
            slide,
            Inches(0.5), Inches(5.3), Inches(12.333), Inches(1.9),
            self.style.primary_color,
        )

        self._add_text_box(
            slide,
            Inches(0.7), Inches(5.5), Inches(12), Inches(0.4),
            "下一步行动 NEXT STEPS",
            font_size=16,
            font_color=(255, 255, 255),
            bold=True,
        )

        next_steps = (
            "1. 成立战略执行委员会，CEO挂帅   "
            "2. 启动5G网络加速项目，确保Q4达95%覆盖   "
            "3. 组建企业解决方案中心   "
            "4. 建立双周KPI追踪机制"
        )
        self._add_text_box(
            slide,
            Inches(0.7), Inches(6.0), Inches(12), Inches(1),
            next_steps,
            font_size=12,
            font_color=(255, 255, 255),
        )

    def _add_closing_slide(self):
        """Add closing slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        # Full red background
        self._add_shape(
            slide, 0, 0, Inches(13.333), Inches(7.5),
            self.style.primary_color
        )

        # Decorative element
        self._add_shape(
            slide, Inches(10), Inches(0), Inches(3.333), Inches(7.5),
            (170, 0, 10),
        )

        # Thank you
        self._add_text_box(
            slide,
            Inches(0.5), Inches(2.3), Inches(9), Inches(1.2),
            "谢谢",
            font_size=80,
            font_color=(255, 255, 255),
            bold=True,
        )

        self._add_text_box(
            slide,
            Inches(0.5), Inches(3.6), Inches(9), Inches(0.7),
            "THANK YOU",
            font_size=36,
            font_color=(255, 200, 200),
        )

        # Footer
        self._add_text_box(
            slide,
            Inches(0.5), Inches(6.3), Inches(9), Inches(0.5),
            f"BLM Strategic Analysis | Generated {datetime.now().strftime('%Y-%m-%d')}",
            font_size=12,
            font_color=(255, 200, 200),
        )

        self._add_text_box(
            slide,
            Inches(0.5), Inches(6.7), Inches(9), Inches(0.4),
            "Powered by BLM Analysis Framework",
            font_size=10,
            font_color=(255, 180, 180),
        )

    # =========================================================================
    # Helper Methods
    # =========================================================================

    def _add_header(self, slide, title: str, subtitle: str = ""):
        """Add standard header to slide."""
        # Red accent bar
        self._add_shape(
            slide, 0, 0, Inches(13.333), Inches(0.1),
            self.style.primary_color
        )

        # Title
        self._add_text_box(
            slide,
            Inches(0.5), Inches(0.25), Inches(10), Inches(0.55),
            title,
            font_size=24,
            font_color=self.style.text_color,
            bold=True,
        )

        # Subtitle
        if subtitle:
            self._add_text_box(
                slide,
                Inches(0.5), Inches(0.75), Inches(10), Inches(0.35),
                subtitle,
                font_size=11,
                font_color=self.style.primary_color,
                bold=True,
            )

        # Page number
        self._add_shape(
            slide,
            Inches(12.5), Inches(7), Inches(0.6), Inches(0.35),
            self.style.primary_color,
        )
        self._add_text_box(
            slide,
            Inches(12.5), Inches(7.02), Inches(0.6), Inches(0.3),
            str(self.slide_num),
            font_size=10,
            font_color=(255, 255, 255),
            bold=True,
            align="center",
        )

    def _add_shape(
        self,
        slide,
        left,
        top,
        width,
        height,
        fill_color: tuple,
        shape_type=MSO_SHAPE.RECTANGLE,
    ):
        """Add a shape to the slide."""
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_color)
        shape.line.fill.background()
        return shape

    def _add_text_box(
        self,
        slide,
        left,
        top,
        width,
        height,
        text: str,
        font_size: int = 14,
        font_color: tuple = (0, 0, 0),
        bold: bool = False,
        align: str = "left",
    ):
        """Add a text box to the slide."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*font_color)
        p.font.bold = bold
        p.font.name = self.style.body_font

        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

        return txBox

    def _sanitize_name(self, name: str) -> str:
        """Sanitize name for filename."""
        return name.lower().replace(" ", "_").replace("/", "_")


def generate_enhanced_blm_ppt(
    five_looks: dict[str, InsightResult],
    three_decisions: dict[str, StrategyResult],
    target_operator: str,
    competitors: list[str] = None,
    financial_data: dict = None,
    competitive_scores: dict = None,
    style: str = "huawei",
    output_dir: str = None,
) -> str:
    """Convenience function to generate enhanced BLM PPT.

    Args:
        five_looks: Five Looks analysis results
        three_decisions: Three Decisions strategy results
        target_operator: Target operator name
        competitors: List of competitor names
        financial_data: Raw financial data for charts
        competitive_scores: Competitive dimension scores
        style: PPT style ("huawei" or "vodafone")
        output_dir: Output directory

    Returns:
        Path to generated PPT file
    """
    generator = BLMPPTGeneratorEnhanced(style=style, output_dir=output_dir)
    return generator.generate(
        five_looks=five_looks,
        three_decisions=three_decisions,
        target_operator=target_operator,
        competitors=competitors,
        financial_data=financial_data,
        competitive_scores=competitive_scores,
    )
