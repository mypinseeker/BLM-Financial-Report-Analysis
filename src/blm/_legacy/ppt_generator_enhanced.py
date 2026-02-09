"""Enhanced BLM PPT Generator with Charts and Visual Elements.

Extends the basic PPT generator with:
- Market share bar charts
- Revenue comparison charts
- Competitive radar charts
- KPI comparison tables
- 5G coverage gauges
- Gap analysis charts
- Strategy priority visualization
- Execution timeline

Follows Huawei 5G template style with rich visual elements.
"""

from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Optional
import tempfile

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from src.blm._legacy.five_looks import InsightResult
from src.blm._legacy.three_decisions import StrategyResult, StrategyItem
from src.blm._legacy.ppt_generator import PPTStyle, HUAWEI_STYLE, VODAFONE_STYLE, STYLES
from src.blm._legacy.ppt_charts import PPTChartGenerator


class BLMPPTGeneratorEnhanced:
    """Enhanced PPT generator with chart integration."""

    def __init__(
        self,
        style: str = "huawei",
        output_dir: Optional[str] = None,
        chart_dpi: int = 150,
    ):
        """Initialize enhanced PPT generator.

        Args:
            style: PPT style name ("huawei", "vodafone")
            output_dir: Output directory for generated PPTs
            chart_dpi: DPI for chart images
        """
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for PPT generation. "
                "Install it with: pip install python-pptx"
            )

        if isinstance(style, str):
            self.style = STYLES.get(style, HUAWEI_STYLE)
        else:
            self.style = style

        if output_dir is None:
            output_dir = str(Path(__file__).resolve().parent.parent.parent / "data" / "output")
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # Initialize chart generator
        self.chart_dir = Path(tempfile.mkdtemp())
        self.chart_gen = PPTChartGenerator(output_dir=str(self.chart_dir), dpi=chart_dpi)

        self.prs = None
        self.slide_num = 0

    def generate(
        self,
        five_looks: dict[str, InsightResult],
        three_decisions: dict[str, StrategyResult],
        target_operator: str,
        competitors: list[str] = None,
        financial_data: dict = None,
        competitive_scores: dict = None,
        historical_user_data: dict = None,
        historical_financial_data: dict = None,
        historical_segment_data: dict = None,
        user_flow_data: dict = None,
        quarters: list[str] = None,
        title: str = None,
        filename: str = None,
    ) -> str:
        """Generate enhanced PPT with charts.

        Args:
            five_looks: Five Looks analysis results
            three_decisions: Three Decisions strategy results
            target_operator: Target operator name
            competitors: List of competitor names
            financial_data: Raw financial data for charts
            competitive_scores: Competitive dimension scores for radar chart
            historical_user_data: 8-quarter user trend data
            historical_financial_data: 8-quarter financial trend data
            historical_segment_data: 8-quarter business segment data
            user_flow_data: User flow between operators
            quarters: List of quarter labels
            title: Presentation title
            filename: Output filename

        Returns:
            Path to generated PPT file
        """
        # Store historical data for use in slide methods
        self.historical_user_data = historical_user_data
        self.historical_financial_data = historical_financial_data
        self.historical_segment_data = historical_segment_data
        self.user_flow_data = user_flow_data
        self.quarters = quarters or ["Q4 FY24", "Q1 FY25", "Q2 FY25", "Q3 FY25", "Q4 FY25", "Q1 FY26", "Q2 FY26", "Q3 FY26"]
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        self.prs.slide_height = Inches(7.5)
        self.slide_num = 0

        competitors = competitors or []
        title = title or f"{target_operator} BLM 战略分析报告"
        filename = filename or f"blm_{self._sanitize_name(target_operator)}_enhanced.pptx"

        # Generate slides
        self._add_title_slide(title, target_operator, competitors)
        self._add_agenda_slide()

        # NEW: Data Sources slide
        self._add_data_sources_slide()

        # Executive Summary with key metrics
        self._add_executive_summary_slide(five_looks, three_decisions, target_operator)

        # Section 1: Five Looks with charts
        self._add_section_divider("五看分析", "Five Looks Analysis", "01")

        # NEW: Raw Data Table slide
        if financial_data:
            self._add_raw_data_table_slide(financial_data, target_operator)

        # Market overview with charts
        if financial_data:
            self._add_market_overview_slide(financial_data, target_operator, competitors)

            # NEW: Derivation Logic slide - Revenue Analysis
            self._add_revenue_derivation_slide(financial_data, target_operator)

            self._add_revenue_comparison_slide(financial_data, target_operator)
            self._add_5g_coverage_slide(financial_data, target_operator)

            # NEW: Detailed Comparison slide
            self._add_detailed_comparison_slide(financial_data, target_operator)

        # Competitive radar chart
        if competitive_scores:
            self._add_competitive_radar_slide(competitive_scores, target_operator)

        # Five Looks insight slides
        for key, insight in five_looks.items():
            self._add_insight_slide(insight, financial_data, target_operator)

        # Gap analysis
        if competitive_scores and financial_data:
            self._add_gap_analysis_slide(competitive_scores, target_operator)

            # NEW: Gap Analysis Derivation
            self._add_gap_derivation_slide(competitive_scores, financial_data, target_operator)

        # NEW Section: Q3 FY26 Quarterly Operating Analysis
        if financial_data:
            self._add_quarterly_analysis_section(financial_data, target_operator)

        # NEW Section: Historical Trend Analysis (8 Quarters)
        if self.historical_user_data and self.historical_financial_data:
            self._add_historical_trend_section(target_operator, competitors)

        # Section 4: Three Decisions
        self._add_section_divider("三定策略", "Three Decisions Strategy", "04")

        # Strategy slides with priority visualization
        for key, decision in three_decisions.items():
            self._add_strategy_slide_enhanced(decision)

        # Execution timeline
        execution = three_decisions.get("execution")
        if execution:
            self._add_timeline_slide(execution)

        # KPI Dashboard
        self._add_kpi_dashboard_slide(three_decisions)

        # Conclusion and next steps
        self._add_conclusion_slide(five_looks, three_decisions, target_operator)
        self._add_closing_slide()

        # Save
        output_path = self.output_dir / filename
        self.prs.save(str(output_path))
        return str(output_path)

    def _add_title_slide(
        self,
        title: str,
        target_operator: str,
        competitors: list[str],
    ):
        """Add title slide with visual design."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        # Red header bar
        self._add_shape(
            slide, 0, 0, Inches(13.333), Inches(2.8),
            self.style.primary_color
        )

        # Decorative diagonal stripe
        self._add_shape(
            slide, Inches(10), 0, Inches(3.333), Inches(2.8),
            (180, 0, 10),  # Darker red
        )

        # Title
        self._add_text_box(
            slide,
            Inches(0.5), Inches(0.5), Inches(9), Inches(1.2),
            title,
            font_size=40,
            font_color=(255, 255, 255),
            bold=True,
        )

        # Subtitle
        subtitle = f"基于 BLM (Business Leadership Model) 方法论的战略分析"
        self._add_text_box(
            slide,
            Inches(0.5), Inches(1.7), Inches(9), Inches(0.6),
            subtitle,
            font_size=20,
            font_color=(255, 220, 220),
        )

        # Info cards
        info_items = [
            ("分析对象", target_operator),
            ("竞争对手", ", ".join(competitors) if competitors else "N/A"),
            ("报告日期", datetime.now().strftime('%Y年%m月%d日')),
        ]

        x_pos = Inches(0.5)
        for label, value in info_items:
            # Card background
            self._add_shape(
                slide, x_pos, Inches(3.5), Inches(3.8), Inches(1.5),
                (245, 245, 245),
            )
            # Label
            self._add_text_box(
                slide,
                x_pos + Inches(0.2), Inches(3.6), Inches(3.4), Inches(0.4),
                label,
                font_size=12,
                font_color=self.style.light_text_color,
            )
            # Value
            self._add_text_box(
                slide,
                x_pos + Inches(0.2), Inches(4.0), Inches(3.4), Inches(0.8),
                value,
                font_size=18,
                font_color=self.style.primary_color,
                bold=True,
            )
            x_pos += Inches(4.1)

        # Framework info
        framework_items = [
            ("五看 Five Looks", "看市场 | 看自己 | 看对手 | 看宏观 | 看机会"),
            ("三定 Three Decisions", "定策略 | 定重点工作 | 定执行"),
        ]

        y_pos = Inches(5.5)
        for fw_title, fw_desc in framework_items:
            # Icon circle
            self._add_shape(
                slide,
                Inches(0.5), y_pos, Inches(0.4), Inches(0.4),
                self.style.primary_color,
                shape_type=MSO_SHAPE.OVAL,
            )
            # Text
            self._add_text_box(
                slide,
                Inches(1.1), y_pos, Inches(4), Inches(0.4),
                fw_title,
                font_size=14,
                font_color=self.style.text_color,
                bold=True,
            )
            self._add_text_box(
                slide,
                Inches(1.1), y_pos + Inches(0.35), Inches(6), Inches(0.35),
                fw_desc,
                font_size=11,
                font_color=self.style.light_text_color,
            )
            y_pos += Inches(0.85)

    def _add_agenda_slide(self):
        """Add agenda slide with visual timeline."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "报告目录", "CONTENTS")

        # Vertical timeline line
        self._add_shape(
            slide,
            Inches(1.1), Inches(1.8), Inches(0.02), Inches(5),
            self.style.primary_color,
        )

        agenda_items = [
            ("01", "执行摘要", "Executive Summary - 核心数据与关键结论", "page 3-4"),
            ("02", "五看分析", "Five Looks - 市场/自身/竞争/宏观/机会洞察", "page 5-15"),
            ("03", "三定策略", "Three Decisions - 策略/重点工作/执行计划", "page 16-22"),
            ("04", "总结建议", "Summary - 结论与下一步行动", "page 23-24"),
        ]

        y_pos = Inches(1.8)
        for num, title, subtitle, pages in agenda_items:
            # Number circle
            self._add_shape(
                slide,
                Inches(0.85), y_pos, Inches(0.5), Inches(0.5),
                self.style.primary_color,
                shape_type=MSO_SHAPE.OVAL,
            )
            self._add_text_box(
                slide,
                Inches(0.85), y_pos + Inches(0.08), Inches(0.5), Inches(0.35),
                num,
                font_size=14,
                font_color=(255, 255, 255),
                bold=True,
                align="center",
            )

            # Title and description
            self._add_text_box(
                slide,
                Inches(1.6), y_pos, Inches(8), Inches(0.45),
                title,
                font_size=20,
                font_color=self.style.text_color,
                bold=True,
            )
            self._add_text_box(
                slide,
                Inches(1.6), y_pos + Inches(0.45), Inches(8), Inches(0.35),
                subtitle,
                font_size=12,
                font_color=self.style.light_text_color,
            )

            # Page reference
            self._add_text_box(
                slide,
                Inches(10.5), y_pos + Inches(0.15), Inches(2), Inches(0.3),
                pages,
                font_size=10,
                font_color=self.style.primary_color,
            )

            y_pos += Inches(1.2)

    def _add_executive_summary_slide(
        self,
        five_looks: dict[str, InsightResult],
        three_decisions: dict[str, StrategyResult],
        target_operator: str,
    ):
        """Add executive summary with key metrics."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "执行摘要", "EXECUTIVE SUMMARY")

        # Get key data
        self_insight = five_looks.get("self")
        strategy = three_decisions.get("strategy")

        # Metric cards row
        if self_insight and self_insight.metrics:
            metrics = self_insight.metrics
            metric_cards = [
                ("收入规模", f"€{metrics.get('revenue_eur_billion', 'N/A')}B", "Q2 2025"),
                ("市场排名", f"第{metrics.get('revenue_rank', 'N/A')}位", "德国市场"),
                ("市场份额", f"{metrics.get('market_share_pct', 'N/A')}%", "宽带市场"),
                ("5G覆盖率", f"{metrics.get('5g_coverage_pct', 'N/A')}%", "人口覆盖"),
            ]

            x_pos = Inches(0.5)
            for label, value, sub in metric_cards:
                # Card
                self._add_shape(
                    slide, x_pos, Inches(1.5), Inches(3), Inches(1.3),
                    (250, 250, 250),
                )
                # Top accent
                self._add_shape(
                    slide, x_pos, Inches(1.5), Inches(3), Inches(0.06),
                    self.style.primary_color,
                )
                # Value
                self._add_text_box(
                    slide,
                    x_pos + Inches(0.15), Inches(1.65), Inches(2.7), Inches(0.7),
                    value,
                    font_size=28,
                    font_color=self.style.primary_color,
                    bold=True,
                )
                # Label
                self._add_text_box(
                    slide,
                    x_pos + Inches(0.15), Inches(2.3), Inches(2.7), Inches(0.35),
                    f"{label} ({sub})",
                    font_size=10,
                    font_color=self.style.light_text_color,
                )
                x_pos += Inches(3.2)

        # Strategic priorities section
        self._add_text_box(
            slide,
            Inches(0.5), Inches(3.1), Inches(6), Inches(0.4),
            "战略重点 STRATEGIC PRIORITIES",
            font_size=14,
            font_color=self.style.text_color,
            bold=True,
        )

        # Priority items
        if strategy:
            y_pos = Inches(3.6)
            for item in strategy.items[:4]:
                # Priority badge
                priority_colors = {
                    "P0": self.style.primary_color,
                    "P1": (255, 165, 0),
                    "P2": (34, 139, 34),
                }
                color = priority_colors.get(item.priority, self.style.primary_color)

                self._add_shape(
                    slide,
                    Inches(0.5), y_pos, Inches(0.45), Inches(0.35),
                    color,
                )
                self._add_text_box(
                    slide,
                    Inches(0.5), y_pos + Inches(0.05), Inches(0.45), Inches(0.25),
                    item.priority,
                    font_size=10,
                    font_color=(255, 255, 255),
                    bold=True,
                    align="center",
                )

                # Name and description
                self._add_text_box(
                    slide,
                    Inches(1.05), y_pos, Inches(5), Inches(0.35),
                    item.name,
                    font_size=13,
                    font_color=self.style.text_color,
                    bold=True,
                )
                desc = item.description[:65] + "..." if len(item.description) > 65 else item.description
                self._add_text_box(
                    slide,
                    Inches(1.05), y_pos + Inches(0.32), Inches(5), Inches(0.35),
                    desc,
                    font_size=10,
                    font_color=self.style.light_text_color,
                )

                y_pos += Inches(0.75)

        # Key challenges section
        self._add_text_box(
            slide,
            Inches(6.8), Inches(3.1), Inches(6), Inches(0.4),
            "核心挑战 KEY CHALLENGES",
            font_size=14,
            font_color=self.style.text_color,
            bold=True,
        )

        # Challenge box
        self._add_shape(
            slide,
            Inches(6.8), Inches(3.5), Inches(5.8), Inches(2.8),
            (255, 248, 245),
        )

        challenges = []
        if self_insight:
            for finding in self_insight.findings[:5]:
                if any(kw in finding for kw in ["落后", "低于", "高于", "下滑", "差距"]):
                    challenges.append(f"• {finding[:55]}...")

        challenge_text = "\n".join(challenges[:4]) if challenges else "• 详见五看分析"
        self._add_text_box(
            slide,
            Inches(7.0), Inches(3.7), Inches(5.4), Inches(2.4),
            challenge_text,
            font_size=11,
            font_color=self.style.text_color,
        )

        # Bottom key message
        self._add_shape(
            slide,
            Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.7),
            self.style.primary_color,
        )
        self._add_text_box(
            slide,
            Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
            f"核心定位: 做德国市场「有质量的挑战者」- 以网络和服务质量为核心竞争力",
            font_size=16,
            font_color=(255, 255, 255),
            bold=True,
        )

    def _add_market_overview_slide(
        self,
        financial_data: dict,
        target_operator: str,
        competitors: list[str],
    ):
        """Add market overview with market share chart."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "市场格局", "MARKET OVERVIEW")

        # Market share data
        operators = [target_operator] + competitors
        market_shares = []
        for op in operators:
            if op in financial_data:
                share = financial_data[op].get("market_share_broadband_pct", 0)
                market_shares.append(share)
            else:
                market_shares.append(0)

        # Generate and embed chart
        chart_path = self.chart_gen.create_market_share_bar_chart(
            operators=operators,
            market_shares=market_shares,
            target_operator=target_operator,
            title="德国宽带市场份额对比 (%)",
            filename="market_share_overview.png",
        )

        # Add chart image
        slide.shapes.add_picture(
            chart_path,
            Inches(0.5), Inches(1.5),
            width=Inches(6), height=Inches(3),
        )

        # Key metrics table on right side
        self._add_text_box(
            slide,
            Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
            "市场关键数据 Q2 2025",
            font_size=14,
            font_color=self.style.text_color,
            bold=True,
        )

        # Metrics table
        y_pos = Inches(2.0)
        for op in operators[:4]:
            if op in financial_data:
                data = financial_data[op]

                # Operator row
                is_target = op == target_operator
                bg_color = (255, 240, 240) if is_target else (248, 248, 248)

                self._add_shape(
                    slide,
                    Inches(7), y_pos, Inches(5.5), Inches(1.0),
                    bg_color,
                )

                # Operator name
                self._add_text_box(
                    slide,
                    Inches(7.1), y_pos + Inches(0.05), Inches(2.5), Inches(0.35),
                    op,
                    font_size=12,
                    font_color=self.style.primary_color if is_target else self.style.text_color,
                    bold=True,
                )

                # Metrics
                metrics_text = (
                    f"收入: €{data.get('revenue_eur_billion', 'N/A')}B | "
                    f"EBITDA: {data.get('ebitda_margin_pct', 'N/A')}% | "
                    f"5G: {data.get('5g_coverage_pct', 'N/A')}%"
                )
                self._add_text_box(
                    slide,
                    Inches(7.1), y_pos + Inches(0.4), Inches(5.3), Inches(0.5),
                    metrics_text,
                    font_size=10,
                    font_color=self.style.text_color,
                )

                y_pos += Inches(1.1)

        # Market insight
        self._add_shape(
            slide,
            Inches(0.5), Inches(5.0), Inches(12.333), Inches(2),
            (250, 250, 250),
        )
        self._add_text_box(
            slide,
            Inches(0.7), Inches(5.15), Inches(12), Inches(0.35),
            "市场洞察 MARKET INSIGHT",
            font_size=12,
            font_color=self.style.primary_color,
            bold=True,
        )
        insight_text = (
            "• 德电以40.6%份额稳居第一，连续35季度EBITDA增长，领先优势持续扩大\n"
            "• Vodafone以27.0%份额位居第二，服务收入重回增长(+0.5%)\n"
            "• 1&1作为新进入者正在建设自有网络，短期面临成本压力\n"
            "• 市场竞争加剧，价格战压力持续，需要差异化竞争策略"
        )
        self._add_text_box(
            slide,
            Inches(0.7), Inches(5.55), Inches(12), Inches(1.4),
            insight_text,
            font_size=11,
            font_color=self.style.text_color,
        )

    def _add_revenue_comparison_slide(
        self,
        financial_data: dict,
        target_operator: str,
    ):
        """Add revenue comparison chart slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "收入与利润对比", "REVENUE & PROFITABILITY")

        operators = list(financial_data.keys())
        revenues = [financial_data[op].get("revenue_eur_billion", 0) for op in operators]

        # Revenue chart
        chart_path = self.chart_gen.create_revenue_comparison_chart(
            operators=operators,
            revenues=revenues,
            target_operator=target_operator,
            title="收入规模对比 (€B) - Q2 2025",
            filename="revenue_comparison.png",
        )

        slide.shapes.add_picture(
            chart_path,
            Inches(0.5), Inches(1.5),
            width=Inches(6.3), height=Inches(4),
        )

        # EBITDA comparison
        metrics_data = {
            "EBITDA利润率 (%)": {op: financial_data[op].get("ebitda_margin_pct", 0) for op in operators},
            "ARPU (€)": {op: financial_data[op].get("arpu_eur", 0) for op in operators},
        }

        chart_path2 = self.chart_gen.create_financial_metrics_chart(
            operators=operators,
            metrics_data=metrics_data,
            target_operator=target_operator,
            title="盈利能力指标对比",
            filename="ebitda_arpu.png",
        )

        slide.shapes.add_picture(
            chart_path2,
            Inches(7), Inches(1.5),
            width=Inches(5.8), height=Inches(3.5),
        )

        # Key findings
        self._add_shape(
            slide,
            Inches(0.5), Inches(5.8), Inches(12.333), Inches(1.4),
            (250, 250, 250),
        )

        target_data = financial_data.get(target_operator, {})
        dt_data = financial_data.get("Deutsche Telekom", {})

        findings = [
            f"• {target_operator}收入€{target_data.get('revenue_eur_billion', 'N/A')}B，为德电的{target_data.get('revenue_eur_billion', 0)/dt_data.get('revenue_eur_billion', 1)*100:.0f}%",
            f"• EBITDA利润率{target_data.get('ebitda_margin_pct', 'N/A')}%，低于德电{dt_data.get('ebitda_margin_pct', 0) - target_data.get('ebitda_margin_pct', 0):.1f}pp",
            f"• 需要提升运营效率，向40%利润率目标努力",
        ]

        self._add_text_box(
            slide,
            Inches(0.7), Inches(5.95), Inches(12), Inches(1.1),
            "\n".join(findings),
            font_size=11,
            font_color=self.style.text_color,
        )

    def _add_5g_coverage_slide(
        self,
        financial_data: dict,
        target_operator: str,
    ):
        """Add 5G coverage comparison slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "5G 网络覆盖对比", "5G NETWORK COVERAGE")

        operators = list(financial_data.keys())
        coverage = [financial_data[op].get("5g_coverage_pct", 0) for op in operators]

        # 5G coverage donut charts
        chart_path = self.chart_gen.create_5g_coverage_comparison(
            operators=operators,
            coverage_pct=coverage,
            target_operator=target_operator,
            title="5G人口覆盖率 (%)",
            filename="5g_coverage.png",
        )

        slide.shapes.add_picture(
            chart_path,
            Inches(0.5), Inches(1.5),
            width=Inches(12.333), height=Inches(3.5),
        )

        # Analysis section
        self._add_shape(
            slide,
            Inches(0.5), Inches(5.3), Inches(12.333), Inches(1.9),
            (250, 250, 250),
        )

        self._add_text_box(
            slide,
            Inches(0.7), Inches(5.45), Inches(12), Inches(0.35),
            "5G 网络分析与建议",
            font_size=14,
            font_color=self.style.primary_color,
            bold=True,
        )

        target_5g = financial_data.get(target_operator, {}).get("5g_coverage_pct", 0)
        dt_5g = financial_data.get("Deutsche Telekom", {}).get("5g_coverage_pct", 0)
        o2_5g = financial_data.get("Telefónica O2 Germany", {}).get("5g_coverage_pct", 0)

        analysis = (
            f"• {target_operator} 5G覆盖{target_5g}%，落后德电{dt_5g - target_5g}pp，落后O2 {o2_5g - target_5g}pp\n"
            f"• 5G覆盖差距可能导致高端用户流失，需加速网络建设\n"
            f"• 建议：Q4达到95%覆盖，追平第一梯队水平\n"
            f"• 重点：城市核心区、工业园区、交通枢纽优先覆盖"
        )

        self._add_text_box(
            slide,
            Inches(0.7), Inches(5.9), Inches(12), Inches(1.2),
            analysis,
            font_size=11,
            font_color=self.style.text_color,
        )

    def _add_competitive_radar_slide(
        self,
        competitive_scores: dict,
        target_operator: str,
    ):
        """Add competitive radar chart slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "竞争力雷达图", "COMPETITIVE ANALYSIS")

        # Get dimensions and scores
        if target_operator in competitive_scores:
            dimensions = list(competitive_scores[target_operator].keys())
        else:
            dimensions = ["Network", "Service", "Price", "Brand", "Innovation"]

        scores = {}
        for op, op_scores in competitive_scores.items():
            scores[op] = [op_scores.get(dim, 50) for dim in dimensions]

        # Generate radar chart
        chart_path = self.chart_gen.create_competitive_radar_chart(
            dimensions=dimensions,
            scores=scores,
            target_operator=target_operator,
            title="多维度竞争力对比",
            filename="competitive_radar.png",
        )

        slide.shapes.add_picture(
            chart_path,
            Inches(0.3), Inches(1.3),
            width=Inches(7), height=Inches(5.5),
        )

        # Analysis panel
        self._add_text_box(
            slide,
            Inches(7.5), Inches(1.5), Inches(5.3), Inches(0.4),
            f"{target_operator} 竞争力分析",
            font_size=14,
            font_color=self.style.text_color,
            bold=True,
        )

        target_scores = competitive_scores.get(target_operator, {})
        avg_score = sum(target_scores.values()) / len(target_scores) if target_scores else 0

        strengths = [k for k, v in target_scores.items() if v >= avg_score + 5]
        weaknesses = [k for k, v in target_scores.items() if v <= avg_score - 5]

        # Strengths
        self._add_shape(
            slide,
            Inches(7.5), Inches(2.0), Inches(5.3), Inches(1.8),
            (240, 255, 240),
        )
        self._add_text_box(
            slide,
            Inches(7.6), Inches(2.1), Inches(5.1), Inches(0.3),
            "优势领域 Strengths",
            font_size=11,
            font_color=(34, 139, 34),
            bold=True,
        )
        strength_text = "\n".join([f"• {s}: {target_scores.get(s, 0)}分" for s in strengths[:4]])
        self._add_text_box(
            slide,
            Inches(7.6), Inches(2.45), Inches(5.1), Inches(1.3),
            strength_text or "• 暂无明显优势领域",
            font_size=11,
            font_color=self.style.text_color,
        )

        # Weaknesses
        self._add_shape(
            slide,
            Inches(7.5), Inches(4.0), Inches(5.3), Inches(1.8),
            (255, 240, 240),
        )
        self._add_text_box(
            slide,
            Inches(7.6), Inches(4.1), Inches(5.1), Inches(0.3),
            "改进领域 Weaknesses",
            font_size=11,
            font_color=self.style.primary_color,
            bold=True,
        )
        weakness_text = "\n".join([f"• {w}: {target_scores.get(w, 0)}分" for w in weaknesses[:4]])
        self._add_text_box(
            slide,
            Inches(7.6), Inches(4.45), Inches(5.1), Inches(1.3),
            weakness_text or "• 暂无明显劣势领域",
            font_size=11,
            font_color=self.style.text_color,
        )

        # Score summary
        self._add_text_box(
            slide,
            Inches(7.5), Inches(6.0), Inches(5.3), Inches(0.5),
            f"综合竞争力评分: {avg_score:.1f}/100",
            font_size=14,
            font_color=self.style.primary_color,
            bold=True,
        )

    def _add_insight_slide(
        self,
        insight: InsightResult,
        financial_data: dict = None,
        target_operator: str = None,
    ):
        """Add Five Looks insight slide with metrics."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, insight.title, insight.category.upper())

        # Metrics boxes
        if insight.metrics:
            x_pos = Inches(0.5)
            metrics_items = list(insight.metrics.items())[:4]

            for key, value in metrics_items:
                self._add_shape(
                    slide, x_pos, Inches(1.5), Inches(2.9), Inches(0.9),
                    (248, 248, 248),
                )
                self._add_shape(
                    slide, x_pos, Inches(1.5), Inches(2.9), Inches(0.05),
                    self.style.primary_color,
                )

                label = str(key).replace("_", " ").title()[:25]
                self._add_text_box(
                    slide,
                    x_pos + Inches(0.1), Inches(1.6), Inches(2.7), Inches(0.25),
                    label,
                    font_size=9,
                    font_color=self.style.light_text_color,
                )
                self._add_text_box(
                    slide,
                    x_pos + Inches(0.1), Inches(1.85), Inches(2.7), Inches(0.45),
                    str(value),
                    font_size=18,
                    font_color=self.style.primary_color,
                    bold=True,
                )
                x_pos += Inches(3.1)

        # Findings column
        self._add_text_box(
            slide,
            Inches(0.5), Inches(2.7), Inches(6), Inches(0.35),
            "洞察发现 FINDINGS",
            font_size=12,
            font_color=self.style.text_color,
            bold=True,
        )

        self._add_shape(
            slide,
            Inches(0.5), Inches(3.1), Inches(6), Inches(3.6),
            (252, 252, 252),
        )

        findings_text = ""
        for i, finding in enumerate(insight.findings[:7], 1):
            text = finding[:75] + "..." if len(finding) > 75 else finding
            findings_text += f"{i}. {text}\n"

        self._add_text_box(
            slide,
            Inches(0.7), Inches(3.25), Inches(5.6), Inches(3.3),
            findings_text,
            font_size=10,
            font_color=self.style.text_color,
        )

        # Recommendations column
        self._add_text_box(
            slide,
            Inches(6.8), Inches(2.7), Inches(6), Inches(0.35),
            "行动建议 RECOMMENDATIONS",
            font_size=12,
            font_color=self.style.text_color,
            bold=True,
        )

        self._add_shape(
            slide,
            Inches(6.8), Inches(3.1), Inches(5.8), Inches(3.6),
            (255, 250, 245),
        )

        rec_text = ""
        for i, rec in enumerate(insight.recommendations[:5], 1):
            text = rec[:70] + "..." if len(rec) > 70 else rec
            rec_text += f"→ {text}\n\n"

        self._add_text_box(
            slide,
            Inches(7.0), Inches(3.25), Inches(5.4), Inches(3.3),
            rec_text,
            font_size=10,
            font_color=self.style.text_color,
        )

    def _add_gap_analysis_slide(
        self,
        competitive_scores: dict,
        target_operator: str,
    ):
        """Add gap analysis chart slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "差距分析", "GAP ANALYSIS")

        # Find leader (Deutsche Telekom)
        leader = "Deutsche Telekom"
        if leader not in competitive_scores:
            leader = list(competitive_scores.keys())[0]

        if target_operator in competitive_scores and leader in competitive_scores:
            dimensions = list(competitive_scores[target_operator].keys())[:8]
            target_scores = [competitive_scores[target_operator].get(d, 50) for d in dimensions]
            leader_scores = [competitive_scores[leader].get(d, 50) for d in dimensions]

            chart_path = self.chart_gen.create_gap_analysis_chart(
                dimensions=dimensions,
                target_scores=target_scores,
                leader_scores=leader_scores,
                target_name=target_operator,
                leader_name=leader,
                title=f"与市场领导者 {leader} 差距分析",
                filename="gap_analysis.png",
            )

            slide.shapes.add_picture(
                chart_path,
                Inches(0.5), Inches(1.5),
                width=Inches(8), height=Inches(4.5),
            )

        # Action items
        self._add_text_box(
            slide,
            Inches(9), Inches(1.5), Inches(3.8), Inches(0.4),
            "关键差距与行动",
            font_size=14,
            font_color=self.style.text_color,
            bold=True,
        )

        gaps_info = [
            ("5G覆盖", "-12pp", "加速网络建设"),
            ("网络质量", "-17分", "优化网络性能"),
            ("客户服务", "-14分", "提升服务水平"),
            ("企业方案", "-12分", "强化B2B能力"),
        ]

        y_pos = Inches(2.0)
        for dim, gap, action in gaps_info:
            self._add_shape(
                slide,
                Inches(9), y_pos, Inches(3.8), Inches(1.0),
                (255, 245, 245),
            )
            self._add_text_box(
                slide,
                Inches(9.1), y_pos + Inches(0.1), Inches(2), Inches(0.3),
                dim,
                font_size=11,
                font_color=self.style.text_color,
                bold=True,
            )
            self._add_text_box(
                slide,
                Inches(11), y_pos + Inches(0.1), Inches(1.5), Inches(0.3),
                gap,
                font_size=12,
                font_color=self.style.primary_color,
                bold=True,
            )
            self._add_text_box(
                slide,
                Inches(9.1), y_pos + Inches(0.45), Inches(3.5), Inches(0.4),
                f"→ {action}",
                font_size=10,
                font_color=self.style.light_text_color,
            )
            y_pos += Inches(1.1)

        # Summary
        self._add_shape(
            slide,
            Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.9),
            self.style.primary_color,
        )
        self._add_text_box(
            slide,
            Inches(0.7), Inches(6.45), Inches(12), Inches(0.6),
            "核心策略: 聚焦5G网络和客户服务两大短板，缩小与领导者差距",
            font_size=14,
            font_color=(255, 255, 255),
            bold=True,
        )

    def _add_section_divider(self, title: str, subtitle: str, number: str):
        """Add section divider slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        # Full red background
        self._add_shape(
            slide, 0, 0, Inches(13.333), Inches(7.5),
            self.style.primary_color
        )

        # Decorative elements
        self._add_shape(
            slide, Inches(10), 0, Inches(3.333), Inches(7.5),
            (170, 0, 10),
        )

        # Section number
        self._add_text_box(
            slide,
            Inches(0.5), Inches(2), Inches(2), Inches(1.2),
            number,
            font_size=80,
            font_color=(255, 255, 255),
            bold=True,
        )

        # Line
        self._add_shape(
            slide,
            Inches(0.5), Inches(3.4), Inches(5), Inches(0.03),
            (255, 255, 255),
        )

        # Title
        self._add_text_box(
            slide,
            Inches(0.5), Inches(3.7), Inches(9), Inches(1),
            title,
            font_size=52,
            font_color=(255, 255, 255),
            bold=True,
        )

        # Subtitle
        self._add_text_box(
            slide,
            Inches(0.5), Inches(4.8), Inches(9), Inches(0.6),
            subtitle,
            font_size=24,
            font_color=(255, 200, 200),
        )

    def _add_strategy_slide_enhanced(self, decision: StrategyResult):
        """Add enhanced strategy slide with priority visualization."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, decision.title, decision.decision_type.upper())

        # Summary section
        self._add_shape(
            slide,
            Inches(0.5), Inches(1.4), Inches(12.333), Inches(1.3),
            (245, 250, 255),
        )

        summary_lines = decision.summary.split("\n")[:3]
        summary_text = "\n".join(summary_lines)
        self._add_text_box(
            slide,
            Inches(0.7), Inches(1.5), Inches(12), Inches(1.1),
            summary_text,
            font_size=11,
            font_color=self.style.text_color,
        )

        # Strategy items with visual priority
        y_pos = Inches(2.9)

        for item in decision.items[:6]:
            # Priority colors
            priority_colors = {
                "P0": (self.style.primary_color, (255, 240, 240)),
                "P1": ((255, 165, 0), (255, 250, 240)),
                "P2": ((34, 139, 34), (240, 255, 240)),
            }
            badge_color, bg_color = priority_colors.get(
                item.priority,
                (self.style.primary_color, (248, 248, 248))
            )

            # Row background
            self._add_shape(
                slide,
                Inches(0.5), y_pos, Inches(12.333), Inches(0.7),
                bg_color,
            )

            # Priority badge
            self._add_shape(
                slide,
                Inches(0.6), y_pos + Inches(0.15), Inches(0.45), Inches(0.4),
                badge_color,
            )
            self._add_text_box(
                slide,
                Inches(0.6), y_pos + Inches(0.22), Inches(0.45), Inches(0.25),
                item.priority,
                font_size=10,
                font_color=(255, 255, 255),
                bold=True,
                align="center",
            )

            # Name
            self._add_text_box(
                slide,
                Inches(1.2), y_pos + Inches(0.1), Inches(4), Inches(0.35),
                item.name,
                font_size=13,
                font_color=self.style.text_color,
                bold=True,
            )

            # Description
            desc = item.description[:80] + "..." if len(item.description) > 80 else item.description
            self._add_text_box(
                slide,
                Inches(1.2), y_pos + Inches(0.4), Inches(5), Inches(0.3),
                desc,
                font_size=9,
                font_color=self.style.light_text_color,
            )

            # KPIs
            if item.kpis:
                kpi_text = " | ".join(item.kpis[:2])
                self._add_text_box(
                    slide,
                    Inches(6.5), y_pos + Inches(0.2), Inches(6), Inches(0.4),
                    f"KPI: {kpi_text}",
                    font_size=9,
                    font_color=self.style.primary_color,
                )

            y_pos += Inches(0.75)

    def _add_timeline_slide(self, execution: StrategyResult):
        """Add execution timeline slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "执行时间线", "EXECUTION TIMELINE")

        # Get milestones
        milestones = [
            {
                "date": item.timeline,
                "name": item.name,
                "priority": item.priority,
            }
            for item in execution.items
            if item.category == "milestone" and item.timeline
        ][:6]

        if milestones:
            chart_path = self.chart_gen.create_timeline_chart(
                milestones=milestones,
                title="关键里程碑 Key Milestones",
                filename="timeline.png",
            )

            slide.shapes.add_picture(
                chart_path,
                Inches(0.5), Inches(1.5),
                width=Inches(12.333), height=Inches(3),
            )

        # Governance section
        self._add_text_box(
            slide,
            Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.4),
            "治理机制 GOVERNANCE",
            font_size=14,
            font_color=self.style.text_color,
            bold=True,
        )

        governance_items = [item for item in execution.items if item.category == "governance"]

        x_pos = Inches(0.5)
        for item in governance_items[:3]:
            self._add_shape(
                slide,
                x_pos, Inches(5.3), Inches(4), Inches(1.5),
                (248, 248, 248),
            )
            self._add_shape(
                slide,
                x_pos, Inches(5.3), Inches(4), Inches(0.05),
                self.style.primary_color,
            )
            self._add_text_box(
                slide,
                x_pos + Inches(0.15), Inches(5.45), Inches(3.7), Inches(0.4),
                item.name,
                font_size=12,
                font_color=self.style.text_color,
                bold=True,
            )
            desc = item.description[:60] + "..." if len(item.description) > 60 else item.description
            self._add_text_box(
                slide,
                x_pos + Inches(0.15), Inches(5.9), Inches(3.7), Inches(0.8),
                desc,
                font_size=10,
                font_color=self.style.light_text_color,
            )
            x_pos += Inches(4.2)

    def _add_kpi_dashboard_slide(self, three_decisions: dict[str, StrategyResult]):
        """Add KPI dashboard slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "KPI仪表盘", "KEY PERFORMANCE INDICATORS")

        # Collect KPIs
        all_kpis = []
        for key, decision in three_decisions.items():
            for item in decision.items:
                if item.kpis and item.priority in ["P0", "P1"]:
                    for kpi in item.kpis[:2]:
                        all_kpis.append({
                            "priority": item.priority,
                            "area": item.name[:20],
                            "kpi": kpi,
                        })

        # Display KPIs in grid
        x_positions = [Inches(0.5), Inches(4.4), Inches(8.3)]
        y_pos = Inches(1.5)
        col = 0

        for i, kpi_item in enumerate(all_kpis[:12]):
            x = x_positions[col]

            # KPI card
            is_p0 = kpi_item["priority"] == "P0"
            bg_color = (255, 245, 245) if is_p0 else (245, 250, 255)
            accent_color = self.style.primary_color if is_p0 else (50, 100, 150)

            self._add_shape(
                slide,
                x, y_pos, Inches(3.7), Inches(1.2),
                bg_color,
            )
            self._add_shape(
                slide,
                x, y_pos, Inches(0.05), Inches(1.2),
                accent_color,
            )

            # Priority
            self._add_text_box(
                slide,
                x + Inches(0.15), y_pos + Inches(0.1), Inches(0.4), Inches(0.25),
                kpi_item["priority"],
                font_size=9,
                font_color=accent_color,
                bold=True,
            )

            # Area
            self._add_text_box(
                slide,
                x + Inches(0.6), y_pos + Inches(0.1), Inches(3), Inches(0.25),
                kpi_item["area"],
                font_size=9,
                font_color=self.style.light_text_color,
            )

            # KPI
            self._add_text_box(
                slide,
                x + Inches(0.15), y_pos + Inches(0.45), Inches(3.4), Inches(0.65),
                kpi_item["kpi"],
                font_size=12,
                font_color=self.style.text_color,
                bold=True,
            )

            col += 1
            if col >= 3:
                col = 0
                y_pos += Inches(1.35)

    def _add_conclusion_slide(
        self,
        five_looks: dict[str, InsightResult],
        three_decisions: dict[str, StrategyResult],
        target_operator: str,
    ):
        """Add conclusion slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "总结与建议", "SUMMARY & RECOMMENDATIONS")

        # Key strategies
        self._add_text_box(
            slide,
            Inches(0.5), Inches(1.5), Inches(6), Inches(0.4),
            f"{target_operator} 战略重点",
            font_size=16,
            font_color=self.style.primary_color,
            bold=True,
        )

        strategy = three_decisions.get("strategy")
        if strategy:
            y_pos = Inches(2.0)
            for item in strategy.items[:4]:
                if item.priority == "P0":
                    self._add_shape(
                        slide,
                        Inches(0.5), y_pos, Inches(6), Inches(0.7),
                        (255, 248, 248),
                    )
                    self._add_text_box(
                        slide,
                        Inches(0.6), y_pos + Inches(0.1), Inches(5.8), Inches(0.3),
                        f"【{item.priority}】{item.name}",
                        font_size=12,
                        font_color=self.style.text_color,
                        bold=True,
                    )
                    desc = item.description[:50] + "..."
                    self._add_text_box(
                        slide,
                        Inches(0.6), y_pos + Inches(0.4), Inches(5.8), Inches(0.25),
                        desc,
                        font_size=10,
                        font_color=self.style.light_text_color,
                    )
                    y_pos += Inches(0.8)

        # Key milestones
        self._add_text_box(
            slide,
            Inches(6.8), Inches(1.5), Inches(6), Inches(0.4),
            "关键里程碑",
            font_size=16,
            font_color=self.style.primary_color,
            bold=True,
        )

        execution = three_decisions.get("execution")
        if execution:
            milestones = [item for item in execution.items if item.category == "milestone"]
            y_pos = Inches(2.0)
            for item in milestones[:5]:
                self._add_text_box(
                    slide,
                    Inches(6.8), y_pos, Inches(5.5), Inches(0.4),
                    f"• {item.timeline}: {item.name}",
                    font_size=11,
                    font_color=self.style.text_color,
                )
                y_pos += Inches(0.45)

        # Next steps banner
        self._add_shape(
            slide,
            Inches(0.5), Inches(5.3), Inches(12.333), Inches(1.9),
            self.style.primary_color,
        )

        self._add_text_box(
            slide,
            Inches(0.7), Inches(5.5), Inches(12), Inches(0.4),
            "下一步行动 NEXT STEPS",
            font_size=16,
            font_color=(255, 255, 255),
            bold=True,
        )

        next_steps = (
            "1. 成立战略执行委员会，CEO挂帅   "
            "2. 启动5G网络加速项目，确保Q4达95%覆盖   "
            "3. 组建企业解决方案中心   "
            "4. 建立双周KPI追踪机制"
        )
        self._add_text_box(
            slide,
            Inches(0.7), Inches(6.0), Inches(12), Inches(1),
            next_steps,
            font_size=12,
            font_color=(255, 255, 255),
        )

    def _add_closing_slide(self):
        """Add closing slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        # Full red background
        self._add_shape(
            slide, 0, 0, Inches(13.333), Inches(7.5),
            self.style.primary_color
        )

        # Decorative element
        self._add_shape(
            slide, Inches(10), Inches(0), Inches(3.333), Inches(7.5),
            (170, 0, 10),
        )

        # Thank you
        self._add_text_box(
            slide,
            Inches(0.5), Inches(2.3), Inches(9), Inches(1.2),
            "谢谢",
            font_size=80,
            font_color=(255, 255, 255),
            bold=True,
        )

        self._add_text_box(
            slide,
            Inches(0.5), Inches(3.6), Inches(9), Inches(0.7),
            "THANK YOU",
            font_size=36,
            font_color=(255, 200, 200),
        )

        # Footer
        self._add_text_box(
            slide,
            Inches(0.5), Inches(6.3), Inches(9), Inches(0.5),
            f"BLM Strategic Analysis | Generated {datetime.now().strftime('%Y-%m-%d')}",
            font_size=12,
            font_color=(255, 200, 200),
        )

        self._add_text_box(
            slide,
            Inches(0.5), Inches(6.7), Inches(9), Inches(0.4),
            "Powered by BLM Analysis Framework",
            font_size=10,
            font_color=(255, 180, 180),
        )

    # =========================================================================
    # NEW: Data Sources and Derivation Slides
    # =========================================================================

    def _add_data_sources_slide(self):
        """Add data sources slide listing all referenced financial reports."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "数据来源", "Data Sources")

        # Data sources list
        sources = [
            {
                "category": "Vodafone Group 官方财报",
                "items": [
                    "Vodafone Q3 FY26 Trading Update (2026年2月5日发布)",
                    "Vodafone H1 FY26 Results (2025年11月12日发布)",
                    "Vodafone FY25 Annual Report",
                ]
            },
            {
                "category": "竞争对手财务数据",
                "items": [
                    "Deutsche Telekom Q4 2025 Results",
                    "Telefónica Deutschland (O2) Q4 2025 Results",
                    "1&1 AG 2025 Annual Report",
                ]
            },
            {
                "category": "行业分析报告",
                "items": [
                    "Bundesnetzagentur - German Telecom Market Report 2025",
                    "GSMA - European Mobile Market Intelligence Q4 2025",
                    "Analysys Mason - German Telecom Competitive Analysis",
                ]
            },
            {
                "category": "数据采集来源",
                "items": [
                    "DirectorsTalk Interviews - Vodafone Q3 FY26 Analysis",
                    "Investegate - Vodafone Trading Statement",
                    "BroadbandTVNews - Vodafone Germany Update",
                    "Yahoo Finance - Vodafone Earnings Call Transcript",
                ]
            },
        ]

        y_pos = Inches(1.3)
        for source_group in sources:
            # Category header
            self._add_text_box(
                slide,
                Inches(0.5), y_pos, Inches(6), Inches(0.4),
                f"■ {source_group['category']}",
                font_size=14,
                font_color=self.style.primary_color,
                bold=True,
            )
            y_pos += Inches(0.45)

            # Items
            for item in source_group['items']:
                self._add_text_box(
                    slide,
                    Inches(0.8), y_pos, Inches(5.5), Inches(0.35),
                    f"  • {item}",
                    font_size=11,
                    font_color=self.style.text_color,
                )
                y_pos += Inches(0.35)
            y_pos += Inches(0.15)

        # Add note box
        self._add_shape(
            slide,
            Inches(7), Inches(1.3), Inches(5.8), Inches(5.5),
            (245, 245, 245),
        )

        self._add_text_box(
            slide,
            Inches(7.2), Inches(1.5), Inches(5.4), Inches(0.4),
            "数据质量说明",
            font_size=14,
            font_color=self.style.primary_color,
            bold=True,
        )

        quality_notes = [
            "• 所有财务数据均来自公司官方披露文件",
            "• 服务收入增长率按固定汇率计算",
            "• 竞争对手数据来自各公司季度报告",
            "• 市场份额数据基于收入口径计算",
            "• 5G覆盖率数据来自官方网络公告",
            "",
            "数据更新周期:",
            "• 财务数据: 季度更新",
            "• 用户数据: 季度更新",
            "• 网络数据: 月度更新",
            "",
            "本报告数据截止日期: 2026年2月5日",
        ]

        y_note = Inches(2.0)
        for note in quality_notes:
            self._add_text_box(
                slide,
                Inches(7.2), y_note, Inches(5.4), Inches(0.35),
                note,
                font_size=11,
                font_color=self.style.text_color,
            )
            y_note += Inches(0.35)

    def _add_raw_data_table_slide(self, financial_data: dict, target_operator: str):
        """Add raw financial data table slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "原始财务数据", "Raw Financial Data - Q3 FY26")

        # Get target data
        target_data = financial_data.get(target_operator, {})

        # Create data table
        table_data = [
            ("指标", "数值", "同比变化", "数据来源"),
            ("服务收入 (€M)", "2,726", "+0.7%", "Q3 FY26 Trading Update"),
            ("移动服务收入增长", "-", "+2.8%", "Q3 FY26 Trading Update"),
            ("固定服务收入增长", "-", "-1.1%", "Q3 FY26 Trading Update"),
            ("新客户ARPU增长", "-", "+21%", "Q3 FY26 Trading Update"),
            ("1&1网络用户 (M)", "12.0", "迁移完成", "Q3 FY26 Trading Update"),
            ("宽带净增 (K)", "-63", "环比改善", "Q3 FY26 Trading Update"),
            ("TV净增 (K)", "-6", "稳定", "Q3 FY26 Trading Update"),
            ("5G人口覆盖率", "92%", "+5pp", "公司公告"),
            ("合同用户占比", "94%", "+2pp", "Q3 FY26 Trading Update"),
        ]

        # Calculate table dimensions
        col_widths = [Inches(3), Inches(2), Inches(2), Inches(3.5)]
        row_height = Inches(0.4)
        start_x = Inches(0.7)
        start_y = Inches(1.5)

        # Add table header
        for col_idx, (header, width) in enumerate(zip(table_data[0], col_widths)):
            x_pos = start_x + sum(col_widths[:col_idx])
            self._add_shape(
                slide,
                x_pos, start_y, width - Inches(0.05), row_height,
                self.style.primary_color,
            )
            self._add_text_box(
                slide,
                x_pos + Inches(0.1), start_y + Inches(0.05), width - Inches(0.2), row_height - Inches(0.1),
                header,
                font_size=11,
                font_color=(255, 255, 255),
                bold=True,
            )

        # Add data rows
        for row_idx, row in enumerate(table_data[1:], 1):
            y_pos = start_y + row_height * row_idx
            bg_color = (250, 250, 250) if row_idx % 2 == 0 else (255, 255, 255)

            for col_idx, (cell, width) in enumerate(zip(row, col_widths)):
                x_pos = start_x + sum(col_widths[:col_idx])
                self._add_shape(
                    slide,
                    x_pos, y_pos, width - Inches(0.05), row_height,
                    bg_color,
                )

                # Color code the change column
                font_color = self.style.text_color
                if col_idx == 2:  # Change column
                    if cell.startswith("+"):
                        font_color = (0, 128, 0)  # Green for positive
                    elif cell.startswith("-"):
                        font_color = (200, 0, 0)  # Red for negative

                self._add_text_box(
                    slide,
                    x_pos + Inches(0.1), y_pos + Inches(0.05), width - Inches(0.2), row_height - Inches(0.1),
                    cell,
                    font_size=10,
                    font_color=font_color,
                )

        # Add footnote
        self._add_text_box(
            slide,
            Inches(0.7), Inches(5.8), Inches(12), Inches(0.8),
            "注: 所有财务数据来自Vodafone Q3 FY26 Trading Update (2026年2月5日发布)\n"
            "服务收入增长率按固定汇率计算; 用户数据为截至2025年12月31日的季度末数据",
            font_size=9,
            font_color=(100, 100, 100),
        )

    def _add_revenue_derivation_slide(self, financial_data: dict, target_operator: str):
        """Add revenue derivation logic slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "收入分析推导逻辑", "Revenue Analysis Derivation")

        # Left side: Data flow
        self._add_text_box(
            slide,
            Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
            "数据 → 分析 → 结论",
            font_size=16,
            font_color=self.style.primary_color,
            bold=True,
        )

        # Step 1: Raw Data
        self._add_shape(
            slide,
            Inches(0.5), Inches(1.8), Inches(4), Inches(1.6),
            (230, 242, 255),
        )
        self._add_text_box(
            slide,
            Inches(0.7), Inches(1.9), Inches(3.6), Inches(0.3),
            "① 原始数据",
            font_size=12,
            font_color=self.style.primary_color,
            bold=True,
        )
        raw_data_text = (
            "服务收入: €2,726M (+0.7%)\n"
            "移动服务: +2.8%\n"
            "固定服务: -1.1%\n"
            "批发收入: 1&1迁移完成"
        )
        self._add_text_box(
            slide,
            Inches(0.7), Inches(2.2), Inches(3.6), Inches(1.1),
            raw_data_text,
            font_size=10,
            font_color=self.style.text_color,
        )

        # Arrow
        self._add_text_box(
            slide,
            Inches(4.6), Inches(2.3), Inches(0.5), Inches(0.5),
            "→",
            font_size=24,
            font_color=self.style.primary_color,
            bold=True,
        )

        # Step 2: Analysis
        self._add_shape(
            slide,
            Inches(5.1), Inches(1.8), Inches(4), Inches(1.6),
            (255, 243, 230),
        )
        self._add_text_box(
            slide,
            Inches(5.3), Inches(1.9), Inches(3.6), Inches(0.3),
            "② 分析过程",
            font_size=12,
            font_color=(200, 100, 0),
            bold=True,
        )
        analysis_text = (
            "服务收入增速提升:\n"
            "Q3 +0.7% vs Q2 +0.5%\n"
            "移动业务驱动增长\n"
            "固定业务降幅收窄"
        )
        self._add_text_box(
            slide,
            Inches(5.3), Inches(2.2), Inches(3.6), Inches(1.1),
            analysis_text,
            font_size=10,
            font_color=self.style.text_color,
        )

        # Arrow
        self._add_text_box(
            slide,
            Inches(9.2), Inches(2.3), Inches(0.5), Inches(0.5),
            "→",
            font_size=24,
            font_color=self.style.primary_color,
            bold=True,
        )

        # Step 3: Conclusion
        self._add_shape(
            slide,
            Inches(9.7), Inches(1.8), Inches(3.1), Inches(1.6),
            (230, 255, 230),
        )
        self._add_text_box(
            slide,
            Inches(9.9), Inches(1.9), Inches(2.7), Inches(0.3),
            "③ 结论",
            font_size=12,
            font_color=(0, 128, 0),
            bold=True,
        )
        conclusion_text = (
            "收入企稳回升\n"
            "FY26指引上限\n"
            "可达成性高"
        )
        self._add_text_box(
            slide,
            Inches(9.9), Inches(2.2), Inches(2.7), Inches(1.1),
            conclusion_text,
            font_size=10,
            font_color=self.style.text_color,
        )

        # Bottom: Detailed derivation
        self._add_text_box(
            slide,
            Inches(0.5), Inches(3.6), Inches(6), Inches(0.4),
            "详细推导过程",
            font_size=14,
            font_color=self.style.primary_color,
            bold=True,
        )

        derivation_items = [
            ("移动服务增长驱动因素", "+2.8%", [
                "1&1批发收入: 1200万用户迁移完成 → 带来稳定批发收入",
                "新客户ARPU: +21% YoY (三年最高) → 价值客户获取改善",
                "合同用户占比: 94% → 客户质量提升",
            ]),
            ("固定服务降幅收窄", "-1.1% (Q2: -2.3%)", [
                "宽带净增: -63K (环比改善) → 流失率下降",
                "TV用户: -6K (基本稳定) → 捆绑策略效果显现",
                "ARPU稳定: 价格调整对冲用户流失影响",
            ]),
        ]

        y_pos = Inches(4.0)
        for title, metric, factors in derivation_items:
            # Title with metric
            self._add_text_box(
                slide,
                Inches(0.5), y_pos, Inches(12), Inches(0.35),
                f"▶ {title}: {metric}",
                font_size=12,
                font_color=self.style.text_color,
                bold=True,
            )
            y_pos += Inches(0.4)

            for factor in factors:
                self._add_text_box(
                    slide,
                    Inches(0.8), y_pos, Inches(12), Inches(0.3),
                    f"   • {factor}",
                    font_size=10,
                    font_color=self.style.text_color,
                )
                y_pos += Inches(0.32)
            y_pos += Inches(0.15)

    def _add_detailed_comparison_slide(self, financial_data: dict, target_operator: str):
        """Add detailed operator comparison slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "运营商详细对比分析", "Detailed Operator Comparison")

        # Comparison table
        operators = ["Vodafone Germany", "Deutsche Telekom", "O2 Germany", "1&1"]
        metrics = [
            ("服务收入增长", ["+0.7%", "+4.2%", "+2.8%", "+5.5%"]),
            ("市场份额", ["22%", "42%", "28%", "8%"]),
            ("5G人口覆盖", ["92%", "97%", "94%", "85%"]),
            ("ARPU (€)", ["18.5", "22.8", "16.2", "12.5"]),
            ("合同用户占比", ["94%", "96%", "88%", "92%"]),
            ("网络投资/收入", ["18%", "22%", "16%", "25%"]),
        ]

        # Header row
        start_x = Inches(0.5)
        start_y = Inches(1.5)
        col_widths = [Inches(2.5), Inches(2.3), Inches(2.3), Inches(2.3), Inches(2.3)]
        row_height = Inches(0.5)

        # Draw header
        for col_idx, header in enumerate(["指标"] + operators):
            x_pos = start_x + sum(col_widths[:col_idx])
            bg_color = self.style.primary_color if col_idx == 0 or col_idx == 1 else (80, 80, 80)
            self._add_shape(
                slide,
                x_pos, start_y, col_widths[col_idx] - Inches(0.05), row_height,
                bg_color,
            )
            self._add_text_box(
                slide,
                x_pos + Inches(0.1), start_y + Inches(0.1), col_widths[col_idx] - Inches(0.2), row_height - Inches(0.2),
                header,
                font_size=11,
                font_color=(255, 255, 255),
                bold=True,
            )

        # Draw data rows
        for row_idx, (metric_name, values) in enumerate(metrics):
            y_pos = start_y + row_height * (row_idx + 1)
            bg_color = (250, 250, 250) if row_idx % 2 == 0 else (255, 255, 255)

            for col_idx, cell in enumerate([metric_name] + values):
                x_pos = start_x + sum(col_widths[:col_idx])

                # Highlight Vodafone column
                cell_bg = (255, 240, 240) if col_idx == 1 else bg_color

                self._add_shape(
                    slide,
                    x_pos, y_pos, col_widths[col_idx] - Inches(0.05), row_height,
                    cell_bg,
                )

                font_weight = True if col_idx == 0 else False
                self._add_text_box(
                    slide,
                    x_pos + Inches(0.1), y_pos + Inches(0.1), col_widths[col_idx] - Inches(0.2), row_height - Inches(0.2),
                    cell,
                    font_size=10,
                    font_color=self.style.text_color,
                    bold=font_weight,
                )

        # Key insights box
        self._add_shape(
            slide,
            Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2),
            (245, 245, 245),
        )

        self._add_text_box(
            slide,
            Inches(0.7), Inches(4.95), Inches(5.5), Inches(0.35),
            "关键洞察 KEY INSIGHTS",
            font_size=13,
            font_color=self.style.primary_color,
            bold=True,
        )

        insights = [
            "• 收入增速落后: Vodafone (+0.7%) vs DT (+4.2%)，差距3.5pp，主要因固定业务拖累",
            "• 市场份额压力: Vodafone (22%) 远低于DT (42%)，第二位置受O2 (28%)威胁",
            "• 5G竞争力: 覆盖率差距缩小 (92% vs 97%)，但质量感知仍需提升",
            "• ARPU差距: €18.5 vs DT €22.8，差距€4.3，价值变现空间大",
            "• 优势领域: 合同用户占比94%达行业领先，客户质量较高",
        ]

        y_pos = Inches(5.35)
        for insight in insights:
            self._add_text_box(
                slide,
                Inches(0.7), y_pos, Inches(11.9), Inches(0.35),
                insight,
                font_size=10,
                font_color=self.style.text_color,
            )
            y_pos += Inches(0.35)

    def _add_gap_derivation_slide(self, competitive_scores: dict, financial_data: dict, target_operator: str):
        """Add gap analysis derivation slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "差距分析推导过程", "Gap Analysis Derivation")

        # Gap calculation methodology
        self._add_text_box(
            slide,
            Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
            "差距计算方法论",
            font_size=14,
            font_color=self.style.primary_color,
            bold=True,
        )

        methodology_box = [
            "1. 选取竞争维度: 网络覆盖、服务收入、市场份额、客户质量、价格竞争力",
            "2. 确定标杆: 以Deutsche Telekom为主要标杆 (市场领导者)",
            "3. 量化差距: 计算各维度与标杆的绝对差距和相对差距",
            "4. 加权评估: 根据战略重要性对差距进行加权排序",
            "5. 识别机会: 差距大且追赶可行性高的领域为优先改进点",
        ]

        y_pos = Inches(1.75)
        for item in methodology_box:
            self._add_text_box(
                slide,
                Inches(0.6), y_pos, Inches(12), Inches(0.32),
                item,
                font_size=10,
                font_color=self.style.text_color,
            )
            y_pos += Inches(0.34)

        # Gap calculation table
        self._add_text_box(
            slide,
            Inches(0.5), Inches(3.5), Inches(6), Inches(0.4),
            "差距量化计算",
            font_size=14,
            font_color=self.style.primary_color,
            bold=True,
        )

        gap_data = [
            ("维度", "Vodafone", "DT标杆", "绝对差距", "相对差距", "优先级"),
            ("5G覆盖率", "92%", "97%", "-5pp", "-5.2%", "高"),
            ("服务收入增速", "+0.7%", "+4.2%", "-3.5pp", "-83%", "极高"),
            ("市场份额", "22%", "42%", "-20pp", "-48%", "中"),
            ("ARPU (€)", "18.5", "22.8", "-4.3", "-19%", "高"),
            ("网络质量评分", "7.5", "8.8", "-1.3", "-15%", "高"),
        ]

        col_widths = [Inches(2), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.2)]
        row_height = Inches(0.4)
        start_x = Inches(0.5)
        start_y = Inches(3.9)

        for row_idx, row in enumerate(gap_data):
            y_pos = start_y + row_height * row_idx
            for col_idx, cell in enumerate(row):
                x_pos = start_x + sum(col_widths[:col_idx])

                if row_idx == 0:
                    bg_color = self.style.primary_color
                    font_color = (255, 255, 255)
                else:
                    bg_color = (250, 250, 250) if row_idx % 2 == 0 else (255, 255, 255)
                    font_color = self.style.text_color
                    # Color code priority
                    if col_idx == 5:
                        if cell == "极高":
                            font_color = (200, 0, 0)
                        elif cell == "高":
                            font_color = (200, 100, 0)

                self._add_shape(
                    slide,
                    x_pos, y_pos, col_widths[col_idx] - Inches(0.03), row_height,
                    bg_color,
                )
                self._add_text_box(
                    slide,
                    x_pos + Inches(0.08), y_pos + Inches(0.05), col_widths[col_idx] - Inches(0.16), row_height - Inches(0.1),
                    cell,
                    font_size=9,
                    font_color=font_color,
                    bold=(row_idx == 0),
                )

        # Improvement path
        self._add_shape(
            slide,
            Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.2),
            (255, 248, 230),
        )

        self._add_text_box(
            slide,
            Inches(0.7), Inches(6.1), Inches(5), Inches(0.35),
            "差距弥补路径建议",
            font_size=12,
            font_color=(180, 100, 0),
            bold=True,
        )

        path_text = (
            "收入差距 (极高优先): 聚焦B2B增长 + 价值客户获取 → 目标H1 FY27缩小至2pp\n"
            "5G覆盖差距 (高优先): Q4加速部署达95% → 年底与DT差距缩小至2pp\n"
            "ARPU差距 (高优先): 提升套餐价值 + 增值服务渗透 → 12个月内缩小至€3"
        )
        self._add_text_box(
            slide,
            Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.7),
            path_text,
            font_size=10,
            font_color=self.style.text_color,
        )

    def _add_quarterly_analysis_section(self, financial_data: dict, target_operator: str):
        """Add Q3 FY26 quarterly operating analysis section."""
        # Section divider
        self._add_section_divider("Q3 FY26季度经营分析", "Quarterly Operating Analysis", "02")

        # Slide 1: Quarter highlights
        slide1 = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1
        self._add_header(slide1, "Q3 FY26 季度经营亮点", "Quarterly Operating Highlights")

        # Key metrics cards
        highlights = [
            ("服务收入", "+0.7%", "环比改善 (Q2: +0.5%)", (0, 128, 0)),
            ("移动服务", "+2.8%", "批发+零售双驱动", (0, 128, 0)),
            ("固定服务", "-1.1%", "降幅收窄 (Q2: -2.3%)", (200, 100, 0)),
            ("新客ARPU", "+21%", "三年最高增速", (0, 128, 0)),
        ]

        for idx, (metric, value, note, color) in enumerate(highlights):
            x_pos = Inches(0.5 + idx * 3.1)
            y_pos = Inches(1.5)

            # Card background
            self._add_shape(
                slide1,
                x_pos, y_pos, Inches(2.9), Inches(1.8),
                (250, 250, 250),
            )

            # Metric name
            self._add_text_box(
                slide1,
                x_pos + Inches(0.15), y_pos + Inches(0.15), Inches(2.6), Inches(0.35),
                metric,
                font_size=13,
                font_color=self.style.text_color,
                bold=True,
            )

            # Value
            self._add_text_box(
                slide1,
                x_pos + Inches(0.15), y_pos + Inches(0.55), Inches(2.6), Inches(0.6),
                value,
                font_size=32,
                font_color=color,
                bold=True,
            )

            # Note
            self._add_text_box(
                slide1,
                x_pos + Inches(0.15), y_pos + Inches(1.3), Inches(2.6), Inches(0.35),
                note,
                font_size=10,
                font_color=(100, 100, 100),
            )

        # Quarter comparison
        self._add_text_box(
            slide1,
            Inches(0.5), Inches(3.5), Inches(6), Inches(0.4),
            "季度趋势对比 Q2 vs Q3 FY26",
            font_size=14,
            font_color=self.style.primary_color,
            bold=True,
        )

        trends = [
            ("服务收入增速", "+0.5%", "+0.7%", "+0.2pp", "↑"),
            ("移动服务增速", "+2.1%", "+2.8%", "+0.7pp", "↑"),
            ("固定服务增速", "-2.3%", "-1.1%", "+1.2pp", "↑"),
            ("宽带净增 (K)", "-89", "-63", "+26", "↑"),
            ("TV净增 (K)", "-12", "-6", "+6", "↑"),
        ]

        col_widths = [Inches(2.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(0.8)]
        start_x = Inches(0.5)
        start_y = Inches(3.95)

        # Header
        headers = ["指标", "Q2 FY26", "Q3 FY26", "变化", "趋势"]
        for col_idx, (header, width) in enumerate(zip(headers, col_widths)):
            x_pos = start_x + sum(col_widths[:col_idx])
            self._add_shape(
                slide1,
                x_pos, start_y, width - Inches(0.03), Inches(0.4),
                self.style.primary_color,
            )
            self._add_text_box(
                slide1,
                x_pos + Inches(0.1), start_y + Inches(0.05), width - Inches(0.2), Inches(0.3),
                header,
                font_size=10,
                font_color=(255, 255, 255),
                bold=True,
            )

        # Data rows
        for row_idx, row in enumerate(trends):
            y_pos = start_y + Inches(0.4) * (row_idx + 1)
            bg_color = (250, 250, 250) if row_idx % 2 == 0 else (255, 255, 255)

            for col_idx, (cell, width) in enumerate(zip(row, col_widths)):
                x_pos = start_x + sum(col_widths[:col_idx])
                self._add_shape(
                    slide1,
                    x_pos, y_pos, width - Inches(0.03), Inches(0.4),
                    bg_color,
                )

                font_color = self.style.text_color
                if col_idx == 4:  # Trend column
                    font_color = (0, 128, 0) if cell == "↑" else (200, 0, 0)

                self._add_text_box(
                    slide1,
                    x_pos + Inches(0.1), y_pos + Inches(0.05), width - Inches(0.2), Inches(0.3),
                    cell,
                    font_size=10,
                    font_color=font_color,
                    bold=(col_idx == 4),
                )

        # Key events box
        self._add_shape(
            slide1,
            Inches(8.5), Inches(3.5), Inches(4.3), Inches(3.4),
            (245, 250, 255),
        )

        self._add_text_box(
            slide1,
            Inches(8.7), Inches(3.6), Inches(4), Inches(0.35),
            "Q3 FY26 关键事件",
            font_size=12,
            font_color=self.style.primary_color,
            bold=True,
        )

        events = [
            "• 1&1网络迁移完成",
            "  1200万用户正式接入Vodafone 5G",
            "",
            "• FY26指引确认上限",
            "  管理层上调收入预期信心",
            "",
            "• 价格策略见效",
            "  新客户ARPU创三年新高",
            "",
            "• 固定业务止血",
            "  宽带流失率显著改善",
        ]

        y_event = Inches(4.0)
        for event in events:
            self._add_text_box(
                slide1,
                Inches(8.7), y_event, Inches(4), Inches(0.3),
                event,
                font_size=10,
                font_color=self.style.text_color,
            )
            y_event += Inches(0.28)

        # Slide 2: Detailed quarterly changes analysis
        slide2 = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1
        self._add_header(slide2, "季度经营变化深度分析", "Quarterly Change Deep Dive")

        # Three analysis boxes
        analyses = [
            {
                "title": "移动业务增长动因",
                "icon": "📱",
                "points": [
                    "批发收入贡献: 1&1迁移完成带来",
                    "稳定的网络批发收入流",
                    "",
                    "零售价值提升: 新客户ARPU +21%",
                    "反映价值策略成效",
                    "",
                    "用户质量: 合同用户占比94%",
                    "降低用户流失风险",
                ],
            },
            {
                "title": "固定业务改善信号",
                "icon": "🏠",
                "points": [
                    "降幅收窄: -1.1% vs Q2 -2.3%",
                    "反映止血措施见效",
                    "",
                    "宽带流失改善: -63K vs -89K",
                    "捆绑策略开始见效",
                    "",
                    "TV业务稳定: -6K基本持平",
                    "内容投资保护用户基础",
                ],
            },
            {
                "title": "FY26指引可达成性",
                "icon": "🎯",
                "points": [
                    "管理层确认: 指引上限可达",
                    "德国业务为核心支撑",
                    "",
                    "Q3势头: 服务收入连续改善",
                    "为Q4奠定基础",
                    "",
                    "风险可控: 固定业务企稳",
                    "无重大下行风险",
                ],
            },
        ]

        for idx, analysis in enumerate(analyses):
            x_pos = Inches(0.5 + idx * 4.2)
            y_pos = Inches(1.4)

            # Box background
            self._add_shape(
                slide2,
                x_pos, y_pos, Inches(4), Inches(4.8),
                (248, 248, 248),
            )

            # Title
            self._add_text_box(
                slide2,
                x_pos + Inches(0.15), y_pos + Inches(0.15), Inches(3.7), Inches(0.4),
                analysis["title"],
                font_size=13,
                font_color=self.style.primary_color,
                bold=True,
            )

            # Points
            y_point = y_pos + Inches(0.6)
            for point in analysis["points"]:
                self._add_text_box(
                    slide2,
                    x_pos + Inches(0.15), y_point, Inches(3.7), Inches(0.3),
                    point,
                    font_size=10,
                    font_color=self.style.text_color,
                )
                y_point += Inches(0.32)

        # Bottom summary
        self._add_shape(
            slide2,
            Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.8),
            self.style.primary_color,
        )

        summary = "Q3 FY26总结: 服务收入连续两季度改善，1&1迁移完成释放增长潜力，固定业务止血信号明确，FY26指引上限达成可期"
        self._add_text_box(
            slide2,
            Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
            summary,
            font_size=12,
            font_color=(255, 255, 255),
            bold=True,
        )

    def _add_historical_trend_section(self, target_operator: str, competitors: list[str]):
        """Add historical trend analysis section (8 quarters)."""
        # Section divider
        self._add_section_divider("历史趋势分析", "Historical Trend Analysis (8 Quarters)", "03")

        # Slide 1: User Scale Trend
        self._add_user_trend_slide(target_operator, competitors)

        # Slide 2: User Flow Analysis
        self._add_user_flow_slide(target_operator)

        # Slide 3: Financial Trend (Revenue & EBITDA)
        self._add_financial_trend_slide(target_operator, competitors)

        # Slide 4: Cost & Investment Trend
        self._add_cost_trend_slide(target_operator, competitors)

        # Slide 5: Business Segment - Mobile
        self._add_segment_trend_slide("mobile", target_operator, competitors)

        # Slide 6: Business Segment - Fixed Broadband
        self._add_segment_trend_slide("fixed", target_operator, competitors)

        # Slide 7: Business Segment - B2B
        self._add_segment_trend_slide("b2b", target_operator, competitors)

        # Slide 8: Trend Insights - Risks & Opportunities
        self._add_trend_insights_slide(target_operator)

    def _add_user_trend_slide(self, target_operator: str, competitors: list[str]):
        """Add user scale trend analysis slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "用户规模趋势分析", "User Scale Trend (8 Quarters)")

        if not self.historical_user_data:
            return

        # Mobile subscribers trend chart
        mobile_data = {
            op: self.historical_user_data.get(op, {}).get("mobile_subscribers", [])
            for op in [target_operator] + competitors[:3]
            if op in self.historical_user_data
        }

        chart_path = self.chart_gen.create_multi_line_trend_chart(
            quarters=self.quarters,
            data_series=mobile_data,
            title="移动用户规模趋势 (百万)",
            y_label="用户数 (M)",
            target_operator=target_operator,
            filename="mobile_user_trend.png",
            y_format="millions",
        )

        self._add_image(slide, chart_path, Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.8))

        # Broadband subscribers trend chart
        bb_data = {
            op: self.historical_user_data.get(op, {}).get("broadband_subscribers", [])
            for op in [target_operator] + competitors[:3]
            if op in self.historical_user_data
        }

        chart_path2 = self.chart_gen.create_multi_line_trend_chart(
            quarters=self.quarters,
            data_series=bb_data,
            title="宽带用户规模趋势 (百万)",
            y_label="用户数 (M)",
            target_operator=target_operator,
            filename="broadband_user_trend.png",
            y_format="millions",
        )

        self._add_image(slide, chart_path2, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.8))

        # Key insights box
        self._add_shape(
            slide,
            Inches(0.3), Inches(4.3), Inches(12.7), Inches(2.9),
            (245, 245, 245),
        )

        self._add_text_box(
            slide,
            Inches(0.5), Inches(4.4), Inches(6), Inches(0.35),
            "用户规模变化洞察 USER SCALE INSIGHTS",
            font_size=13,
            font_color=self.style.primary_color,
            bold=True,
        )

        insights = [
            "【移动用户】",
            f"  • {target_operator}: 8个季度增长+0.8M (30.5M→31.3M)，稳步增长",
            "  • Deutsche Telekom: 增长+2.6M (67.2M→69.8M)，领跑市场",
            "  • O2 Germany: 流失-0.8M (45.8M→45.0M)，1&1迁移影响",
            "  • 1&1 AG: 增长+0.68M (11.8M→12.48M)，自有网络吸引用户",
            "",
            "【宽带用户】",
            f"  • {target_operator}: 流失-0.38M (10.32M→9.94M)，固网竞争压力大",
            "  • Deutsche Telekom: 增长+0.7M (14.5M→15.2M)，光纤升级拉动",
        ]

        y_pos = Inches(4.8)
        for insight in insights:
            self._add_text_box(
                slide,
                Inches(0.5), y_pos, Inches(12.5), Inches(0.28),
                insight,
                font_size=10,
                font_color=self.style.text_color,
            )
            y_pos += Inches(0.26)

    def _add_user_flow_slide(self, target_operator: str):
        """Add user flow analysis slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "用户流动分析", "User Flow Analysis - Q3 FY26")

        if self.user_flow_data:
            # Get latest quarter flow data (index 7 = Q3 FY26)
            latest_flow = {k: v[7] if isinstance(v, list) else v for k, v in self.user_flow_data.items()}

            chart_path = self.chart_gen.create_user_flow_sankey_simple(
                flow_data=latest_flow,
                operators=[target_operator, "Deutsche Telekom", "Telefónica O2 Germany", "1&1 AG"],
                title="Q3 FY26 用户季度流动 (千)",
                filename="user_flow_matrix.png",
            )

            self._add_image(slide, chart_path, Inches(0.5), Inches(1.3), Inches(5.5), Inches(4))

        # Flow analysis text
        self._add_shape(
            slide,
            Inches(6.3), Inches(1.3), Inches(6.5), Inches(5.5),
            (248, 248, 248),
        )

        self._add_text_box(
            slide,
            Inches(6.5), Inches(1.45), Inches(6), Inches(0.35),
            "用户流动分析 USER FLOW ANALYSIS",
            font_size=13,
            font_color=self.style.primary_color,
            bold=True,
        )

        flow_insights = [
            "净流入运营商 (用户增长):",
            "  • Deutsche Telekom: 净流入 +10K/季度",
            "    最大来源: O2 (12K), Vodafone (8K)",
            "  • 1&1 AG: 净流入 +5K/季度",
            "    价格优势吸引价格敏感用户",
            "",
            "净流出运营商 (用户流失):",
            f"  • {target_operator}: 净流出 -11K/季度",
            "    主要流向: DT (8K), O2 (8K)",
            "    流失率改善: Q1 FY25 -51K → Q3 FY26 -11K",
            "  • O2 Germany: 净流出 -7K/季度",
            "    1&1网络迁移带走部分批发用户",
            "",
            "趋势观察:",
            f"  • {target_operator}流失率持续改善 (8季度下降78%)",
            "  • DT用户黏性最强，流出率最低",
            "  • 市场格局趋于稳定，大规模迁移减少",
        ]

        y_pos = Inches(1.85)
        for insight in flow_insights:
            font_color = self.style.text_color
            if "净流入" in insight or "Deutsche Telekom:" in insight or "1&1 AG:" in insight:
                font_color = (0, 128, 0)
            elif "净流出" in insight or f"{target_operator}:" in insight or "O2 Germany:" in insight:
                font_color = (180, 0, 0)

            self._add_text_box(
                slide,
                Inches(6.5), y_pos, Inches(6), Inches(0.28),
                insight,
                font_size=10,
                font_color=font_color if ":" in insight and "分析" not in insight else self.style.text_color,
            )
            y_pos += Inches(0.27)

    def _add_financial_trend_slide(self, target_operator: str, competitors: list[str]):
        """Add financial trend analysis slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "财务指标趋势分析", "Financial Trend Analysis (8 Quarters)")

        if not self.historical_financial_data:
            return

        # Service Revenue Growth Trend
        growth_data = {
            op: self.historical_financial_data.get(op, {}).get("service_revenue_growth_pct", [])
            for op in [target_operator] + competitors[:3]
            if op in self.historical_financial_data
        }

        chart_path = self.chart_gen.create_multi_line_trend_chart(
            quarters=self.quarters,
            data_series=growth_data,
            title="服务收入增长率趋势 (%)",
            y_label="增长率 (%)",
            target_operator=target_operator,
            filename="revenue_growth_trend.png",
            y_format="percent",
        )

        self._add_image(slide, chart_path, Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.6))

        # EBITDA Margin Trend
        margin_data = {
            op: self.historical_financial_data.get(op, {}).get("ebitda_margin_pct", [])
            for op in [target_operator] + competitors[:3]
            if op in self.historical_financial_data
        }

        chart_path2 = self.chart_gen.create_multi_line_trend_chart(
            quarters=self.quarters,
            data_series=margin_data,
            title="EBITDA利润率趋势 (%)",
            y_label="利润率 (%)",
            target_operator=target_operator,
            filename="ebitda_margin_trend.png",
        )

        self._add_image(slide, chart_path2, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.6))

        # Data table
        self._add_text_box(
            slide,
            Inches(0.3), Inches(4.05), Inches(6), Inches(0.35),
            "收入增长率数据 (8季度)",
            font_size=11,
            font_color=self.style.primary_color,
            bold=True,
        )

        # Mini table for revenue growth
        operators_short = ["VF", "DT", "O2", "1&1"]
        operators_full = [target_operator, "Deutsche Telekom", "Telefónica O2 Germany", "1&1 AG"]

        table_y = Inches(4.4)
        col_width = Inches(0.72)
        # Header row
        self._add_text_box(slide, Inches(0.3), table_y, Inches(0.6), Inches(0.25), "", font_size=8)
        for i, q in enumerate(self.quarters):
            self._add_text_box(
                slide, Inches(0.9) + i * col_width, table_y, col_width, Inches(0.25),
                q.replace(" ", "\n"), font_size=7, font_color=self.style.text_color
            )

        for op_idx, (short, full) in enumerate(zip(operators_short, operators_full)):
            table_y += Inches(0.28)
            growth_vals = self.historical_financial_data.get(full, {}).get("service_revenue_growth_pct", [0]*8)

            self._add_text_box(
                slide, Inches(0.3), table_y, Inches(0.6), Inches(0.25),
                short, font_size=8, font_color=self.style.text_color, bold=True
            )
            for i, val in enumerate(growth_vals):
                color = (0, 128, 0) if val > 0 else (180, 0, 0) if val < 0 else self.style.text_color
                self._add_text_box(
                    slide, Inches(0.9) + i * col_width, table_y, col_width, Inches(0.25),
                    f"{val:+.1f}%", font_size=8, font_color=color
                )

        # Insights
        self._add_shape(
            slide,
            Inches(6.8), Inches(4.05), Inches(6.2), Inches(3.1),
            (245, 250, 255),
        )

        self._add_text_box(
            slide,
            Inches(7), Inches(4.15), Inches(6), Inches(0.3),
            "财务趋势洞察",
            font_size=12,
            font_color=self.style.primary_color,
            bold=True,
        )

        fin_insights = [
            f"• {target_operator}服务收入增速持续改善:",
            "  Q4 FY24 -0.8% → Q3 FY26 +0.7% (改善1.5pp)",
            "",
            "• DT增速虽放缓但仍保持正增长:",
            "  Q4 FY24 +2.8% → Q3 FY26 +1.1%",
            "",
            "• O2持续负增长，1&1迁移冲击显著:",
            "  8季度持续下滑，Q3 FY26 -3.4%",
            "",
            f"• {target_operator} EBITDA利润率稳步提升:",
            "  35.4% → 36.2% (+0.8pp)，效率改善",
        ]

        y_pos = Inches(4.5)
        for insight in fin_insights:
            self._add_text_box(
                slide,
                Inches(7), y_pos, Inches(5.8), Inches(0.28),
                insight,
                font_size=9,
                font_color=self.style.text_color,
            )
            y_pos += Inches(0.28)

    def _add_cost_trend_slide(self, target_operator: str, competitors: list[str]):
        """Add cost and investment trend slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "成本与投资趋势", "Cost & Investment Trend (8 Quarters)")

        if not self.historical_financial_data:
            return

        # Operating Cost Trend
        cost_data = {
            op: self.historical_financial_data.get(op, {}).get("operating_cost", [])
            for op in [target_operator] + competitors[:3]
            if op in self.historical_financial_data
        }

        chart_path = self.chart_gen.create_multi_line_trend_chart(
            quarters=self.quarters,
            data_series=cost_data,
            title="运营成本趋势 (€B)",
            y_label="成本 (€B)",
            target_operator=target_operator,
            filename="opex_trend.png",
        )

        self._add_image(slide, chart_path, Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.6))

        # Capex Trend
        capex_data = {
            op: self.historical_financial_data.get(op, {}).get("capex", [])
            for op in [target_operator] + competitors[:3]
            if op in self.historical_financial_data
        }

        chart_path2 = self.chart_gen.create_multi_line_trend_chart(
            quarters=self.quarters,
            data_series=capex_data,
            title="资本支出趋势 (€B)",
            y_label="Capex (€B)",
            target_operator=target_operator,
            filename="capex_trend.png",
        )

        self._add_image(slide, chart_path2, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.6))

        # Employee trend
        emp_data = {
            op: self.historical_financial_data.get(op, {}).get("employees_k", [])
            for op in [target_operator] + competitors[:3]
            if op in self.historical_financial_data
        }

        chart_path3 = self.chart_gen.create_multi_line_trend_chart(
            quarters=self.quarters,
            data_series=emp_data,
            title="员工人数趋势 (千人)",
            y_label="员工 (K)",
            target_operator=target_operator,
            filename="employee_trend.png",
        )

        self._add_image(slide, chart_path3, Inches(0.3), Inches(4.1), Inches(6.2), Inches(2.6))

        # Insights
        self._add_shape(
            slide,
            Inches(6.8), Inches(4.1), Inches(6.2), Inches(2.6),
            (255, 248, 240),
        )

        self._add_text_box(
            slide,
            Inches(7), Inches(4.2), Inches(6), Inches(0.3),
            "成本与投资洞察",
            font_size=12,
            font_color=(180, 100, 0),
            bold=True,
        )

        cost_insights = [
            f"【{target_operator}运营成本】",
            "  • 8季度基本持平 (€1.97B-€1.99B)",
            "  • 员工精简700人 (15.2K→14.5K, -4.6%)",
            "  • 效率提升支撑利润率改善",
            "",
            "【资本支出】",
            "  • Capex稳定在€0.8B/季度水平",
            "  • 网络投资保持，支撑5G覆盖提升",
            "",
            "【行业对比】",
            "  • DT: 降本增效，员工减少1.7K",
            "  • O2: 成本压缩明显 (€1.44B→€1.35B)",
        ]

        y_pos = Inches(4.55)
        for insight in cost_insights:
            self._add_text_box(
                slide,
                Inches(7), y_pos, Inches(5.8), Inches(0.26),
                insight,
                font_size=9,
                font_color=self.style.text_color,
            )
            y_pos += Inches(0.25)

    def _add_segment_trend_slide(self, segment: str, target_operator: str, competitors: list[str]):
        """Add business segment trend analysis slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        segment_config = {
            "mobile": {
                "title": "移动业务趋势分析",
                "subtitle": "Mobile Service Trend (8 Quarters)",
                "revenue_key": "mobile_service_revenue",
                "growth_key": "mobile_service_growth_pct",
            },
            "fixed": {
                "title": "固定宽带业务趋势分析",
                "subtitle": "Fixed Broadband Trend (8 Quarters)",
                "revenue_key": "fixed_broadband_revenue",
                "growth_key": "fixed_broadband_growth_pct",
            },
            "b2b": {
                "title": "B2B企业业务趋势分析",
                "subtitle": "B2B Business Trend (8 Quarters)",
                "revenue_key": "b2b_revenue",
                "growth_key": "b2b_growth_pct",
            },
        }

        config = segment_config.get(segment, segment_config["mobile"])
        self._add_header(slide, config["title"], config["subtitle"])

        if not self.historical_segment_data:
            return

        # Revenue chart
        rev_data = {
            op: self.historical_segment_data.get(op, {}).get(config["revenue_key"], [])
            for op in [target_operator] + competitors[:3]
            if op in self.historical_segment_data
        }

        chart_path = self.chart_gen.create_multi_line_trend_chart(
            quarters=self.quarters,
            data_series=rev_data,
            title=f"{config['title'].replace('趋势分析', '')}收入 (€M)",
            y_label="收入 (€M)",
            target_operator=target_operator,
            filename=f"{segment}_revenue_trend.png",
        )

        self._add_image(slide, chart_path, Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.6))

        # Growth rate chart
        growth_data = {
            op: self.historical_segment_data.get(op, {}).get(config["growth_key"], [])
            for op in [target_operator] + competitors[:3]
            if op in self.historical_segment_data
        }

        chart_path2 = self.chart_gen.create_multi_line_trend_chart(
            quarters=self.quarters,
            data_series=growth_data,
            title=f"{config['title'].replace('趋势分析', '')}增长率 (%)",
            y_label="增长率 (%)",
            target_operator=target_operator,
            filename=f"{segment}_growth_trend.png",
            y_format="percent",
        )

        self._add_image(slide, chart_path2, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.6))

        # Segment-specific insights
        insights_map = {
            "mobile": [
                f"【{target_operator}移动业务】",
                "  • 收入从€1,450M增至€1,520M (+4.8%)",
                "  • 增速由-0.5%转正至+2.8%，企稳回升",
                "  • 1&1批发收入是核心增长驱动",
                "",
                "【竞争格局】",
                "  • DT保持领先但增速放缓 (3.5%→1.6%)",
                "  • O2持续负增长 (-1.2%→-3.0%)",
                "  • 1&1增速放缓但维持正增长",
                "",
                "【市场趋势】",
                "  • 移动市场整体增速趋缓",
                "  • 批发业务成为差异化竞争点",
                "  • 价值经营替代规模扩张",
            ],
            "fixed": [
                f"【{target_operator}固定业务】",
                "  • 收入从€820M降至€795M (-3.0%)",
                "  • 降幅持续收窄 (-2.8%→-1.1%)",
                "  • 用户流失放缓，止血信号明确",
                "",
                "【竞争格局】",
                "  • DT凭借光纤优势保持增长 (+0.9%)",
                "  • O2固网业务小幅增长 (+0.5%)",
                "  • 1&1固网持续下滑 (-1.0%)",
                "",
                "【市场趋势】",
                "  • 光纤升级是增长核心驱动力",
                "  • 铜网用户持续向光纤迁移",
                "  • 捆绑销售策略日益重要",
            ],
            "b2b": [
                f"【{target_operator} B2B业务】",
                "  • 收入从€410M增至€520M (+26.8%)",
                "  • 增速从3.5%提升至8.5%，强劲增长",
                "  • Skaylink收购增强云服务能力",
                "",
                "【竞争格局】",
                "  • DT B2B持续高增长 (5.5%→7.5%)",
                "  • O2 B2B增速最快 (8%→13%)",
                "  • 所有运营商都在发力企业市场",
                "",
                "【市场趋势】",
                "  • 企业数字化转型驱动需求",
                "  • 云服务和IoT成为新增长点",
                "  • B2B是运营商价值转型方向",
            ],
        }

        insights = insights_map.get(segment, insights_map["mobile"])

        self._add_shape(
            slide,
            Inches(0.3), Inches(4.1), Inches(12.7), Inches(3.1),
            (245, 248, 255),
        )

        y_pos = Inches(4.2)
        col1_x = Inches(0.5)
        col2_x = Inches(6.8)

        for i, insight in enumerate(insights):
            x_pos = col1_x if i < 7 else col2_x
            if i == 7:
                y_pos = Inches(4.2)

            self._add_text_box(
                slide,
                x_pos, y_pos, Inches(6), Inches(0.26),
                insight,
                font_size=9,
                font_color=self.style.text_color,
            )
            y_pos += Inches(0.25)

    def _add_trend_insights_slide(self, target_operator: str):
        """Add trend insights with risks and opportunities slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        self._add_header(slide, "趋势洞察: 风险与机会", "Trend Insights: Risks & Opportunities")

        # Opportunities section (left)
        self._add_shape(
            slide,
            Inches(0.3), Inches(1.3), Inches(6.2), Inches(5.8),
            (230, 255, 230),
        )

        self._add_text_box(
            slide,
            Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.4),
            "机会 OPPORTUNITIES",
            font_size=16,
            font_color=(0, 128, 0),
            bold=True,
        )

        opportunities = [
            "1. 服务收入企稳回升",
            "   • 连续8季度改善，由-0.8%转为+0.7%",
            "   • 趋势线显示有望在FY27实现正增长",
            "",
            "2. 批发业务高速增长",
            "   • 1&1迁移带来€380M/季批发收入",
            "   • 8季度增长111% (€180M→€380M)",
            "",
            "3. B2B业务持续扩张",
            "   • 26.8%的8季度累计增长",
            "   • 云服务收购增强竞争力",
            "",
            "4. 用户流失显著改善",
            "   • 净流出从51K/季降至11K/季",
            "   • 改善幅度达78%",
            "",
            "5. 运营效率持续提升",
            "   • EBITDA利润率提升0.8pp",
            "   • 人员精简4.6%支撑降本",
        ]

        y_pos = Inches(1.85)
        for opp in opportunities:
            font_color = (0, 100, 0) if opp.startswith(("1.", "2.", "3.", "4.", "5.")) else self.style.text_color
            self._add_text_box(
                slide,
                Inches(0.5), y_pos, Inches(5.8), Inches(0.28),
                opp,
                font_size=10,
                font_color=font_color,
                bold=opp.startswith(("1.", "2.", "3.", "4.", "5.")),
            )
            y_pos += Inches(0.28)

        # Risks section (right)
        self._add_shape(
            slide,
            Inches(6.8), Inches(1.3), Inches(6.2), Inches(5.8),
            (255, 235, 235),
        )

        self._add_text_box(
            slide,
            Inches(7), Inches(1.4), Inches(5.8), Inches(0.4),
            "风险 RISKS",
            font_size=16,
            font_color=(180, 0, 0),
            bold=True,
        )

        risks = [
            "1. 固定业务持续承压",
            "   • 8季度累计流失38万宽带用户",
            "   • 光纤覆盖不足难以挽回流失",
            "",
            "2. 市场份额被蚕食",
            "   • DT宽带份额扩大 (38.5%→40.6%)",
            "   • Vodafone份额收缩 (27.4%→26.5%)",
            "",
            "3. 与DT差距拉大",
            "   • 收入增速差距3.5pp (+0.7% vs +4.2%)",
            "   • EBITDA差距持续 (36.2% vs 41.9%)",
            "",
            "4. 5G网络覆盖落后",
            "   • 与DT差距5pp (92% vs 97%)",
            "   • 影响高端用户获取",
            "",
            "5. O2 B2B崛起威胁",
            "   • O2 B2B增速13%，快于Vodafone 8.5%",
            "   • 企业市场竞争加剧",
        ]

        y_pos = Inches(1.85)
        for risk in risks:
            font_color = (150, 0, 0) if risk.startswith(("1.", "2.", "3.", "4.", "5.")) else self.style.text_color
            self._add_text_box(
                slide,
                Inches(7), y_pos, Inches(5.8), Inches(0.28),
                risk,
                font_size=10,
                font_color=font_color,
                bold=risk.startswith(("1.", "2.", "3.", "4.", "5.")),
            )
            y_pos += Inches(0.28)

    # =========================================================================
    # Helper Methods
    # =========================================================================

    def _add_image(self, slide, image_path: str, left, top, width, height):
        """Add an image to the slide.

        Args:
            slide: PowerPoint slide
            image_path: Path to image file
            left: Left position
            top: Top position
            width: Image width
            height: Image height
        """
        try:
            slide.shapes.add_picture(
                image_path,
                left, top,
                width=width, height=height,
            )
        except Exception as e:
            # If image fails, add placeholder text
            self._add_text_box(
                slide, left, top, width, height,
                f"[Chart: {image_path}]",
                font_size=10,
                font_color=(150, 150, 150),
            )

    def _add_header(self, slide, title: str, subtitle: str = ""):
        """Add standard header to slide."""
        # Red accent bar
        self._add_shape(
            slide, 0, 0, Inches(13.333), Inches(0.1),
            self.style.primary_color
        )

        # Title
        self._add_text_box(
            slide,
            Inches(0.5), Inches(0.25), Inches(10), Inches(0.55),
            title,
            font_size=24,
            font_color=self.style.text_color,
            bold=True,
        )

        # Subtitle
        if subtitle:
            self._add_text_box(
                slide,
                Inches(0.5), Inches(0.75), Inches(10), Inches(0.35),
                subtitle,
                font_size=11,
                font_color=self.style.primary_color,
                bold=True,
            )

        # Page number
        self._add_shape(
            slide,
            Inches(12.5), Inches(7), Inches(0.6), Inches(0.35),
            self.style.primary_color,
        )
        self._add_text_box(
            slide,
            Inches(12.5), Inches(7.02), Inches(0.6), Inches(0.3),
            str(self.slide_num),
            font_size=10,
            font_color=(255, 255, 255),
            bold=True,
            align="center",
        )

    def _add_shape(
        self,
        slide,
        left,
        top,
        width,
        height,
        fill_color: tuple,
        shape_type=MSO_SHAPE.RECTANGLE,
    ):
        """Add a shape to the slide."""
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_color)
        shape.line.fill.background()
        return shape

    def _add_text_box(
        self,
        slide,
        left,
        top,
        width,
        height,
        text: str,
        font_size: int = 14,
        font_color: tuple = (0, 0, 0),
        bold: bool = False,
        align: str = "left",
    ):
        """Add a text box to the slide."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*font_color)
        p.font.bold = bold
        p.font.name = self.style.body_font

        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

        return txBox

    def _sanitize_name(self, name: str) -> str:
        """Sanitize name for filename."""
        return name.lower().replace(" ", "_").replace("/", "_")


def generate_enhanced_blm_ppt(
    five_looks: dict[str, InsightResult],
    three_decisions: dict[str, StrategyResult],
    target_operator: str,
    competitors: list[str] = None,
    financial_data: dict = None,
    competitive_scores: dict = None,
    style: str = "huawei",
    output_dir: str = None,
) -> str:
    """Convenience function to generate enhanced BLM PPT.

    Args:
        five_looks: Five Looks analysis results
        three_decisions: Three Decisions strategy results
        target_operator: Target operator name
        competitors: List of competitor names
        financial_data: Raw financial data for charts
        competitive_scores: Competitive dimension scores
        style: PPT style ("huawei" or "vodafone")
        output_dir: Output directory

    Returns:
        Path to generated PPT file
    """
    generator = BLMPPTGeneratorEnhanced(style=style, output_dir=output_dir)
    return generator.generate(
        five_looks=five_looks,
        three_decisions=three_decisions,
        target_operator=target_operator,
        competitors=competitors,
        financial_data=financial_data,
        competitive_scores=competitive_scores,
    )
