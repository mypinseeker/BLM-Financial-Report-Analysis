"""综合分析PPT生成模块

基于germany_market_comprehensive_data.py数据生成详细分析PPT
包含: 经营分析、业务分析、战略分析、组织分析、网络分析、资费对比
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
import tempfile

from .ppt_charts import PPTChartGenerator
from . import germany_market_comprehensive_data as data


class ComprehensiveAnalysisPPT:
    """生成综合分析PPT"""

    def __init__(self, output_dir: str = None):
        self.output_dir = Path(output_dir) if output_dir else Path(tempfile.mkdtemp())
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # 初始化
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        self.chart_gen = PPTChartGenerator(str(self.output_dir))
        self.slide_num = 0

        # 华为配色
        self.primary_color = (199, 0, 11)  # Huawei Red
        self.text_color = (51, 51, 51)
        self.light_text = (102, 102, 102)

        # 运营商
        self.operators = ["Vodafone Germany", "Deutsche Telekom Germany",
                         "Telefónica O2 Germany", "1&1 AG"]
        self.target = "Vodafone Germany"

    def generate(self, output_filename: str = "comprehensive_analysis.pptx") -> str:
        """生成完整的综合分析PPT"""

        # 封面
        self._add_cover_slide()

        # 第一部分: 经营分析
        self._add_section_divider("经营层面分析", "Operating Analysis", "01")
        self._add_revenue_analysis_slides()
        self._add_profitability_analysis_slides()
        self._add_investment_analysis_slides()

        # 第二部分: 业务分析
        self._add_section_divider("业务层面分析", "Business Segment Analysis", "02")
        self._add_mobile_business_slides()
        self._add_fixed_business_slides()
        self._add_b2b_business_slides()
        self._add_tv_fmc_slides()

        # 第三部分: 战略与组织
        self._add_section_divider("战略与组织分析", "Strategy & Organization", "03")
        self._add_strategy_analysis_slides()
        self._add_executive_changes_slide()

        # 第四部分: 网络与资费
        self._add_section_divider("网络与资费分析", "Network & Tariff Analysis", "04")
        self._add_network_analysis_slides()
        self._add_tariff_comparison_slides()

        # 第五部分: 总结
        self._add_summary_slide()

        # 保存
        output_path = self.output_dir / output_filename
        self.prs.save(str(output_path))
        return str(output_path)

    # ==========================================================================
    # 基础方法
    # ==========================================================================

    def _add_shape(self, slide, left, top, width, height, color):
        """添加形状"""
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        if isinstance(color, tuple):
            shape.fill.fore_color.rgb = RGBColor(*color)
        shape.line.fill.background()
        return shape

    def _add_text_box(self, slide, left, top, width, height, text,
                      font_size=12, font_color=None, bold=False, align="left"):
        """添加文本框"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        if font_color:
            if isinstance(font_color, tuple):
                p.font.color.rgb = RGBColor(*font_color)
        p.font.bold = bold
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        return txBox

    def _add_header(self, slide, title: str, subtitle: str = ""):
        """添加页眉"""
        # 红色顶条
        self._add_shape(slide, 0, 0, Inches(13.333), Inches(0.08), self.primary_color)

        # 标题
        self._add_text_box(slide, Inches(0.5), Inches(0.25), Inches(10), Inches(0.5),
                          title, font_size=24, font_color=self.text_color, bold=True)

        if subtitle:
            self._add_text_box(slide, Inches(0.5), Inches(0.7), Inches(10), Inches(0.3),
                              subtitle, font_size=12, font_color=self.light_text)

        # 页码
        self._add_text_box(slide, Inches(12), Inches(0.25), Inches(1), Inches(0.3),
                          str(self.slide_num), font_size=14,
                          font_color=self.primary_color, bold=True, align="right")

    def _add_image(self, slide, image_path, left, top, width, height):
        """添加图片"""
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)

    # ==========================================================================
    # 封面和分隔页
    # ==========================================================================

    def _add_cover_slide(self):
        """添加封面"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        # 红色背景
        self._add_shape(slide, 0, 0, Inches(13.333), Inches(7.5), self.primary_color)

        # 装饰
        self._add_shape(slide, Inches(10), 0, Inches(3.333), Inches(7.5), (170, 0, 10))

        # 标题
        self._add_text_box(slide, Inches(0.8), Inches(2.2), Inches(9), Inches(1.2),
                          "德国电信市场综合分析",
                          font_size=48, font_color=(255, 255, 255), bold=True)

        self._add_text_box(slide, Inches(0.8), Inches(3.5), Inches(9), Inches(0.6),
                          "Germany Telecom Market Comprehensive Analysis",
                          font_size=24, font_color=(255, 200, 200))

        # 分隔线
        self._add_shape(slide, Inches(0.8), Inches(4.3), Inches(5), Inches(0.03), (255, 255, 255))

        # 副标题
        self._add_text_box(slide, Inches(0.8), Inches(4.6), Inches(8), Inches(0.8),
                          "8季度趋势分析 | 经营/业务/战略/网络/资费全维度\n"
                          "数据范围: Q4 FY24 - Q3 FY26",
                          font_size=16, font_color=(255, 220, 220))

        # 日期
        self._add_text_box(slide, Inches(0.8), Inches(6.5), Inches(4), Inches(0.4),
                          "2026年2月",
                          font_size=14, font_color=(255, 200, 200))

    def _add_section_divider(self, title: str, subtitle: str, number: str):
        """添加章节分隔页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1

        # 红色背景
        self._add_shape(slide, 0, 0, Inches(13.333), Inches(7.5), self.primary_color)
        self._add_shape(slide, Inches(10), 0, Inches(3.333), Inches(7.5), (170, 0, 10))

        # 章节号
        self._add_text_box(slide, Inches(0.5), Inches(2), Inches(2), Inches(1.2),
                          number, font_size=80, font_color=(255, 255, 255), bold=True)

        # 分隔线
        self._add_shape(slide, Inches(0.5), Inches(3.4), Inches(5), Inches(0.03), (255, 255, 255))

        # 标题
        self._add_text_box(slide, Inches(0.5), Inches(3.7), Inches(9), Inches(1),
                          title, font_size=52, font_color=(255, 255, 255), bold=True)

        # 副标题
        self._add_text_box(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(0.6),
                          subtitle, font_size=24, font_color=(255, 200, 200))

    # ==========================================================================
    # 第一部分: 经营分析
    # ==========================================================================

    def _add_revenue_analysis_slides(self):
        """添加收入分析页面"""
        # 页面1: 总收入和服务收入对比
        self._add_revenue_overview_slide()
        # 页面2: 收入增长率趋势
        self._add_revenue_growth_slide()
        # 页面3: 收入结构分析
        self._add_revenue_structure_slide()

    def _add_revenue_overview_slide(self):
        """收入规模对比"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1
        self._add_header(slide, "收入规模对比分析", "Revenue Scale Comparison (8 Quarters)")

        # 总收入趋势图
        revenue_data = {}
        for op in self.operators:
            revenue_data[op] = [r/1000 for r in data.REVENUE_DATA_8Q[op]["total_revenue"]]

        chart_path = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=revenue_data,
            title="总收入趋势 (€B)",
            y_label="收入 (€B)",
            target_operator=self.target,
            filename="revenue_total_trend.png",
        )
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2), Inches(6.2), Inches(2.8))

        # 服务收入趋势图
        svc_data = {}
        for op in self.operators:
            svc_data[op] = [r/1000 for r in data.REVENUE_DATA_8Q[op]["service_revenue"]]

        chart_path2 = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=svc_data,
            title="服务收入趋势 (€B)",
            y_label="收入 (€B)",
            target_operator=self.target,
            filename="revenue_service_trend.png",
        )
        self._add_image(slide, chart_path2, Inches(6.8), Inches(1.2), Inches(6.2), Inches(2.8))

        # 关键洞察
        self._add_shape(slide, Inches(0.3), Inches(4.2), Inches(12.7), Inches(3), (245, 245, 250))

        self._add_text_box(slide, Inches(0.5), Inches(4.3), Inches(6), Inches(0.35),
                          "收入洞察 REVENUE INSIGHTS", font_size=13,
                          font_color=self.primary_color, bold=True)

        # Q3 FY26数据
        vf_rev = data.REVENUE_DATA_8Q["Vodafone Germany"]["total_revenue"][-1]
        dt_rev = data.REVENUE_DATA_8Q["Deutsche Telekom Germany"]["total_revenue"][-1]
        o2_rev = data.REVENUE_DATA_8Q["Telefónica O2 Germany"]["total_revenue"][-1]

        insights = [
            f"【市场格局】DT: €{dt_rev/1000:.1f}B (51%) > VF: €{vf_rev/1000:.1f}B (25%) > O2: €{o2_rev/1000:.1f}B (16%) > 1&1: €1.0B (8%)",
            "",
            "【Vodafone Germany表现】",
            f"  • 总收入: €{vf_rev}M，8季度增长+1.4% (€3050M→€3092M)",
            "  • 服务收入增速: Q4 FY24 -0.8% → Q3 FY26 +0.7%，扭亏为盈",
            "  • 移动服务收入持续增长，固网服务承压",
            "",
            "【竞争对手动态】",
            "  • DT Germany: 收入增速放缓但仍保持正增长 (+1.1%)",
            "  • O2 Germany: 1&1批发用户流失导致收入下滑 (-3.4%)",
            "  • 1&1 AG: 自建网络投资期，收入增长乏力 (+0.1%)",
        ]

        y_pos = Inches(4.7)
        for insight in insights:
            self._add_text_box(slide, Inches(0.5), y_pos, Inches(12.5), Inches(0.28),
                              insight, font_size=10, font_color=self.text_color)
            y_pos += Inches(0.24)

    def _add_revenue_growth_slide(self):
        """收入增长率趋势"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1
        self._add_header(slide, "收入增长率趋势分析", "Revenue Growth Rate Trend")

        # 服务收入增长率
        growth_data = {op: data.REVENUE_DATA_8Q[op]["service_revenue_growth_pct"] for op in self.operators}

        chart_path = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=growth_data,
            title="服务收入YoY增长率 (%)",
            y_label="增长率 (%)",
            target_operator=self.target,
            filename="service_revenue_growth.png",
            y_format="percent",
        )
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2), Inches(6.2), Inches(2.8))

        # 移动收入增长率
        mobile_growth = {op: data.REVENUE_DATA_8Q[op]["mobile_service_growth_pct"] for op in self.operators}

        chart_path2 = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=mobile_growth,
            title="移动服务收入增长率 (%)",
            y_label="增长率 (%)",
            target_operator=self.target,
            filename="mobile_revenue_growth.png",
            y_format="percent",
        )
        self._add_image(slide, chart_path2, Inches(6.8), Inches(1.2), Inches(6.2), Inches(2.8))

        # 数据表格
        self._add_text_box(slide, Inches(0.3), Inches(4.2), Inches(6), Inches(0.3),
                          "服务收入增长率 (季度数据)", font_size=11,
                          font_color=self.primary_color, bold=True)

        # 表格绘制
        self._draw_growth_table(slide, Inches(0.3), Inches(4.5), growth_data)

        # 洞察框
        self._add_shape(slide, Inches(6.8), Inches(4.2), Inches(6.2), Inches(3), (255, 248, 240))

        self._add_text_box(slide, Inches(7), Inches(4.3), Inches(6), Inches(0.3),
                          "增长趋势洞察", font_size=12, font_color=(180, 100, 0), bold=True)

        insights = [
            "Vodafone Germany增长轨迹:",
            "  • 服务收入: -0.8% → +0.7% (改善1.5pp)",
            "  • 移动收入: -0.5% → +2.8% (改善3.3pp)",
            "  • 固网收入: -2.8% → -1.1% (改善1.7pp)",
            "",
            "行业对比:",
            "  • VF移动增速领先O2 (+5.8pp)",
            "  • VF固网仍落后DT (-2.0pp)",
            "  • B2B增速最快 (+8.5%)",
        ]

        y_pos = Inches(4.65)
        for insight in insights:
            self._add_text_box(slide, Inches(7), y_pos, Inches(6), Inches(0.26),
                              insight, font_size=10, font_color=self.text_color)
            y_pos += Inches(0.26)

    def _draw_growth_table(self, slide, left, top, growth_data):
        """绘制增长率表格"""
        col_width = Inches(0.72)
        row_height = Inches(0.26)

        # 表头
        headers = [""] + [q.replace(" ", "\n") for q in data.QUARTERS_8Q]
        for i, h in enumerate(headers):
            x = left + i * col_width
            self._add_text_box(slide, x, top, col_width, row_height,
                              h, font_size=7, font_color=self.text_color, bold=True)

        # 数据行
        short_names = {"Vodafone Germany": "VF", "Deutsche Telekom Germany": "DT",
                      "Telefónica O2 Germany": "O2", "1&1 AG": "1&1"}

        for row_idx, op in enumerate(self.operators):
            y = top + (row_idx + 1) * row_height
            self._add_text_box(slide, left, y, col_width, row_height,
                              short_names[op], font_size=8, font_color=self.text_color, bold=True)

            for col_idx, val in enumerate(growth_data[op]):
                x = left + (col_idx + 1) * col_width
                color = (0, 128, 0) if val > 0 else (180, 0, 0) if val < 0 else self.text_color
                self._add_text_box(slide, x, y, col_width, row_height,
                                  f"{val:+.1f}%", font_size=8, font_color=color)

    def _add_revenue_structure_slide(self):
        """收入结构分析"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1
        self._add_header(slide, "收入结构对比分析", "Revenue Structure Comparison")

        # Q3 FY26收入结构对比
        self._add_text_box(slide, Inches(0.3), Inches(1.2), Inches(6), Inches(0.3),
                          "Q3 FY26 收入结构 (€M)", font_size=13,
                          font_color=self.primary_color, bold=True)

        # 收入结构表格
        y_base = Inches(1.6)
        headers = ["运营商", "总收入", "移动服务", "固网服务", "B2B", "TV", "批发"]
        col_widths = [Inches(1.8), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2), Inches(0.9), Inches(0.9)]

        # 表头行
        x = Inches(0.3)
        for i, h in enumerate(headers):
            self._add_shape(slide, x, y_base, col_widths[i], Inches(0.35), (220, 220, 225))
            self._add_text_box(slide, x + Inches(0.05), y_base + Inches(0.05),
                              col_widths[i], Inches(0.25), h,
                              font_size=10, font_color=self.text_color, bold=True)
            x += col_widths[i]

        # 数据行
        for row_idx, op in enumerate(self.operators):
            y = y_base + Inches(0.4) + row_idx * Inches(0.45)
            rev = data.REVENUE_DATA_8Q[op]
            row_data = [
                op.replace(" Germany", "").replace("Telefónica ", ""),
                f"€{rev['total_revenue'][-1]:,}",
                f"€{rev['mobile_service_revenue'][-1]:,}",
                f"€{rev['fixed_service_revenue'][-1]:,}",
                f"€{rev['b2b_revenue'][-1]:,}",
                f"€{rev['tv_revenue'][-1]:,}",
                f"€{rev['wholesale_revenue'][-1]:,}",
            ]

            bg_color = (255, 245, 245) if op == self.target else (248, 248, 248)
            x = Inches(0.3)
            for i, val in enumerate(row_data):
                self._add_shape(slide, x, y, col_widths[i], Inches(0.4), bg_color)
                font_weight = i == 0
                self._add_text_box(slide, x + Inches(0.05), y + Inches(0.08),
                                  col_widths[i], Inches(0.28), val,
                                  font_size=10, font_color=self.text_color, bold=font_weight)
                x += col_widths[i]

        # 收入结构饼图说明
        self._add_text_box(slide, Inches(0.3), Inches(3.8), Inches(6), Inches(0.3),
                          "Vodafone Germany收入结构变化", font_size=12,
                          font_color=self.primary_color, bold=True)

        # 收入结构变化分析
        self._add_shape(slide, Inches(0.3), Inches(4.15), Inches(6.2), Inches(3.0), (245, 250, 255))

        structure_insights = [
            "Q4 FY24 → Q3 FY26 结构变化:",
            "",
            "【移动服务】 47.5% → 49.2% (+1.7pp)",
            "  增长驱动: 后付费用户增长, ARPU提升",
            "",
            "【固网服务】 26.9% → 25.7% (-1.2pp)",
            "  下滑原因: 用户流失, Cable竞争",
            "",
            "【B2B业务】 13.4% → 16.8% (+3.4pp)",
            "  增长亮点: 数字化服务快速增长",
            "",
            "【批发收入】 5.9% → 12.3% (+6.4pp)",
            "  1&1 MVNO批发贡献显著增长",
        ]

        y_pos = Inches(4.25)
        for insight in structure_insights:
            self._add_text_box(slide, Inches(0.5), y_pos, Inches(6), Inches(0.25),
                              insight, font_size=10, font_color=self.text_color)
            y_pos += Inches(0.23)

        # 右侧: B2B增长分析
        self._add_text_box(slide, Inches(6.8), Inches(3.8), Inches(6), Inches(0.3),
                          "B2B业务高增长分析", font_size=12,
                          font_color=self.primary_color, bold=True)

        self._add_shape(slide, Inches(6.8), Inches(4.15), Inches(6.2), Inches(3.0), (240, 255, 240))

        b2b_insights = [
            "【Vodafone Germany B2B增长 +8.5%】",
            "  • 连接服务: 企业移动、固定连接",
            "  • 数字化服务: IoT、云、安全、边缘计算",
            "  • 行业方案: 制造业4.0、智慧城市",
            "",
            "【竞争格局】",
            f"  • DT B2B: €{data.REVENUE_DATA_8Q['Deutsche Telekom Germany']['b2b_revenue'][-1]}M (+7.5%)",
            f"  • VF B2B: €{data.REVENUE_DATA_8Q['Vodafone Germany']['b2b_revenue'][-1]}M (+8.5%)",
            f"  • O2 B2B: €{data.REVENUE_DATA_8Q['Telefónica O2 Germany']['b2b_revenue'][-1]}M (+13.0%)",
            "",
            "【机会】",
            "  • O2 B2B增速最快但基数小",
            "  • 数字化转型需求强劲",
            "  • 5G专网/边缘计算是新增长点",
        ]

        y_pos = Inches(4.25)
        for insight in b2b_insights:
            self._add_text_box(slide, Inches(7), y_pos, Inches(6), Inches(0.25),
                              insight, font_size=10, font_color=self.text_color)
            y_pos += Inches(0.23)

    def _add_profitability_analysis_slides(self):
        """盈利分析页面"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1
        self._add_header(slide, "盈利能力分析", "Profitability Analysis (EBITDA)")

        # EBITDA绝对值趋势
        ebitda_data = {}
        for op in self.operators:
            ebitda_data[op] = [e/1000 for e in data.PROFITABILITY_DATA_8Q[op]["ebitda"]]

        chart_path = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=ebitda_data,
            title="EBITDA趋势 (€B)",
            y_label="EBITDA (€B)",
            target_operator=self.target,
            filename="ebitda_trend.png",
        )
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2), Inches(6.2), Inches(2.6))

        # EBITDA Margin趋势
        margin_data = {op: data.PROFITABILITY_DATA_8Q[op]["ebitda_margin_pct"] for op in self.operators}

        chart_path2 = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=margin_data,
            title="EBITDA Margin趋势 (%)",
            y_label="Margin (%)",
            target_operator=self.target,
            filename="ebitda_margin_trend.png",
        )
        self._add_image(slide, chart_path2, Inches(6.8), Inches(1.2), Inches(6.2), Inches(2.6))

        # EBITDA增长率趋势
        growth_data = {op: data.PROFITABILITY_DATA_8Q[op]["ebitda_growth_pct"] for op in self.operators}

        chart_path3 = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=growth_data,
            title="EBITDA YoY增长率 (%)",
            y_label="增长率 (%)",
            target_operator=self.target,
            filename="ebitda_growth_trend.png",
            y_format="percent",
        )
        self._add_image(slide, chart_path3, Inches(0.3), Inches(4.0), Inches(6.2), Inches(2.6))

        # 盈利洞察
        self._add_shape(slide, Inches(6.8), Inches(4.0), Inches(6.2), Inches(2.6), (245, 250, 255))

        self._add_text_box(slide, Inches(7), Inches(4.1), Inches(6), Inches(0.3),
                          "盈利能力洞察", font_size=12, font_color=self.primary_color, bold=True)

        vf_margin = data.PROFITABILITY_DATA_8Q["Vodafone Germany"]["ebitda_margin_pct"]
        dt_margin = data.PROFITABILITY_DATA_8Q["Deutsche Telekom Germany"]["ebitda_margin_pct"]

        insights = [
            f"【EBITDA Margin对比 Q3 FY26】",
            f"  • DT Germany: {dt_margin[-1]}% (行业领先)",
            f"  • VF Germany: {vf_margin[-1]}% (差距: -{dt_margin[-1]-vf_margin[-1]:.1f}pp)",
            f"  • O2 Germany: {data.PROFITABILITY_DATA_8Q['Telefónica O2 Germany']['ebitda_margin_pct'][-1]}%",
            f"  • 1&1 AG: {data.PROFITABILITY_DATA_8Q['1&1 AG']['ebitda_margin_pct'][-1]}% (网络投资期)",
            "",
            f"【VF Germany盈利改善】",
            f"  • Margin: {vf_margin[0]}% → {vf_margin[-1]}% (+{vf_margin[-1]-vf_margin[0]:.1f}pp)",
            f"  • EBITDA增长: -2.0% → +1.5% (扭亏)",
            "  • 效率提升: 员工减少700人，成本优化",
        ]

        y_pos = Inches(4.45)
        for insight in insights:
            self._add_text_box(slide, Inches(7), y_pos, Inches(6), Inches(0.24),
                              insight, font_size=10, font_color=self.text_color)
            y_pos += Inches(0.23)

    def _add_investment_analysis_slides(self):
        """投资分析页面 (OPEX/CAPEX)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1
        self._add_header(slide, "投资与成本分析", "Investment Analysis (OPEX & CAPEX)")

        # OPEX/Revenue比例趋势
        opex_ratio = {op: data.INVESTMENT_DATA_8Q[op]["opex_to_revenue_pct"] for op in self.operators}

        chart_path = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=opex_ratio,
            title="OPEX/Revenue比例 (%)",
            y_label="占比 (%)",
            target_operator=self.target,
            filename="opex_ratio_trend.png",
        )
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2), Inches(6.2), Inches(2.6))

        # CAPEX/Revenue比例趋势
        capex_ratio = {op: data.INVESTMENT_DATA_8Q[op]["capex_to_revenue_pct"] for op in self.operators}

        chart_path2 = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=capex_ratio,
            title="CAPEX/Revenue比例 (%)",
            y_label="占比 (%)",
            target_operator=self.target,
            filename="capex_ratio_trend.png",
        )
        self._add_image(slide, chart_path2, Inches(6.8), Inches(1.2), Inches(6.2), Inches(2.6))

        # 员工数趋势
        emp_data = {op: data.INVESTMENT_DATA_8Q[op]["employees_k"] for op in self.operators}

        chart_path3 = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=emp_data,
            title="员工人数趋势 (千人)",
            y_label="员工 (K)",
            target_operator=self.target,
            filename="employees_trend.png",
        )
        self._add_image(slide, chart_path3, Inches(0.3), Inches(4.0), Inches(6.2), Inches(2.6))

        # 投资洞察
        self._add_shape(slide, Inches(6.8), Inches(4.0), Inches(6.2), Inches(2.6), (255, 248, 240))

        self._add_text_box(slide, Inches(7), Inches(4.1), Inches(6), Inches(0.3),
                          "投资效率洞察", font_size=12, font_color=(180, 100, 0), bold=True)

        vf_opex = data.INVESTMENT_DATA_8Q["Vodafone Germany"]["opex_to_revenue_pct"]
        vf_capex = data.INVESTMENT_DATA_8Q["Vodafone Germany"]["capex_to_revenue_pct"]
        dt_capex = data.INVESTMENT_DATA_8Q["Deutsche Telekom Germany"]["capex_to_revenue_pct"]

        insights = [
            "【OPEX效率】",
            f"  • VF OPEX占比: {vf_opex[0]}% → {vf_opex[-1]}% (优化{vf_opex[0]-vf_opex[-1]:.1f}pp)",
            "  • 员工优化: 15.2K → 14.5K (-700人)",
            "  • DT OPEX最优: 57.9% (规模效应)",
            "",
            "【CAPEX投资】",
            f"  • VF: {vf_capex[-1]}% (保持较高投资)",
            f"  • DT: {dt_capex[-1]}% (投资效率更高)",
            "  • 1&1: 38.6% (自建网络高投入期)",
            "",
            "【效率差距】",
            "  • VF成本结构仍待优化",
            "  • 数字化/自动化是关键",
        ]

        y_pos = Inches(4.45)
        for insight in insights:
            self._add_text_box(slide, Inches(7), y_pos, Inches(6), Inches(0.23),
                              insight, font_size=10, font_color=self.text_color)
            y_pos += Inches(0.22)

    # ==========================================================================
    # 第二部分: 业务分析
    # ==========================================================================

    def _add_mobile_business_slides(self):
        """移动业务分析"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1
        self._add_header(slide, "移动业务分析", "Mobile Business Analysis")

        # 消费者移动用户趋势 (不含IoT)
        mobile_subs = {op: data.MOBILE_BUSINESS_DATA_8Q[op]["consumer_mobile_subs_m"] for op in self.operators}

        chart_path = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=mobile_subs,
            title="消费者移动用户趋势 (不含IoT, 百万)",
            y_label="用户 (M)",
            target_operator=self.target,
            filename="consumer_mobile_trend.png",
        )
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2), Inches(6.2), Inches(2.6))

        # ARPU趋势
        arpu_data = {op: data.MOBILE_BUSINESS_DATA_8Q[op]["mobile_arpu_eur"] for op in self.operators}

        chart_path2 = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=arpu_data,
            title="移动ARPU趋势 (€/月)",
            y_label="ARPU (€)",
            target_operator=self.target,
            filename="mobile_arpu_trend.png",
        )
        self._add_image(slide, chart_path2, Inches(6.8), Inches(1.2), Inches(6.2), Inches(2.6))

        # 详细数据表格
        self._add_text_box(slide, Inches(0.3), Inches(4.0), Inches(6), Inches(0.3),
                          "Q3 FY26 移动业务KPI对比", font_size=12,
                          font_color=self.primary_color, bold=True)

        # KPI表格
        self._add_shape(slide, Inches(0.3), Inches(4.35), Inches(6.2), Inches(2.8), (248, 248, 248))

        mobile_kpis = [
            ("指标", "VF", "DT", "O2", "1&1"),
            ("消费者用户(M)", "29.5", "46.5", "34.0", "12.2"),
            ("IoT连接(M)", "3.0", "3.5", "8.0", "0.3"),
            ("后付费占比", "78.6%", "82.0%", "68.5%", "75.0%"),
            ("ARPU(€)", "12.8", "14.2", "10.5", "9.8"),
            ("新客ARPU增长", "+21%", "+15%", "+12%", "+8%"),
            ("月度流失率", "1.05%", "0.85%", "1.20%", "1.35%"),
        ]

        y_pos = Inches(4.4)
        for row in mobile_kpis:
            x_pos = Inches(0.4)
            for i, val in enumerate(row):
                width = Inches(1.4) if i == 0 else Inches(1.1)
                bold = i == 0 or row == mobile_kpis[0]
                self._add_text_box(slide, x_pos, y_pos, width, Inches(0.35),
                                  val, font_size=10, font_color=self.text_color, bold=bold)
                x_pos += width
            y_pos += Inches(0.38)

        # 移动业务洞察
        self._add_shape(slide, Inches(6.8), Inches(4.0), Inches(6.2), Inches(2.8), (240, 248, 255))

        self._add_text_box(slide, Inches(7), Inches(4.1), Inches(6), Inches(0.3),
                          "移动业务洞察", font_size=12, font_color=self.primary_color, bold=True)

        insights = [
            "【用户规模】",
            "  • 市场格局: DT(46.5M) > O2(34M) > VF(29.5M) > 1&1(12.2M)",
            "  • VF 8季度净增+1.3M消费者用户",
            "  • IoT连接: O2领先(8M)，工业IoT优势",
            "",
            "【价值提升】",
            "  • VF ARPU €11.8 → €12.8 (+8.5%)",
            "  • 新客户ARPU增长+21% (价值导向策略)",
            "  • 后付费占比提升: 77.3% → 78.6%",
            "",
            "【流失控制】",
            "  • VF流失率改善: 1.15% → 1.05%",
            "  • 仍高于DT (0.85%)，需继续优化",
        ]

        y_pos = Inches(4.45)
        for insight in insights:
            self._add_text_box(slide, Inches(7), y_pos, Inches(6), Inches(0.23),
                              insight, font_size=10, font_color=self.text_color)
            y_pos += Inches(0.22)

    def _add_fixed_business_slides(self):
        """固网业务分析"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1
        self._add_header(slide, "固网宽带业务分析", "Fixed Broadband Business Analysis")

        # 宽带用户趋势
        bb_subs = {op: data.FIXED_BROADBAND_DATA_8Q[op]["broadband_subs_m"] for op in self.operators}

        chart_path = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=bb_subs,
            title="宽带用户趋势 (百万)",
            y_label="用户 (M)",
            target_operator=self.target,
            filename="broadband_subs_trend.png",
        )
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2), Inches(6.2), Inches(2.6))

        # 固网收入增长趋势
        fixed_growth = {op: data.REVENUE_DATA_8Q[op]["fixed_service_growth_pct"] for op in self.operators}

        chart_path2 = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=fixed_growth,
            title="固网服务收入增长率 (%)",
            y_label="增长率 (%)",
            target_operator=self.target,
            filename="fixed_revenue_growth.png",
            y_format="percent",
        )
        self._add_image(slide, chart_path2, Inches(6.8), Inches(1.2), Inches(6.2), Inches(2.6))

        # 网络类型占比分析
        self._add_text_box(slide, Inches(0.3), Inches(4.0), Inches(6), Inches(0.3),
                          "Q3 FY26 宽带接入类型占比", font_size=12,
                          font_color=self.primary_color, bold=True)

        self._add_shape(slide, Inches(0.3), Inches(4.35), Inches(6.2), Inches(2.8), (245, 250, 255))

        # 接入类型表格
        access_data = [
            ("运营商", "FTTH/B", "Cable/HFC", "DSL", "总用户(M)"),
            ("VF Germany", "8%", "72%", "20%", "9.9"),
            ("DT Germany", "35%", "0%", "65%", "15.2"),
            ("O2 Germany", "25%", "0%", "75%", "2.4"),
            ("1&1 AG", "15%", "0%", "85%", "4.2"),
        ]

        y_pos = Inches(4.45)
        for row in access_data:
            x_pos = Inches(0.4)
            for i, val in enumerate(row):
                width = Inches(1.4) if i == 0 else Inches(1.1)
                bold = i == 0 or row == access_data[0]
                self._add_text_box(slide, x_pos, y_pos, width, Inches(0.35),
                                  val, font_size=10, font_color=self.text_color, bold=bold)
                x_pos += width
            y_pos += Inches(0.42)

        # 固网洞察
        self._add_shape(slide, Inches(6.8), Inches(4.0), Inches(6.2), Inches(2.8), (255, 245, 240))

        self._add_text_box(slide, Inches(7), Inches(4.1), Inches(6), Inches(0.3),
                          "固网业务洞察", font_size=12, font_color=(180, 80, 0), bold=True)

        insights = [
            "【Vodafone挑战】",
            "  • 用户流失: 10.32M → 9.94M (-380K)",
            "  • Cable/HFC网络老化，升级成本高",
            "  • FTTH覆盖仅8%，远低于DT(35%)",
            "",
            "【DT优势】",
            "  • 光纤战略领先，FTTH持续扩容",
            "  • 收入正增长(+0.9%)，VF负增长(-1.1%)",
            "",
            "【策略建议】",
            "  • 加速FTTH/DOCSIS 4.0升级",
            "  • 推动FMC融合提升粘性",
            "  • 聚焦高价值存量用户",
        ]

        y_pos = Inches(4.45)
        for insight in insights:
            self._add_text_box(slide, Inches(7), y_pos, Inches(6), Inches(0.22),
                              insight, font_size=10, font_color=self.text_color)
            y_pos += Inches(0.22)

    def _add_b2b_business_slides(self):
        """B2B业务分析"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1
        self._add_header(slide, "B2B业务分析", "B2B Business Analysis")

        # B2B收入趋势
        b2b_rev = {}
        for op in self.operators:
            b2b_rev[op] = [r/1000 for r in data.REVENUE_DATA_8Q[op]["b2b_revenue"]]

        chart_path = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=b2b_rev,
            title="B2B收入趋势 (€B)",
            y_label="收入 (€B)",
            target_operator=self.target,
            filename="b2b_revenue_trend.png",
        )
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2), Inches(6.2), Inches(2.6))

        # B2B增长率趋势
        b2b_growth = {op: data.REVENUE_DATA_8Q[op]["b2b_growth_pct"] for op in self.operators}

        chart_path2 = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=b2b_growth,
            title="B2B收入增长率 (%)",
            y_label="增长率 (%)",
            target_operator=self.target,
            filename="b2b_growth_trend.png",
            y_format="percent",
        )
        self._add_image(slide, chart_path2, Inches(6.8), Inches(1.2), Inches(6.2), Inches(2.6))

        # B2B业务结构
        self._add_text_box(slide, Inches(0.3), Inches(4.0), Inches(6), Inches(0.3),
                          "Vodafone B2B业务构成 (Q3 FY26)", font_size=12,
                          font_color=self.primary_color, bold=True)

        self._add_shape(slide, Inches(0.3), Inches(4.35), Inches(6.2), Inches(2.8), (240, 248, 255))

        b2b_structure = [
            ("业务类型", "收入(€M)", "占比", "增长率"),
            ("企业连接服务", "280", "53.8%", "+5.0%"),
            ("数字化服务合计", "240", "46.2%", "+12.5%"),
            ("  - IoT/M2M", "95", "18.3%", "+18%"),
            ("  - 云服务", "72", "13.8%", "+22%"),
            ("  - 安全服务", "48", "9.2%", "+15%"),
            ("  - 边缘计算", "25", "4.8%", "+35%"),
        ]

        y_pos = Inches(4.45)
        for row in b2b_structure:
            x_pos = Inches(0.4)
            for i, val in enumerate(row):
                width = Inches(2.0) if i == 0 else Inches(1.3)
                bold = row == b2b_structure[0] or row == b2b_structure[1] or row == b2b_structure[2]
                self._add_text_box(slide, x_pos, y_pos, width, Inches(0.32),
                                  val, font_size=10, font_color=self.text_color, bold=bold)
                x_pos += width
            y_pos += Inches(0.35)

        # B2B洞察
        self._add_shape(slide, Inches(6.8), Inches(4.0), Inches(6.2), Inches(2.8), (240, 255, 245))

        self._add_text_box(slide, Inches(7), Inches(4.1), Inches(6), Inches(0.3),
                          "B2B业务洞察", font_size=12, font_color=(0, 128, 80), bold=True)

        insights = [
            "【增长亮点】",
            "  • B2B整体增速+8.5%，高于消费者业务",
            "  • 数字化服务增长+12.5%，是核心驱动力",
            "  • 边缘计算增长最快(+35%)，5G赋能",
            "",
            "【竞争格局】",
            f"  • DT B2B: €{data.REVENUE_DATA_8Q['Deutsche Telekom Germany']['b2b_revenue'][-1]/1000:.2f}B (市场第一)",
            f"  • VF B2B: €{data.REVENUE_DATA_8Q['Vodafone Germany']['b2b_revenue'][-1]/1000:.2f}B (第二)",
            "  • O2增速最快(+13%)但基数较小",
            "",
            "【策略方向】",
            "  • 深耕制造业/汽车行业",
            "  • 5G专网+边缘计算组合方案",
            "  • 强化云/安全生态合作",
        ]

        y_pos = Inches(4.45)
        for insight in insights:
            self._add_text_box(slide, Inches(7), y_pos, Inches(6), Inches(0.22),
                              insight, font_size=10, font_color=self.text_color)
            y_pos += Inches(0.22)

    def _add_tv_fmc_slides(self):
        """TV和FMC融合业务分析"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1
        self._add_header(slide, "TV与FMC融合业务分析", "TV & FMC (Fixed-Mobile Convergence) Analysis")

        # TV用户趋势
        tv_subs = {op: data.TV_FMC_DATA_8Q[op]["tv_subs_m"] for op in self.operators}

        chart_path = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=tv_subs,
            title="TV用户趋势 (百万)",
            y_label="用户 (M)",
            target_operator=self.target,
            filename="tv_subs_trend.png",
        )
        self._add_image(slide, chart_path, Inches(0.3), Inches(1.2), Inches(6.2), Inches(2.6))

        # FMC用户占比趋势
        fmc_ratio = {op: data.TV_FMC_DATA_8Q[op]["fmc_penetration_pct"] for op in self.operators}

        chart_path2 = self.chart_gen.create_multi_line_trend_chart(
            quarters=data.QUARTERS_8Q,
            data_series=fmc_ratio,
            title="FMC融合用户占比 (%)",
            y_label="占比 (%)",
            target_operator=self.target,
            filename="fmc_ratio_trend.png",
        )
        self._add_image(slide, chart_path2, Inches(6.8), Inches(1.2), Inches(6.2), Inches(2.6))

        # TV/FMC KPI对比
        self._add_text_box(slide, Inches(0.3), Inches(4.0), Inches(6), Inches(0.3),
                          "Q3 FY26 TV & FMC KPI对比", font_size=12,
                          font_color=self.primary_color, bold=True)

        self._add_shape(slide, Inches(0.3), Inches(4.35), Inches(6.2), Inches(2.8), (248, 245, 255))

        tv_kpis = [
            ("指标", "VF", "DT", "O2", "1&1"),
            ("TV用户(M)", "9.2", "4.8", "0.5", "0.3"),
            ("IPTV占比", "45%", "100%", "100%", "100%"),
            ("Cable TV占比", "55%", "0%", "0%", "0%"),
            ("FMC融合率", "32%", "45%", "18%", "25%"),
            ("融合ARPU提升", "+18%", "+22%", "+15%", "+12%"),
        ]

        y_pos = Inches(4.45)
        for row in tv_kpis:
            x_pos = Inches(0.4)
            for i, val in enumerate(row):
                width = Inches(1.5) if i == 0 else Inches(1.1)
                bold = row == tv_kpis[0]
                self._add_text_box(slide, x_pos, y_pos, width, Inches(0.35),
                                  val, font_size=10, font_color=self.text_color, bold=bold)
                x_pos += width
            y_pos += Inches(0.4)

        # TV/FMC洞察
        self._add_shape(slide, Inches(6.8), Inches(4.0), Inches(6.2), Inches(2.8), (255, 248, 255))

        self._add_text_box(slide, Inches(7), Inches(4.1), Inches(6), Inches(0.3),
                          "TV & FMC洞察", font_size=12, font_color=(128, 0, 128), bold=True)

        insights = [
            "【TV业务】",
            "  • VF TV用户最多(9.2M)，Cable资产优势",
            "  • 流媒体冲击，传统TV承压",
            "  • GigaTV融合战略是差异化",
            "",
            "【FMC融合】",
            "  • DT融合率最高(45%)，One套餐策略",
            "  • VF融合率提升: 28% → 32% (+4pp)",
            "  • 融合用户ARPU+18%，流失率-30%",
            "",
            "【策略重点】",
            "  • 提升FMC渗透率(目标: 40%+)",
            "  • 内容差异化(体育/独家)",
            "  • 超融合产品(移动+固网+TV+IoT)",
        ]

        y_pos = Inches(4.45)
        for insight in insights:
            self._add_text_box(slide, Inches(7), y_pos, Inches(6), Inches(0.22),
                              insight, font_size=10, font_color=self.text_color)
            y_pos += Inches(0.22)

    # ==========================================================================
    # 第三部分: 战略与组织
    # ==========================================================================

    def _add_strategy_analysis_slides(self):
        """战略分析页面"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1
        self._add_header(slide, "客户战略对比分析", "Customer Strategy Comparison")

        # 各运营商战略对比
        y_base = Inches(1.2)

        for i, op in enumerate(self.operators[:4]):
            strategy = data.STRATEGY_DATA.get(op, {})

            # 运营商名称框
            x = Inches(0.3) + i * Inches(3.2)
            color = self.primary_color if op == self.target else (80, 80, 120)

            self._add_shape(slide, x, y_base, Inches(3.0), Inches(0.4), color)
            short_name = op.replace(" Germany", "").replace("Telefónica ", "")
            self._add_text_box(slide, x + Inches(0.1), y_base + Inches(0.05),
                              Inches(2.8), Inches(0.3), short_name,
                              font_size=12, font_color=(255, 255, 255), bold=True)

            # 战略内容
            self._add_shape(slide, x, y_base + Inches(0.45), Inches(3.0), Inches(5.3), (248, 248, 252))

            content = [
                "【战略定位】",
                strategy.get("positioning", "N/A")[:40],
                "",
                "【核心策略】",
            ]
            priorities = strategy.get("priorities", [])[:4]
            for p in priorities:
                content.append(f"• {p[:35]}")

            content.extend([
                "",
                "【关键举措】",
            ])
            initiatives = strategy.get("initiatives", [])[:3]
            for init in initiatives:
                content.append(f"• {init[:35]}")

            y_pos = y_base + Inches(0.55)
            for line in content:
                bold = "【" in line
                self._add_text_box(slide, x + Inches(0.1), y_pos, Inches(2.8), Inches(0.25),
                                  line, font_size=9, font_color=self.text_color, bold=bold)
                y_pos += Inches(0.28)

    def _add_executive_changes_slide(self):
        """高管变动分析"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1
        self._add_header(slide, "组织与高管变动分析", "Organization & Executive Changes")

        y_base = Inches(1.2)

        for i, op in enumerate(self.operators[:4]):
            changes = data.EXECUTIVE_CHANGES.get(op, {})

            x = Inches(0.3) + i * Inches(3.2)
            color = self.primary_color if op == self.target else (60, 80, 100)

            self._add_shape(slide, x, y_base, Inches(3.0), Inches(0.4), color)
            short_name = op.replace(" Germany", "").replace("Telefónica ", "")
            self._add_text_box(slide, x + Inches(0.1), y_base + Inches(0.05),
                              Inches(2.8), Inches(0.3), short_name,
                              font_size=12, font_color=(255, 255, 255), bold=True)

            self._add_shape(slide, x, y_base + Inches(0.45), Inches(3.0), Inches(5.3), (245, 248, 252))

            content = []

            # CEO
            ceo = changes.get("ceo", {})
            content.append("【CEO】")
            content.append(f"{ceo.get('name', 'N/A')}")
            content.append(f"任期: {ceo.get('tenure', 'N/A')}")
            content.append("")

            # Recent changes
            recent = changes.get("recent_changes", [])
            if recent:
                content.append("【近期变动】")
                for change in recent[:3]:
                    content.append(f"• {change[:35]}")
                content.append("")

            # Org structure
            org = changes.get("org_focus", [])
            if org:
                content.append("【组织重点】")
                for o in org[:3]:
                    content.append(f"• {o[:35]}")

            y_pos = y_base + Inches(0.55)
            for line in content:
                bold = "【" in line
                self._add_text_box(slide, x + Inches(0.1), y_pos, Inches(2.8), Inches(0.25),
                                  line, font_size=9, font_color=self.text_color, bold=bold)
                y_pos += Inches(0.28)

    # ==========================================================================
    # 第四部分: 网络与资费
    # ==========================================================================

    def _add_network_analysis_slides(self):
        """网络分析页面"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1
        self._add_header(slide, "网络基础设施对比分析", "Network Infrastructure Comparison")

        # 移动网络对比
        self._add_text_box(slide, Inches(0.3), Inches(1.2), Inches(6), Inches(0.3),
                          "移动网络覆盖与质量", font_size=13,
                          font_color=self.primary_color, bold=True)

        self._add_shape(slide, Inches(0.3), Inches(1.55), Inches(6.2), Inches(2.6), (245, 250, 255))

        mobile_net = [
            ("指标", "VF", "DT", "O2", "1&1"),
            ("5G覆盖率(人口)", "92%", "97%", "95%", "50%"),
            ("5G基站数", "18K", "28K", "22K", "6K"),
            ("4G覆盖率", "99%", "99%", "99%", "95%"),
            ("平均下载速度", "85Mbps", "120Mbps", "95Mbps", "65Mbps"),
            ("网络质量评分", "78/100", "92/100", "82/100", "68/100"),
        ]

        y_pos = Inches(1.65)
        for row in mobile_net:
            x_pos = Inches(0.4)
            for i, val in enumerate(row):
                width = Inches(1.5) if i == 0 else Inches(1.1)
                bold = row == mobile_net[0]
                self._add_text_box(slide, x_pos, y_pos, width, Inches(0.35),
                                  val, font_size=10, font_color=self.text_color, bold=bold)
                x_pos += width
            y_pos += Inches(0.38)

        # 固网网络对比
        self._add_text_box(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(0.3),
                          "固定网络覆盖与技术", font_size=13,
                          font_color=self.primary_color, bold=True)

        self._add_shape(slide, Inches(6.8), Inches(1.55), Inches(6.2), Inches(2.6), (250, 248, 245))

        fixed_net = [
            ("指标", "VF", "DT", "O2", "1&1"),
            ("FTTH覆盖(户)", "2.5M", "8.5M", "0.8M", "0.3M"),
            ("Cable/HFC覆盖", "24M", "0", "0", "0"),
            ("最高速率", "1Gbps", "1Gbps", "250Mbps", "250Mbps"),
            ("DOCSIS 4.0计划", "2026启动", "N/A", "N/A", "N/A"),
            ("网络投资占比", "26%", "19%", "25%", "39%"),
        ]

        y_pos = Inches(1.65)
        for row in fixed_net:
            x_pos = Inches(6.9)
            for i, val in enumerate(row):
                width = Inches(1.5) if i == 0 else Inches(1.1)
                bold = row == fixed_net[0]
                self._add_text_box(slide, x_pos, y_pos, width, Inches(0.35),
                                  val, font_size=10, font_color=self.text_color, bold=bold)
                x_pos += width
            y_pos += Inches(0.38)

        # 网络洞察
        self._add_shape(slide, Inches(0.3), Inches(4.4), Inches(12.7), Inches(2.8), (248, 252, 248))

        self._add_text_box(slide, Inches(0.5), Inches(4.5), Inches(6), Inches(0.3),
                          "网络洞察与建议", font_size=12, font_color=self.primary_color, bold=True)

        # 左侧洞察
        left_insights = [
            "【VF移动网络差距】",
            "  • 5G覆盖率: 92% vs DT 97% (-5pp)",
            "  • 网络质量评分: 78 vs DT 92 (-14分)",
            "  • 需加速5G部署和网络优化",
            "",
            "【1&1网络挑战】",
            "  • 自建网络进度缓慢(50%覆盖)",
            "  • 仍依赖O2批发接入",
        ]

        # 右侧洞察
        right_insights = [
            "【VF固网优势与挑战】",
            "  • Cable/HFC 24M覆盖是核心资产",
            "  • FTTH仅2.5M，落后DT(8.5M)",
            "  • DOCSIS 4.0升级是关键机会",
            "",
            "【策略建议】",
            "  • 移动: 聚焦5G关键城市覆盖",
            "  • 固网: Cable升级+选择性FTTH",
        ]

        y_pos = Inches(4.85)
        for insight in left_insights:
            self._add_text_box(slide, Inches(0.5), y_pos, Inches(6), Inches(0.24),
                              insight, font_size=10, font_color=self.text_color)
            y_pos += Inches(0.24)

        y_pos = Inches(4.85)
        for insight in right_insights:
            self._add_text_box(slide, Inches(6.8), y_pos, Inches(6), Inches(0.24),
                              insight, font_size=10, font_color=self.text_color)
            y_pos += Inches(0.24)

    def _add_tariff_comparison_slides(self):
        """资费对比分析"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1
        self._add_header(slide, "套餐资费对比分析", "Tariff & Package Comparison")

        # 移动套餐对比
        self._add_text_box(slide, Inches(0.3), Inches(1.2), Inches(6), Inches(0.3),
                          "移动主力套餐对比 (无限流量)", font_size=13,
                          font_color=self.primary_color, bold=True)

        self._add_shape(slide, Inches(0.3), Inches(1.55), Inches(6.2), Inches(2.4), (248, 248, 255))

        mobile_tariff = [
            ("运营商", "套餐名", "月费", "流量", "5G", "特点"),
            ("VF", "GigaMobil XL", "€79.99", "无限", "是", "120Mbps"),
            ("DT", "MagentaM. XL", "€84.95", "无限", "是", "5G优先"),
            ("O2", "Free Unlim.", "€49.99", "无限", "是", "最低价"),
            ("1&1", "All-Net XL", "€49.99", "无限", "是", "前6月优惠"),
        ]

        y_pos = Inches(1.65)
        col_widths = [Inches(0.6), Inches(1.4), Inches(0.9), Inches(0.8), Inches(0.6), Inches(1.2)]
        for row in mobile_tariff:
            x_pos = Inches(0.4)
            for i, val in enumerate(row):
                bold = row == mobile_tariff[0]
                self._add_text_box(slide, x_pos, y_pos, col_widths[i], Inches(0.35),
                                  val, font_size=9, font_color=self.text_color, bold=bold)
                x_pos += col_widths[i]
            y_pos += Inches(0.38)

        # 固网套餐对比
        self._add_text_box(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(0.3),
                          "固网宽带套餐对比 (1Gbps)", font_size=13,
                          font_color=self.primary_color, bold=True)

        self._add_shape(slide, Inches(6.8), Inches(1.55), Inches(6.2), Inches(2.4), (255, 248, 245))

        fixed_tariff = [
            ("运营商", "套餐名", "月费", "速率", "技术", "TV"),
            ("VF", "Cable Max", "€49.99", "1Gbps", "Cable", "可选"),
            ("DT", "MagentaZu. XL", "€59.95", "1Gbps", "FTTH", "MagentaTV"),
            ("O2", "my Home XL", "€44.99", "250Mbps", "DSL", "无"),
            ("1&1", "1&1 Glasfaser", "€49.99", "1Gbps", "FTTH", "可选"),
        ]

        y_pos = Inches(1.65)
        for row in fixed_tariff:
            x_pos = Inches(6.9)
            for i, val in enumerate(row):
                bold = row == fixed_tariff[0]
                self._add_text_box(slide, x_pos, y_pos, col_widths[i], Inches(0.35),
                                  val, font_size=9, font_color=self.text_color, bold=bold)
                x_pos += col_widths[i]
            y_pos += Inches(0.38)

        # 融合套餐对比
        self._add_text_box(slide, Inches(0.3), Inches(4.2), Inches(12.7), Inches(0.3),
                          "FMC融合套餐策略对比", font_size=13,
                          font_color=self.primary_color, bold=True)

        self._add_shape(slide, Inches(0.3), Inches(4.55), Inches(12.7), Inches(2.6), (252, 250, 248))

        # 融合策略表格
        fmc_data = [
            ("运营商", "融合品牌", "核心优惠", "融合折扣", "差异化"),
            ("Vodafone", "GigaKombi", "移动+固网+TV", "€10-15/月", "Cable+TV内容"),
            ("DT", "MagentaEINS", "移动+固网+TV", "€10/月+StreamOn", "品牌一体化"),
            ("O2", "O2 my Home", "移动+固网", "€5-10/月", "低价策略"),
            ("1&1", "All-in-One", "移动+固网", "€5/月", "价格灵活"),
        ]

        y_pos = Inches(4.7)
        col_widths2 = [Inches(1.4), Inches(1.6), Inches(2.5), Inches(1.5), Inches(3.0)]
        for row in fmc_data:
            x_pos = Inches(0.5)
            for i, val in enumerate(row):
                bold = row == fmc_data[0]
                self._add_text_box(slide, x_pos, y_pos, col_widths2[i], Inches(0.35),
                                  val, font_size=10, font_color=self.text_color, bold=bold)
                x_pos += col_widths2[i]
            y_pos += Inches(0.42)

        # 资费洞察
        tariff_insight = "【资费策略洞察】VF定价处于中高端，需通过网络质量和服务体验证明价值溢价；O2低价策略抢占价格敏感用户"
        self._add_text_box(slide, Inches(0.5), Inches(6.9), Inches(12.5), Inches(0.35),
                          tariff_insight, font_size=11, font_color=self.text_color)

    # ==========================================================================
    # 第五部分: 总结
    # ==========================================================================

    def _add_summary_slide(self):
        """添加总结页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_num += 1
        self._add_header(slide, "分析总结与战略建议", "Summary & Strategic Recommendations")

        # 左侧: 关键发现
        self._add_text_box(slide, Inches(0.3), Inches(1.2), Inches(6), Inches(0.35),
                          "关键发现 KEY FINDINGS", font_size=14,
                          font_color=self.primary_color, bold=True)

        self._add_shape(slide, Inches(0.3), Inches(1.6), Inches(6.2), Inches(5.5), (248, 250, 255))

        findings = [
            "【经营层面】",
            "  • 收入企稳回升: 服务收入增速 -0.8% → +0.7%",
            "  • 盈利改善: EBITDA Margin 35.4% → 36.2%",
            "  • 成本优化: OPEX占比下降0.8pp",
            "",
            "【业务层面】",
            "  • 移动业务: 用户增长+ARPU提升双轮驱动",
            "  • 固网业务: 用户流失压力，Cable升级紧迫",
            "  • B2B业务: 增速最快(+8.5%)，数字化是引擎",
            "",
            "【竞争格局】",
            "  • DT继续领跑，但增速放缓",
            "  • O2因1&1迁移持续承压",
            "  • 1&1自建网络仍在投入期",
            "",
            "【核心差距】",
            "  • 5G覆盖率: 92% vs DT 97%",
            "  • FTTH规模: 2.5M vs DT 8.5M",
            "  • 网络质量评分: 78 vs DT 92",
        ]

        y_pos = Inches(1.7)
        for finding in findings:
            bold = "【" in finding
            self._add_text_box(slide, Inches(0.5), y_pos, Inches(6), Inches(0.26),
                              finding, font_size=10, font_color=self.text_color, bold=bold)
            y_pos += Inches(0.27)

        # 右侧: 战略建议
        self._add_text_box(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(0.35),
                          "战略建议 RECOMMENDATIONS", font_size=14,
                          font_color=self.primary_color, bold=True)

        self._add_shape(slide, Inches(6.8), Inches(1.6), Inches(6.2), Inches(5.5), (255, 248, 245))

        recommendations = [
            "【P0 - 网络质量提升】",
            "  • 5G覆盖加速: 目标2026年达95%",
            "  • Cable升级: DOCSIS 4.0路线图",
            "  • 网络体验: 关键城市NPS改善",
            "",
            "【P1 - 价值增长战略】",
            "  • 移动: 后付费渗透+ARPU提升",
            "  • 固网: FMC融合率达40%+",
            "  • B2B: 数字化服务占比达50%+",
            "",
            "【P1 - 客户体验优化】",
            "  • 数字化渠道: App/自助服务提升",
            "  • 服务响应: 首次解决率提升至85%",
            "",
            "【P2 - 成本效率】",
            "  • 运营自动化: AI/RPA全面应用",
            "  • 组织精简: 人均效能提升15%",
            "",
            "【执行优先级】",
            "  网络 > 价值增长 > 客户体验 > 成本",
        ]

        y_pos = Inches(1.7)
        for rec in recommendations:
            bold = "【" in rec
            color = self.primary_color if "P0" in rec else self.text_color
            self._add_text_box(slide, Inches(7), y_pos, Inches(6), Inches(0.26),
                              rec, font_size=10, font_color=color, bold=bold)
            y_pos += Inches(0.27)


def generate_comprehensive_ppt(output_dir: str = None) -> str:
    """生成综合分析PPT的便捷函数"""
    generator = ComprehensiveAnalysisPPT(output_dir)
    return generator.generate()


if __name__ == "__main__":
    output_path = generate_comprehensive_ppt()
    print(f"Generated: {output_path}")
