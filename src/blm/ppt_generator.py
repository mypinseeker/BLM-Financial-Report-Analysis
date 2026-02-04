"""BLM Strategic Analysis PPT Generator.

Generates professional PowerPoint presentations from BLM analysis results.
Supports Huawei-style template with clean, professional design.

Requires: python-pptx (pip install python-pptx)

Usage:
    from src.blm.ppt_generator import BLMPPTGenerator

    generator = BLMPPTGenerator(style="huawei")
    ppt_path = generator.generate(five_looks, three_decisions, "Vodafone Germany")
"""

import io
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.oxml.ns import nsmap
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from src.blm.five_looks import InsightResult
from src.blm.three_decisions import StrategyResult, StrategyItem


# =============================================================================
# Huawei Style Configuration
# =============================================================================

@dataclass
class PPTStyle:
    """PPT style configuration."""
    name: str
    primary_color: tuple  # RGB
    secondary_color: tuple
    accent_color: tuple
    text_color: tuple
    light_text_color: tuple
    background_color: tuple
    title_font: str
    body_font: str
    title_size: int  # Pt
    subtitle_size: int
    heading_size: int
    body_size: int
    small_size: int


HUAWEI_STYLE = PPTStyle(
    name="Huawei",
    primary_color=(199, 0, 11),      # Huawei Red #C7000B
    secondary_color=(0, 0, 0),        # Black
    accent_color=(229, 32, 50),       # Lighter Red #E52032
    text_color=(51, 51, 51),          # Dark Gray #333333
    light_text_color=(128, 128, 128), # Gray #808080
    background_color=(255, 255, 255), # White
    title_font="Arial",
    body_font="Arial",
    title_size=36,
    subtitle_size=24,
    heading_size=20,
    body_size=14,
    small_size=11,
)

VODAFONE_STYLE = PPTStyle(
    name="Vodafone",
    primary_color=(230, 0, 0),        # Vodafone Red #E60000
    secondary_color=(51, 51, 51),     # Dark Gray
    accent_color=(255, 255, 255),     # White
    text_color=(51, 51, 51),
    light_text_color=(128, 128, 128),
    background_color=(255, 255, 255),
    title_font="Arial",
    body_font="Arial",
    title_size=36,
    subtitle_size=24,
    heading_size=20,
    body_size=14,
    small_size=11,
)

STYLES = {
    "huawei": HUAWEI_STYLE,
    "vodafone": VODAFONE_STYLE,
}


class BLMPPTGenerator:
    """Generate PowerPoint presentations from BLM analysis results."""

    def __init__(
        self,
        style: str = "huawei",
        output_dir: Optional[str] = None,
    ):
        """Initialize PPT generator.

        Args:
            style: PPT style name ("huawei", "vodafone", or custom PPTStyle)
            output_dir: Output directory for generated PPTs
        """
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for PPT generation. "
                "Install it with: pip install python-pptx"
            )

        if isinstance(style, str):
            self.style = STYLES.get(style, HUAWEI_STYLE)
        else:
            self.style = style

        if output_dir is None:
            output_dir = str(Path(__file__).resolve().parent.parent.parent / "data" / "output")
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        self.prs = None

    def generate(
        self,
        five_looks: dict[str, InsightResult],
        three_decisions: dict[str, StrategyResult],
        target_operator: str,
        competitors: list[str] = None,
        title: str = None,
        filename: str = None,
    ) -> str:
        """Generate PPT from BLM analysis results.

        Args:
            five_looks: Five Looks analysis results
            three_decisions: Three Decisions strategy results
            target_operator: Target operator name
            competitors: List of competitor names
            title: Presentation title
            filename: Output filename

        Returns:
            Path to generated PPT file
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        self.prs.slide_height = Inches(7.5)

        competitors = competitors or []
        title = title or f"{target_operator} BLM 战略分析报告"
        filename = filename or f"blm_{self._sanitize_name(target_operator)}_strategy.pptx"

        # Generate slides
        self._add_title_slide(title, target_operator, competitors)
        self._add_agenda_slide()
        self._add_executive_summary_slide(five_looks, three_decisions, target_operator)

        # Five Looks section
        self._add_section_divider("五看分析", "Five Looks Analysis", "01")
        for key, insight in five_looks.items():
            self._add_insight_slide(insight)

        # Three Decisions section
        self._add_section_divider("三定策略", "Three Decisions Strategy", "02")
        for key, decision in three_decisions.items():
            self._add_strategy_slide(decision)

        # Conclusion
        self._add_conclusion_slide(five_looks, three_decisions, target_operator)
        self._add_closing_slide()

        # Save
        output_path = self.output_dir / filename
        self.prs.save(str(output_path))
        return str(output_path)

    def _add_title_slide(
        self,
        title: str,
        target_operator: str,
        competitors: list[str],
    ):
        """Add title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Red header bar
        self._add_shape(
            slide, 0, 0, Inches(13.333), Inches(2.5),
            self.style.primary_color
        )

        # Title
        self._add_text_box(
            slide,
            Inches(0.5), Inches(0.6), Inches(12.333), Inches(1),
            title,
            font_size=self.style.title_size,
            font_color=(255, 255, 255),
            bold=True,
        )

        # Subtitle
        subtitle = f"基于 BLM (Business Leadership Model) 方法论"
        self._add_text_box(
            slide,
            Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.6),
            subtitle,
            font_size=self.style.subtitle_size,
            font_color=(255, 255, 255),
            bold=False,
        )

        # Target operator info
        info_text = f"分析对象: {target_operator}"
        if competitors:
            info_text += f"\n竞争对手: {', '.join(competitors)}"
        info_text += f"\n报告日期: {datetime.now().strftime('%Y年%m月%d日')}"

        self._add_text_box(
            slide,
            Inches(0.5), Inches(3.2), Inches(6), Inches(2),
            info_text,
            font_size=self.style.heading_size,
            font_color=self.style.text_color,
        )

        # Framework description
        framework_text = """BLM 战略分析框架:

• 五看 (Five Looks): 看市场、看自己、看对手、看宏观、看机会
• 三定 (Three Decisions): 定策略、定重点工作、定执行"""

        self._add_text_box(
            slide,
            Inches(7), Inches(3.2), Inches(5.8), Inches(2.5),
            framework_text,
            font_size=self.style.body_size,
            font_color=self.style.text_color,
        )

    def _add_agenda_slide(self):
        """Add agenda/contents slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Header
        self._add_header(slide, "报告目录", "AGENDA")

        # Agenda items
        agenda_items = [
            ("01", "执行摘要", "Executive Summary"),
            ("02", "五看分析", "看市场 | 看自己 | 看对手 | 看宏观 | 看机会"),
            ("03", "三定策略", "定策略 | 定重点工作 | 定执行"),
            ("04", "总结与建议", "Summary & Recommendations"),
        ]

        y_pos = Inches(1.8)
        for num, title, subtitle in agenda_items:
            # Number circle
            self._add_shape(
                slide,
                Inches(0.8), y_pos, Inches(0.6), Inches(0.6),
                self.style.primary_color,
                shape_type=MSO_SHAPE.OVAL,
            )
            self._add_text_box(
                slide,
                Inches(0.8), y_pos + Inches(0.1), Inches(0.6), Inches(0.4),
                num,
                font_size=16,
                font_color=(255, 255, 255),
                bold=True,
                align="center",
            )

            # Title and subtitle
            self._add_text_box(
                slide,
                Inches(1.7), y_pos, Inches(10), Inches(0.5),
                title,
                font_size=self.style.heading_size,
                font_color=self.style.text_color,
                bold=True,
            )
            self._add_text_box(
                slide,
                Inches(1.7), y_pos + Inches(0.45), Inches(10), Inches(0.4),
                subtitle,
                font_size=self.style.body_size,
                font_color=self.style.light_text_color,
            )

            y_pos += Inches(1.3)

    def _add_executive_summary_slide(
        self,
        five_looks: dict[str, InsightResult],
        three_decisions: dict[str, StrategyResult],
        target_operator: str,
    ):
        """Add executive summary slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self._add_header(slide, "执行摘要", "EXECUTIVE SUMMARY")

        # Key metrics from Five Looks
        market = five_looks.get("market")
        self_insight = five_looks.get("self")

        # Left column - Market position
        self._add_text_box(
            slide,
            Inches(0.5), Inches(1.6), Inches(4), Inches(0.4),
            "市场地位",
            font_size=16,
            font_color=self.style.primary_color,
            bold=True,
        )

        market_text = ""
        if self_insight and self_insight.metrics:
            metrics = self_insight.metrics
            market_text = f"""• 收入规模: €{metrics.get('revenue_eur_billion', 'N/A')}B
• 市场排名: 第{metrics.get('revenue_rank', 'N/A')}位
• 市场份额: {metrics.get('market_share_pct', 'N/A')}%
• 5G覆盖率: {metrics.get('5g_coverage_pct', 'N/A')}%"""

        self._add_text_box(
            slide,
            Inches(0.5), Inches(2.1), Inches(4), Inches(2.5),
            market_text,
            font_size=self.style.body_size,
            font_color=self.style.text_color,
        )

        # Middle column - Key challenges
        self._add_text_box(
            slide,
            Inches(4.8), Inches(1.6), Inches(4), Inches(0.4),
            "核心痛点",
            font_size=16,
            font_color=self.style.primary_color,
            bold=True,
        )

        challenges = []
        if self_insight:
            for finding in self_insight.findings[:4]:
                if "落后" in finding or "低于" in finding or "高于" in finding:
                    challenges.append(f"• {finding[:50]}...")

        self._add_text_box(
            slide,
            Inches(4.8), Inches(2.1), Inches(4), Inches(2.5),
            "\n".join(challenges) if challenges else "• 详见五看分析",
            font_size=self.style.body_size,
            font_color=self.style.text_color,
        )

        # Right column - Strategic priorities
        self._add_text_box(
            slide,
            Inches(9.1), Inches(1.6), Inches(3.8), Inches(0.4),
            "战略重点",
            font_size=16,
            font_color=self.style.primary_color,
            bold=True,
        )

        strategy = three_decisions.get("strategy")
        strategy_text = ""
        if strategy:
            p0_items = [i for i in strategy.items if i.priority == "P0"][:3]
            for item in p0_items:
                strategy_text += f"• [P0] {item.name}\n"

        self._add_text_box(
            slide,
            Inches(9.1), Inches(2.1), Inches(3.8), Inches(2.5),
            strategy_text or "• 详见三定策略",
            font_size=self.style.body_size,
            font_color=self.style.text_color,
        )

        # Bottom banner
        self._add_shape(
            slide,
            Inches(0.5), Inches(5.2), Inches(12.333), Inches(1.8),
            (245, 245, 245),  # Light gray background
        )

        # Key recommendation
        rec_title = "核心定位: 做德国市场「有质量的挑战者」"
        self._add_text_box(
            slide,
            Inches(0.8), Inches(5.4), Inches(11.8), Inches(0.5),
            rec_title,
            font_size=18,
            font_color=self.style.primary_color,
            bold=True,
        )

        rec_text = "以网络和服务质量为核心竞争力，聚焦B2B市场差异化，避免C端价格战"
        self._add_text_box(
            slide,
            Inches(0.8), Inches(5.9), Inches(11.8), Inches(0.8),
            rec_text,
            font_size=self.style.body_size,
            font_color=self.style.text_color,
        )

    def _add_section_divider(self, title: str, subtitle: str, number: str):
        """Add section divider slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Full red background
        self._add_shape(
            slide, 0, 0, Inches(13.333), Inches(7.5),
            self.style.primary_color
        )

        # Section number
        self._add_text_box(
            slide,
            Inches(0.5), Inches(2), Inches(2), Inches(1),
            number,
            font_size=72,
            font_color=(255, 255, 255),
            bold=True,
        )

        # Divider line
        self._add_shape(
            slide,
            Inches(0.5), Inches(3.3), Inches(5), Inches(0.03),
            (255, 255, 255),
        )

        # Title
        self._add_text_box(
            slide,
            Inches(0.5), Inches(3.6), Inches(12), Inches(1),
            title,
            font_size=48,
            font_color=(255, 255, 255),
            bold=True,
        )

        # Subtitle
        self._add_text_box(
            slide,
            Inches(0.5), Inches(4.6), Inches(12), Inches(0.6),
            subtitle,
            font_size=24,
            font_color=(255, 200, 200),
        )

    def _add_insight_slide(self, insight: InsightResult):
        """Add Five Looks insight slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self._add_header(slide, insight.title, insight.category.upper())

        # Metrics boxes (top row)
        if insight.metrics:
            x_pos = Inches(0.5)
            metrics_items = list(insight.metrics.items())[:4]
            box_width = Inches(3)

            for key, value in metrics_items:
                # Metric box
                self._add_shape(
                    slide, x_pos, Inches(1.6), box_width, Inches(0.9),
                    (245, 245, 245),
                )
                # Label
                label = key.replace("_", " ").title()[:20]
                self._add_text_box(
                    slide,
                    x_pos + Inches(0.1), Inches(1.65), box_width - Inches(0.2), Inches(0.3),
                    label,
                    font_size=10,
                    font_color=self.style.light_text_color,
                )
                # Value
                self._add_text_box(
                    slide,
                    x_pos + Inches(0.1), Inches(1.95), box_width - Inches(0.2), Inches(0.5),
                    str(value),
                    font_size=20,
                    font_color=self.style.primary_color,
                    bold=True,
                )
                x_pos += box_width + Inches(0.2)

        # Findings (left column)
        self._add_text_box(
            slide,
            Inches(0.5), Inches(2.8), Inches(6), Inches(0.4),
            "洞察发现",
            font_size=14,
            font_color=self.style.primary_color,
            bold=True,
        )

        findings_text = ""
        for finding in insight.findings[:6]:
            # Truncate long findings
            text = finding[:80] + "..." if len(finding) > 80 else finding
            findings_text += f"• {text}\n"

        self._add_text_box(
            slide,
            Inches(0.5), Inches(3.2), Inches(6), Inches(3.5),
            findings_text,
            font_size=self.style.small_size,
            font_color=self.style.text_color,
        )

        # Recommendations (right column)
        self._add_text_box(
            slide,
            Inches(6.8), Inches(2.8), Inches(6), Inches(0.4),
            "行动建议",
            font_size=14,
            font_color=self.style.primary_color,
            bold=True,
        )

        # Recommendation box
        self._add_shape(
            slide,
            Inches(6.8), Inches(3.2), Inches(5.8), Inches(3.5),
            (255, 250, 240),  # Light orange background
        )

        rec_text = ""
        for rec in insight.recommendations[:4]:
            text = rec[:70] + "..." if len(rec) > 70 else rec
            rec_text += f"→ {text}\n\n"

        self._add_text_box(
            slide,
            Inches(7), Inches(3.4), Inches(5.4), Inches(3.2),
            rec_text,
            font_size=self.style.small_size,
            font_color=self.style.text_color,
        )

    def _add_strategy_slide(self, decision: StrategyResult):
        """Add Three Decisions strategy slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self._add_header(slide, decision.title, decision.decision_type.upper())

        # Summary box
        self._add_shape(
            slide,
            Inches(0.5), Inches(1.5), Inches(12.333), Inches(1.5),
            (240, 248, 255),  # Light blue background
        )

        # Take first part of summary
        summary_lines = decision.summary.split("\n")[:4]
        summary_text = "\n".join(summary_lines)

        self._add_text_box(
            slide,
            Inches(0.7), Inches(1.6), Inches(12), Inches(1.3),
            summary_text,
            font_size=self.style.small_size,
            font_color=self.style.text_color,
        )

        # Strategy items
        y_pos = Inches(3.2)

        for item in decision.items[:5]:
            # Priority badge
            priority_colors = {
                "P0": (199, 0, 11),    # Red
                "P1": (255, 165, 0),   # Orange
                "P2": (34, 139, 34),   # Green
            }
            color = priority_colors.get(item.priority, self.style.primary_color)

            self._add_shape(
                slide,
                Inches(0.5), y_pos, Inches(0.5), Inches(0.35),
                color,
            )
            self._add_text_box(
                slide,
                Inches(0.5), y_pos + Inches(0.05), Inches(0.5), Inches(0.25),
                item.priority,
                font_size=10,
                font_color=(255, 255, 255),
                bold=True,
                align="center",
            )

            # Item name and description
            self._add_text_box(
                slide,
                Inches(1.1), y_pos, Inches(5), Inches(0.35),
                item.name,
                font_size=14,
                font_color=self.style.text_color,
                bold=True,
            )

            desc_text = item.description[:100] + "..." if len(item.description) > 100 else item.description
            self._add_text_box(
                slide,
                Inches(1.1), y_pos + Inches(0.35), Inches(5), Inches(0.4),
                desc_text,
                font_size=self.style.small_size,
                font_color=self.style.light_text_color,
            )

            # KPIs
            if item.kpis:
                kpi_text = " | ".join(item.kpis[:2])
                self._add_text_box(
                    slide,
                    Inches(6.5), y_pos + Inches(0.1), Inches(6), Inches(0.5),
                    f"KPI: {kpi_text}",
                    font_size=10,
                    font_color=self.style.primary_color,
                )

            y_pos += Inches(0.85)

    def _add_conclusion_slide(
        self,
        five_looks: dict[str, InsightResult],
        three_decisions: dict[str, StrategyResult],
        target_operator: str,
    ):
        """Add conclusion slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self._add_header(slide, "总结与建议", "SUMMARY")

        # Key takeaways
        self._add_text_box(
            slide,
            Inches(0.5), Inches(1.6), Inches(6), Inches(0.4),
            f"{target_operator} 战略重点",
            font_size=18,
            font_color=self.style.primary_color,
            bold=True,
        )

        # Get P0 strategies
        strategy = three_decisions.get("strategy")
        key_tasks = three_decisions.get("key_tasks")

        takeaways = []
        if strategy:
            for item in strategy.items:
                if item.priority == "P0":
                    takeaways.append(f"【{item.priority}】{item.name}")

        takeaway_text = "\n".join(takeaways[:4])
        self._add_text_box(
            slide,
            Inches(0.5), Inches(2.1), Inches(6), Inches(2),
            takeaway_text,
            font_size=self.style.body_size,
            font_color=self.style.text_color,
        )

        # Key milestones
        self._add_text_box(
            slide,
            Inches(6.8), Inches(1.6), Inches(6), Inches(0.4),
            "关键里程碑",
            font_size=18,
            font_color=self.style.primary_color,
            bold=True,
        )

        execution = three_decisions.get("execution")
        milestones = []
        if execution:
            for item in execution.items:
                if item.category == "milestone" and item.timeline:
                    milestones.append(f"• {item.timeline}: {item.name}")

        milestone_text = "\n".join(milestones[:5])
        self._add_text_box(
            slide,
            Inches(6.8), Inches(2.1), Inches(6), Inches(2),
            milestone_text,
            font_size=self.style.body_size,
            font_color=self.style.text_color,
        )

        # Bottom message
        self._add_shape(
            slide,
            Inches(0.5), Inches(5), Inches(12.333), Inches(2),
            self.style.primary_color,
        )

        self._add_text_box(
            slide,
            Inches(0.8), Inches(5.3), Inches(11.8), Inches(0.6),
            "下一步行动",
            font_size=20,
            font_color=(255, 255, 255),
            bold=True,
        )

        next_steps = "1. 成立战略执行委员会  2. 启动5G网络加速项目  3. 组建企业解决方案中心  4. 双周跟踪KPI进展"
        self._add_text_box(
            slide,
            Inches(0.8), Inches(5.9), Inches(11.8), Inches(0.8),
            next_steps,
            font_size=self.style.body_size,
            font_color=(255, 255, 255),
        )

    def _add_closing_slide(self):
        """Add closing/thank you slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Full red background
        self._add_shape(
            slide, 0, 0, Inches(13.333), Inches(7.5),
            self.style.primary_color
        )

        # Thank you text
        self._add_text_box(
            slide,
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.2),
            "谢谢",
            font_size=72,
            font_color=(255, 255, 255),
            bold=True,
            align="center",
        )

        self._add_text_box(
            slide,
            Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.6),
            "THANK YOU",
            font_size=32,
            font_color=(255, 200, 200),
            align="center",
        )

        # Footer
        self._add_text_box(
            slide,
            Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5),
            f"BLM Strategic Analysis | Generated {datetime.now().strftime('%Y-%m-%d')}",
            font_size=12,
            font_color=(255, 200, 200),
            align="center",
        )

    # =========================================================================
    # Helper Methods
    # =========================================================================

    def _add_header(self, slide, title: str, subtitle: str = ""):
        """Add standard header to slide."""
        # Red accent bar
        self._add_shape(
            slide, 0, 0, Inches(13.333), Inches(0.08),
            self.style.primary_color
        )

        # Title
        self._add_text_box(
            slide,
            Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
            title,
            font_size=self.style.subtitle_size,
            font_color=self.style.text_color,
            bold=True,
        )

        # Subtitle/category
        if subtitle:
            self._add_text_box(
                slide,
                Inches(0.5), Inches(0.85), Inches(10), Inches(0.4),
                subtitle,
                font_size=12,
                font_color=self.style.primary_color,
                bold=True,
            )

        # Page number area (bottom right)
        self._add_shape(
            slide,
            Inches(12.5), Inches(7), Inches(0.6), Inches(0.35),
            self.style.primary_color,
        )

    def _add_shape(
        self,
        slide,
        left,
        top,
        width,
        height,
        fill_color: tuple,
        shape_type=MSO_SHAPE.RECTANGLE,
    ):
        """Add a shape to the slide."""
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_color)
        shape.line.fill.background()  # No border
        return shape

    def _add_text_box(
        self,
        slide,
        left,
        top,
        width,
        height,
        text: str,
        font_size: int = 14,
        font_color: tuple = (0, 0, 0),
        bold: bool = False,
        align: str = "left",
    ):
        """Add a text box to the slide."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*font_color)
        p.font.bold = bold
        p.font.name = self.style.body_font

        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

        return txBox

    def _sanitize_name(self, name: str) -> str:
        """Sanitize name for filename."""
        return name.lower().replace(" ", "_").replace("/", "_")


def generate_blm_ppt(
    five_looks: dict[str, InsightResult],
    three_decisions: dict[str, StrategyResult],
    target_operator: str,
    competitors: list[str] = None,
    style: str = "huawei",
    output_dir: str = None,
) -> str:
    """Convenience function to generate BLM PPT.

    Args:
        five_looks: Five Looks analysis results
        three_decisions: Three Decisions strategy results
        target_operator: Target operator name
        competitors: List of competitor names
        style: PPT style ("huawei" or "vodafone")
        output_dir: Output directory

    Returns:
        Path to generated PPT file
    """
    generator = BLMPPTGenerator(style=style, output_dir=output_dir)
    return generator.generate(
        five_looks=five_looks,
        three_decisions=three_decisions,
        target_operator=target_operator,
        competitors=competitors,
    )
