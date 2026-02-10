"""Tests for BLM chart engine and PPT generator."""

import os
import tempfile
from pathlib import Path

import pytest

from src.output.ppt_styles import PPTStyle, get_style, VODAFONE_STYLE, DEFAULT_STYLE
from src.output.ppt_charts import BLMChartGenerator, _rgb_to_hex, _rgb_to_mpl
from src.blm.engine import FiveLooksResult
from src.models.provenance import ProvenanceStore, SourceReference, SourceType, Confidence
from src.models.trend import TrendAnalysis, PESTAnalysis, PESTFactor
from src.models.market import (
    MarketCustomerInsight, MarketChange, CustomerSegment, APPEALSAssessment
)
from src.models.competition import (
    CompetitionInsight, PorterForce, CompetitorDeepDive, CompetitorImplication
)
from src.models.self_analysis import (
    SelfInsight, SegmentAnalysis, NetworkAnalysis, BMCCanvas,
    ExposurePoint, SegmentChange, ChangeAttribution
)
from src.models.swot import SWOTAnalysis
from src.models.opportunity import (
    OpportunityInsight, SPANPosition, OpportunityItem
)


# =============================================================================
# Test fixtures
# =============================================================================

@pytest.fixture
def tmp_dir():
    with tempfile.TemporaryDirectory() as d:
        yield d


@pytest.fixture
def style():
    return VODAFONE_STYLE


@pytest.fixture
def chart_gen(tmp_dir, style):
    return BLMChartGenerator(style=style, output_dir=tmp_dir, dpi=72)


@pytest.fixture
def mock_pest():
    return PESTAnalysis(
        political_factors=[
            PESTFactor(dimension="P", dimension_name="Political",
                       factor_name="Gigabit Strategy", impact_type="opportunity",
                       severity="high"),
            PESTFactor(dimension="P", dimension_name="Political",
                       factor_name="Spectrum Regulation", impact_type="both",
                       severity="medium"),
        ],
        economic_factors=[
            PESTFactor(dimension="E", dimension_name="Economic",
                       factor_name="GDP Growth", impact_type="neutral",
                       severity="medium"),
        ],
        society_factors=[
            PESTFactor(dimension="S", dimension_name="Society",
                       factor_name="Digital Adoption", impact_type="opportunity",
                       severity="high"),
        ],
        technology_factors=[
            PESTFactor(dimension="T", dimension_name="Technology",
                       factor_name="5G Rollout", impact_type="opportunity",
                       severity="high"),
            PESTFactor(dimension="T", dimension_name="Technology",
                       factor_name="AI/ML in Networks", impact_type="opportunity",
                       severity="medium"),
        ],
        key_message="Favorable macro environment for 5G investment",
    )


@pytest.fixture
def mock_trends(mock_pest):
    return TrendAnalysis(
        pest=mock_pest,
        industry_market_size="EUR 60B",
        industry_growth_rate="2.5%",
        industry_profit_trend="Stable margins",
        industry_concentration="CR4 = 95%",
        industry_lifecycle_stage="mature",
        new_business_models=["FWA as broadband substitute", "IoT bundling"],
        technology_revolution=["5G SA deployment", "Open RAN"],
        key_success_factors=["Network quality", "Convergence", "B2B growth"],
        value_transfer_trends=["Shift from voice to data revenue"],
        key_message="Mature market with 5G growth opportunity",
    )


@pytest.fixture
def mock_appeals():
    return [
        APPEALSAssessment(dimension="$", dimension_name="Price",
                          our_score=3.5, competitor_scores={"dt": 3.0, "o2": 4.0},
                          customer_priority="critical"),
        APPEALSAssessment(dimension="A1", dimension_name="Availability",
                          our_score=4.0, competitor_scores={"dt": 4.5, "o2": 3.5},
                          customer_priority="important"),
        APPEALSAssessment(dimension="P1", dimension_name="Packaging",
                          our_score=3.8, competitor_scores={"dt": 3.5, "o2": 3.0},
                          customer_priority="important"),
        APPEALSAssessment(dimension="P2", dimension_name="Performance",
                          our_score=4.2, competitor_scores={"dt": 4.8, "o2": 3.0},
                          customer_priority="critical"),
        APPEALSAssessment(dimension="E", dimension_name="Ease",
                          our_score=3.0, competitor_scores={"dt": 3.5, "o2": 3.5},
                          customer_priority="important"),
        APPEALSAssessment(dimension="A2", dimension_name="Assurances",
                          our_score=3.5, competitor_scores={"dt": 4.0, "o2": 3.0},
                          customer_priority="nice_to_have"),
        APPEALSAssessment(dimension="L", dimension_name="Lifecycle",
                          our_score=3.2, competitor_scores={"dt": 3.0, "o2": 3.5},
                          customer_priority="nice_to_have"),
        APPEALSAssessment(dimension="S", dimension_name="Social",
                          our_score=3.0, competitor_scores={"dt": 3.5, "o2": 2.5},
                          customer_priority="nice_to_have"),
    ]


@pytest.fixture
def mock_market_customer(mock_appeals):
    return MarketCustomerInsight(
        market_snapshot={"total_revenue": "EUR 55B", "total_subscribers": "130M",
                         "avg_arpu": "EUR 12.50"},
        changes=[
            MarketChange(change_type="pricing", description="Price war in postpaid",
                         impact_type="threat"),
            MarketChange(change_type="technology", description="FWA adoption growing",
                         impact_type="opportunity"),
        ],
        opportunities=[
            MarketChange(change_type="technology", description="5G enterprise services"),
        ],
        threats=[
            MarketChange(change_type="new_entrant", description="1&1 network launch"),
        ],
        customer_segments=[
            CustomerSegment(segment_name="Value-Conscious", growth_trend="growing",
                            size_estimate="40%", our_share="25%",
                            unmet_needs=["Low-cost 5G"]),
            CustomerSegment(segment_name="Premium Users", growth_trend="stable",
                            size_estimate="20%", our_share="35%"),
            CustomerSegment(segment_name="Enterprise", growth_trend="growing",
                            segment_type="enterprise"),
        ],
        appeals_assessment=mock_appeals,
        customer_value_migration="Shift toward bundled services and convergence",
        key_message="Market favors convergence; price pressure from new entrants",
    )


@pytest.fixture
def mock_competition():
    return CompetitionInsight(
        five_forces={
            "existing_competitors": PorterForce(
                force_name="existing_competitors", force_level="high",
                key_factors=[{"name": "Price war", "description": "Aggressive pricing"}]),
            "new_entrants": PorterForce(
                force_name="new_entrants", force_level="medium",
                key_factors=[{"name": "1&1", "description": "Building own network"}]),
            "substitutes": PorterForce(force_name="substitutes", force_level="low"),
            "supplier_power": PorterForce(force_name="supplier_power", force_level="medium"),
            "buyer_power": PorterForce(force_name="buyer_power", force_level="high"),
        },
        overall_competition_intensity="high",
        competitor_analyses={
            "deutsche_telekom": CompetitorDeepDive(
                operator="deutsche_telekom",
                financial_health={"revenue": "EUR 25B", "margin": "32%"},
                strengths=["Largest network", "Brand trust"],
                weaknesses=["Premium pricing", "Legacy systems"],
                growth_strategy="Fiber and convergence",
            ),
            "telefonica_o2": CompetitorDeepDive(
                operator="telefonica_o2",
                financial_health={"revenue": "EUR 8B"},
                strengths=["Price leadership"],
                weaknesses=["Network quality"],
            ),
        },
        comparison_table={
            "Revenue": {"vodafone": "EUR 12B", "dt": "EUR 25B", "o2": "EUR 8B"},
            "Subscribers": {"vodafone": "30M", "dt": "50M", "o2": "45M"},
        },
        competitive_landscape="Highly competitive 4-player market",
        key_message="High competition intensity; DT leads on network, O2 on price",
    )


@pytest.fixture
def mock_self():
    return SelfInsight(
        financial_health={"revenue": "EUR 12.1B", "ebitda_margin": "32%",
                          "capex_ratio": "18%", "free_cash_flow": "EUR 2.3B"},
        revenue_breakdown={"mobile": "55%", "fixed": "30%", "b2b": "15%"},
        market_positions={"mobile": "#2 (25%)", "fixed": "#3 (20%)"},
        segment_analyses=[
            SegmentAnalysis(
                segment_name="Mobile", segment_id="mobile",
                key_metrics={"revenue": "EUR 6.5B", "subscribers": "22M"},
                health_status="stable",
                changes=[SegmentChange(metric="revenue", direction="stable")],
                attributions=[ChangeAttribution(
                    attribution_type="tariff_competition",
                    description="Price erosion from O2 unlimited plans")],
                key_message="Mobile stable but under pricing pressure"),
            SegmentAnalysis(
                segment_name="Fixed Broadband", segment_id="fixed",
                key_metrics={"revenue": "EUR 3.6B"},
                health_status="weakening",
                key_message="Fixed under pressure from fiber migration"),
        ],
        network=NetworkAnalysis(
            technology_mix={"cable": 70, "fiber": 5, "dsl_resale": 25},
            coverage={"5g": 90, "4g": 99.5},
            investment_direction="Cable to fiber migration",
        ),
        bmc=BMCCanvas(
            key_partners=["Ericsson", "Huawei", "Open Fiber"],
            key_activities=["Network operation", "Customer acquisition"],
            key_resources=["Cable network", "Spectrum licenses"],
            value_propositions=["Converged bundles", "5G speed"],
            customer_relationships=["Retail stores", "Digital channels"],
            channels=["Online", "Retail", "B2B sales"],
            customer_segments=["Consumer", "SME", "Enterprise"],
            cost_structure=["Network opex", "Spectrum", "Content"],
            revenue_streams=["Subscriptions", "B2B services"],
        ),
        strengths=["Strong cable network", "Brand recognition", "5G spectrum"],
        weaknesses=["Limited fiber", "Premium pricing", "Legacy IT"],
        exposure_points=[
            ExposurePoint(trigger_action="1&1 migrating 11M users",
                          side_effect="Network load surge",
                          severity="high"),
        ],
        health_rating="stable",
        key_message="Stable but facing cable-to-fiber transition challenge",
    )


@pytest.fixture
def mock_swot():
    return SWOTAnalysis(
        strengths=["Cable infrastructure", "5G spectrum", "Brand"],
        weaknesses=["Limited FTTH", "Price perception", "Legacy IT"],
        opportunities=["5G enterprise", "FWA", "IoT"],
        threats=["1&1 network launch", "DT fiber expansion", "Price war"],
        so_strategies=["Leverage 5G for enterprise services"],
        wo_strategies=["Partner for fiber access"],
        st_strategies=["Defend base with convergence bundles"],
        wt_strategies=["Selective retreat from unprofitable segments"],
        key_message="Use cable+5G strengths while addressing fiber gap",
    )


@pytest.fixture
def mock_opportunities():
    return OpportunityInsight(
        span_positions=[
            SPANPosition(opportunity_name="5G Enterprise",
                         market_attractiveness=8.5, competitive_position=7.0,
                         bubble_size=3.0, quadrant="grow_invest"),
            SPANPosition(opportunity_name="FWA Broadband",
                         market_attractiveness=7.0, competitive_position=8.0,
                         bubble_size=2.0, quadrant="grow_invest"),
            SPANPosition(opportunity_name="IoT/M2M",
                         market_attractiveness=6.0, competitive_position=4.0,
                         bubble_size=1.5, quadrant="acquire_skills"),
            SPANPosition(opportunity_name="Content Bundling",
                         market_attractiveness=3.0, competitive_position=5.0,
                         bubble_size=1.0, quadrant="harvest"),
        ],
        grow_invest=["5G Enterprise", "FWA Broadband"],
        acquire_skills=["IoT/M2M"],
        harvest=["Content Bundling"],
        avoid_exit=[],
        opportunities=[
            OpportunityItem(name="5G Enterprise", priority="P0",
                            description="B2B 5G services",
                            addressable_market="EUR 2B",
                            time_window="immediate"),
            OpportunityItem(name="FWA Broadband", priority="P0",
                            description="Fixed wireless access",
                            addressable_market="EUR 1.5B"),
            OpportunityItem(name="IoT/M2M", priority="P1",
                            description="Industrial IoT",
                            addressable_market="N/A"),
        ],
        key_message="Focus on 5G enterprise and FWA as primary growth engines",
    )


@pytest.fixture
def mock_provenance():
    prov = ProvenanceStore()
    src = SourceReference(
        source_type=SourceType.FINANCIAL_REPORT_PDF,
        document_name="Vodafone Q3 FY26 Results",
        publisher="Vodafone Group",
        confidence=Confidence.HIGH,
    )
    prov.register_source(src)
    prov.track(12.1, "revenue", operator="vodafone", period="CQ4_2025",
               source=src, unit="EUR B")
    prov.track(32, "ebitda_margin", operator="vodafone", period="CQ4_2025",
               source=src, unit="%")
    prov.track(22, "mobile_subscribers", operator="vodafone", period="CQ4_2025",
               source=src, unit="M")
    return prov


@pytest.fixture
def mock_result(mock_trends, mock_market_customer, mock_competition,
                mock_self, mock_swot, mock_opportunities, mock_provenance):
    return FiveLooksResult(
        target_operator="vodafone_germany",
        market="germany",
        analysis_period="CQ4_2025",
        trends=mock_trends,
        market_customer=mock_market_customer,
        competition=mock_competition,
        self_analysis=mock_self,
        swot=mock_swot,
        opportunities=mock_opportunities,
        provenance=mock_provenance,
    )


# =============================================================================
# Color utility tests
# =============================================================================

class TestColorUtils:
    def test_rgb_to_hex(self):
        assert _rgb_to_hex((255, 0, 0)) == "#FF0000"
        assert _rgb_to_hex((0, 0, 0)) == "#000000"
        assert _rgb_to_hex((230, 0, 0)) == "#E60000"

    def test_rgb_to_mpl(self):
        r, g, b = _rgb_to_mpl((255, 128, 0))
        assert abs(r - 1.0) < 0.01
        assert abs(g - 0.502) < 0.01
        assert abs(b - 0.0) < 0.01


# =============================================================================
# BLMChartGenerator tests
# =============================================================================

class TestBLMChartGenerator:
    def test_init_default_style(self, tmp_dir):
        gen = BLMChartGenerator(output_dir=tmp_dir)
        assert gen.style == DEFAULT_STYLE

    def test_init_with_style(self, chart_gen, style):
        assert chart_gen.style == style

    def test_create_bar_chart(self, chart_gen):
        path = chart_gen.create_bar_chart(
            ["A", "B", "C"], [10, 20, 15], target_category="B",
            title="Test", filename="test_bar.png")
        assert Path(path).exists()
        assert Path(path).stat().st_size > 0

    def test_create_horizontal_bar_chart(self, chart_gen):
        path = chart_gen.create_horizontal_bar_chart(
            ["DT", "Vodafone", "O2"], [50, 25, 20], target_category="Vodafone",
            title="Market Share", value_suffix="%", filename="test_hbar.png")
        assert Path(path).exists()

    def test_create_radar_chart(self, chart_gen):
        path = chart_gen.create_radar_chart(
            ["Price", "Network", "Brand", "Innovation", "Service"],
            {"Vodafone": [70, 80, 75, 65, 60], "DT": [60, 90, 85, 70, 75]},
            target_operator="Vodafone", filename="test_radar.png")
        assert Path(path).exists()

    def test_create_multi_line_trend(self, chart_gen):
        path = chart_gen.create_multi_line_trend(
            ["Q1", "Q2", "Q3", "Q4"],
            {"Vodafone": [10, 11, 12, 13], "DT": [15, 16, 17, 18]},
            target_operator="Vodafone", filename="test_trend.png")
        assert Path(path).exists()

    def test_create_multi_line_trend_with_none(self, chart_gen):
        path = chart_gen.create_multi_line_trend(
            ["Q1", "Q2", "Q3", "Q4"],
            {"Vodafone": [10, None, 12, 13]},
            filename="test_trend_none.png")
        assert Path(path).exists()

    def test_create_stacked_bar(self, chart_gen):
        path = chart_gen.create_stacked_bar(
            ["Q1", "Q2"], {"Mobile": [5, 6], "Fixed": [3, 3]},
            filename="test_stacked.png")
        assert Path(path).exists()

    def test_create_kpi_table_chart(self, chart_gen):
        path = chart_gen.create_kpi_table_chart(
            ["Revenue", "EBITDA"], {"Vodafone": ["12B", "4B"], "DT": ["25B", "8B"]},
            target_operator="Vodafone", filename="test_kpi.png")
        assert Path(path).exists()

    def test_create_donut_gauges(self, chart_gen):
        path = chart_gen.create_donut_gauges(
            ["Vodafone", "DT", "O2"], [90, 95, 80],
            target_label="Vodafone", title="5G Coverage", filename="test_donut.png")
        assert Path(path).exists()

    def test_create_gap_analysis_chart(self, chart_gen):
        path = chart_gen.create_gap_analysis_chart(
            ["Price", "Network", "Brand"], [70, 80, 75], [60, 90, 85],
            target_name="Vodafone", leader_name="DT", filename="test_gap.png")
        assert Path(path).exists()

    def test_create_heatmap(self, chart_gen):
        path = chart_gen.create_heatmap(
            ["VF", "DT"], ["VF", "DT"], [[0, 50], [30, 0]],
            filename="test_heatmap.png")
        assert Path(path).exists()

    def test_create_segment_comparison(self, chart_gen):
        path = chart_gen.create_segment_comparison(
            ["Q1", "Q2"], {"Vodafone": [5, 6], "DT": [8, 9]},
            target_operator="Vodafone", filename="test_segcomp.png")
        assert Path(path).exists()

    def test_create_priority_chart(self, chart_gen):
        path = chart_gen.create_priority_chart(
            ["5G Enterprise", "FWA", "IoT"], ["P0", "P0", "P1"],
            filename="test_priority.png")
        assert Path(path).exists()

    def test_create_timeline_chart(self, chart_gen):
        path = chart_gen.create_timeline_chart(
            [{"date": "Q1", "name": "Launch", "priority": "P0"},
             {"date": "Q2", "name": "Scale", "priority": "P1"}],
            filename="test_timeline.png")
        assert Path(path).exists()

    # --- New chart types ---

    def test_create_span_bubble_chart(self, chart_gen, mock_opportunities):
        path = chart_gen.create_span_bubble_chart(
            mock_opportunities.span_positions, filename="test_span.png")
        assert Path(path).exists()

    def test_create_span_bubble_chart_with_dicts(self, chart_gen):
        positions = [
            {"opportunity_name": "Test", "competitive_position": 7,
             "market_attractiveness": 8, "bubble_size": 2, "quadrant": "grow_invest"},
        ]
        path = chart_gen.create_span_bubble_chart(positions, filename="test_span_dict.png")
        assert Path(path).exists()

    def test_create_porter_five_forces(self, chart_gen, mock_competition):
        path = chart_gen.create_porter_five_forces(
            mock_competition.five_forces, filename="test_porter.png")
        assert Path(path).exists()

    def test_create_porter_five_forces_with_strings(self, chart_gen):
        forces = {"existing_competitors": "high", "new_entrants": "medium",
                  "substitutes": "low", "supplier_power": "medium",
                  "buyer_power": "high"}
        path = chart_gen.create_porter_five_forces(forces, filename="test_porter_str.png")
        assert Path(path).exists()

    def test_create_swot_matrix(self, chart_gen, mock_swot):
        path = chart_gen.create_swot_matrix(mock_swot, filename="test_swot.png")
        assert Path(path).exists()

    def test_create_swot_matrix_with_dict(self, chart_gen):
        path = chart_gen.create_swot_matrix(
            {"strengths": ["A"], "weaknesses": ["B"],
             "opportunities": ["C"], "threats": ["D"]},
            filename="test_swot_dict.png")
        assert Path(path).exists()

    def test_create_appeals_radar(self, chart_gen, mock_appeals):
        path = chart_gen.create_appeals_radar(
            mock_appeals, target_operator="vodafone",
            filename="test_appeals.png")
        assert Path(path).exists()

    def test_create_appeals_radar_empty(self, chart_gen):
        path = chart_gen.create_appeals_radar([], filename="test_appeals_empty.png")
        assert Path(path).exists()

    def test_create_bmc_canvas(self, chart_gen, mock_self):
        path = chart_gen.create_bmc_canvas(mock_self.bmc, filename="test_bmc.png")
        assert Path(path).exists()

    def test_create_pest_dashboard(self, chart_gen, mock_pest):
        path = chart_gen.create_pest_dashboard(mock_pest, filename="test_pest.png")
        assert Path(path).exists()

    def test_create_pest_dashboard_with_dict(self, chart_gen):
        pest = {"political_factors": [], "economic_factors": [],
                "society_factors": [], "technology_factors": []}
        path = chart_gen.create_pest_dashboard(pest, filename="test_pest_dict.png")
        assert Path(path).exists()


# =============================================================================
# BLMPPTGenerator tests
# =============================================================================

class TestBLMPPTGenerator:
    @pytest.fixture
    def ppt_gen(self, tmp_dir, style):
        from src.output.ppt_generator import BLMPPTGenerator
        return BLMPPTGenerator(
            style=style, operator_id="vodafone_germany",
            output_dir=tmp_dir, chart_dpi=72)

    def test_generate_draft_mode(self, ppt_gen, mock_result):
        path = ppt_gen.generate(mock_result, mode="draft")
        assert Path(path).exists()
        assert path.endswith(".pptx")
        assert Path(path).stat().st_size > 10000  # Non-trivial file

    def test_generate_custom_filename(self, ppt_gen, mock_result):
        path = ppt_gen.generate(mock_result, filename="custom_name.pptx")
        assert Path(path).name == "custom_name.pptx"

    def test_draft_slide_count(self, ppt_gen, mock_result):
        path = ppt_gen.generate(mock_result, mode="draft")
        from pptx import Presentation
        prs = Presentation(path)
        slide_count = len(prs.slides)
        # Expect 30+ slides with all sections populated
        assert slide_count >= 25, f"Expected >= 25 slides, got {slide_count}"
        assert slide_count <= 50, f"Expected <= 50 slides, got {slide_count}"

    def test_final_mode_filters_slides(self, ppt_gen, mock_result):
        # First generate draft to get slide count
        draft_path = ppt_gen.generate(mock_result, mode="draft",
                                      filename="draft.pptx")
        from pptx import Presentation
        draft_prs = Presentation(draft_path)
        draft_count = len(draft_prs.slides)

        # Now generate final with some slides removed
        decisions = {
            "trend_deep_dive": "remove",
            "market_deep_dive": "remove",
            "opp_deep_dive": "remove",
            "new_entrants": "remove",
        }
        final_path = ppt_gen.generate(mock_result, mode="final",
                                      user_decisions=decisions,
                                      filename="final.pptx")
        final_prs = Presentation(final_path)
        final_count = len(final_prs.slides)

        assert final_count < draft_count

    def test_key_message_bar_present(self, ppt_gen, mock_result):
        """Verify that key message bars are present on content slides."""
        path = ppt_gen.generate(mock_result, mode="draft")
        from pptx import Presentation
        prs = Presentation(path)

        key_message_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text and "Key Message:" in para.text:
                            key_message_count += 1

        # Most content slides should have a key message
        assert key_message_count >= 10, \
            f"Expected >= 10 key messages, found {key_message_count}"

    def test_cover_slide_has_operator_info(self, ppt_gen, mock_result):
        path = ppt_gen.generate(mock_result, mode="draft")
        from pptx import Presentation
        prs = Presentation(path)

        cover = prs.slides[0]
        texts = [shape.text_frame.text for shape in cover.shapes if shape.has_text_frame]
        all_text = " ".join(texts)
        assert "Vodafone" in all_text
        assert "CQ4_2025" in all_text

    def test_style_applied(self, ppt_gen, mock_result):
        """Verify the PPT uses the correct brand style."""
        path = ppt_gen.generate(mock_result, mode="draft")
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        prs = Presentation(path)

        # Check first slide has brand color
        cover = prs.slides[0]
        found_brand = False
        expected = RGBColor(*VODAFONE_STYLE.primary_color)
        for shape in cover.shapes:
            if shape.shape_type == 1:  # Auto shape
                if hasattr(shape, 'fill') and shape.fill.type is not None:
                    try:
                        if shape.fill.fore_color.rgb == expected:
                            found_brand = True
                    except Exception:
                        pass
        assert found_brand, "Brand color not found on cover slide"

    def test_provenance_appendix_exists(self, ppt_gen, mock_result):
        path = ppt_gen.generate(mock_result, mode="draft")
        from pptx import Presentation
        prs = Presentation(path)

        found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and "Data Provenance" in shape.text_frame.text:
                    found = True
                    break
        assert found, "Provenance appendix slide not found"

    def test_generate_with_none_sections(self, tmp_dir, style):
        """Test graceful handling of None analysis sections."""
        from src.output.ppt_generator import BLMPPTGenerator
        gen = BLMPPTGenerator(style=style, output_dir=tmp_dir, chart_dpi=72)
        result = FiveLooksResult(
            target_operator="test_op",
            market="test_market",
            analysis_period="CQ1_2026",
        )
        path = gen.generate(result, mode="draft")
        assert Path(path).exists()
        # Should still have cover, TOC, summary, back cover at minimum
        from pptx import Presentation
        prs = Presentation(path)
        assert len(prs.slides) >= 5

    def test_section_dividers_present(self, ppt_gen, mock_result):
        path = ppt_gen.generate(mock_result, mode="draft")
        from pptx import Presentation
        prs = Presentation(path)

        divider_texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if ("Look at Trends" in text or "Look at Market" in text
                            or "Look at Competition" in text or "Look at Self" in text
                            or "Look at Opportunities" in text or "SWOT Synthesis" in text):
                        divider_texts.append(text)
        assert len(divider_texts) >= 5, f"Expected >= 5 section dividers, found {len(divider_texts)}"

    def test_slide_width_is_widescreen(self, ppt_gen, mock_result):
        path = ppt_gen.generate(mock_result, mode="draft")
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation(path)
        # 16:9 = 13.333 x 7.5
        assert abs(prs.slide_width - Inches(13.333)) < Inches(0.01)
        assert abs(prs.slide_height - Inches(7.5)) < Inches(0.01)
